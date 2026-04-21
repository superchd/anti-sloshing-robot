#!/usr/bin/env python3
"""
create_sbsfc_figure.py
Recreate Choi et al. 2024 Fig. 1 — Overview of the proposed slosh-free control
scheme for food serving robots — as an editable PowerPoint slide.

Run:    python3 create_sbsfc_figure.py
Output: results/sbsfc_overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os, math

# ─── Colors ──────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

BLACK     = rgb(0x1A, 0x1A, 0x1A)
WHITE     = rgb(0xFF, 0xFF, 0xFF)
BLUE      = rgb(0x26, 0x66, 0xCF)   # Upper body controller border
BLUE_BG   = rgb(0xEB, 0xF3, 0xFF)
BLUE_FG   = rgb(0x1A, 0x50, 0xAA)
YELLOW    = rgb(0xCC, 0x99, 0x00)    # Lower body controller border
YELLOW_BG = rgb(0xFF, 0xFB, 0xE6)
YELLOW_FG = rgb(0x99, 0x77, 0x00)
RED       = rgb(0xCC, 0x22, 0x22)    # SBSFC outer border
RED_DASH  = rgb(0xDD, 0x44, 0x44)
GRAY      = rgb(0x55, 0x55, 0x55)
LTGRAY    = rgb(0xAA, 0xAA, 0xAA)
GREEN     = rgb(0x22, 0x88, 0x22)
ORANGE    = rgb(0xDD, 0x88, 0x00)
BOX_BG    = rgb(0xFA, 0xFA, 0xFA)   # Default box fill
BOX_BDR   = rgb(0x33, 0x33, 0x33)   # Default box border
SIG_COLOR = rgb(0x22, 0x22, 0x22)   # Signal arrow color
BLUE_SIG  = rgb(0x22, 0x55, 0xCC)   # Blue signal lines
YEL_SIG   = rgb(0xBB, 0x88, 0x00)   # Yellow signal lines

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT
RGT = PP_ALIGN.RIGHT

# ─── Presentation setup ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# White background
bg = slide.background
bf = bg.fill
bf.solid()
bf.fore_color.rgb = WHITE

# ─── Helpers ─────────────────────────────────────────────────────────────────

def box(l, t, w, h, fill=BOX_BG, border=BOX_BDR, bpt=1.5, dash=False, rounded=False):
    kind = 5 if rounded else 1
    s = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(bpt)
    if dash:
        ln = s.line._ln
        for old in ln.findall(qn('a:prstDash')): ln.remove(old)
        pd = etree.SubElement(ln, qn('a:prstDash'))
        pd.set('val', 'lgDash')
    return s

def nobox(l, t, w, h):
    """Invisible rectangle (no fill, no border) for text."""
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.background()
    s.line.fill.background()
    return s

def write(shape, rows, anchor='ctr'):
    tf = shape.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', anchor)
        bp.set('lIns', '36000')  # reduced margin
        bp.set('rIns', '36000')
        bp.set('tIns', '18000')
        bp.set('bIns', '18000')
    tf.clear()
    for i, (s, sz, bd, co, al) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if not s: continue
        r = p.add_run()
        r.text = s
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = co

def label(l, t, w, h, text, sz, bold, color, al=CTR):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = al
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def italic_label(l, t, w, h, text, sz, color, al=CTR):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = al
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.italic = True
    r.font.bold = True
    r.font.color.rgb = color
    return tb

def fwd_arrow(x1, y1, x2, y2, color=SIG_COLOR, pt=2.0):
    """Arrow from (x1,y1) → (x2,y2), arrowhead at destination."""
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    ln = c.line._ln
    for tag in [qn('a:headEnd'), qn('a:tailEnd')]:
        for el in ln.findall(tag): ln.remove(el)
    etree.SubElement(ln, qn('a:headEnd')).set('type', 'none')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')

def line(x1, y1, x2, y2, color=SIG_COLOR, pt=1.5):
    """Plain line, no arrowhead."""
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(pt)

def summing_junction(cx, cy, r=0.14):
    """Small circle with + inside (summing junction)."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r), Inches(2*r), Inches(2*r))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BOX_BDR
    s.line.width = Pt(1.5)
    return s

def plus_minus_label(cx, cy, text, sz=8, dx=0, dy=-0.18, color=BLACK):
    """Small +/- label near a summing junction."""
    label(cx + dx - 0.1, cy + dy, 0.2, 0.16, text, sz, True, color)


# ═══════════════════════════════════════════════════════════════════════════════
#  TOP SECTION: SBSFC Block Diagram
# ═══════════════════════════════════════════════════════════════════════════════

# --- SBSFC outer boundary (red dashed) ---
sbsfc = box(0.3, 0.45, 8.6, 3.65, WHITE, RED_DASH, bpt=2.5, dash=True)
label(0.35, 3.7, 4.5, 0.35,
      "[Self-balancing slosh-free control (SBSFC)]", 10, True, BLACK, LFT)

# --- Upper body Controller group (blue dashed) ---
box(0.5, 0.75, 5.85, 1.95, BLUE_BG, BLUE, bpt=2.0, dash=True)
italic_label(1.9, 0.55, 2.8, 0.25, "Upper body Controller", 10, BLUE_FG)

# --- Lower body Controller group (yellow dashed) ---
box(4.95, 2.95, 2.65, 0.95, YELLOW_BG, YELLOW, bpt=2.0, dash=True)
label(5.8, 3.62, 2.0, 0.25, "Lower body Controller", 9, True, YELLOW_FG)

# ─── Blocks ──────────────────────────────────────────────────────────────────

# Reference Shaping
rs = box(1.1, 1.20, 1.35, 0.50, BOX_BG, BOX_BDR, 1.5)
write(rs, [("Reference", 9, True, BLACK, CTR),
           ("Shaping",   9, True, BLACK, CTR)])

# Upper body controller
uc = box(3.1, 1.20, 1.60, 0.50, BOX_BG, BOX_BDR, 1.5)
write(uc, [("Upper body", 9, True, BLACK, CTR),
           ("controller", 9, True, BLACK, CTR)])

# Disturbance observer
do = box(3.2, 2.10, 1.60, 0.50, BOX_BG, BOX_BDR, 1.5)
write(do, [("Disturbance", 9, True, BLACK, CTR),
           ("observer",    9, True, BLACK, CTR)])

# Compensator (Lower body)
comp = box(5.15, 3.08, 1.70, 0.65, BOX_BG, YELLOW, 1.5)
write(comp, [("Compensator",  9, True, BLACK, CTR),
             ("sgn(\u03c8\u1d38\u209c)", 10, True, YELLOW_FG, CTR)])

# ─── State vector q_t display ────────────────────────────────────────────────
qbox = nobox(1.75, 1.85, 1.10, 0.90)
write(qbox, [
    ("q\u209c =",            10, True, BLACK, CTR),
    ("\u03c8\u209c",          9, False, BLACK, CTR),
    ("\u03c8\u0307\u209c",    9, False, BLACK, CTR),
    ("x\u209c",              9, False, BLACK, CTR),
    ("x\u0307\u209c",        9, False, BLACK, CTR),
], anchor='ctr')
# Left bracket decoration
line(1.85, 1.98, 1.85, 2.60, BLACK, 1.5)
line(1.85, 1.98, 1.92, 1.98, BLACK, 1.5)
line(1.85, 2.60, 1.92, 2.60, BLACK, 1.5)
# Right bracket
line(2.68, 1.98, 2.68, 2.60, BLACK, 1.5)
line(2.68, 1.98, 2.61, 1.98, BLACK, 1.5)
line(2.68, 2.60, 2.61, 2.60, BLACK, 1.5)

# ─── Signal labels ───────────────────────────────────────────────────────────

# v_ref input
label(0.08, 1.22, 0.65, 0.28, "v\u1d63\u2091\u2092", 11, True, BLACK)
# v_d between blocks
label(2.45, 0.98, 0.55, 0.25, "v\u1d48", 11, True, BLACK)
# u_t^u from upper body controller
label(5.0, 0.98, 0.55, 0.25, "u\u209c\u1d58", 11, True, BLUE_FG)
# d_hat_t from disturbance observer
label(5.15, 2.15, 0.55, 0.25, "d\u0302\u209c", 11, True, BLUE_FG)
# u_t^L from compensator
label(6.95, 3.18, 0.55, 0.25, "u\u209c\u1d38", 11, True, YELLOW_FG)
# d_t external disturbance
label(7.65, 0.30, 0.55, 0.25, "d\u209c", 12, True, RED)
# u_t = x_ddot output
label(8.35, 1.20, 1.0, 0.28, "u\u209c = x\u0308\u209c", 11, True, BLACK)

# ─── Summing junctions ──────────────────────────────────────────────────────

# Sum1: u_t^u + (-d_hat) + u_t^L  at (6.6, 1.45)
sj1 = summing_junction(6.6, 1.45, 0.14)
# Sum2: + d_t  at (7.85, 1.45)
sj2 = summing_junction(7.85, 1.45, 0.14)

# +/- labels on summing junctions
plus_minus_label(6.40, 1.22, "+", 8, color=BLACK)     # u_t^u arrives from left
plus_minus_label(6.40, 1.65, "\u2013", 8, color=RED)   # -d_hat from below
plus_minus_label(7.65, 1.22, "+", 8, color=BLACK)      # from sum1
plus_minus_label(8.00, 1.22, "+", 8, color=RED)        # d_t from above

# ─── Signal flow arrows ─────────────────────────────────────────────────────

# v_ref → Reference Shaping
fwd_arrow(0.65, 1.45, 1.10, 1.45, BLACK, 2.0)

# Reference Shaping → Upper body controller
fwd_arrow(2.45, 1.45, 3.10, 1.45, BLACK, 2.0)

# Upper body controller → u_t^u → Sum1
fwd_arrow(4.70, 1.45, 6.46, 1.45, BLUE_SIG, 2.0)

# q_t feedback down from v_d line
line(2.70, 1.45, 2.70, 1.85, BLACK, 1.5)           # vertical down to q_t
fwd_arrow(2.70, 1.72, 2.70, 1.85, BLACK, 1.5)

# q_t → Upper body controller (up-right)
fwd_arrow(2.70, 1.85, 3.10, 1.60, GRAY, 1.2)

# q_t → Disturbance observer (right)
fwd_arrow(2.85, 2.35, 3.20, 2.35, BLACK, 1.5)

# Disturbance observer → d_hat → Sum1 (up and right path)
fwd_arrow(4.80, 2.35, 6.60, 2.35, BLUE_SIG, 1.5)   # horizontal right
line(6.60, 2.35, 6.60, 1.59, BLUE_SIG, 1.5)         # vertical up to sum1

# Compensator → u_t^L → Sum1 area (up from lower body)
fwd_arrow(6.85, 3.40, 7.30, 3.40, YEL_SIG, 1.5)     # horizontal right
line(7.30, 3.40, 7.30, 1.59, YEL_SIG, 1.5)           # vertical up
fwd_arrow(7.30, 1.80, 7.30, 1.59, YEL_SIG, 1.5)      # arrow into junction area

# Compensator ← ψ_t^L feedback (from right side, loops back)
label(5.10, 3.65, 0.50, 0.22, "\u03c8\u209c\u1d38", 9, True, YELLOW_FG)

# Sum1 → Sum2
fwd_arrow(6.74, 1.45, 7.71, 1.45, BLACK, 2.0)

# d_t → Sum2 (vertical down from above)
fwd_arrow(7.85, 0.60, 7.85, 1.31, RED, 2.0)

# Sum2 → u_t output (right)
fwd_arrow(7.99, 1.45, 8.90, 1.45, BLACK, 2.5)

# Output arrow to robot (continues right)
fwd_arrow(8.90, 1.45, 9.30, 1.45, BLACK, 2.5)

# ═══════════════════════════════════════════════════════════════════════════════
#  TOP-RIGHT: Robot Illustration Placeholder
# ═══════════════════════════════════════════════════════════════════════════════

rp = box(9.5, 0.5, 3.2, 3.2, rgb(0xF8, 0xF8, 0xF8), LTGRAY, 1.0, dash=True)
write(rp, [
    ("[Robot Illustration]", 11, True, GRAY, CTR),
    ("", 6, False, GRAY, CTR),
    ("Insert photo/drawing", 9, False, LTGRAY, CTR),
    ("of food-serving robot", 9, False, LTGRAY, CTR),
    ("(Choi et al. Fig.1 right)", 8, False, LTGRAY, CTR),
])

# ═══════════════════════════════════════════════════════════════════════════════
#  BOTTOM-LEFT: 1D Moving Space (schematic)
# ═══════════════════════════════════════════════════════════════════════════════

label(0.3, 4.25, 4.2, 0.30,
      "[1-dimensional moving space omitting steering]", 10, True, BLACK, LFT)

# Robot icon (simplified top-view rounded rect)
robot_top = box(0.8, 4.75, 0.60, 1.10, rgb(0xE8, 0xE8, 0xE8), GRAY, 1.5, rounded=True)
write(robot_top, [("", 1, False, BLACK, CTR)])

# ω_t rotation arrow label
label(1.10, 4.55, 0.45, 0.25, "\u03c9\u209c", 11, True, RED)

# x_t position label
label(0.75, 5.95, 0.55, 0.25, "x\u209c", 11, True, BLACK)

# Arrow from robot to destination
fwd_arrow(1.50, 5.30, 3.10, 5.30, BLACK, 2.0)

# Destination box
dest = box(3.10, 5.05, 1.10, 0.50, rgb(0xE8, 0xF5, 0xFF), BLUE, 1.5)
write(dest, [("Destination", 9, True, BLUE_FG, CTR)])

# x_d label
label(3.20, 5.55, 0.55, 0.25, "x\u1d48", 11, True, BLUE_FG)

# Cup on robot (small oval)
cup = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(0.90), Inches(4.85), Inches(0.40), Inches(0.40))
cup.fill.solid()
cup.fill.fore_color.rgb = rgb(0xCC, 0xDD, 0xFF)
cup.line.color.rgb = BLUE
cup.line.width = Pt(1)

# ═══════════════════════════════════════════════════════════════════════════════
#  BOTTOM-CENTER: Comparison of Liquids (placeholder)
# ═══════════════════════════════════════════════════════════════════════════════

label(0.5, 6.35, 3.5, 0.30,
      "[Comparison of liquids]", 10, True, BLACK, LFT)

# Two placeholder boxes for photos
ph_slosh = box(0.6, 6.65, 1.40, 0.65,
               rgb(0xFF, 0xF0, 0xE0), ORANGE, 1.0, rounded=True)
write(ph_slosh, [("Photo:", 7, False, GRAY, CTR),
                 ("sloshing", 8, True, GRAY, CTR)])

ph_calm = box(2.15, 6.65, 1.40, 0.65,
              rgb(0xE6, 0xF7, 0xED), GREEN, 1.0, rounded=True)
write(ph_calm, [("Photo:", 7, False, GRAY, CTR),
                ("no sloshing", 8, True, GRAY, CTR)])

label(0.60, 7.28, 1.40, 0.22, "With sloshing", 8, True, ORANGE, CTR)
label(2.15, 7.28, 1.40, 0.22, "Without sloshing", 8, True, GREEN, CTR)

# ═══════════════════════════════════════════════════════════════════════════════
#  BOTTOM-RIGHT: Physical Robot Schematic
# ═══════════════════════════════════════════════════════════════════════════════

# Main robot body outline (simplified side view)
body = box(6.3, 4.65, 1.40, 2.10,
           rgb(0xF5, 0xF5, 0xF5), GRAY, 1.5, rounded=True)

# Tray on top of robot
tray = box(5.9, 4.55, 2.20, 0.18, rgb(0xCC, 0xDD, 0xCC), GRAY, 1.0)

# IMUs label
imus = box(5.2, 4.95, 0.70, 0.35, rgb(0xFF, 0xFB, 0xE6), YELLOW, 1.5)
write(imus, [("IMUs", 8, True, BLACK, CTR)])

# Push & Pull label (top center)
pp = box(5.0, 4.30, 0.85, 0.50, rgb(0xFF, 0xEE, 0xDD), ORANGE, 1.5, rounded=True)
write(pp, [("Push", 8, True, BLACK, CTR),
           ("& Pull", 8, True, BLACK, CTR)])

# Arrow from Push & Pull to robot
fwd_arrow(5.45, 4.80, 6.30, 5.20, ORANGE, 1.5)

# Tilt angles label
label(8.0, 4.55, 1.80, 0.30,
      "[Tilt angles", 9, True, BLACK, LFT)
label(8.0, 4.75, 1.80, 0.30,
      "measured by IMUs]", 9, True, BLACK, LFT)

# ψ_t angle label (upper body tilt)
label(7.8, 5.10, 0.55, 0.25, "\u03c8\u209c", 12, True, BLUE_FG)

# ψ_t^L angle label (lower body tilt)
label(7.0, 6.10, 0.55, 0.25, "\u03c8\u209c\u1d38", 12, True, YELLOW_FG)

# Dashed angle arcs (ψ_t — represented as dashed line from vertical)
line(7.00, 4.75, 7.00, 6.00, BLUE, 1.0)      # vertical reference
line(7.00, 4.75, 7.30, 5.80, BLUE, 1.5)       # tilted body line

# Lower angle dashed lines
line(6.80, 5.80, 6.80, 6.80, YELLOW, 1.0)     # vertical reference lower
line(6.80, 5.80, 7.10, 6.70, YELLOW, 1.5)     # tilted lower body

# Pendulum length label
label(8.5, 5.80, 0.40, 0.25, "l", 12, True, BLACK)
line(8.10, 4.75, 8.10, 6.30, BLACK, 1.5)       # pendulum length line
# Small tick marks
line(8.00, 4.75, 8.20, 4.75, BLACK, 1.0)
line(8.00, 6.30, 8.20, 6.30, BLACK, 1.0)

# Rigid connectors label
label(9.0, 5.40, 1.80, 0.60,
      "Rigid\nconnectors\nbetween\nupper body\nand lower body",
      8, True, BLACK, LFT)

# Front and rear casters label
label(5.0, 6.50, 1.30, 0.45,
      "Front\nand rear\ncasters",
      8, True, BLACK, CTR)

# Ground line with hatching
line(4.8, 6.95, 10.5, 6.95, GRAY, 2.0)
for i in range(20):
    x = 4.9 + i * 0.28
    line(x, 6.95, x - 0.12, 7.10, GRAY, 0.8)

label(5.5, 7.05, 3.0, 0.25,
      "Slightly uneven ground", 9, True, GRAY, CTR)

# Wheel outlines (simplified)
wh1 = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(6.20), Inches(6.60), Inches(0.30), Inches(0.35))
wh1.fill.solid(); wh1.fill.fore_color.rgb = rgb(0x44, 0x44, 0x44)
wh1.line.color.rgb = BLACK; wh1.line.width = Pt(1.5)

wh2 = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(7.50), Inches(6.60), Inches(0.30), Inches(0.35))
wh2.fill.solid(); wh2.fill.fore_color.rgb = rgb(0x44, 0x44, 0x44)
wh2.line.color.rgb = BLACK; wh2.line.width = Pt(1.5)

# ═══════════════════════════════════════════════════════════════════════════════
#  Caption at bottom
# ═══════════════════════════════════════════════════════════════════════════════

label(0.3, 7.15, 12.5, 0.35,
      "Fig. 1.  Overview of the proposed slosh-free control scheme for food serving robots. "
      "(One-dimensional schematic view omitting steering operation).",
      9, False, GRAY, CTR)

# Source reference
label(9.5, 7.15, 3.5, 0.30,
      "Adapted from Choi et al. (2024)", 8, False, LTGRAY, RGT)

# ═══════════════════════════════════════════════════════════════════════════════
#  Save
# ═══════════════════════════════════════════════════════════════════════════════

out_dir = os.path.join(os.path.dirname(__file__), 'results')
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'sbsfc_overview.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
print("Open in PowerPoint and edit any element — all shapes are editable.")
print("Replace placeholder boxes with actual photos/drawings from the paper.")
