#!/usr/bin/env python3
"""
create_sbsfc_combined.py
Single slide that combines:
  (top)    SBSFC architecture strip — 3 feedback blocks color-coded
  (bottom) "Why three blocks?" explanation columns, color-matched

Run:    python3 create_sbsfc_combined.py
Output: results/sbsfc_combined.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

def rgb(r, g, b): return RGBColor(r, g, b)

BLACK  = rgb(0x1A, 0x1A, 0x1A)
WHITE  = rgb(0xFF, 0xFF, 0xFF)
GRAY   = rgb(0x66, 0x66, 0x66)
LTGRAY = rgb(0xCC, 0xCC, 0xCC)

BLUE_FG  = rgb(0x1A, 0x50, 0xAA)
BLUE_BG  = rgb(0xEB, 0xF3, 0xFF)
BLUE_BRD = rgb(0x26, 0x66, 0xCF)

GRN_FG   = rgb(0x1E, 0x6E, 0x2E)
GRN_BG   = rgb(0xEB, 0xF7, 0xEE)
GRN_BRD  = rgb(0x2E, 0x99, 0x4A)

RED_FG   = rgb(0xA8, 0x1E, 0x1E)
RED_BG   = rgb(0xFC, 0xEC, 0xEC)
RED_BRD  = rgb(0xCC, 0x33, 0x33)

GOLD_FG  = rgb(0x99, 0x66, 0x00)
GOLD_BG  = rgb(0xFF, 0xF7, 0xE0)
GOLD_BRD = rgb(0xCC, 0x99, 0x00)

PROB_BG  = rgb(0xFF, 0xEC, 0xEC)
PROB_BR  = rgb(0xCC, 0x66, 0x66)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
#  TITLE
# ============================================================
tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.5))
tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
p = tf.paragraphs[0]; p.alignment = CTR
r = p.add_run(); r.text = "SBSFC Architecture  —  Why We Need Three Feedback Blocks"
r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = BLACK

tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.62), Inches(12.5), Inches(0.3))
tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
p = tf.paragraphs[0]; p.alignment = CTR
r = p.add_run(); r.text = "Each colored block in the diagram maps to one column below"
r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = GRAY

# ============================================================
#  TOP:  ARCHITECTURE STRIP
# ============================================================
DIAG_Y = 1.00
DIAG_H = 2.35


def box(x, y, w, h, fg, bg, brd, text, font_size=11, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = bg
    shp.line.color.rgb = brd; shp.line.width = Pt(2)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = text
    r.font.size = Pt(font_size); r.font.bold = bold; r.font.color.rgb = fg
    return shp


def circle(x, y, d, text, color=BLACK, fsize=14):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = color; shp.line.width = Pt(2)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = text
    r.font.size = Pt(fsize); r.font.bold = True; r.font.color.rgb = color
    return shp


def arrow(x1, y1, x2, y2, color=BLACK, w=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                   Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    line = c.line._get_or_add_ln()
    tail = etree.SubElement(line, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return c


def text_at(x, y, w, h, text, fsize=10, color=BLACK, bold=False, italic=False,
            align=CTR):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fsize); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color


# ---- Architecture outline (dashed container) ----
outline = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.4), Inches(DIAG_Y),
                                 Inches(10.1), Inches(DIAG_H))
outline.fill.background()
outline.line.color.rgb = LTGRAY; outline.line.width = Pt(1.2)
outline.line.dash_style = 7  # dash
outline.shadow.inherit = False
text_at(0.55, DIAG_Y + 0.03, 4.0, 0.25,
        "Self-balancing slosh-free controller (SBSFC)",
        fsize=9, italic=True, color=GRAY, align=LFT)

# ---- Reference shaping (feedforward, gold) ----
box(0.65, DIAG_Y + 0.55, 1.5, 0.65, GOLD_FG, GOLD_BG, GOLD_BRD,
    "Reference\nShaping  F_c(s)", font_size=10)
text_at(0.65, DIAG_Y + 1.25, 1.5, 0.22, "(feedforward)",
        fsize=8, italic=True, color=GOLD_FG, align=CTR)
text_at(0.10, DIAG_Y + 0.75, 0.6, 0.3, "v_ref", fsize=10, italic=True,
        color=GRAY, align=CTR)
arrow(0.55, DIAG_Y + 0.88, 0.65, DIAG_Y + 0.88, GRAY)

# ---- LQT (blue) ----
box(2.65, DIAG_Y + 0.55, 1.7, 0.65, BLUE_FG, BLUE_BG, BLUE_BRD,
    "LQT Controller", font_size=11)
arrow(2.15, DIAG_Y + 0.88, 2.65, DIAG_Y + 0.88, GRAY)
text_at(2.25, DIAG_Y + 0.42, 0.6, 0.25, "v_d", fsize=9, italic=True,
        color=GRAY, align=CTR)

# ---- DOB (red) ----
box(2.65, DIAG_Y + 1.45, 1.7, 0.65, RED_FG, RED_BG, RED_BRD,
    "Disturbance\nObserver", font_size=10)

# ---- Compensator (green) ----
box(4.85, DIAG_Y + 1.75, 1.7, 0.45, GRN_FG, GRN_BG, GRN_BRD,
    "Compensator  sign(ψ̇_L)", font_size=10)

# ---- Summation junction ----
circle(5.05, DIAG_Y + 0.68, 0.4, "Σ", color=BLACK, fsize=13)
# arrows into sum
arrow(4.35, DIAG_Y + 0.88, 5.05, DIAG_Y + 0.88, BLUE_BRD, w=2)
arrow(4.35, DIAG_Y + 1.78, 5.25, DIAG_Y + 1.08, RED_BRD, w=1.5)
arrow(4.85, DIAG_Y + 1.95, 5.25, DIAG_Y + 1.08, GRN_BRD, w=1.5)
# u_total out
arrow(5.45, DIAG_Y + 0.88, 6.4, DIAG_Y + 0.88, BLACK, w=2)
text_at(5.5, DIAG_Y + 0.45, 1.0, 0.25, "u_total", fsize=10, bold=True,
        color=BLACK, align=CTR)

# ---- Plant (robot) ----
box(6.4, DIAG_Y + 0.55, 1.85, 0.65, BLACK, WHITE, BLACK,
    "Robot + Liquid\n(two-wheel + pendulum)", font_size=10)
arrow(8.25, DIAG_Y + 0.88, 9.3, DIAG_Y + 0.88, BLACK, w=2)
text_at(8.3, DIAG_Y + 0.45, 1.0, 0.25, "ψ, x, θ", fsize=10, italic=True,
        color=BLACK, align=CTR)

# ---- Sensors feedback line ----
arrow(9.3, DIAG_Y + 1.3, 9.3, DIAG_Y + 2.05, GRAY, w=1.2)
arrow(9.3, DIAG_Y + 2.05, 3.5, DIAG_Y + 2.05, GRAY, w=1.2)
arrow(3.5, DIAG_Y + 2.05, 3.5, DIAG_Y + 1.22, GRAY, w=1.2)
text_at(6.0, DIAG_Y + 2.07, 2.5, 0.2, "IMU  +  wheel encoders",
        fsize=9, italic=True, color=GRAY, align=CTR)

# ============================================================
#  RIGHT SIDE:  small callout "How they interact"
# ============================================================
CO_X = 10.7
CO_Y = DIAG_Y
CO_W = 2.5
CO_H = DIAG_H
panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(CO_X), Inches(CO_Y),
                               Inches(CO_W), Inches(CO_H))
panel.fill.solid(); panel.fill.fore_color.rgb = rgb(0xFA, 0xFA, 0xFA)
panel.line.color.rgb = LTGRAY; panel.line.width = Pt(1)
panel.shadow.inherit = False

text_at(CO_X + 0.1, CO_Y + 0.08, CO_W - 0.2, 0.3,
        "How the blocks cooperate",
        fsize=12, bold=True, color=BLACK, align=CTR)

bullets = [
    ("Feedforward", "shapes v_ref to avoid exciting sloshing (no feedback needed)",
     GOLD_FG),
    ("LQT", "balances the body & tracks v_d", BLUE_FG),
    ("Compensator", "fixes slow nonlinear pitch drift", GRN_FG),
    ("DOB", "cancels pushes / bumps the model didn't predict", RED_FG),
]
yc = CO_Y + 0.45
for name, desc, col in bullets:
    tb = slide.shapes.add_textbox(Inches(CO_X + 0.1), Inches(yc),
                                  Inches(CO_W - 0.2), Inches(0.45))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = "■ " + name + "  "
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = col
    r2 = p.add_run(); r2.text = desc
    r2.font.size = Pt(10); r2.font.color.rgb = BLACK
    yc += 0.45

# ============================================================
#  BOTTOM:  THREE "WHY" COLUMNS
# ============================================================
COL_Y = 3.60
COL_H = 3.80
COL_W = 4.10
GAP   = 0.14
START = 0.40


def why_column(x_in, order, title, subtitle, problem, eq, fix, fg, bg, brd):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x_in), Inches(COL_Y),
                                   Inches(COL_W), Inches(COL_H))
    panel.fill.solid(); panel.fill.fore_color.rgb = bg
    panel.line.color.rgb = brd; panel.line.width = Pt(2)
    panel.shadow.inherit = False

    # order badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(x_in + 0.12), Inches(COL_Y + 0.12),
                                   Inches(0.45), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = fg
    badge.line.color.rgb = fg
    badge.shadow.inherit = False
    tf = badge.text_frame; tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = str(order)
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

    # title
    text_at(x_in + 0.65, COL_Y + 0.10, COL_W - 0.8, 0.35,
            title, fsize=18, bold=True, color=fg, align=LFT)
    text_at(x_in + 0.65, COL_Y + 0.45, COL_W - 0.8, 0.25,
            subtitle, fsize=10, italic=True, color=GRAY, align=LFT)

    # problem
    text_at(x_in + 0.25, COL_Y + 0.78, COL_W - 0.5, 0.25,
            "⚠  Problem",
            fsize=10, bold=True, color=rgb(0xB0, 0x44, 0x44), align=LFT)
    prob_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x_in + 0.25), Inches(COL_Y + 1.05),
                                      Inches(COL_W - 0.5), Inches(0.95))
    prob_box.fill.solid(); prob_box.fill.fore_color.rgb = PROB_BG
    prob_box.line.color.rgb = PROB_BR; prob_box.line.width = Pt(1)
    prob_box.shadow.inherit = False
    tf = prob_box.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = problem
    r.font.size = Pt(10); r.font.color.rgb = BLACK

    # fix label
    text_at(x_in + 0.25, COL_Y + 2.10, COL_W - 0.5, 0.25,
            "✓  How it fixes",
            fsize=10, bold=True, color=fg, align=LFT)

    # equation
    eq_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x_in + 0.25), Inches(COL_Y + 2.37),
                                    Inches(COL_W - 0.5), Inches(0.42))
    eq_box.fill.solid(); eq_box.fill.fore_color.rgb = WHITE
    eq_box.line.color.rgb = brd; eq_box.line.width = Pt(1)
    eq_box.shadow.inherit = False
    tf = eq_box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = eq
    r.font.size = Pt(12); r.font.bold = True; r.font.name = "Courier New"
    r.font.color.rgb = BLACK

    # fix text
    text_at(x_in + 0.25, COL_Y + 2.85, COL_W - 0.5, 0.90,
            fix, fsize=10, color=BLACK, align=LFT)


why_column(
    x_in=START, order=1,
    title="LQT",
    subtitle="Linear Quadratic Tracker",
    problem="Without control the two-wheel robot is open-loop UNSTABLE — it falls over. We also need to track v_d.",
    eq="u_lqt = K·q − K_track·q_d",
    fix="Uses q = [ψ, ψ̇, x, ẋ] from IMU + encoders. Gain K is optimal for the linear model, balancing body and tracking together.",
    fg=BLUE_FG, bg=BLUE_BG, brd=BLUE_BRD,
)

why_column(
    x_in=START + COL_W + GAP, order=2,
    title="Compensator",
    subtitle="Nonlinear auxiliary term",
    problem="LQT is based on the LINEAR model — sin(ψ) ≠ ψ at real angles. That mismatch leaves a slow pitch drift LQT can't remove (no integral action).",
    eq="u_aux = K_c · sign(ψ̇_L,filt)",
    fix="Adds a tiny nonlinear nudge (K_c = 0.001) driven by the low-passed pitch-rate trend. Pushes body back toward upright when drift appears.",
    fg=GRN_FG, bg=GRN_BG, brd=GRN_BRD,
)

why_column(
    x_in=START + 2 * (COL_W + GAP), order=3,
    title="DOB",
    subtitle="Disturbance Observer",
    problem="External forces (pushes, bumps, terrain) are NOT in the model. Pure feedback reacts only AFTER the tilt builds up.",
    eq="u_dob = − d̂ / m",
    fix="Compares measured ψ̈ with model's prediction — mismatch = external force d̂. Feeds −d̂/m forward to cancel it BEFORE tilt grows.",
    fg=RED_FG, bg=RED_BG, brd=RED_BRD,
)

# ============================================================
#  CONNECTOR LINES  (top diagram block → bottom column)
# ============================================================
# LQT box center ≈ (3.5, DIAG_Y + 0.88); column 1 top center ≈ (START + COL_W/2, COL_Y)
def dashed(x1, y1, x2, y2, color):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                   Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(1.2)
    c.line.dash_style = 7

dashed(3.5, DIAG_Y + 1.22, START + COL_W/2, COL_Y, BLUE_BRD)
dashed(5.70, DIAG_Y + 2.22, START + COL_W + GAP + COL_W/2, COL_Y, GRN_BRD)
dashed(3.5, DIAG_Y + 2.10, START + 2*(COL_W + GAP) + COL_W/2, COL_Y, RED_BRD)

# ============================================================
#  SAVE
# ============================================================
os.makedirs("results", exist_ok=True)
out = "results/sbsfc_combined.pptx"
prs.save(out)
print(f"Saved: {out}")
