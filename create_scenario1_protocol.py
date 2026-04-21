#!/usr/bin/env python3
"""
create_scenario1_protocol.py
Single-slide explainer for the "Sudden Start and Stop" test protocol
(Scenario 1).

Run:    python3 create_scenario1_protocol.py
Output: results/scenario1_protocol.pptx
        results/scenario1_vref.png   (embedded in the slide)
"""

import os
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def rgb(r, g, b): return RGBColor(r, g, b)

BLACK = rgb(0x1A, 0x1A, 0x1A)
WHITE = rgb(0xFF, 0xFF, 0xFF)
GRAY  = rgb(0x66, 0x66, 0x66)
LTGRAY = rgb(0xDD, 0xDD, 0xDD)

BLUE_FG = rgb(0x1A, 0x50, 0xAA)
BLUE_BG = rgb(0xEB, 0xF3, 0xFF)
BLUE_BRD = rgb(0x26, 0x66, 0xCF)

ORG_FG = rgb(0xB5, 0x55, 0x00)
ORG_BG = rgb(0xFF, 0xF2, 0xE0)
ORG_BRD = rgb(0xE0, 0x85, 0x20)

GRN_FG = rgb(0x1E, 0x6E, 0x2E)
GRN_BG = rgb(0xEB, 0xF7, 0xEE)
GRN_BRD = rgb(0x2E, 0x99, 0x4A)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT

# ============================================================
# 1. Generate v_ref(t) plot
# ============================================================
os.makedirs("results", exist_ok=True)

dt = 0.001
T  = 30.0
t  = np.arange(0, T, dt)
v  = np.zeros_like(t)
v[(t >= 2)  & (t < 8)]  = 0.6
v[(t >= 20) & (t < 28)] = 0.3

fig, ax = plt.subplots(figsize=(8, 3.6), dpi=160)
ax.plot(t, v, color="#1a50aa", linewidth=2.5)
ax.fill_between(t, 0, v, alpha=0.15, color="#1a50aa")

# Event markers
events = [
    (2.0,  0.60, "Sudden\nSTART",  "#b55500"),
    (8.0,  0.60, "Sudden\nSTOP",   "#a81e1e"),
    (20.0, 0.30, "Start",          "#b55500"),
    (28.0, 0.30, "Stop",           "#a81e1e"),
]
for te, ve, lbl, col in events:
    ax.axvline(te, color=col, linestyle="--", linewidth=1.0, alpha=0.6)
    ax.annotate(lbl, xy=(te, ve), xytext=(te, ve + 0.15),
                color=col, fontsize=9, fontweight="bold",
                ha="center",
                arrowprops=dict(arrowstyle="->", color=col, lw=1.2))

ax.set_xlabel("Time t [s]", fontsize=11)
ax.set_ylabel(r"$v_{\mathrm{ref}}$  [m/s]", fontsize=11)
ax.set_title("Scenario 1 — Commanded Velocity Profile", fontsize=12,
             fontweight="bold")
ax.set_xlim(0, T)
ax.set_ylim(-0.05, 0.9)
ax.grid(True, alpha=0.3)
ax.set_xticks([0, 2, 5, 8, 10, 15, 20, 25, 28, 30])

plt.tight_layout()
vref_png = "results/scenario1_vref.png"
plt.savefig(vref_png, dpi=160, bbox_inches="tight")
plt.close()
print(f"Saved plot: {vref_png}")

# ============================================================
# 2. Build the slide
# ============================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def text(x, y, w, h, s, fsize=11, color=BLACK, bold=False, italic=False,
         align=LFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.size = Pt(fsize); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color


def panel(x, y, w, h, fg, bg, brd, title):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = bg
    shp.line.color.rgb = brd; shp.line.width = Pt(1.8)
    shp.shadow.inherit = False
    # Title strip
    text(x + 0.15, y + 0.08, w - 0.3, 0.35, title,
         fsize=14, bold=True, color=fg, align=LFT)


# --- Title ---
text(0.4, 0.15, 12.5, 0.55,
     "Test Protocol  —  Scenario 1: Sudden Start and Stop",
     fsize=26, bold=True, color=BLACK, align=CTR)
text(0.4, 0.70, 12.5, 0.35,
     "Step-shaped velocity command to excite sloshing across all frequencies",
     fsize=13, italic=True, color=GRAY, align=CTR)

# --- Left: Velocity profile plot ---
slide.shapes.add_picture(vref_png, Inches(0.35), Inches(1.2),
                         width=Inches(7.3))

# --- Right: Event timeline panel ---
panel(7.85, 1.15, 5.1, 3.2, BLUE_FG, BLUE_BG, BLUE_BRD,
      "Event Timeline")

events_text = [
    ("0 – 2 s",   "Settle at rest (warm-up)",            BLACK),
    ("t = 2 s",   "STEP ↑ to 0.6 m/s (sharp launch)",    ORG_FG),
    ("2 – 8 s",   "Cruise at 0.6 m/s",                   BLACK),
    ("t = 8 s",   "STEP ↓ to 0 (hard brake)",            rgb(0xA8, 0x1E, 0x1E)),
    ("8 – 20 s",  "Idle — observe residual ringing",     BLACK),
    ("t = 20 s",  "STEP ↑ to 0.3 m/s (gentler)",         ORG_FG),
    ("20 – 28 s", "Cruise at 0.3 m/s",                   BLACK),
    ("t = 28 s",  "STEP ↓ to 0",                         rgb(0xA8, 0x1E, 0x1E)),
]

yy = 1.70
for tm, desc, col in events_text:
    text(8.05, yy, 1.50, 0.30, tm, fsize=11, bold=True, color=col)
    text(9.55, yy, 3.35, 0.30, desc, fsize=11, color=BLACK)
    yy += 0.32

# --- Bottom-left: What is measured ---
panel(0.35, 4.55, 6.2, 2.75, GRN_FG, GRN_BG, GRN_BRD, "What we measure")
items = [
    ("θ(t)",         "sloshing angle — time series"),
    ("|θ|_mean",     "average sloshing severity"),
    ("|θ|_peak",     "worst-case wave (spill risk)"),
    ("σ²(θ)",        "oscillation energy / variance"),
    ("ψ(t)",         "body tilt — shows controller smoothness"),
    ("x(t) vs ∫v_d", "tracking error (did we arrive on time?)"),
]
yy = 5.15
for sym, desc in items:
    text(0.60, yy, 2.0, 0.30, sym,
         fsize=12, bold=True, color=GRN_FG, italic=True)
    text(2.60, yy, 3.9, 0.30, desc, fsize=11, color=BLACK)
    yy += 0.33

# --- Bottom-right: Why this test ---
panel(6.75, 4.55, 6.2, 2.75, ORG_FG, ORG_BG, ORG_BRD, "Why this design?")
bullets = [
    ("Deterministic",
     "no random disturbance — LPF vs SBSFC difference is purely due to control."),
    ("Broadband excitation",
     "a step demands infinite acceleration → kicks all frequencies, including ω_f."),
    ("Realistic failure mode",
     "start / stop / start / stop is the canonical serving-robot workload (stop at tables, obstacles)."),
    ("Two amplitudes",
     "0.6 m/s + 0.3 m/s → test whether benefit scales with speed."),
    ("Matches Choi",
     "mirrors the step-response test in Figure 6 of the paper."),
]
yy = 5.15
for hd, desc in bullets:
    text(7.00, yy, 1.7, 0.30, "• " + hd,
         fsize=11, bold=True, color=ORG_FG)
    text(8.70, yy, 4.1, 0.30, desc, fsize=10.5, color=BLACK)
    yy += 0.42

# --- Save ---
out = "results/scenario1_protocol.pptx"
prs.save(out)
print(f"Saved slide: {out}")
