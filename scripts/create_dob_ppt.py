#!/usr/bin/env python3
"""
create_dob_ppt.py
Teaching deck for Section 4.5 of Choi et al. (2024):
"Disturbance observer for robustifying a robot"

Walks students through:
   what the lumped disturbance d_t actually contains       (Eq. 5)
   -> the DOB chasing law       d_hat_dot = eta (d - d_hat) (Eq. 6)
   -> derivative-noise problem -> ksi = d_hat - eta * q     (Eq. 7)
   -> recovery   d_hat = ksi + eta * q
   -> role-split with the auxiliary compensator
   -> mapping to controller.m

Run:    python3 scripts/create_dob_ppt.py
Output: results/dob_explained.pptx
"""
import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)
PURPLE_BG= rgb(0xF5, 0xEC, 0xFB)
TEAL     = rgb(0x00, 0x7A, 0x86)
TEAL_BG  = rgb(0xE0, 0xF2, 0xF1)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def line_(slide, x1, y1, x2, y2, color=BLACK, weight=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = conn.line._get_or_add_ln()
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', 'dash')
    return conn


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE,
        bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.1, 11.7, 1.1,
         "Disturbance Observer (DOB)",
         size=46, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.3, 11.7, 1.4,
         "Estimate what the model can't see — then cancel it\n"
         "d̂_t  →  d_t   over time",
         size=22, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 5.5, 11.7, 0.5,
         "Choi et al. 2024 · Section 4.5",
         size=20, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.5, 11.7, 0.5,
         "MECE 6397   ·   Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2 — Roadmap
# =============================================================================
s = blank_slide()
title_bar(s, "Roadmap: 5 steps to understand the DOB",
          "Each box becomes one slide")

steps = [
    (1, "Why we need DOB at all",                            BLUE),
    (2, "Eq. (5):  what  d_t  contains",                     ORANGE),
    (3, "Eq. (6):  the chasing law   d̂̇ = η (d − d̂)",       RED),
    (4, "Eq. (7):  remove derivative noise via  ξ",          PURPLE),
    (5, "Role split with aux compensator + code mapping",    GREEN),
]
y = 1.7
for i, (n, lab, c) in enumerate(steps):
    yi = y + i * 0.95
    box(s, 1.4, yi, 10.5, 0.75, fill=WHITE, border=c, bpt=2.5)
    box(s, 1.4, yi, 0.75, 0.75, fill=c, border=c, bpt=0)
    text_box(s, 1.4, yi + 0.10, 0.75, 0.55, str(n),
             size=22, bold=True, color=WHITE, align=CTR)
    text_box(s, 2.4, yi + 0.18, 9.0, 0.6, lab,
             size=18, bold=True, color=c, align=LFT)
    if i < len(steps) - 1:
        arrow(s, 6.65, yi + 0.78, 6.65, yi + 0.92,
              color=GRAY, weight=1.5)

text_box(s, 0.6, 6.85, 12.1, 0.5,
         "Goal of the DOB:  make the dirty real robot LOOK like the "
         "clean nominal model so LQT works.",
         size=13, color=GRAY, align=CTR)


# =============================================================================
# SLIDE 3 — Why we need it
# =============================================================================
s = blank_slide()
title_bar(s, "Step 1: Why a DOB on top of LQT and aux compensator?",
          "LQT trusts the model; the real world does not")

# Top message
box(s, 0.6, 1.3, 12.1, 1.05, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "The nominal model says one thing — the real robot does another.",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.80, 11.8, 0.5,
         "Anything not in the model becomes a 'disturbance'. "
         "If we can ESTIMATE it, we can CANCEL it.",
         size=13, color=BLACK)

# what's missing
box(s, 0.5, 2.55, 6.2, 4.4, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 2.65, 5.8, 0.5,
         "What the nominal model has",
         size=15, bold=True, color=BLUE)
text_box(s, 0.7, 3.10, 5.8, 3.7,
         ["q̇  =  A q  +  B u",
          "",
          "  • inertia,  geometry",
          "  • gravity   (M g L term)",
          "  • motor command  u",
          "",
          "Used by:",
          "  • LQT to design  K"],
         size=13, color=BLACK, font='Consolas')

box(s, 6.9, 2.55, 6.0, 4.4, fill=RED_BG, border=RED)
text_box(s, 7.1, 2.65, 5.6, 0.5,
         "What it does NOT have",
         size=15, bold=True, color=RED)
text_box(s, 7.1, 3.10, 5.6, 3.7,
         ["• caster wheel kicks",
          "• floor friction / slope",
          "• modeling error (mass off, inertia off)",
          "• motor saturation, dead-zone",
          "• payload changes (food added/removed)",
          "",
          "All of that  →  lumped into  d_t.",
          "DOB's job: figure out what  d_t  is."],
         size=13, color=BLACK)


# =============================================================================
# SLIDE 4 — Eq. (5):  what is d_t
# =============================================================================
s = blank_slide()
title_bar(s, "Step 2: Eq. (5) — what the lumped disturbance contains",
          "d_t is the bucket of everything LQT does not know")

# Big formula
box(s, 0.6, 1.3, 12.1, 1.5, fill=ORANGE_BG, border=ORANGE, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Eq. (5):  decomposition of  d_t",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.85, 11.8, 0.7,
         "d_t   =   B K_c · sgn(d ψ^L/dt)   +   d_{a,t}   +   d_{c,t}   +   ΔU",
         size=18, bold=True, color=BLACK, font='Consolas', align=CTR)

# 4 cards explaining each term
y = 3.0
cards = [
    ("B K_c sgn(d ψ^L/dt)", "the aux compensator term itself\n"
                            "(from Sec 4.4)", BLUE),
    ("d_{a,t}",             "abrupt impulse-like\ndisturbances "
                            "(caster kicks)", ORANGE),
    ("d_{c,t}",             "continuous\ndisturbances\n"
                            "(friction, slope)", PURPLE),
    ("ΔU",                  "system uncertainty\n(model error,\n"
                            "wrong masses)",  RED),
]
w = 2.95
for i, (eq, desc, c) in enumerate(cards):
    x = 0.6 + i * (w + 0.1)
    box(s, x, y, w, 3.6, fill=WHITE, border=c, bpt=2.5)
    box(s, x, y, w, 0.6, fill=c, border=c, bpt=0)
    text_box(s, x, y + 0.07, w, 0.5, eq,
             size=14, bold=True, color=WHITE, font='Consolas', align=CTR)
    text_box(s, x + 0.2, y + 0.85, w - 0.4, 2.6, desc,
             size=13, color=BLACK, align=CTR)


# =============================================================================
# SLIDE 5 — Eq. (6): chasing law
# =============================================================================
s = blank_slide()
title_bar(s, "Step 3: Eq. (6) — DOB chases the disturbance",
          "Low-pass filter that drives  d̂  toward  d")

# Big formula
box(s, 0.6, 1.3, 12.1, 1.5, fill=RED_BG, border=RED, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Eq. (6):  estimator update law",
         size=14, bold=True, color=RED)
text_box(s, 0.8, 1.85, 11.8, 0.7,
         "d̂̇_t   =   η  ( d_t   −   d̂_t )",
         size=24, bold=True, color=BLACK, font='Consolas', align=CTR)

# left — interpretation
box(s, 0.5, 2.95, 6.2, 4.0, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 3.05, 5.8, 0.5,
         "How to read it",
         size=15, bold=True, color=BLUE)
text_box(s, 0.7, 3.55, 5.8, 3.4,
         ["• d̂  =  our estimate of d",
          "• d̂̇  =  rate of change of estimate",
          "• (d − d̂)  =  error",
          "• η  =  bandwidth (how fast we chase)",
          "",
          "Behavior:",
          "  • If d̂ < d  →  d̂̇ > 0  →  d̂ rises",
          "  • If d̂ > d  →  d̂̇ < 0  →  d̂ falls",
          "  • In steady-state:  d̂  →  d  ✅"],
         size=12, color=BLACK)

# right — but how do we know d?
box(s, 6.9, 2.95, 6.0, 4.0, fill=YELLOW_BG, border=ORANGE)
text_box(s, 7.1, 3.05, 5.6, 0.5,
         "But we never measure  d  directly!",
         size=15, bold=True, color=ORANGE)
text_box(s, 7.1, 3.55, 5.6, 3.4,
         ["From the dynamics:",
          "  q̇  =  A q  +  B u  +  d",
          "Rearrange:",
          "  d  =  q̇  −  A q  −  B u",
          "",
          "(this is the underbraced part of Eq. 6)",
          "",
          "Plug in:  estimate = LPF(observed gap)",
          "    d̂̇  =  η  (q̇ − A q − B u  −  d̂)"],
         size=12, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 6 — Eq. (7): the trick
# =============================================================================
s = blank_slide()
title_bar(s, "Step 4: Eq. (7) — kill the noisy derivative",
          "Trade  q̇  for an integrator on  ξ  =  d̂  −  η q")

# Problem statement
box(s, 0.6, 1.3, 12.1, 0.95, fill=RED_BG, border=RED)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Problem with Eq. (6):  it needs  q̇.",
         size=14, bold=True, color=RED)
text_box(s, 0.8, 1.80, 11.8, 0.4,
         "Sensors give q (smooth); computing q̇ = (q − q_prev)/dt amplifies noise. "
         "Garbage q̇ → garbage d̂.",
         size=12, color=BLACK)

# trick
box(s, 0.6, 2.4, 12.1, 0.9, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 2.50, 11.8, 0.4,
         "Trick: define a helper variable",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 2.90, 11.8, 0.4,
         "ξ_t   ≡   d̂_t   −   η q_t",
         size=16, bold=True, color=BLACK, font='Consolas', align=CTR)

# math chain
box(s, 0.6, 3.45, 12.1, 0.7, fill=WHITE, border=PURPLE, bpt=1.5)
text_box(s, 0.8, 3.55, 11.8, 0.5,
         "Differentiate:   ξ̇ = d̂̇ − η q̇",
         size=14, color=BLACK, font='Consolas')

box(s, 0.6, 4.20, 12.1, 0.7, fill=WHITE, border=PURPLE, bpt=1.5)
text_box(s, 0.8, 4.30, 11.8, 0.5,
         "Substitute Eq. (6):  ξ̇ = η(q̇ − A q − B u − d̂) − η q̇",
         size=14, color=BLACK, font='Consolas')

box(s, 0.6, 4.95, 12.1, 0.85, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 0.8, 5.05, 11.8, 0.4,
         "η q̇  cancels — only q remains.   Eq. (7):",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 5.40, 11.8, 0.4,
         "ξ̇   =   η ( − A q − B u )   −   η d̂",
         size=16, bold=True, color=GREEN, font='Consolas', align=CTR)

# recover
box(s, 0.6, 5.95, 12.1, 1.4, fill=ORANGE_BG, border=ORANGE)
text_box(s, 0.8, 6.05, 11.8, 0.4,
         "How to use it online (each timestep):",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 6.45, 11.8, 0.95,
         ["    1.  integrate    ξ_new = ξ_old + dt · ( η(−A q − B u) − η d̂ )",
          "    2.  recover     d̂  =  ξ + η q          ← uses q only, never q̇"],
         size=14, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 7 — convergence visual
# =============================================================================
s = blank_slide()
title_bar(s, "How d̂ catches up to d  —  what η controls",
          "η is the bandwidth knob: speed vs. noise rejection")

# 3 panels: small η, medium η, large η
def panel(slide, x0, y0, w, h, eta_label, eta_speed, color):
    box(slide, x0, y0, w, h, fill=WHITE, border=color, bpt=2.0)
    text_box(slide, x0, y0 + 0.08, w, 0.4, eta_label,
             size=14, bold=True, color=color, align=CTR)
    # axes
    ax_l = x0 + 0.3
    ax_t = y0 + 0.55
    ax_w = w - 0.6
    ax_h = h - 0.85
    line_(slide, ax_l, ax_t + ax_h - 0.1, ax_l + ax_w, ax_t + ax_h - 0.1,
          color=BLACK, weight=1)
    line_(slide, ax_l, ax_t + 0.1, ax_l, ax_t + ax_h - 0.1,
          color=BLACK, weight=1)
    # true d (step at t=0.2)
    base_y = ax_t + ax_h - 0.6
    step_y = ax_t + 0.4
    sx = ax_l + 0.20 * ax_w
    line_(slide, ax_l, base_y, sx, base_y, color=BLACK, weight=2.5)
    line_(slide, sx, base_y, sx, step_y, color=BLACK, weight=2.5)
    line_(slide, sx, step_y, ax_l + ax_w, step_y, color=BLACK, weight=2.5)
    # estimate: exponential rise after step
    n = 24
    prev = (sx, base_y)
    for i in range(1, n + 1):
        frac = i / n
        x = sx + frac * (ax_l + ax_w - sx)
        # exponential approach with given rate
        val = 1 - math.exp(-eta_speed * frac * 4.0)
        y = base_y - val * (base_y - step_y)
        line_(slide, prev[0], prev[1], x, y, color=color, weight=2.5)
        prev = (x, y)
    # legend
    text_box(slide, x0 + 0.2, y0 + h - 0.45, w - 0.4, 0.35,
             "—— true d   —— d̂",
             size=10, color=GRAY, align=CTR)

w = 4.0; h = 4.4
panel(s, 0.6,  1.4, w, h, "small η (slow, smooth)",   0.4, BLUE)
panel(s, 4.7,  1.4, w, h, "medium η  (sweet spot)",   1.5, GREEN)
panel(s, 8.8,  1.4, w, h, "large η (fast, noisy)",    4.0, RED)

# bottom takeaway
box(s, 0.6, 5.95, 12.1, 1.35, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 6.05, 11.8, 0.4,
         "Tradeoff:",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 6.45, 11.8, 0.95,
         ["• small η  →  slow chase, smooth d̂  (rejects noise but lags real disturbance)",
          "• large η  →  fast chase, but amplifies sensor noise (chatter)",
          "• tune η so DOB is faster than disturbance changes but slower than noise"],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 8 — role split + block diagram
# =============================================================================
s = blank_slide()
title_bar(s, "Role split: aux compensator + DOB cover what LQT can't",
          "Two robustness blocks, two timescales")

# left — role table
box(s, 0.5, 1.3, 6.2, 5.6, fill=GREEN_BG, border=GREEN)
text_box(s, 0.7, 1.4, 5.8, 0.5,
         "Who handles what?",
         size=16, bold=True, color=GREEN)
text_box(s, 0.7, 1.95, 5.8, 4.0,
         ["LQT  (Sec 4.3)",
          "   • smooth nominal balance",
          "   • velocity tracking",
          "",
          "Aux compensator  (Sec 4.4)",
          "   • caster impulses (fast, sharp)",
          "   • bounded ±K_c shove",
          "",
          "DOB  (Sec 4.5)",
          "   • residual continuous biases",
          "   • model error,  friction,  slope",
          "   • slow but accurate"],
         size=13, color=BLACK)
text_box(s, 0.7, 5.95, 5.8, 0.9,
         "Together: every disturbance band is\ncovered "
         "without LQT having to do it.",
         size=12, color=GRAY)

# right — full block diagram
box(s, 6.9, 1.3, 6.0, 5.6, fill=BLUE_BG, border=BLUE)
text_box(s, 7.1, 1.4, 5.6, 0.5,
         "Where the DOB sits",
         size=16, bold=True, color=BLUE)

# blocks
yb = 2.05
box(s, 7.1, yb, 1.3, 0.55, fill=WHITE, border=PURPLE)
text_box(s, 7.1, yb + 0.05, 1.3, 0.45, "F_e(s)",
         size=12, bold=True, color=PURPLE, align=CTR)
arrow(s, 8.4, yb + 0.27, 8.7, yb + 0.27, color=GRAY, weight=2)
box(s, 8.7, yb, 1.3, 0.55, fill=WHITE, border=BLUE)
text_box(s, 8.7, yb + 0.05, 1.3, 0.45, "LQT",
         size=12, bold=True, color=BLUE, align=CTR)
arrow(s, 10.0, yb + 0.27, 10.3, yb + 0.27, color=GRAY, weight=2)
box(s, 10.3, yb, 0.7, 0.55, fill=WHITE, border=BLACK)
text_box(s, 10.3, yb + 0.02, 0.7, 0.45, "+",
         size=18, bold=True, color=BLACK, align=CTR)
arrow(s, 11.0, yb + 0.27, 11.3, yb + 0.27, color=GRAY, weight=2)
box(s, 11.3, yb, 1.4, 0.55, fill=WHITE, border=ORANGE)
text_box(s, 11.3, yb + 0.05, 1.4, 0.45, "Plant",
         size=12, bold=True, color=ORANGE, align=CTR)

# u_aux feeding +
box(s, 7.6, yb + 1.3, 2.2, 0.65, fill=GREEN_BG, border=GREEN)
text_box(s, 7.6, yb + 1.4, 2.2, 0.5, "u_aux",
         size=12, bold=True, color=GREEN, align=CTR, font='Consolas')
arrow(s, 9.8, yb + 1.55, 10.6, yb + 0.55, color=GREEN, weight=2)

# u_dob feeding +
box(s, 9.6, yb + 1.3, 2.2, 0.65, fill=RED_BG, border=RED)
text_box(s, 9.6, yb + 1.4, 2.2, 0.5, "u_dob = − d̂ / m",
         size=11, bold=True, color=RED, align=CTR, font='Consolas')
arrow(s, 10.7, yb + 1.3, 10.7, yb + 0.55, color=RED, weight=2)

# DOB box itself, fed by q,u
box(s, 7.6, yb + 2.2, 4.5, 1.0, fill=PURPLE_BG, border=PURPLE)
text_box(s, 7.6, yb + 2.30, 4.5, 0.4,
         "DOB  (Eq. 7)",
         size=13, bold=True, color=PURPLE, align=CTR)
text_box(s, 7.6, yb + 2.65, 4.5, 0.4,
         "ξ̇ = η(−Aq − Bu) − η d̂  →  d̂ = ξ + η q",
         size=11, color=BLACK, font='Consolas', align=CTR)
arrow(s, 9.85, yb + 2.2, 10.7, yb + 1.95, color=PURPLE, weight=2)

# feedback path
arrow(s, 12.0, yb + 0.55, 12.0, yb + 4.0, color=GRAY, weight=1.5)
arrow(s, 12.0, yb + 4.0, 7.8, yb + 4.0, color=GRAY, weight=1.5)
arrow(s, 7.8, yb + 4.0, 7.8, yb + 3.2, color=GRAY, weight=1.5)
text_box(s, 9.5, yb + 4.05, 2.0, 0.35, "q feedback",
         size=10, color=GRAY)


# =============================================================================
# SLIDE 9 — Code mapping
# =============================================================================
s = blank_slide()
title_bar(s, "How the DOB shows up in your controller.m",
          "controller.m, lines ~67-79")

# code block
box(s, 0.6, 1.3, 12.1, 4.0, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "DOB block (case 'sbsfc'):",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.85, 11.8, 3.4,
         ["  q_dot_est = ctrl_state.A * q + ctrl_state.B * u_lqt;",
          "  residual  = (q(2) - ctrl_state.q2_prev) / p.dt - q_dot_est(2);",
          "  ctrl_state.q2_prev = q(2);",
          "",
          "  alpha_dob = p.eta * p.dt / (1 + p.eta * p.dt);",
          "  ctrl_state.d_hat_scalar = ...",
          "      (1 - alpha_dob) * ctrl_state.d_hat_scalar + ...",
          "      alpha_dob * (p.I_p * residual);",
          "",
          "  u_dob = -ctrl_state.d_hat_scalar / p.m;",
          "  u_dob = max(min(u_dob, 5), -5);   % saturation"],
         size=13, color=BLACK, font='Consolas')

# annotation column
box(s, 0.6, 5.45, 5.95, 1.85, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 5.55, 5.6, 0.4, "Match with the paper",
         size=14, bold=True, color=BLUE)
text_box(s, 0.8, 5.95, 5.6, 1.4,
         ["• residual          = q̇ − A q − B u",
          "• alpha_dob         = discrete LPF (= η)",
          "• d_hat_scalar      = d̂",
          "• u_dob             = − d̂ / m"],
         size=12, color=BLACK, font='Consolas')

# tuning column
box(s, 6.75, 5.45, 5.95, 1.85, fill=GREEN_BG, border=GREEN)
text_box(s, 6.95, 5.55, 5.6, 0.4, "Knobs to tune", size=14, bold=True,
         color=GREEN)
text_box(s, 6.95, 5.95, 5.6, 1.4,
         ["• p.eta   →  bandwidth η",
          "    (high → fast but noisy,",
          "     low  → slow but smooth)",
          "• ±5 saturation guards against",
          "  bad transient estimates."],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 10 — Summary
# =============================================================================
s = blank_slide()
title_bar(s, "Summary — DOB on one page",
          "Estimate. Smooth. Cancel.")

# 5 bullets
y = 1.4
items = [
    ("Why",       "make the dirty real robot look like the clean nominal model",  BLUE),
    ("Eq. (5)",   "d_t  =  aux + d_{a,t} + d_{c,t} + ΔU   (the bucket)",          ORANGE),
    ("Eq. (6)",   "d̂̇  =  η ( d − d̂ )   (chase the truth)",                      RED),
    ("Eq. (7)",   "ξ ≡ d̂ − η q   →   integrate ξ, never differentiate q",         PURPLE),
    ("Cancel",    "u_dob  =  − d̂ / m   (subtract the estimate from input)",      GREEN),
]
for i, (k, v, c) in enumerate(items):
    yi = y + i * 0.85
    box(s, 0.6, yi, 12.1, 0.7, fill=WHITE, border=c, bpt=2.0)
    box(s, 0.6, yi, 1.0, 0.7, fill=c, border=c, bpt=0)
    text_box(s, 0.6, yi + 0.13, 1.0, 0.5, k,
             size=14, bold=True, color=WHITE, align=CTR)
    text_box(s, 1.8, yi + 0.13, 10.7, 0.5, v,
             size=14, color=BLACK, font='Consolas')

# closing
box(s, 0.6, 6.05, 12.1, 1.25, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 6.15, 11.8, 0.5,
         "Big picture",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 6.55, 11.8, 0.7,
         "DOB is the slow, accurate complement to the aux compensator. "
         "Together they take care of everything LQT was never told about — "
         "letting the optimal balancing law work on a clean nominal system.",
         size=13, color=BLACK)


# ----- save ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'dob_explained.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  {len(prs.slides)} slides')
