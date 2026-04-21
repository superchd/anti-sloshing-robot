#!/usr/bin/env python3
"""
create_pipeline_ppt.py
Generate an editable PowerPoint pipeline diagram for the Anti-Sloshing Robot project.
Run: python3 create_pipeline_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─── Color palette ───────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

BLUE      = rgb(0x26, 0x66, 0xBF);  BLUE_BG   = rgb(0xE8, 0xF2, 0xFF)
TEAL      = rgb(0x1A, 0x7A, 0x85);  TEAL_BG   = rgb(0xE0, 0xF5, 0xF7)
ORANGE    = rgb(0xCC, 0x6A, 0x10);  ORANGE_BG = rgb(0xFF, 0xF0, 0xE0)
GREEN     = rgb(0x14, 0x80, 0x38);  GREEN_BG  = rgb(0xE6, 0xF7, 0xED)
RED       = rgb(0xBF, 0x20, 0x20);  RED_BG    = rgb(0xFF, 0xEC, 0xEC)
PURPLE    = rgb(0x75, 0x28, 0xAD);  PURPLE_BG = rgb(0xF3, 0xE8, 0xFF)
GRAY      = rgb(0x60, 0x60, 0x60);  GRAY_BG   = rgb(0xF4, 0xF4, 0xF4)
GOLD      = rgb(0x99, 0x77, 0x00);  GOLD_BG   = rgb(0xFF, 0xFD, 0xE0)
WHITE     = rgb(0xFF, 0xFF, 0xFF)
BLACK     = rgb(0x1A, 0x1A, 0x1A)
LTGRAY    = rgb(0x99, 0x99, 0x99)

CTR = PP_ALIGN.CENTER

# ─── Presentation setup ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ─── Helpers ─────────────────────────────────────────────────────────────────

def box(l, t, w, h, fill, border, bpt=1.5, dash=False, rounded=True):
    """Add a rounded (or plain) rectangle."""
    kind = 5 if rounded else 1   # 5 = RoundedRect, 1 = Rect
    s = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(bpt)
    if dash:
        ln = s.line._ln
        for old in ln.findall(qn('a:prstDash')): ln.remove(old)
        pd = etree.SubElement(ln, qn('a:prstDash'))
        pd.set('val', 'lgDash')
    return s


def write(shape, rows, anchor='ctr'):
    """
    Write text rows into a shape.
    rows: list of (text, size_pt, bold, color, align)
    anchor: 'ctr' | 't' | 'b'
    """
    tf = shape.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', anchor)
    tf.clear()
    for i, (s, sz, bd, co, al) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if not s:
            continue
        r = p.add_run()
        r.text = s
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = co


def label(l, t, w, h, text, sz, bold, color, al=CTR):
    """Floating text box (no background)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = al
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def arrow(x1, y1, x2, y2, color, pt=2.2):
    """Straight connector with arrowhead at end."""
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    ln = c.line._ln
    for tag in [qn('a:headEnd'), qn('a:tailEnd')]:
        for el in ln.findall(tag): ln.remove(el)
    etree.SubElement(ln, qn('a:tailEnd')).set('type', 'none')
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')


# ─── Layout constants ────────────────────────────────────────────────────────
# Four stage columns:          x     width
S1X, S1W = 0.10, 2.65   # Problem & Modeling
S2X, S2W = 3.05, 3.80   # Control Design
S3X, S3W = 7.10, 2.90   # Simulation Engine
S4X, S4W = 10.25, 2.95  # Analysis & Results
# Common stage y / height
SGY, SGH = 0.90, 5.10

# ═══════════════════════════════════════════════════════════════════════════
#  TITLE
# ═══════════════════════════════════════════════════════════════════════════
label(0.15, 0.04, 10.3, 0.50,
      'Integrated Pipeline: Anti-Sloshing Control for Food-Serving Robots',
      17, True, BLACK)
label(0.15, 0.50, 10.3, 0.30,
      'Classical Control (SBSFC)  vs  Reinforcement Learning'
      '   \u2014   MECE 6397, University of Houston',
      9, False, LTGRAY)

# Tools banner (top-right)
b = box(10.85, 0.06, 2.35, 0.62, GRAY_BG, GRAY, 1.0)
write(b, [('MATLAB  |  RL Toolbox', 8, True, GRAY, CTR)])

# ═══════════════════════════════════════════════════════════════════════════
#  STAGE GROUP BORDERS  (dashed)
# ═══════════════════════════════════════════════════════════════════════════
stage_defs = [
    (S1X, S1W, BLUE,   rgb(0xF4, 0xF8, 0xFF), 'Problem & Modeling'),
    (S2X, S2W, TEAL,   rgb(0xF2, 0xFB, 0xFC), 'Control Design'),
    (S3X, S3W, GREEN,  rgb(0xF2, 0xFB, 0xF5), 'Simulation Engine'),
    (S4X, S4W, RED,    rgb(0xFD, 0xF4, 0xF4), 'Analysis & Results'),
]
for (sx, sw, col, bg, ttl) in stage_defs:
    box(sx, SGY, sw, SGH, bg, col, 1.8, dash=True, rounded=False)
    label(sx+0.05, SGY+0.05, sw-0.10, 0.32, ttl, 11, True, col)

# ═══════════════════════════════════════════════════════════════════════════
#  STAGE 1 — Problem & Modeling
# ═══════════════════════════════════════════════════════════════════════════
CX, CW = S1X+0.13, S1W-0.26

# Robot model
s = box(CX, 1.34, CW, 1.30, BLUE_BG, BLUE, 1.5)
write(s, [
    ('Serving Robot',                    11, True,  BLUE,   CTR),
    ('Wheeled inverted pendulum',          8, False, GRAY,   CTR),
    ('q = [\u03c8, \u03c8\u0307, x, x\u0307]', 8.5, False, BLUE, CTR),
])

# Sloshing proxy
s = box(CX, 2.76, CW, 1.30, ORANGE_BG, ORANGE, 1.5)
write(s, [
    ('Liquid Sloshing Proxy',            11, True,  ORANGE, CTR),
    ('Nonlinear pendulum on tray',         8, False, GRAY,   CTR),
    ('\u03b8\u0308 = \u2212(g/l)sin\u03b8 + (a/l)cos\u03b8',
                                          8.5, False, ORANGE, CTR),
])

# Parameters
s = box(CX, 4.18, CW, 0.82, GRAY_BG, GRAY, 1.2)
write(s, [
    ('Parameters',                       10, True,  rgb(0x33,0x33,0x33), CTR),
    ('parameters.m  |  serving_robot.urdf', 7.5, False, GRAY, CTR),
])

# Novel contribution callout
s = box(CX, 5.12, CW, 0.75, ORANGE_BG, ORANGE, 1.8)
write(s, [
    ('\u2605  Novel Contribution',         9, True,  ORANGE, CTR),
    ('RL handles 2D turning + bumps',     7.5, False, GRAY,   CTR),
    ('better than classical 1D control',  7,   False, GRAY,   CTR),
])

# ═══════════════════════════════════════════════════════════════════════════
#  STAGE 2 — Control Design
# ═══════════════════════════════════════════════════════════════════════════
# Left track (Classical): x = S2X+0.12, width = 1.68
LX, LW = S2X+0.12, 1.68
# Right track (RL):       x = LX+LW+0.25
RX = LX + LW + 0.25
RW = S2X + S2W - RX - 0.12

# ── Classical SBSFC ──
s = box(LX, 1.34, LW, 1.05, TEAL_BG, TEAL, 1.5)
write(s, [
    ('Classical: SBSFC',                 10, True,  TEAL, CTR),
    ('Input Shaping + LQT + DOB',         8, False, GRAY, CTR),
])

# SBSFC sub-boxes (4)
sub_names = ['IS', 'LQT', 'DOB', 'Aux']
sw = (LW - 0.05) / 4
for i, nm in enumerate(sub_names):
    sx = LX + 0.01 + i*(sw + 0.01)
    s = box(sx, 2.52, sw, 0.68, TEAL_BG, TEAL, 0.8)
    write(s, [(nm, 7, True, TEAL, CTR)])

# SBSFC detail card
s = box(LX, 3.33, LW, 1.60, TEAL_BG, TEAL, 1.0)
write(s, [
    ('Component Roles',                  8.5, True,  TEAL, CTR),
    ('IS \u2192 suppress slosh freq',     7.5, False, GRAY, CTR),
    ('LQT \u2192 track velocity',         7.5, False, GRAY, CTR),
    ('DOB \u2192 reject disturbance',     7.5, False, GRAY, CTR),
    ('Aux \u2192 damp residual pitch',    7.5, False, GRAY, CTR),
    ('',                                  2,   False, GRAY, CTR),
    ('controller.m  |  design_lqt.m',    7,   False, GRAY, CTR),
], anchor='t')

# vs label
label(LX+LW+0.03, 1.72, 0.22, 0.40, 'vs', 13, True, LTGRAY)

# ── RL DDPG / TD3 ──
s = box(RX, 1.34, RW, 1.05, PURPLE_BG, PURPLE, 1.5)
write(s, [
    ('RL: DDPG / TD3',                   10, True,  PURPLE, CTR),
    ('Learn policy from reward',           8, False, GRAY,   CTR),
])

# RL sub-boxes (2)
sw_rl = (RW - 0.03) / 2
for i, nm in enumerate(['Env', 'Agent']):
    sx = RX + 0.01 + i*(sw_rl + 0.01)
    s = box(sx, 2.52, sw_rl, 0.68, PURPLE_BG, PURPLE, 0.8)
    write(s, [(nm, 7, True, PURPLE, CTR)])

# RL detail card
s = box(RX, 3.33, RW, 1.60, PURPLE_BG, PURPLE, 1.0)
write(s, [
    ('Reward Signal',                     8.5, True,  PURPLE, CTR),
    ('r = \u2212|\u03b8| \u2212 \u03b1|u|', 9, False, PURPLE, CTR),
    ('',                                   2,  False, GRAY,   CTR),
    ('State:   q + \u03b8',               7.5, False, GRAY,   CTR),
    ('Action:  acceleration u',           7.5, False, GRAY,   CTR),
    ('',                                   2,  False, GRAY,   CTR),
    ('controller.m  (to implement)',       7,  False, GRAY,   CTR),
], anchor='t')

# ═══════════════════════════════════════════════════════════════════════════
#  STAGE 3 — Simulation Engine
# ═══════════════════════════════════════════════════════════════════════════
CX, CW = S3X+0.12, S3W-0.24

# Test Scenarios
s = box(CX, 1.34, CW, 1.05, GREEN_BG, GREEN, 1.5)
write(s, [
    ('Test Scenarios',                   10, True,  GREEN, CTR),
    ('5 disturbance profiles',            8, False, GRAY,  CTR),
    ('scenarios.m',                      7.5, False, GRAY,  CTR),
])

# Closed-Loop Simulation
s = box(CX, 2.52, CW, 1.05, GREEN_BG, GREEN, 1.5)
write(s, [
    ('Closed-Loop Simulation',           10, True,  GREEN, CTR),
    ('q\u0307 = Aq + Bu + d',            8.5, False, GREEN, CTR),
    ('simulate_system.m',                7.5, False, GRAY,  CTR),
])

# Mode selector
s = box(CX, 3.70, CW, 0.75, GOLD_BG, GOLD, 1.5)
write(s, [('Mode:  none | lpf | sbsfc | rl', 8, True, GOLD, CTR)])

# Feedback
s = box(CX, 4.57, CW, 0.70, GREEN_BG, GREEN, 1.0)
write(s, [
    ('Feedback States',                  8.5, True,  GREEN, CTR),
    ('q = [\u03c8, \u03c8\u0307, x, x\u0307]   \u03b8', 7.5, False, GREEN, CTR),
])

# ═══════════════════════════════════════════════════════════════════════════
#  STAGE 4 — Analysis & Results
# ═══════════════════════════════════════════════════════════════════════════
CX, CW = S4X+0.12, S4W-0.24

# Performance Metrics
s = box(CX, 1.34, CW, 1.05, RED_BG, RED, 1.5)
write(s, [
    ('Performance Metrics',              10, True,  RED,  CTR),
    ('|\u03b8|\u2098\u2091\u2090\u2099  '
     '|\u03b8|\u2098\u2090\u2093  |\u03c8|  energy',
                                          8, False, GRAY, CTR),
    ('print_comparison_table.m',          7, False, GRAY, CTR),
])

# Comparison Plots
s = box(CX, 2.52, CW, 0.82, RED_BG, RED, 1.5)
write(s, [
    ('Comparison Plots',                 10, True,  RED,  CTR),
    ('plot_results.m',                   7.5, False, GRAY, CTR),
])

# 2D Animation
s = box(CX, 3.47, CW, 0.82, RED_BG, RED, 1.5)
write(s, [
    ('2D Animation',                     10, True,  RED,  CTR),
    ('animate_2d.m',                     7.5, False, GRAY, CTR),
])

# Conclusion
s = box(CX, 4.42, CW, 0.82, GOLD_BG, GOLD, 1.8)
write(s, [
    ('\u2605  Conclusion',               10, True,  GOLD, CTR),
    ('SBSFC vs RL \u2014 which wins?',    8, False, GRAY, CTR),
])

# ═══════════════════════════════════════════════════════════════════════════
#  CONNECTING ARROWS between stages
# ═══════════════════════════════════════════════════════════════════════════
AY = 2.87   # arrow vertical center
arrow(S1X+S1W+0.02, AY, S2X-0.02, AY, BLUE,  2.5)
arrow(S2X+S2W+0.02, AY, S3X-0.02, AY, TEAL,  2.5)
arrow(S3X+S3W+0.02, AY, S4X-0.02, AY, GREEN, 2.5)

label(S1X+S1W,      AY-0.28, 0.42, 0.25, 'A, B',   7, True, BLUE)
label(S2X+S2W,      AY-0.28, 0.42, 0.25, 'u(q)',   7, True, TEAL)
label(S3X+S3W,      AY-0.28, 0.50, 0.25, 'states', 7, True, GREEN)

# ═══════════════════════════════════════════════════════════════════════════
#  BOTTOM STRIP — Disturbance Scenarios
# ═══════════════════════════════════════════════════════════════════════════
SY, SH = 6.13, 1.22
box(0.10, SY, 13.10, SH, GRAY_BG, GRAY, 1.2, dash=True, rounded=False)
label(0.10, SY+0.05, 13.10, 0.28,
      'Disturbance Scenarios  (scenarios.m)', 9, True, rgb(0x44,0x44,0x44))

scenarios = [
    ('Sudden Start/Stop', 'v_ref = step'),
    ('Speed Bumps',       'impulse  F'),
    ('Accel / Decel',     'repeated ramp'),
    ('External Push',     'F_push pulse'),
    ('Rough Terrain',     'random  F(t)'),
]
gap   = 0.15
sc_w  = (13.10 - 0.24 - 4*gap) / 5
x0    = 0.22
for i, (name, icon) in enumerate(scenarios):
    sx = x0 + i*(sc_w + gap)
    s = box(sx, SY+0.36, sc_w, 0.72, WHITE, GRAY, 1.0)
    write(s, [
        (f'{i+1}. {name}',  8,   True,  BLACK, CTR),
        (icon,              7.5, False, GRAY,  CTR),
    ])

# ═══════════════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   'results', 'project_pipeline.pptx')
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'Saved → {out}')
