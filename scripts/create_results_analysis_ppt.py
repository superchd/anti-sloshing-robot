"""
Teaching PPT: Results Analysis, Limitations & Next Steps
For MECE 6397 anti-sloshing robot project (Choi et al. 2024)

Numbers (Scenario 1, 30 s, 30,000 steps):
    RL    : mean|θ| = 0.2377°,  var = 0.1934°²,  max = 2.0809°
    SBSFC : mean|θ| = 0.3167°,  var = 0.1931°²,  max = 1.7120°
    LPF   : mean|θ| = 5.0659°,  var = 35.4129°², max = 11.4124°

Slide map:
  1. Title
  2. Headline finding + 3-column summary
  3. Three layers of insight
  4. The key trade-off — RL vs SBSFC
  5. Fairness check — ψ (robot tilt) comparison
  6. Limitations — top 4 critical
  7. Limitations — remaining 4
  8. What we can / cannot claim
  9. Next steps — 3 tiers
 10. Bottom-line conclusion
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY   = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF4, 0xF1, 0xEC)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=NAVY, line_w=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, 0.4, 0.08, 12.5, 0.55, title,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.4, 0.78, 12.5, 0.4, subtitle,
                 size=14, color=GRAY)


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------------- Slide 1 — Title ----------------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

add_text(s, 0.6, 1.9, 12, 1.4,
         "Results Analysis",
         size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.3, 12, 0.8,
         "Findings, Limitations & Next Steps",
         size=24, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.4, 12, 0.6,
         "RL vs. SBSFC vs. LPF on Scenario 1 (sudden start/stop)",
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 5.6, 12, 0.5,
         "MECE 6397 — Anti-Sloshing Robot Project",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ---------------- Slide 2 — Headline finding ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Headline Finding",
           "What the numbers tell us in one sentence")

# Big take-away box
add_box(s, 0.6, 1.0, 12.1, 1.5, fill=ORANGE, line=ORANGE)
add_text(s, 0.85, 1.10, 11.6, 0.4,
         "THE TAKE-AWAY",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.85, 1.50, 11.6, 0.95,
         "Both SBSFC and RL eliminate sloshing  (>93 % reduction vs LPF baseline).\n"
         "RL wins on average performance.   SBSFC wins on worst-case peak.",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Compact 3-column results table (no fill, paper style)
add_text(s, 0.6, 2.85, 12.0, 0.4,
         "Scenario 1 — sudden start and stop, T = 30 s",
         size=14, bold=True, color=NAVY)

table_x, table_w = 1.0, 11.3
# Top thick rule
ln = s.shapes.add_connector(1, Inches(table_x), Inches(3.30),
                            Inches(table_x + table_w), Inches(3.30))
ln.line.color.rgb = BLACK; ln.line.width = Pt(2.0)

add_text(s, table_x,        3.40, 4.0, 0.40, "Value",
         size=14, color=BLACK, align=PP_ALIGN.LEFT)
add_text(s, table_x + 4.0,  3.40, 2.4, 0.40, "RL",
         size=14, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 6.4,  3.40, 2.4, 0.40, "With SBSFC",
         size=14, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 8.8,  3.40, 2.5, 0.40, "Without SBSFC",
         size=14, color=BLACK, align=PP_ALIGN.CENTER)

# Header rule
ln = s.shapes.add_connector(1, Inches(table_x), Inches(3.85),
                            Inches(table_x + table_w), Inches(3.85))
ln.line.color.rgb = BLACK; ln.line.width = Pt(1.0)

rows = [
    ("Mean ( |θ|  [degree] )",        "0.2377",  "0.3167",  "5.0659"),
    ("Variance ( θ  [degree²] )",     "0.1934",  "0.1931",  "35.4129"),
    ("Max ( |θ|  [degree] )",         "2.0809",  "1.7120",  "11.4124"),
]
y = 3.95
for v, c1, c2, c3 in rows:
    add_text(s, table_x,        y, 4.0, 0.40, v,
             size=13, color=BLACK, align=PP_ALIGN.LEFT)
    add_text(s, table_x + 4.0,  y, 2.4, 0.40, c1,
             size=13, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s, table_x + 6.4,  y, 2.4, 0.40, c2,
             size=13, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s, table_x + 8.8,  y, 2.5, 0.40, c3,
             size=13, color=BLACK, align=PP_ALIGN.CENTER)
    y += 0.45

# Bottom rule
ln = s.shapes.add_connector(1, Inches(table_x), Inches(y + 0.05),
                            Inches(table_x + table_w), Inches(y + 0.05))
ln.line.color.rgb = BLACK; ln.line.width = Pt(2.0)

# Reduction summary
add_box(s, 0.6, 5.85, 12.1, 1.30, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 5.95, 11.6, 0.40,
         "Reduction vs. Without SBSFC (LPF baseline):",
         size=14, bold=True, color=GREEN)
add_text(s, 0.85, 6.35, 11.6, 0.40,
         "•  SBSFC :   Mean = 93.7 %     Variance = 99.5 %",
         size=13, bold=True, color=GREEN, font="Consolas")
add_text(s, 0.85, 6.75, 11.6, 0.40,
         "•  RL    :   Mean = 95.3 %     Variance = 99.5 %",
         size=13, bold=True, color=GREEN, font="Consolas")


# ---------------- Slide 3 — Three layers of insight ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Three Layers of Insight",
           "Reading the numbers from shallow to deep")

layers = [
    ("Layer 1", "LPF is fundamentally inadequate", ORANGE,
     "LPF only smooths the input command. It does NOTHING to actively damp the slosh mode.\n"
     "5° mean and 11° peak  →  spilled liquid. Active control is mandatory, not optional."),
    ("Layer 2", "RL ≈ SBSFC  —  this is the real result", GREEN,
     "A learned policy MATCHES a physics-derived controller (Choi 2024).\n"
     "RL had no model of the slosh dynamics — it discovered the strategy from rewards alone."),
    ("Layer 3", "The trade-off is structurally meaningful", PURPLE,
     "RL wins steady-state (mean) — gradient descent fine-tunes across many states.\n"
     "SBSFC wins peaks (max) — its shaper is analytically optimal for the slosh frequency.\n"
     "This is physics, not bad tuning."),
]

y = 1.15
for tag, title, color, body in layers:
    add_box(s, 0.6, y, 12.1, 1.85, fill=LIGHT, line=color)
    add_text(s, 0.85, y + 0.10, 2.0, 0.45, tag,
             size=20, bold=True, color=color)
    add_text(s, 2.85, y + 0.10, 9.7, 0.45, title,
             size=18, bold=True, color=NAVY)
    add_text(s, 0.85, y + 0.65, 11.6, 1.20, body,
             size=13, color=NAVY)
    y += 2.05


# ---------------- Slide 4 — The key trade-off ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "The Key Trade-off — RL vs. SBSFC",
           "Each method wins where its design philosophy is strongest")

# Left half — RL wins mean
add_box(s, 0.6, 1.10, 5.95, 4.4, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 1.20, 5.7, 0.45, "RL wins  →  Mean",
         size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, 0.85, 1.75, 5.7, 0.5,
         "0.2377°   vs   0.3167°",
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 0.85, 2.40, 5.7, 0.4,
         "(RL is 25 % lower)",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 0.85, 3.0, 5.7, 0.4, "WHY?",
         size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, 0.85, 3.45, 5.7, 1.95,
         "•  Gradient descent fine-tunes\n"
         "    across millions of states\n"
         "•  Implicitly learns the dynamics\n"
         "•  Steady-state behavior is what\n"
         "    most reward signals capture",
         size=12, color=NAVY)

# Right half — SBSFC wins max
add_box(s, 6.75, 1.10, 5.95, 4.4, fill=LIGHT, line=ORANGE)
add_text(s, 7.0, 1.20, 5.7, 0.45, "SBSFC wins  →  Max",
         size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 7.0, 1.75, 5.7, 0.5,
         "1.7120°   vs   2.0809°",
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 7.0, 2.40, 5.7, 0.4,
         "(SBSFC is 18 % lower)",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 7.0, 3.0, 5.7, 0.4, "WHY?",
         size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 7.0, 3.45, 5.7, 1.95,
         "•  Input shaper places zeros\n"
         "    EXACTLY at the slosh poles\n"
         "•  T_e = π/ω_d  (Choi 2024)\n"
         "•  Mathematically optimal for\n"
         "    the dominant slosh frequency",
         size=12, color=NAVY)

# Bottom synthesis box
add_box(s, 0.6, 5.65, 12.1, 1.55, fill=NAVY, line=NAVY)
add_text(s, 0.85, 5.75, 11.6, 0.45,
         "Synthesis:",
         size=15, bold=True, color=WHITE)
add_text(s, 0.85, 6.20, 11.6, 0.95,
         "RL is empirically tuned across all states  →  better average.\n"
         "SBSFC is structurally locked to the slosh frequency  →  better peak.\n"
         "A hybrid (SBSFC shaper + RL residual) could combine both strengths.",
         size=13, color=WHITE)


# ---------------- Slide 5 — Fairness check (ψ) ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Fairness Check — Robot Tilt ψ",
           "Are the controllers doing equal work on the easy job?")

add_box(s, 0.6, 1.10, 12.1, 1.85, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 1.25, 11.6, 0.45,
         "ψ (robot tilt) comparison:",
         size=16, bold=True, color=NAVY)
add_text(s, 0.85, 1.75, 11.6, 0.40,
         "•  SBSFC :  mean |ψ| = 3.0956°       •  LPF :  mean |ψ| = 3.2138°",
         size=14, color=NAVY, font="Consolas")
add_text(s, 0.85, 2.20, 11.6, 0.65,
         "→  Nearly identical (~ 4 % difference)",
         size=14, bold=True, color=GREEN)

# Interpretation
add_box(s, 0.6, 3.10, 12.1, 4.05, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 3.20, 11.6, 0.45,
         "What this means — the comparison is FAIR:",
         size=15, bold=True, color=GREEN)

bullets = [
    "Stabilizing the unstable robot (LQT base) is the EASY job — both methods do it equally well.",
    "Suppressing slosh is the HARD job — only SBSFC's auxiliary compensator + DOB and RL's learned policy earn their keep here.",
    "The differences in performance appear ONLY in the sloshing channel θ, not in the stabilization channel ψ.",
    "→  This validates that we are comparing the slosh-suppression mechanisms head-to-head, not confounded by stabilization differences.",
]
for i, b in enumerate(bullets):
    add_text(s, 0.85, 3.70 + i * 0.78, 11.6, 0.75, "•  " + b,
             size=13, color=NAVY)


# ---------------- Slide 6 — Limitations Top 4 ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Limitations — Top 4 Critical",
           "Be honest with the audience about what we did NOT show")

lims = [
    ("A",
     "Tested on only 1 of 5 scenarios",
     "Choi 2024 used 5 disturbance types. RL might generalize — or fail — on scenarios 2–5. We don't know.",
     RED),
    ("B",
     "Simulation only — no hardware",
     "Plant is linearized; real slosh is nonlinear. Sim-to-real gap unknown. RL is especially vulnerable here.",
     RED),
    ("C",
     "RL is a black box (no stability proof)",
     "67 K weights, no Lyapunov certificate. SBSFC has explicit pole placement (eigs −71.7, −3.16, −0.16, −0.16). Safety-critical apps need certifiability.",
     ORANGE),
    ("D",
     "Reward weights tuned by trial-and-error",
     "50 : 20 : 2 : 0.05 chosen empirically. No formal sensitivity analysis on these ratios.",
     ORANGE),
]

y = 1.10
for tag, title, body, color in lims:
    add_box(s, 0.6, y, 12.1, 1.45, fill=LIGHT, line=color)
    add_text(s, 0.85, y + 0.08, 0.7, 0.5, tag,
             size=22, bold=True, color=color)
    add_text(s, 1.55, y + 0.08, 11.0, 0.45, title,
             size=16, bold=True, color=NAVY)
    add_text(s, 1.55, y + 0.55, 11.0, 0.85, body,
             size=12, color=NAVY)
    y += 1.55


# ---------------- Slide 7 — Other limitations ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Other Limitations",
           "Smaller but worth acknowledging")

others = [
    ("E",
     "Only one randomized plant",
     "Domain randomization covered ±20 % on L, m, damping. Real-world variation could be larger. Untested on different liquids."),
    ("F",
     "RL training cost",
     "600 episodes × 10 s = ~100 min simulated  →  hours of wall clock. SBSFC: one-time analytical design + LQT solve."),
    ("G",
     "Worst-case slosh not characterized",
     "Peak 2.08° (RL) — is that below spill threshold? Depends on liquid level, viscosity, cup geometry. Not tested."),
    ("H",
     "No hybrid controller tested",
     "SBSFC shaper (good for transients) + RL residual (good for steady state) — likely best-of-both-worlds, untried."),
]

y = 1.10
for tag, title, body in others:
    add_box(s, 0.6, y, 12.1, 1.45, fill=LIGHT, line=GRAY)
    add_text(s, 0.85, y + 0.08, 0.7, 0.5, tag,
             size=22, bold=True, color=GRAY)
    add_text(s, 1.55, y + 0.08, 11.0, 0.45, title,
             size=16, bold=True, color=NAVY)
    add_text(s, 1.55, y + 0.55, 11.0, 0.85, body,
             size=12, color=NAVY)
    y += 1.55


# ---------------- Slide 8 — What we can / cannot claim ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "What We Can Claim — and What We Cannot",
           "The honest scope of our results")

# Left — defensible claims
add_box(s, 0.6, 1.10, 6.05, 5.95, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 1.20, 5.7, 0.45, "✓  Defensible claims",
         size=18, bold=True, color=GREEN)

defensible = [
    "On Scenario 1, SAC matches the Choi 2024 SBSFC controller within 0.08° mean slosh.",
    "Model-free RL is VIABLE for slosh suppression on this class of systems.",
    "Both controllers reduce slosh by >93 % vs. naive LPF.",
    "Stabilization (ψ) is equally effective in both methods — the comparison is fair.",
    "RL favors steady-state; SBSFC favors transient peaks — predictable from design.",
]
for i, b in enumerate(defensible):
    add_text(s, 0.85, 1.75 + i * 1.05, 5.7, 1.0, "•  " + b,
             size=12, color=NAVY)

# Right — overclaims
add_box(s, 6.85, 1.10, 6.05, 5.95, fill=LIGHT, line=RED)
add_text(s, 7.10, 1.20, 5.7, 0.45, "✗  Cannot claim (yet)",
         size=18, bold=True, color=RED)

overclaims = [
    "✗  \"RL is better than SBSFC.\"  — RL wins mean, SBSFC wins max.",
    "✗  \"RL generalizes to all scenarios.\"  — only Scenario 1 tested.",
    "✗  \"RL works on hardware.\"  — simulation only.",
    "✗  \"Stable for any plant variation.\"  — only ±20 % randomization tested.",
    "✗  \"RL is the future for slosh control.\"  — too strong without robustness data.",
]
for i, b in enumerate(overclaims):
    add_text(s, 7.10, 1.75 + i * 1.05, 5.7, 1.0, b,
             size=12, color=NAVY)


# ---------------- Slide 9 — Next steps ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Next Steps — Three Priority Tiers",
           "Ordered by value-per-week of work")

# Tier 1 — must-do
add_box(s, 0.6, 1.05, 12.1, 1.85, fill=LIGHT, line=ORANGE)
add_text(s, 0.85, 1.15, 11.6, 0.40,
         "Tier 1 — Must-do  (each ≤ 1 week)",
         size=16, bold=True, color=ORANGE)
add_text(s, 0.85, 1.55, 11.6, 1.30,
         "1.  Test RL on all 5 scenarios.    ← single highest-value experiment remaining\n"
         "2.  Compute |ψ| metric for RL  (currently shown as “—” in our table).\n"
         "3.  Reward-weight sensitivity sweep  (3 alternative ratios).",
         size=13, color=NAVY)

# Tier 2 — strong additions
add_box(s, 0.6, 3.05, 12.1, 1.85, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 3.15, 11.6, 0.40,
         "Tier 2 — Strong additions  (each ≤ 2 weeks)",
         size=16, bold=True, color=GREEN)
add_text(s, 0.85, 3.55, 11.6, 1.30,
         "4.  Hybrid controller experiment  →  SBSFC shaper + RL residual correction.\n"
         "5.  Robustness test  →  evaluate trained agent on plant params OUTSIDE the randomization range.\n"
         "6.  Training curves and wall-clock time  →  show convergence is stable.",
         size=13, color=NAVY)

# Tier 3 — stretch
add_box(s, 0.6, 5.05, 12.1, 2.10, fill=LIGHT, line=PURPLE)
add_text(s, 0.85, 5.15, 11.6, 0.40,
         "Tier 3 — Stretch goals  (research-paper level)",
         size=16, bold=True, color=PURPLE)
add_text(s, 0.85, 5.55, 11.6, 1.55,
         "7.  Hardware deployment  →  cart + cup of water rig.\n"
         "8.  Different liquids  →  viscosity sensitivity (water vs. honey).\n"
         "9.  4-method comparison  →  add MPC  (mpc_pendcartImplicitMPC.slxc already in matlab/).\n"
         "10. Lyapunov-style stability analysis for the trained RL policy.",
         size=13, color=NAVY)


# ---------------- Slide 10 — Bottom-line conclusion ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Conclusion",
           "The honest summary of what this project showed")

# Big quote box
add_box(s, 0.6, 1.05, 12.1, 2.4, fill=NAVY, line=NAVY)
add_text(s, 0.85, 1.20, 11.6, 0.45,
         "The bottom-line message:",
         size=15, bold=True, color=ORANGE)
add_text(s, 0.85, 1.70, 11.6, 1.65,
         "\"We did not aim to beat SBSFC.\n"
         "  We aimed to show that a model-free policy can reach the same performance\n"
         "  as a model-based design  —  which is exactly what we demonstrated.\n"
         "  The next phase is testing whether that equivalence holds across all\n"
         "  5 scenarios and on hardware.\"",
         size=15, italic=True, color=WHITE)

# Three remember-this boxes
boxes = [
    ("≥ 93 %",
     "Slosh reduction",
     "Both SBSFC and RL crush the LPF baseline",
     GREEN),
    ("0.08°",
     "Mean gap",
     "RL beats SBSFC mean by only 0.08° — they're effectively tied",
     ORANGE),
    ("1 of 5",
     "Scenarios tested",
     "Generalization to scenarios 2–5 is the next priority",
     RED),
]

y = 3.65
x = 0.6
w = 4.0
gap = 0.05
for big, label, desc, color in boxes:
    add_box(s, x, y, w, 3.4, fill=LIGHT, line=color)
    add_text(s, x, y + 0.20, w, 0.9, big,
             size=36, bold=True, color=color, align=PP_ALIGN.CENTER, font="Consolas")
    add_text(s, x, y + 1.20, w, 0.4, label,
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + 1.75, w, 1.5, desc,
             size=12, color=NAVY, align=PP_ALIGN.CENTER)
    x += w + gap


# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/results_analysis.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
