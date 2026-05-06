"""
Teaching PPT: Full Critic Pipeline — All 3 Layers
For MECE 6397 anti-sloshing robot project (Choi et al. 2024)

Critic input: [obs (7) ; action (1)] = 8-dim
Critic output: Q-value (single scalar — expected return)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=NAVY, line_w=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=ORANGE, weight=3):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = color
    return line


def header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, 0.4, 0.08, 12.5, 0.55, title,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.4, 0.78, 12.5, 0.4, subtitle,
                 size=14, color=GRAY)


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------------- Slide 1 — Title ----------------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

add_text(s, 0.6, 2.2, 12, 1.4,
         "Full Critic Pipeline",
         size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.7, 12, 0.7,
         "All 3 layers of the SAC critic Q-network",
         size=24, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.5, 12, 0.6,
         "[obs (7) ; action (1)]  →  256  →  256  →  Q-value (1)",
         size=20, color=LIGHT, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 0.6, 5.4, 12, 0.6,
         "MECE 6397 — Anti-Sloshing Robot Project",
         size=16, color=LIGHT, align=PP_ALIGN.CENTER)


# ---------------- Slide 2 — Pipeline overview (with W1 matrix in Layer 1) ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Full Critic Pipeline — Overview",
           "Input → Layer 1 (with W₁ matrix) → Layer 2 → Layer 3 → Output")

# INPUT box (8-dim: obs + action)
y_in = 0.80
add_box(s, 3.5, y_in, 6.3, 0.50, fill=LIGHT, line=NAVY)
add_text(s, 3.5, y_in + 0.05, 6.3, 0.25, "INPUT",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 3.5, y_in + 0.27, 6.3, 0.22,
         "[obs ; action] = [ψ, ψ̇, x, ẋ, θ, θ̇, v_ref, u]   shape: (8,)",
         size=10, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_arrow(s, 6.65, y_in + 0.52, 6.65, y_in + 0.65)

# LAYER 1  (taller — includes W1 matrix structure, 256 x 8)
y1 = 1.45
add_box(s, 1.0, y1, 11.3, 2.50, fill=LIGHT, line=ORANGE)
add_text(s, 1.2, y1 + 0.04, 11.0, 0.30, "LAYER 1   —   feature extraction",
         size=14, bold=True, color=ORANGE)
add_text(s, 1.2, y1 + 0.35, 11.0, 0.25,
         "z₁ = W₁ · [obs;u] + b₁         W₁:(256, 8)   b₁:(256)   total: 2,304 weights",
         size=11, color=NAVY, font="Consolas")

# Matrix sub-box inside Layer 1 (W1: 256 x 8)
add_box(s, 1.4, y1 + 0.66, 10.5, 1.45, fill=WHITE, line=NAVY, line_w=0.6)
add_text(s, 1.55, y1 + 0.72, 10.3, 1.35,
         "         [  w₁,₁     w₁,₂   ...   w₁,₈ ]   ← row 1   (8 weights for output #1)\n"
         "         [  w₂,₁     w₂,₂   ...   w₂,₈ ]   ← row 2\n"
         "W₁  =   [    .         .       .     .   ]\n"
         "         [    .         .       .     .   ]\n"
         "         [w₂₅₆,₁  w₂₅₆,₂   ...  w₂₅₆,₈]   ← row 256",
         size=10, color=NAVY, font="Consolas")

add_text(s, 1.2, y1 + 2.15, 11.0, 0.25,
         "h₁ = ReLU(z₁)        shape: (256)        Purpose: extract 256 features from (state, action)",
         size=10, color=GRAY, font="Consolas")
add_arrow(s, 6.65, y1 + 2.52, 6.65, y1 + 2.62)

# LAYER 2
y2 = 4.10
add_box(s, 1.0, y2, 11.3, 1.05, fill=LIGHT, line=ORANGE)
add_text(s, 1.2, y2 + 0.04, 11.0, 0.30, "LAYER 2   —   feature combination",
         size=13, bold=True, color=ORANGE)
add_text(s, 1.2, y2 + 0.34, 11.0, 0.25,
         "z₂ = W₂ · h₁ + b₂          W₂:(256, 256)   b₂:(256)   total: 65,792 weights",
         size=11, color=NAVY, font="Consolas")
add_text(s, 1.2, y2 + 0.62, 11.0, 0.25,
         "h₂ = ReLU(z₂)        shape: (256)        Purpose: combine features → higher abstractions",
         size=10, color=GRAY, font="Consolas")
add_arrow(s, 6.65, y2 + 1.07, 6.65, y2 + 1.18, color=GREEN)

# LAYER 3 (output — Q-value, single scalar, no activation)
y3 = 5.30
add_box(s, 1.0, y3, 11.3, 1.05, fill=LIGHT, line=GREEN)
add_text(s, 1.2, y3 + 0.04, 11.0, 0.30, "LAYER 3   —   output (no activation)",
         size=13, bold=True, color=GREEN)
add_text(s, 1.2, y3 + 0.34, 11.0, 0.25,
         "z₃ = W₃ · h₂ + b₃          W₃:(1, 256)    b₃:(1)     total: 257 weights",
         size=11, color=NAVY, font="Consolas")
add_text(s, 1.2, y3 + 0.62, 11.0, 0.25,
         "Q(s, a) = z₃    (no ReLU — Q-value can be ±, not clipped)",
         size=10, color=GRAY, font="Consolas")
add_arrow(s, 6.65, y3 + 1.07, 6.65, y3 + 1.18, color=NAVY)

# OUTPUT
y_out = 6.50
add_box(s, 3.5, y_out, 6.3, 0.50, fill=ORANGE, line=ORANGE)
add_text(s, 3.5, y_out + 0.04, 6.3, 0.25, "OUTPUT",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 3.5, y_out + 0.27, 6.3, 0.22,
         "Q(s, a)   (e.g., −12.4 — expected return)",
         size=10, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")


# ---------------- Slide 3 — INPUT ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Input — What the Critic Receives",
           "An 8-dim vector: state (7) + action (1)")

add_box(s, 1.0, 1.3, 11.3, 1.5, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.45, 11.3, 0.4, "Input vector (state-action pair):",
         size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 1.85, 11.3, 0.6,
         "[ ψ ,  ψ̇ ,  x ,  ẋ ,  θ ,  θ̇ ,  v_ref ,  u ]",
         size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 2.5, 11.3, 0.3,
         "shape: (8,)",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

items = [
    ("ψ",     "robot tilt angle (rad)"),
    ("ψ̇",     "tilt angular velocity (rad/s)"),
    ("x",     "cart position (m)"),
    ("ẋ",     "cart velocity (m/s)"),
    ("θ",     "slosh angle (rad)"),
    ("θ̇",     "slosh angular velocity (rad/s)"),
    ("v_ref", "commanded velocity reference (m/s)"),
    ("u",     "action chosen by actor (m/s²)  ← extra input vs. actor"),
]

add_text(s, 0.8, 3.05, 12, 0.4,
         "Each entry — what it means:",
         size=15, bold=True, color=GREEN)

y = 3.5
for i, (var, desc) in enumerate(items):
    add_box(s, 0.8, y + i * 0.45, 11.7, 0.40, fill=LIGHT, line=NAVY, line_w=0.6)
    add_text(s, 1.0, y + i * 0.45 + 0.06, 1.5, 0.32, var,
             size=13, bold=True, color=ORANGE, font="Consolas")
    add_text(s, 2.6, y + i * 0.45 + 0.06, 9.5, 0.32, desc,
             size=12, color=NAVY)


# ---------------- Slide 4 — LAYER 1 ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer 1 — Feature Extraction",
           "Transform 8 raw inputs (state + action) into 256 learned features")

add_box(s, 0.6, 1.1, 12.1, 6.2, fill=LIGHT, line=ORANGE)

add_text(s, 0.85, 1.18, 11.6, 0.4, "LAYER 1",
         size=18, bold=True, color=ORANGE)

# Math operation
add_text(s, 0.85, 1.55, 11.6, 0.35, "Math operation:",
         size=12, bold=True, color=NAVY)
add_box(s, 1.5, 1.88, 10.3, 0.5, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 1.95, 10.3, 0.4,
         "z₁  =  W₁ · [obs ; u]  +  b₁",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

# Explicit matrix structure (256 x 8)
add_text(s, 0.85, 2.55, 11.6, 0.35, "Weight matrix W₁ (size 256 × 8):",
         size=12, bold=True, color=NAVY)

add_box(s, 0.85, 2.95, 11.6, 1.95, fill=WHITE, line=NAVY, line_w=0.8)
add_text(s, 1.0, 3.05, 11.4, 1.85,
         "         [  w₁,₁     w₁,₂   ...   w₁,₈ ]   ← row 1   (8 weights for output #1)\n"
         "         [  w₂,₁     w₂,₂   ...   w₂,₈ ]   ← row 2\n"
         "W₁  =   [    .         .       .     .   ]\n"
         "         [    .         .       .     .   ]\n"
         "         [w₂₅₆,₁  w₂₅₆,₂   ...  w₂₅₆,₈]   ← row 256",
         size=11, color=NAVY, font="Consolas")

# Shapes/totals
add_text(s, 0.85, 5.05, 11.6, 0.35, "Shapes & totals:",
         size=12, bold=True, color=NAVY)
add_text(s, 1.2, 5.4, 11.0, 0.32,
         "•  W₁: (256, 8)     b₁: (256,)     Total: 256 × 8 + 256 = 2,304 weights",
         size=11, color=NAVY, font="Consolas")

# Activation
add_text(s, 0.85, 5.78, 11.6, 0.35, "Activation:",
         size=12, bold=True, color=NAVY)
add_box(s, 1.5, 6.1, 10.3, 0.5, fill=GREEN, line=GREEN)
add_text(s, 1.5, 6.18, 10.3, 0.4,
         "h₁  =  ReLU(z₁)        (negatives → 0)        shape: (256,)",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 6.75, 11.6, 0.4,
         "Purpose: extract 256 'features' from the (state, action) pair.",
         size=12, bold=True, color=GREEN)


# ---------------- Slide 5 — LAYER 2 ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer 2 — Feature Combination",
           "Combine 256 features into 256 higher-level abstractions")

add_box(s, 0.6, 1.2, 12.1, 5.8, fill=LIGHT, line=ORANGE)

add_text(s, 0.85, 1.3, 11.6, 0.4, "LAYER 2",
         size=20, bold=True, color=ORANGE)
add_text(s, 0.85, 1.7, 11.6, 0.4, "Math operation:",
         size=14, bold=True, color=NAVY)

add_box(s, 1.5, 2.1, 10.3, 0.65, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 2.18, 10.3, 0.5,
         "z₂  =  W₂ · h₁  +  b₂",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 2.95, 11.6, 0.4,
         "Shapes:",
         size=14, bold=True, color=NAVY)
add_text(s, 1.2, 3.3, 11.0, 0.4,
         "•  W₂ shape:  (256, 256)       — 256 rows × 256 columns",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 3.65, 11.0, 0.4,
         "•  b₂ shape:  (256,)            — one bias per output neuron",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 4.0, 11.0, 0.4,
         "•  Total weights:  256 × 256 + 256  =  65,792    ← largest layer!",
         size=13, bold=True, color=ORANGE, font="Consolas")

add_text(s, 0.85, 4.55, 11.6, 0.4,
         "Activation:",
         size=14, bold=True, color=NAVY)
add_box(s, 1.5, 4.95, 10.3, 0.65, fill=GREEN, line=GREEN)
add_text(s, 1.5, 5.05, 10.3, 0.45,
         "h₂  =  ReLU(z₂)",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 5.85, 11.6, 0.4,
         "Output shape:  (256,)",
         size=13, color=NAVY)
add_text(s, 0.85, 6.25, 11.6, 0.4,
         "Purpose:",
         size=14, bold=True, color=GREEN)
add_text(s, 0.85, 6.6, 11.6, 0.4,
         "Combine Layer 1 features into higher-level abstractions used to score the (state, action) pair.",
         size=13, color=NAVY)


# ---------------- Slide 6 — LAYER 3 ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer 3 — Output (No Activation)",
           "Collapse 256 features into a single Q-value")

add_box(s, 0.6, 1.2, 12.1, 5.8, fill=LIGHT, line=GREEN)

add_text(s, 0.85, 1.3, 11.6, 0.4, "LAYER 3 (output)",
         size=20, bold=True, color=GREEN)
add_text(s, 0.85, 1.7, 11.6, 0.4, "Math operation:",
         size=14, bold=True, color=NAVY)

add_box(s, 1.5, 2.1, 10.3, 0.65, fill=GREEN, line=GREEN)
add_text(s, 1.5, 2.18, 10.3, 0.5,
         "z₃  =  W₃ · h₂  +  b₃",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 2.95, 11.6, 0.4,
         "Shapes:",
         size=14, bold=True, color=NAVY)
add_text(s, 1.2, 3.3, 11.0, 0.4,
         "•  W₃ shape:  (1, 256)         — one row, 256 weights",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 3.65, 11.0, 0.4,
         "•  b₃ shape:  (1,)              — one bias",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 4.0, 11.0, 0.4,
         "•  Total weights:  256 × 1 + 1  =  257",
         size=13, bold=True, color=GREEN, font="Consolas")

add_text(s, 0.85, 4.55, 11.6, 0.4,
         "Activation:   NONE",
         size=14, bold=True, color=NAVY)
add_box(s, 1.5, 4.95, 10.3, 0.65, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 5.05, 10.3, 0.45,
         "Q(s, a)  =  z₃           (Q-value, can be ±)",
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 5.85, 11.6, 0.4,
         "Not clipped — Q can be any real number (typically negative for our reward shape).",
         size=13, color=NAVY)
add_text(s, 0.85, 6.25, 11.6, 0.4,
         "Purpose:",
         size=14, bold=True, color=GREEN)
add_text(s, 0.85, 6.6, 11.6, 0.4,
         "Collapse 256 features into a single number = expected discounted return from (s, a).",
         size=13, color=NAVY)


# ---------------- Slide 7 — Output / What the Q-value means ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Output — A Single Q-Value",
           "What the critic produces and how it's used")

add_box(s, 1.0, 1.3, 11.3, 2.3, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 1.45, 11.3, 0.5, "OUTPUT",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 2.0, 11.3, 0.7,
         "Q(s, a)",
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 2.85, 11.3, 0.5,
         "(a single scalar — the predicted total reward starting from state s if action a is taken)",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_box(s, 0.8, 3.85, 12.0, 3.2, fill=LIGHT, line=GREEN)
add_text(s, 1.0, 3.95, 11.6, 0.4, "What this single number represents and how it's used:",
         size=15, bold=True, color=GREEN)

bullets = [
    "Q(s, a) ≈ expected discounted return Σ γᵏ · r(t+k)  if we follow the policy from (s, a).",
    "More negative Q  →  worse state-action pair (lots of slosh / control penalty ahead).",
    "The actor is updated to maximize Q  —  so the critic teaches the actor what's good.",
    "Two critics Q₁ and Q₂ are trained; SAC takes  min(Q₁, Q₂)  to fight overestimation.",
    "Targets via Bellman:  y = r + γ · min(Q̄₁, Q̄₂)(s', a') − α · log π(a'|s').",
    "Recomputed every training step from the replay buffer minibatch (size 256).",
]

for i, b in enumerate(bullets):
    add_text(s, 1.0, 4.4 + i * 0.42, 11.6, 0.4, "•  " + b,
             size=12, color=NAVY)


# ---------------- Slide 8 — Number summary ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer-by-Layer Numbers — One Critic",
           "Shapes and weight counts (input is 8-dim, not 7)")

add_box(s, 0.7, 1.2, 11.9, 0.55, fill=NAVY, line=NAVY)
hdrs = ["", "Layer 1", "Layer 2", "Layer 3"]
xs = [0.7, 4.0, 7.0, 10.0]
ws = [3.3, 3.0, 3.0, 2.6]
for i, hd in enumerate(hdrs):
    add_text(s, xs[i], 1.3, ws[i], 0.4, hd,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Input dim",          "8 (obs+u)",   "256",         "256"),
    ("Output dim",         "256",         "256",         "1"),
    ("W shape",            "(256, 8)",    "(256, 256)",  "(1, 256)"),
    ("b shape",            "(256,)",      "(256,)",      "(1,)"),
    ("Weights count",      "2,048",       "65,536",      "256"),
    ("Biases count",       "256",         "256",         "1"),
    ("TOTAL params",       "2,304",       "65,792",      "257"),
    ("Activation",         "ReLU",        "ReLU",        "none"),
    ("Purpose",            "extract",     "combine",     "score"),
]

for r, row in enumerate(rows):
    yy = 1.78 + r * 0.45
    fill = LIGHT if r % 2 == 0 else WHITE
    add_box(s, 0.7, yy, 11.9, 0.42, fill=fill, line=NAVY, line_w=0.5)
    is_total = (row[0] == "TOTAL params")
    for c, cell in enumerate(row):
        col = ORANGE if (is_total and c > 0) else NAVY
        bold = is_total or (c == 0)
        font = "Consolas" if c > 0 else "Calibri"
        add_text(s, xs[c], yy + 0.05, ws[c], 0.32, cell,
                 size=12, bold=bold, color=col,
                 align=PP_ALIGN.CENTER, font=font)

add_box(s, 1.0, 6.1, 11.3, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 6.2, 11.3, 0.4,
         "Total per critic:",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.55, 11.3, 0.5,
         "2,304  +  65,792  +  257  =  68,353  ≈  67K   (×2 critics ≈ 134K)",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")


# ---------------- Slide 9 — Actor vs Critic comparison ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Actor vs. Critic — What's Different?",
           "Same 3-layer MLP shape, different inputs and outputs")

add_box(s, 0.7, 1.2, 11.9, 0.55, fill=NAVY, line=NAVY)
hdrs = ["", "ACTOR  π(a|s)", "CRITIC  Q(s,a)"]
xs = [0.7, 4.5, 9.0]
ws = [3.7, 4.4, 3.5]
for i, hd in enumerate(hdrs):
    add_text(s, xs[i], 1.3, ws[i], 0.4, hd,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Input",              "obs (7)",                  "[obs ; action]  (8)"),
    ("Output",             "action u (1, clipped)",    "Q-value (1, unclipped)"),
    ("Output meaning",     "what to do",               "how good (s,a) is"),
    ("Hidden layers",      "256, 256",                 "256, 256"),
    ("Layer 1 W shape",    "(256, 7)",                 "(256, 8)"),
    ("# of networks",      "1",                        "2  (Q₁, Q₂) + 2 targets"),
    ("Trained to",         "maximize Q(s, π(s))",      "match Bellman target y"),
    ("Total params",       "≈ 67K",                    "≈ 67K each  →  ≈ 134K"),
]

for r, row in enumerate(rows):
    yy = 1.78 + r * 0.50
    fill = LIGHT if r % 2 == 0 else WHITE
    add_box(s, 0.7, yy, 11.9, 0.47, fill=fill, line=NAVY, line_w=0.5)
    for c, cell in enumerate(row):
        bold = (c == 0)
        col = NAVY if c == 0 else (ORANGE if c == 1 else GREEN)
        font = "Calibri" if c == 0 else "Consolas"
        add_text(s, xs[c], yy + 0.08, ws[c], 0.35, cell,
                 size=12, bold=bold, color=col,
                 align=PP_ALIGN.CENTER, font=font)


# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/critic_pipeline_explained.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
