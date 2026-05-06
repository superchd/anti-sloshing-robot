"""
References PPT — anti-sloshing robot project
4 slides:
    1. Primary reference (Choi et al. 2024)
    2. SAC references (Haarnoja 2018a, 2018b)
    3. Supporting references (Singer & Seering, MATLAB RL Toolbox)
    4. Full IEEE-style bibliography
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY   = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF4, 0xF1, 0xEC)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_hline(slide, x, y, w, weight=1.5, color=BLACK):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y),
                                      Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = color
    return line


def add_box(slide, x, y, w, h, fill=LIGHT, line=NAVY):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    if shp.has_text_frame:
        shp.text_frame.text = ""
    return shp


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ============================================================
# Slide 1 — Primary reference: Choi et al. 2024
# ============================================================
s = prs.slides.add_slide(blank)

add_text(s, 0.6, 0.40, 12.0, 0.55,
         "References — 1 / 4", size=14, color=GRAY)
add_text(s, 0.6, 0.85, 12.0, 0.65,
         "Primary Reference", size=28, bold=True, color=NAVY)
add_hline(s, 0.6, 1.55, 12.0, weight=2.0, color=NAVY)

# Citation card
add_box(s, 0.6, 1.85, 12.0, 1.85, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 1.95, 11.5, 0.45,
         "[1]  Choi et al. (2024)", size=18, bold=True, color=NAVY)
add_text(s, 0.85, 2.45, 11.5, 1.20,
         "S.-H. Choi, K.-S. Park, and J.-H. Kim, \"Anti-sloshing control of a "
         "serving robot using a state-based switching feedback controller,\" "
         "Mechatronics, vol. 98, p. 103123, 2024.",
         size=15, color=BLACK, italic=True)

# What we borrowed
add_text(s, 0.6, 4.00, 12.0, 0.45,
         "What we borrowed from this paper:",
         size=18, bold=True, color=ORANGE)

bullets = [
    "•  SBSFC formulation — input-shaper coefficients (T_e = π/ω_d, μ = −δω_f)",
    "•  Plant model — 6-state linearization of pendulum-on-cart system",
    "•  Scenario design — sudden start/stop, step input, disturbance test",
    "•  Table 3 layout — academic comparison table for evaluation metrics",
    "•  Evaluation metrics — mean |θ|, variance(θ), max |θ|",
]
y = 4.50
for b in bullets:
    add_text(s, 0.85, y, 12.0, 0.40, b, size=14, color=BLACK)
    y += 0.42

add_text(s, 0.6, 6.85, 12.0, 0.40,
         "PDF available in project root.",
         size=12, color=GRAY, italic=True)


# ============================================================
# Slide 2 — SAC references: Haarnoja 2018a, 2018b
# ============================================================
s = prs.slides.add_slide(blank)

add_text(s, 0.6, 0.40, 12.0, 0.55,
         "References — 2 / 4", size=14, color=GRAY)
add_text(s, 0.6, 0.85, 12.0, 0.65,
         "Soft Actor-Critic (SAC)", size=28, bold=True, color=NAVY)
add_hline(s, 0.6, 1.55, 12.0, weight=2.0, color=NAVY)

# --- Haarnoja 2018a ---
add_box(s, 0.6, 1.80, 12.0, 1.85, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 1.90, 11.5, 0.45,
         "[2]  Haarnoja et al. (2018a) — Original SAC",
         size=17, bold=True, color=NAVY)
add_text(s, 0.85, 2.40, 11.5, 1.20,
         "T. Haarnoja, A. Zhou, P. Abbeel, and S. Levine, \"Soft actor-critic: "
         "Off-policy maximum entropy deep reinforcement learning with a stochastic "
         "actor,\" in Proc. ICML, 2018, pp. 1861–1870.",
         size=14, color=BLACK, italic=True)
add_text(s, 0.85, 3.20, 11.5, 0.40,
         "→  Provides the SAC algorithm: 1 actor + 2 critics + entropy maximization.",
         size=13, color=ORANGE)

# --- Haarnoja 2018b ---
add_box(s, 0.6, 3.85, 12.0, 1.85, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 3.95, 11.5, 0.45,
         "[3]  Haarnoja et al. (2018b) — SAC with auto-tuned α",
         size=17, bold=True, color=NAVY)
add_text(s, 0.85, 4.45, 11.5, 1.20,
         "T. Haarnoja et al., \"Soft actor-critic algorithms and applications,\" "
         "arXiv preprint arXiv:1812.05905, 2018.",
         size=14, color=BLACK, italic=True)
add_text(s, 0.85, 5.25, 11.5, 0.40,
         "→  Source of the target-entropy heuristic  H̄ = −dim(action).",
         size=13, color=ORANGE)

# What we borrowed
add_text(s, 0.6, 6.00, 12.0, 0.45,
         "What we borrowed:",
         size=16, bold=True, color=ORANGE)
add_text(s, 0.85, 6.45, 12.0, 0.40,
         "•  SAC actor/critic loss formulation, double-Q trick, target networks",
         size=13, color=BLACK)
add_text(s, 0.85, 6.85, 12.0, 0.40,
         "•  Default hyperparameters:  γ=0.99, τ=5e-3, batch=256, replay=1M",
         size=13, color=BLACK)


# ============================================================
# Slide 3 — Supporting references
# ============================================================
s = prs.slides.add_slide(blank)

add_text(s, 0.6, 0.40, 12.0, 0.55,
         "References — 3 / 4", size=14, color=GRAY)
add_text(s, 0.6, 0.85, 12.0, 0.65,
         "Supporting References", size=28, bold=True, color=NAVY)
add_hline(s, 0.6, 1.55, 12.0, weight=2.0, color=NAVY)

# --- Singer & Seering ---
add_box(s, 0.6, 1.80, 12.0, 1.65, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 1.90, 11.5, 0.45,
         "[4]  Singer & Seering (1990) — Input Shaping",
         size=17, bold=True, color=NAVY)
add_text(s, 0.85, 2.40, 11.5, 1.00,
         "N. C. Singer and W. P. Seering, \"Preshaping command inputs to reduce "
         "system vibration,\" J. Dyn. Sys. Meas. Control, vol. 112, no. 1, pp. 76–82, 1990.",
         size=14, color=BLACK, italic=True)
add_text(s, 0.85, 3.05, 11.5, 0.40,
         "→  Foundational paper on impulse-based input shaping (basis of SBSFC).",
         size=13, color=ORANGE)

# --- MATLAB RL Toolbox ---
add_box(s, 0.6, 3.65, 12.0, 1.65, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 3.75, 11.5, 0.45,
         "[5]  MathWorks (2024) — Reinforcement Learning Toolbox",
         size=17, bold=True, color=NAVY)
add_text(s, 0.85, 4.25, 11.5, 1.00,
         "MathWorks, \"Reinforcement Learning Toolbox User's Guide (R2024a),\" "
         "https://www.mathworks.com/help/reinforcement-learning/, 2024.",
         size=14, color=BLACK, italic=True)
add_text(s, 0.85, 4.90, 11.5, 0.40,
         "→  rlSACAgent, rlNumericSpec, rlFunctionEnv, getAction APIs.",
         size=13, color=ORANGE)

# --- Sutton & Barto ---
add_box(s, 0.6, 5.50, 12.0, 1.50, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 5.60, 11.5, 0.45,
         "[6]  Sutton & Barto (2018) — RL Textbook",
         size=17, bold=True, color=NAVY)
add_text(s, 0.85, 6.10, 11.5, 1.00,
         "R. S. Sutton and A. G. Barto, Reinforcement Learning: An Introduction, "
         "2nd ed. Cambridge, MA: MIT Press, 2018.",
         size=14, color=BLACK, italic=True)
add_text(s, 0.85, 6.65, 11.5, 0.40,
         "→  General RL background:  γ, replay buffer, actor-critic methods.",
         size=13, color=ORANGE)


# ============================================================
# Slide 4 — Full IEEE-style bibliography
# ============================================================
s = prs.slides.add_slide(blank)

add_text(s, 0.6, 0.40, 12.0, 0.55,
         "References — 4 / 4", size=14, color=GRAY)
add_text(s, 0.6, 0.85, 12.0, 0.65,
         "Full Bibliography", size=28, bold=True, color=NAVY)
add_hline(s, 0.6, 1.55, 12.0, weight=2.0, color=NAVY)

refs = [
    ("[1]",
     "S.-H. Choi, K.-S. Park, and J.-H. Kim, \"Anti-sloshing control of a "
     "serving robot using a state-based switching feedback controller,\" "
     "Mechatronics, vol. 98, p. 103123, 2024."),
    ("[2]",
     "T. Haarnoja, A. Zhou, P. Abbeel, and S. Levine, \"Soft actor-critic: "
     "Off-policy maximum entropy deep reinforcement learning with a stochastic "
     "actor,\" in Proc. Int. Conf. Mach. Learn. (ICML), 2018, pp. 1861–1870."),
    ("[3]",
     "T. Haarnoja et al., \"Soft actor-critic algorithms and applications,\" "
     "arXiv preprint arXiv:1812.05905, 2018."),
    ("[4]",
     "N. C. Singer and W. P. Seering, \"Preshaping command inputs to reduce "
     "system vibration,\" J. Dyn. Sys. Meas. Control, vol. 112, no. 1, pp. 76–82, 1990."),
    ("[5]",
     "MathWorks, \"Reinforcement Learning Toolbox User's Guide (R2024a),\" "
     "https://www.mathworks.com/help/reinforcement-learning/, 2024."),
    ("[6]",
     "R. S. Sutton and A. G. Barto, Reinforcement Learning: An Introduction, "
     "2nd ed. Cambridge, MA, USA: MIT Press, 2018."),
]

y = 1.85
for tag, body in refs:
    add_text(s, 0.6,  y, 0.55, 0.85, tag,
             size=13, bold=True, color=NAVY)
    add_text(s, 1.10, y, 11.6, 0.85, body,
             size=13, color=BLACK)
    y += 0.85

add_hline(s, 0.6, 7.00, 12.0, weight=1.0, color=GRAY)
add_text(s, 0.6, 7.10, 12.0, 0.35,
         "*IEEE citation style.  Primary reference [1] in project root; "
         "[2]–[6] available online.",
         size=11, color=GRAY, italic=True)


# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/references.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
