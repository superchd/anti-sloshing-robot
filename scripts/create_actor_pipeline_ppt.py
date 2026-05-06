"""
Teaching PPT: Full Actor Pipeline — All 3 Layers
For MECE 6397 anti-sloshing robot project (Choi et al. 2024)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=NAVY, line_w=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=ORANGE, weight=3):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = color
    return line


def header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, 0.4, 0.08, 12.5, 0.55, title,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.4, 0.78, 12.5, 0.4, subtitle,
                 size=14, color=GRAY)


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------------- Slide 1 — Title ----------------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

add_text(s, 0.6, 2.2, 12, 1.4,
         "Full Actor Pipeline",
         size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.7, 12, 0.7,
         "All 3 layers of the SAC actor network",
         size=24, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.5, 12, 0.6,
         "obs (7)  →  256  →  256  →  action (1)",
         size=20, color=LIGHT, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 0.6, 5.4, 12, 0.6,
         "MECE 6397 — Anti-Sloshing Robot Project",
         size=16, color=LIGHT, align=PP_ALIGN.CENTER)


# ---------------- Slide 2 — Pipeline overview (with W1 matrix in Layer 1) ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Full Actor Pipeline — Overview",
           "Input → Layer 1 (with W₁ matrix) → Layer 2 → Layer 3 → Output")

# INPUT box
y_in = 0.80
add_box(s, 3.5, y_in, 6.3, 0.50, fill=LIGHT, line=NAVY)
add_text(s, 3.5, y_in + 0.05, 6.3, 0.25, "INPUT",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 3.5, y_in + 0.27, 6.3, 0.22,
         "obs = [ψ, ψ̇, x, ẋ, θ, θ̇, v_ref]   shape: (7,)",
         size=10, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_arrow(s, 6.65, y_in + 0.52, 6.65, y_in + 0.65)

# LAYER 1  (taller — includes W1 matrix structure)
y1 = 1.45
add_box(s, 1.0, y1, 11.3, 2.50, fill=LIGHT, line=ORANGE)
add_text(s, 1.2, y1 + 0.04, 11.0, 0.30, "LAYER 1   —   feature extraction",
         size=14, bold=True, color=ORANGE)
add_text(s, 1.2, y1 + 0.35, 11.0, 0.25,
         "z₁ = W₁ · obs + b₁          W₁:(256, 7)   b₁:(256)   total: 2,048 weights",
         size=11, color=NAVY, font="Consolas")

# Matrix sub-box inside Layer 1
add_box(s, 1.4, y1 + 0.66, 10.5, 1.45, fill=WHITE, line=NAVY, line_w=0.6)
add_text(s, 1.55, y1 + 0.72, 10.3, 1.35,
         "         [  w₁,₁     w₁,₂   ...   w₁,₇ ]   ← row 1   (7 weights for output #1)\n"
         "         [  w₂,₁     w₂,₂   ...   w₂,₇ ]   ← row 2\n"
         "W₁  =   [    .         .       .     .   ]\n"
         "         [    .         .       .     .   ]\n"
         "         [w₂₅₆,₁  w₂₅₆,₂   ...  w₂₅₆,₇]   ← row 256",
         size=10, color=NAVY, font="Consolas")

add_text(s, 1.2, y1 + 2.15, 11.0, 0.25,
         "h₁ = ReLU(z₁)        shape: (256)        Purpose: extract 256 features from raw obs",
         size=10, color=GRAY, font="Consolas")
add_arrow(s, 6.65, y1 + 2.52, 6.65, y1 + 2.62)

# LAYER 2
y2 = 4.10
add_box(s, 1.0, y2, 11.3, 1.05, fill=LIGHT, line=ORANGE)
add_text(s, 1.2, y2 + 0.04, 11.0, 0.30, "LAYER 2   —   feature combination",
         size=13, bold=True, color=ORANGE)
add_text(s, 1.2, y2 + 0.34, 11.0, 0.25,
         "z₂ = W₂ · h₁ + b₂          W₂:(256, 256)   b₂:(256)   total: 65,792 weights",
         size=11, color=NAVY, font="Consolas")
add_text(s, 1.2, y2 + 0.62, 11.0, 0.25,
         "h₂ = ReLU(z₂)        shape: (256)        Purpose: combine features → higher abstractions",
         size=10, color=GRAY, font="Consolas")
add_arrow(s, 6.65, y2 + 1.07, 6.65, y2 + 1.18, color=GREEN)

# LAYER 3
y3 = 5.30
add_box(s, 1.0, y3, 11.3, 1.05, fill=LIGHT, line=GREEN)
add_text(s, 1.2, y3 + 0.04, 11.0, 0.30, "LAYER 3   —   output (no activation)",
         size=13, bold=True, color=GREEN)
add_text(s, 1.2, y3 + 0.34, 11.0, 0.25,
         "z₃ = W₃ · h₂ + b₃          W₃:(1, 256)    b₃:(1)     total: 257 weights",
         size=11, color=NAVY, font="Consolas")
add_text(s, 1.2, y3 + 0.62, 11.0, 0.25,
         "action = z₃    (no ReLU — output can be ±)    clipped to [−4, +4] m/s²",
         size=10, color=GRAY, font="Consolas")
add_arrow(s, 6.65, y3 + 1.07, 6.65, y3 + 1.18, color=NAVY)

# OUTPUT
y_out = 6.50
add_box(s, 3.5, y_out, 6.3, 0.50, fill=ORANGE, line=ORANGE)
add_text(s, 3.5, y_out + 0.04, 6.3, 0.25, "OUTPUT",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 3.5, y_out + 0.27, 6.3, 0.22,
         "action = u   (e.g., 1.5 m/s²)",
         size=10, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")


# ---------------- Slide 3 — INPUT ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Input — What the Actor Receives",
           "A 7-dimensional observation vector")

add_box(s, 1.0, 1.3, 11.3, 1.5, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.45, 11.3, 0.4, "Input vector:",
         size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 1.85, 11.3, 0.6,
         "obs  =  [ ψ ,  ψ̇ ,  x ,  ẋ ,  θ ,  θ̇ ,  v_ref ]",
         size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 2.5, 11.3, 0.3,
         "shape: (7,)",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Variable explanations
items = [
    ("ψ",     "robot tilt angle (rad)"),
    ("ψ̇",     "tilt angular velocity (rad/s)"),
    ("x",     "cart position (m)"),
    ("ẋ",     "cart velocity (m/s)"),
    ("θ",     "slosh angle (rad)"),
    ("θ̇",     "slosh angular velocity (rad/s)"),
    ("v_ref", "commanded velocity reference (m/s)"),
]

add_text(s, 0.8, 3.05, 12, 0.4,
         "Each entry — what it means:",
         size=15, bold=True, color=GREEN)

y = 3.5
for i, (var, desc) in enumerate(items):
    add_box(s, 0.8, y + i * 0.5, 11.7, 0.45, fill=LIGHT, line=NAVY, line_w=0.6)
    add_text(s, 1.0, y + i * 0.5 + 0.08, 1.5, 0.35, var,
             size=14, bold=True, color=ORANGE, font="Consolas")
    add_text(s, 2.6, y + i * 0.5 + 0.08, 9.5, 0.35, desc,
             size=12, color=NAVY)


# ---------------- Slide 4 — LAYER 1 ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer 1 — Feature Extraction",
           "Transform 7 raw inputs into 256 learned features")

add_box(s, 0.6, 1.1, 12.1, 6.2, fill=LIGHT, line=ORANGE)

add_text(s, 0.85, 1.18, 11.6, 0.4, "LAYER 1",
         size=18, bold=True, color=ORANGE)

# Math operation
add_text(s, 0.85, 1.55, 11.6, 0.35, "Math operation:",
         size=12, bold=True, color=NAVY)
add_box(s, 1.5, 1.88, 10.3, 0.5, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 1.95, 10.3, 0.4,
         "z₁  =  W₁ · obs  +  b₁",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

# Explicit matrix structure
add_text(s, 0.85, 2.55, 11.6, 0.35, "Weight matrix W₁ (size 256 × 7):",
         size=12, bold=True, color=NAVY)

add_box(s, 0.85, 2.95, 11.6, 1.95, fill=WHITE, line=NAVY, line_w=0.8)
add_text(s, 1.0, 3.05, 11.4, 1.85,
         "         [  w₁,₁     w₁,₂   ...   w₁,₇ ]   ← row 1   (7 weights for output #1)\n"
         "         [  w₂,₁     w₂,₂   ...   w₂,₇ ]   ← row 2\n"
         "W₁  =   [    .         .       .     .   ]\n"
         "         [    .         .       .     .   ]\n"
         "         [w₂₅₆,₁  w₂₅₆,₂   ...  w₂₅₆,₇]   ← row 256",
         size=11, color=NAVY, font="Consolas")

# Shapes/totals — compact
add_text(s, 0.85, 5.05, 11.6, 0.35, "Shapes & totals:",
         size=12, bold=True, color=NAVY)
add_text(s, 1.2, 5.4, 11.0, 0.32,
         "•  W₁: (256, 7)     b₁: (256,)     Total: 256 × 7 + 256 = 2,048 weights",
         size=11, color=NAVY, font="Consolas")

# Activation
add_text(s, 0.85, 5.78, 11.6, 0.35, "Activation:",
         size=12, bold=True, color=NAVY)
add_box(s, 1.5, 6.1, 10.3, 0.5, fill=GREEN, line=GREEN)
add_text(s, 1.5, 6.18, 10.3, 0.4,
         "h₁  =  ReLU(z₁)        (negatives → 0)        shape: (256,)",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

# Purpose
add_text(s, 0.85, 6.75, 11.6, 0.4,
         "Purpose: extract 256 'features' from the raw 7-dim obs.",
         size=12, bold=True, color=GREEN)


# ---------------- Slide 5 — LAYER 2 ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer 2 — Feature Combination",
           "Combine 256 features into 256 higher-level abstractions")

add_box(s, 0.6, 1.2, 12.1, 5.8, fill=LIGHT, line=ORANGE)

add_text(s, 0.85, 1.3, 11.6, 0.4, "LAYER 2",
         size=20, bold=True, color=ORANGE)
add_text(s, 0.85, 1.7, 11.6, 0.4, "Math operation:",
         size=14, bold=True, color=NAVY)

add_box(s, 1.5, 2.1, 10.3, 0.65, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 2.18, 10.3, 0.5,
         "z₂  =  W₂ · h₁  +  b₂",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 2.95, 11.6, 0.4,
         "Shapes:",
         size=14, bold=True, color=NAVY)
add_text(s, 1.2, 3.3, 11.0, 0.4,
         "•  W₂ shape:  (256, 256)       — 256 rows × 256 columns",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 3.65, 11.0, 0.4,
         "•  b₂ shape:  (256,)            — one bias per output neuron",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 4.0, 11.0, 0.4,
         "•  Total weights:  256 × 256 + 256  =  65,792    ← largest layer!",
         size=13, bold=True, color=ORANGE, font="Consolas")

add_text(s, 0.85, 4.55, 11.6, 0.4,
         "Activation:",
         size=14, bold=True, color=NAVY)
add_box(s, 1.5, 4.95, 10.3, 0.65, fill=GREEN, line=GREEN)
add_text(s, 1.5, 5.05, 10.3, 0.45,
         "h₂  =  ReLU(z₂)",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 5.85, 11.6, 0.4,
         "Output shape:  (256,)",
         size=13, color=NAVY)
add_text(s, 0.85, 6.25, 11.6, 0.4,
         "Purpose:",
         size=14, bold=True, color=GREEN)
add_text(s, 0.85, 6.6, 11.6, 0.4,
         "Combine Layer 1 features into higher-level abstractions (the 'deep thinking' layer).",
         size=13, color=NAVY)


# ---------------- Slide 6 — LAYER 3 ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer 3 — Output (No Activation)",
           "Collapse 256 features into a single action value")

add_box(s, 0.6, 1.2, 12.1, 5.8, fill=LIGHT, line=GREEN)

add_text(s, 0.85, 1.3, 11.6, 0.4, "LAYER 3 (output)",
         size=20, bold=True, color=GREEN)
add_text(s, 0.85, 1.7, 11.6, 0.4, "Math operation:",
         size=14, bold=True, color=NAVY)

add_box(s, 1.5, 2.1, 10.3, 0.65, fill=GREEN, line=GREEN)
add_text(s, 1.5, 2.18, 10.3, 0.5,
         "z₃  =  W₃ · h₂  +  b₃",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 2.95, 11.6, 0.4,
         "Shapes:",
         size=14, bold=True, color=NAVY)
add_text(s, 1.2, 3.3, 11.0, 0.4,
         "•  W₃ shape:  (1, 256)         — one row, 256 weights",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 3.65, 11.0, 0.4,
         "•  b₃ shape:  (1,)              — one bias",
         size=13, color=NAVY, font="Consolas")
add_text(s, 1.2, 4.0, 11.0, 0.4,
         "•  Total weights:  256 × 1 + 1  =  257",
         size=13, bold=True, color=GREEN, font="Consolas")

add_text(s, 0.85, 4.55, 11.6, 0.4,
         "Activation:   NONE",
         size=14, bold=True, color=NAVY)
add_box(s, 1.5, 4.95, 10.3, 0.65, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 5.05, 10.3, 0.45,
         "action  =  z₃           (output can be ±)",
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.85, 5.85, 11.6, 0.4,
         "Then clipped to [−4, +4] m/s² (motor saturation limit).",
         size=13, color=NAVY)
add_text(s, 0.85, 6.25, 11.6, 0.4,
         "Purpose:",
         size=14, bold=True, color=GREEN)
add_text(s, 0.85, 6.6, 11.6, 0.4,
         "Collapse 256 high-level features into a single action value u.",
         size=13, color=NAVY)


# ---------------- Slide 7 — Output ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Output — A Single Action",
           "What the actor produces")

add_box(s, 1.0, 1.5, 11.3, 2.5, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 1.7, 11.3, 0.5, "OUTPUT",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 2.4, 11.3, 0.7,
         "action  =  u",
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 3.2, 11.3, 0.5,
         "(a single scalar, e.g., u = 1.5 m/s²)",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)

add_box(s, 0.8, 4.3, 12.0, 2.7, fill=LIGHT, line=GREEN)
add_text(s, 1.0, 4.4, 11.6, 0.4, "What this single number represents:",
         size=15, bold=True, color=GREEN)

bullets = [
    "It's the cart-acceleration command applied to the simulator next.",
    "Bounded to [−4, +4] m/s² (motor's physical capability).",
    "Computed FROM all 7 obs THROUGH all 3 layers in microseconds.",
    "Re-computed every Δt = 0.01 seconds (100 Hz control rate).",
    "During training: also drives gradient updates that improve the network.",
]

for i, b in enumerate(bullets):
    add_text(s, 1.0, 4.85 + i * 0.4, 11.6, 0.4, "•  " + b,
             size=13, color=NAVY)


# ---------------- Slide 8 — Number summary ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Layer-by-Layer Numbers",
           "Shapes and weight counts")

# table header
add_box(s, 0.7, 1.2, 11.9, 0.55, fill=NAVY, line=NAVY)
hdrs = ["", "Layer 1", "Layer 2", "Layer 3"]
xs = [0.7, 4.0, 7.0, 10.0]
ws = [3.3, 3.0, 3.0, 2.6]
for i, hd in enumerate(hdrs):
    add_text(s, xs[i], 1.3, ws[i], 0.4, hd,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Input dim",          "7",          "256",         "256"),
    ("Output dim",         "256",        "256",         "1"),
    ("W shape",            "(256, 7)",   "(256, 256)",  "(1, 256)"),
    ("b shape",            "(256,)",     "(256,)",      "(1,)"),
    ("Weights count",      "1,792",      "65,536",      "256"),
    ("Biases count",       "256",        "256",         "1"),
    ("TOTAL params",       "2,048",      "65,792",      "257"),
    ("Activation",         "ReLU",       "ReLU",        "none"),
    ("Purpose",            "extract",    "combine",     "collapse"),
]

for r, row in enumerate(rows):
    yy = 1.78 + r * 0.45
    fill = LIGHT if r % 2 == 0 else WHITE
    add_box(s, 0.7, yy, 11.9, 0.42, fill=fill, line=NAVY, line_w=0.5)
    is_total = (row[0] == "TOTAL params")
    for c, cell in enumerate(row):
        col = ORANGE if (is_total and c > 0) else NAVY
        bold = is_total or (c == 0)
        font = "Consolas" if c > 0 else "Calibri"
        add_text(s, xs[c], yy + 0.05, ws[c], 0.32, cell,
                 size=12, bold=bold, color=col,
                 align=PP_ALIGN.CENTER, font=font)

add_box(s, 1.0, 6.1, 11.3, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 6.2, 11.3, 0.4,
         "Total actor weights:",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.55, 11.3, 0.5,
         "2,048  +  65,792  +  257  =  68,097  ≈  67K",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")


# ---------------- Slide 9 — Summary ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Summary — The 3-Layer Pipeline",
           "obs (7) → 256 → 256 → action (1)")

steps = [
    ("Layer 1 — Feature extraction",
     "z₁ = W₁·obs + b₁ ,  h₁ = ReLU(z₁)",
     "From 7 raw inputs, build 256 different feature signals.",
     ORANGE),
    ("Layer 2 — Feature combination",
     "z₂ = W₂·h₁ + b₂ ,  h₂ = ReLU(z₂)",
     "Combine the 256 Layer-1 features into 256 higher-level abstractions.",
     ORANGE),
    ("Layer 3 — Output",
     "action = W₃·h₂ + b₃   (no ReLU, clipped to [−4, +4])",
     "Vote across all 256 features → produce a single action value u.",
     GREEN),
]

y = 1.2
for title, eq, body, color in steps:
    add_box(s, 0.6, y, 12.1, 1.4, fill=LIGHT, line=color)
    add_text(s, 0.85, y + 0.05, 11.7, 0.4, title,
             size=16, bold=True, color=color)
    add_text(s, 0.85, y + 0.45, 11.7, 0.4, eq,
             size=14, bold=True, color=NAVY, font="Consolas")
    add_text(s, 0.85, y + 0.85, 11.7, 0.5, body,
             size=12, color=GRAY)
    y += 1.5

add_box(s, 1.0, 6.1, 11.3, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 6.2, 11.3, 0.4,
         "The Big Idea",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.6, 11.3, 0.4,
         "Stack 3 simple operations (linear + ReLU) → emergent decision-making policy.",
         size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/actor_pipeline_explained.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
