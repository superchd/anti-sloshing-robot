#!/usr/bin/env python3
"""
create_lqt_progress_ppt.py
Teaching deck for Section 4.3 of Choi et al. (2024):
"Optimal balancing control" — i.e. LQR / LQT.

Walks students through:
   plant model q-dot = A q + B u
   -> cost function J = int( q^T Q q + u^T R u ) dt
   -> Riccati equation -> gain K  (LQR)
   -> tracking term  -B^+ (A + B K) q_d
   -> final LQT control law  u = K q  -  B^+ (A + B K) q_d   (Eq. 3)

Run:    python3 scripts/create_lqt_progress_ppt.py
Output: results/lqt_explained.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)
PURPLE_BG= rgb(0xF5, 0xEC, 0xFB)
TEAL     = rgb(0x00, 0x7A, 0x86)
TEAL_BG  = rgb(0xE0, 0xF2, 0xF1)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def line_(slide, x1, y1, x2, y2, color=BLACK, weight=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = conn.line._get_or_add_ln()
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', 'dash')
    return conn


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE,
        bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.1, 11.7, 1.1,
         "Optimal Balancing Control",
         size=48, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.3, 11.7, 1.4,
         "From the cost function  J  →  the LQT control law\n"
         "u = K q  −  B⁺ (A + B K) q_d",
         size=22, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 5.5, 11.7, 0.5,
         "Choi et al. 2024 · Section 4.3",
         size=20, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.5, 11.7, 0.5,
         "MECE 6397   ·   Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2 — Roadmap
# =============================================================================
s = blank_slide()
title_bar(s, "Roadmap: 5 steps from cost function to control law",
          "Each box becomes one slide")

steps = [
    (1, "Recap the plant   q̇ = A q + B u",                BLUE),
    (2, "Define the cost  J = ∫ (qᵀQq + uᵀRu) dt",        ORANGE),
    (3, "Solve Riccati  →  feedback gain  K   (LQR)",      RED),
    (4, "Add tracking term  − B⁺ (A + B K) q_d",          PURPLE),
    (5, "Final LQT law  u = K q − B⁺ (A + B K) q_d",       GREEN),
]

y = 1.7
for i, (n, lab, c) in enumerate(steps):
    yi = y + i * 0.95
    box(s, 1.4, yi, 10.5, 0.75, fill=WHITE, border=c, bpt=2.5)
    box(s, 1.4, yi, 0.75, 0.75, fill=c, border=c, bpt=0)
    text_box(s, 1.4, yi + 0.10, 0.75, 0.55, str(n),
             size=22, bold=True, color=WHITE, align=CTR)
    text_box(s, 2.4, yi + 0.18, 9.0, 0.6, lab,
             size=18, bold=True, color=c, align=LFT, font='Consolas')
    if i < len(steps) - 1:
        arrow(s, 6.65, yi + 0.78, 6.65, yi + 0.92,
              color=GRAY, weight=1.5)

text_box(s, 0.6, 6.85, 12.1, 0.5,
         "We follow the textbook order: cost  →  gain  →  tracking.",
         size=13, color=GRAY, align=CTR)


# =============================================================================
# SLIDE 3 — Recap of the plant
# =============================================================================
s = blank_slide()
title_bar(s, "Step 1: Plant recap — what the controller is acting on",
          "Linear, time-invariant, but UNSTABLE (inverted pendulum)")

# Top: the equation
box(s, 0.6, 1.3, 12.1, 1.1, fill=BLUE_BG, border=BLUE, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Nominal disturbance-free plant (eq. 1)",
         size=14, bold=True, color=BLUE)
text_box(s, 0.8, 1.85, 11.8, 0.5,
         "q̇  =  A q  +  B uᵒ",
         size=22, bold=True, color=BLACK, font='Consolas', align=CTR)

# Two columns: state | matrices
box(s, 0.5, 2.55, 6.2, 4.4, fill=ORANGE_BG, border=ORANGE)
text_box(s, 0.7, 2.65, 5.8, 0.5,
         "State vector q   (4 components)",
         size=15, bold=True, color=ORANGE)
text_box(s, 0.7, 3.15, 5.8, 3.3,
         ["q  =  [ ψ ,  ψ̇ ,  x ,  ẋ ]ᵀ",
          "",
          "  ψ   =  body tilt angle  [rad]",
          "  ψ̇   =  tilt rate         [rad/s]",
          "  x   =  position           [m]",
          "  ẋ   =  velocity           [m/s]",
          "",
          "uᵒ =  acceleration command  [m/s²]"],
         size=14, color=BLACK, font='Consolas')

box(s, 6.9, 2.55, 6.0, 4.4, fill=RED_BG, border=RED)
text_box(s, 7.1, 2.65, 5.6, 0.5,
         "Why we need a controller",
         size=15, bold=True, color=RED)
text_box(s, 7.1, 3.15, 5.6, 3.5,
         ["• A has a positive eigenvalue",
          "  (gravity tips the body further when ψ ≠ 0).",
          "",
          "• Open-loop:  robot falls over.",
          "",
          "• Need feedback  uᵒ = K q  to push back",
          "  faster than gravity tips.",
          "",
          "• On top of balancing, we also want",
          "  to TRACK a desired velocity  v_d."],
         size=13, color=BLACK)


# =============================================================================
# SLIDE 4 — Cost function
# =============================================================================
s = blank_slide()
title_bar(s, "Step 2: Define the cost function  J",
          "Tells the optimizer what 'good behavior' means")

# Big formula
box(s, 0.6, 1.3, 12.1, 1.4, fill=ORANGE_BG, border=ORANGE, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Quadratic (LQ) cost over all future time",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.85, 11.8, 0.7,
         "J   =   ∫₀^∞   (  qᵀ Q q   +   uᵀ R u  )  dt",
         size=24, bold=True, color=BLACK, font='Consolas', align=CTR)

# left — Q meaning
box(s, 0.5, 2.85, 6.2, 4.1, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 2.95, 5.8, 0.5,
         "Q   →  penalty on STATE error",
         size=15, bold=True, color=BLUE)
text_box(s, 0.7, 3.45, 5.8, 3.4,
         ["• Big Q on ψ      → keep body upright",
          "• Big Q on ẋ      → match desired speed",
          "• Big Q on x      → don't drift in position",
          "",
          "Form (4×4 diagonal in practice):",
          "",
          "       ⌈ q_ψ   0   0   0   ⌉",
          " Q  =  |  0  q_ψ̇  0   0   |",
          "       |  0   0  q_x  0   |",
          "       ⌊  0   0   0  q_ẋ  ⌋",
          "",
          " Q > 0  (positive-definite, symmetric)"],
         size=12, color=BLACK, font='Consolas')

# right — R meaning
box(s, 6.9, 2.85, 6.0, 4.1, fill=RED_BG, border=RED)
text_box(s, 7.1, 2.95, 5.6, 0.5,
         "R   →  penalty on CONTROL effort",
         size=15, bold=True, color=RED)
text_box(s, 7.1, 3.45, 5.6, 3.4,
         ["• Big R   → save energy, gentle u",
          "• Small R → aggressive u, fast response",
          "",
          "For SISO control input (1 actuator):",
          "",
          "    R   =  scalar  > 0",
          "",
          "Tradeoff knob:",
          "    Q / R   ratio decides",
          "    'how aggressive' the controller is.",
          "",
          " R > 0  (positive-definite)"],
         size=12, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 5 — Riccati & K
# =============================================================================
s = blank_slide()
title_bar(s, "Step 3: Solve the Riccati equation  →  feedback gain  K",
          "MATLAB does this automatically:  K = lqr(A, B, Q, R)")

# Concept box
box(s, 0.6, 1.3, 12.1, 1.0, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Idea: minimize J subject to the dynamics  q̇ = A q + B u.",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 1.80, 11.8, 0.4,
         "Optimal-control theory says the answer is linear feedback "
         "u = K q with a specific K.",
         size=13, color=BLACK)

# Riccati
box(s, 0.6, 2.45, 12.1, 1.5, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 2.55, 11.8, 0.4,
         "Algebraic Riccati Equation (ARE)",
         size=14, bold=True, color=BLUE)
text_box(s, 0.8, 3.00, 11.8, 0.6,
         "Aᵀ P  +  P A  −  P B R⁻¹ Bᵀ P  +  Q   =   0",
         size=20, bold=True, color=BLACK, font='Consolas', align=CTR)
text_box(s, 0.8, 3.65, 11.8, 0.4,
         "Solve for the symmetric positive-definite matrix  P.",
         size=12, color=GRAY, align=CTR)

# K formula
box(s, 0.6, 4.10, 12.1, 1.2, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 0.8, 4.20, 11.8, 0.4,
         "Optimal feedback gain  K   (1 × 4 row vector)",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 4.65, 11.8, 0.6,
         "K   =   − R⁻¹ Bᵀ P",
         size=22, bold=True, color=GREEN, font='Consolas', align=CTR)

# In code
box(s, 0.6, 5.50, 12.1, 1.5, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 5.60, 11.8, 0.4,
         "In MATLAB (your design_lqt.m):",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 6.00, 11.8, 0.95,
         ["    [K_lqr, ~, ~] = lqr(A, B, p.Q, p.R);",
          "    K = -K_lqr;                  % flip sign convention",
          "    % K is now ready to use as a row vector."],
         size=14, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 6 — Why pure LQR is not enough
# =============================================================================
s = blank_slide()
title_bar(s, "Step 4a: But pure LQR only REGULATES — not enough!",
          "Pure  u = K q   drives every state to zero — including ẋ!")

# top message
box(s, 0.6, 1.3, 12.1, 1.0, fill=RED_BG, border=RED)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Problem with pure LQR  u = K q :",
         size=14, bold=True, color=RED)
text_box(s, 0.8, 1.80, 11.8, 0.4,
         "It tries to drive q → 0.   But we WANT  ẋ → v_d  (≠ 0). "
         "So we have to bias the controller toward q_d, not 0.",
         size=13, color=BLACK)

# desired state
box(s, 0.5, 2.45, 6.2, 4.4, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.7, 2.55, 5.8, 0.5,
         "Desired state vector  q_d",
         size=16, bold=True, color=PURPLE)
text_box(s, 0.7, 3.05, 5.8, 0.6,
         "q_d   =   [ x_d ,  v_d ,  0 ,  0 ]ᵀ",
         size=18, bold=True, color=BLACK, font='Consolas')
text_box(s, 0.7, 3.80, 5.8, 3.0,
         ["• x_d = ∫ v_d dt  (target position)",
          "• v_d = velocity command (from F_e shaper!)",
          "• body tilt and rate desired = 0",
          "  (we want the robot upright)",
          "",
          "Tracking error:",
          "    e_t  =  q_d  −  q",
          "",
          "We want  e_t  → 0 ."],
         size=13, color=BLACK)

# right — fix
box(s, 6.9, 2.45, 6.0, 4.4, fill=GREEN_BG, border=GREEN)
text_box(s, 7.1, 2.55, 5.6, 0.5,
         "Fix: add a feedforward term",
         size=16, bold=True, color=GREEN)
text_box(s, 7.1, 3.05, 5.6, 0.6,
         "u   =   K q   +   u_ff(q_d)",
         size=18, bold=True, color=GREEN, font='Consolas')

text_box(s, 7.1, 3.85, 5.6, 0.5,
         "What u_ff has to do:",
         size=14, bold=True, color=GREEN)
text_box(s, 7.1, 4.30, 5.6, 2.5,
         ["• When q  ≡  q_d  (perfect tracking),",
          "  the closed-loop equation must satisfy:",
          "",
          "     q̇_d  =  A q_d  +  B u",
          "",
          "  Solving for u (and using u = K q_d + u_ff)",
          "  gives the feedforward we need:",
          "",
          "     u_ff   =   −  B⁺ (A + B K) q_d"],
         size=12, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 7 — Step 4b: derive the tracking term
# =============================================================================
s = blank_slide()
title_bar(s, "Step 4b: Derive the tracking term  − B⁺ (A + B K) q_d",
          "Algebra you can do on a napkin")

# 1
box(s, 0.6, 1.3, 12.1, 1.0, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Goal:  pick u so that  q  STAYS at  q_d (steady tracking).",
         size=14, bold=True, color=BLUE)
text_box(s, 0.8, 1.80, 11.8, 0.4,
         "If q = q_d at all times, the dynamics must produce  q̇ = q̇_d  "
         "automatically.",
         size=13, color=BLACK)

# 2
box(s, 0.6, 2.45, 12.1, 0.85, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 2.55, 11.8, 0.6,
         "Substitute q = q_d into  q̇ = A q + B u :",
         size=15, color=BLACK)
text_box(s, 0.8, 2.95, 11.8, 0.4,
         "q̇_d   =   A q_d   +   B u",
         size=18, bold=True, color=BLACK, font='Consolas', align=CTR)

# 3
box(s, 0.6, 3.45, 12.1, 0.85, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 3.55, 11.8, 0.4,
         "Use the LQR feedback structure  u = K q_d + u_ff :",
         size=15, color=BLACK)
text_box(s, 0.8, 3.95, 11.8, 0.4,
         "q̇_d   =   (A + B K) q_d   +   B u_ff",
         size=18, bold=True, color=BLACK, font='Consolas', align=CTR)

# 4
box(s, 0.6, 4.45, 12.1, 0.85, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 4.55, 11.8, 0.4,
         "Rearrange for u_ff   (B might not be square — use pseudo-inverse B⁺):",
         size=15, color=BLACK)
text_box(s, 0.8, 4.95, 11.8, 0.4,
         "B u_ff   =   q̇_d   −   (A + B K) q_d",
         size=18, bold=True, color=BLACK, font='Consolas', align=CTR)

# 5
box(s, 0.6, 5.45, 12.1, 1.5, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 0.8, 5.55, 11.8, 0.4,
         "If we set q̇_d = 0  (steady tracking on a slow ramp):",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 6.00, 11.8, 0.6,
         "u_ff   =   −  B⁺ (A + B K) q_d",
         size=22, bold=True, color=GREEN, font='Consolas', align=CTR)
text_box(s, 0.8, 6.65, 11.8, 0.3,
         "= the feedforward needed to hold the desired state.",
         size=12, color=GRAY, align=CTR)


# =============================================================================
# SLIDE 8 — Final law (Eq. 3)
# =============================================================================
s = blank_slide()
title_bar(s, "Step 5: Final LQT control law  (paper Eq. 3)",
          "Combine feedback (regulator) and feedforward (tracker)")

# Big formula at top
box(s, 0.6, 1.3, 12.1, 1.5, fill=GREEN_BG, border=GREEN, bpt=3)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "LQT control law",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 1.85, 11.8, 0.7,
         "u   =   K q     −     B⁺ (A + B K) q_d",
         size=28, bold=True, color=GREEN, font='Consolas', align=CTR)

# decompose into two halves
box(s, 0.5, 3.0, 6.2, 3.4, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 3.10, 5.8, 0.5,
         "Term 1   K q       (feedback)",
         size=16, bold=True, color=BLUE)
text_box(s, 0.7, 3.65, 5.8, 2.6,
         ["• Drives q toward 0",
          "  (and stabilizes the robot — keeps it",
          "   upright against gravity).",
          "",
          "• K is the LQR gain from Step 3.",
          "",
          "• In code:  K * q",
          "",
          "• Without this term, the body tips."],
         size=12, color=BLACK)

box(s, 6.9, 3.0, 6.0, 3.4, fill=PURPLE_BG, border=PURPLE)
text_box(s, 7.1, 3.10, 5.6, 0.5,
         "Term 2   − B⁺ (A + B K) q_d   (feedforward)",
         size=14, bold=True, color=PURPLE)
text_box(s, 7.1, 3.65, 5.6, 2.6,
         ["• Biases the loop so that q TRACKS q_d.",
          "",
          "• q_d  =  [x_d, v_d, 0, 0]ᵀ.",
          "",
          "• v_d comes from the F_e shaper.",
          "",
          "• Designers call this 'precompensation'.",
          "",
          "• Without this term: robot stands still",
          "  instead of moving at v_d."],
         size=12, color=BLACK)

# block diagram
box(s, 0.6, 6.55, 12.1, 0.8, fill=WHITE, border=BLUE, bpt=1.5)
text_box(s, 0.8, 6.55, 1.4, 0.7, "v_d",
         size=14, bold=True, color=BLACK, align=CTR)
text_box(s, 1.6, 6.55, 1.6, 0.7, "→ shape →",
         size=12, color=GRAY, align=CTR)
text_box(s, 3.2, 6.55, 1.5, 0.7, "F_e(s)",
         size=14, bold=True, color=RED, align=CTR)
text_box(s, 4.7, 6.55, 1.6, 0.7, "→ q_d →",
         size=12, color=GRAY, align=CTR)
text_box(s, 6.3, 6.55, 2.2, 0.7, "LQT (Eq. 3)",
         size=14, bold=True, color=GREEN, align=CTR)
text_box(s, 8.5, 6.55, 1.6, 0.7, "→ u →",
         size=12, color=GRAY, align=CTR)
text_box(s, 10.1, 6.55, 1.5, 0.7, "Plant",
         size=14, bold=True, color=BLUE, align=CTR)
text_box(s, 11.6, 6.55, 1.5, 0.7, "→ q",
         size=14, bold=True, color=BLACK, align=CTR)


# =============================================================================
# SLIDE 9 — Summary
# =============================================================================
s = blank_slide()
title_bar(s, "Summary: the whole LQT story on one page",
          "Five steps, one final equation")

# Steps with one-liner summaries
y = 1.3
items = [
    ("1.", "Plant",       "q̇ = A q + B u  (linear, unstable)",         BLUE),
    ("2.", "Cost",        "J = ∫(qᵀ Q q + uᵀ R u) dt",                  ORANGE),
    ("3.", "Riccati → K", "K = −R⁻¹ Bᵀ P    (LQR feedback gain)",       RED),
    ("4.", "Tracking",    "u_ff = −B⁺ (A + B K) q_d",                  PURPLE),
    ("5.", "LQT law",     "u = K q  −  B⁺ (A + B K) q_d   (Eq. 3)",     GREEN),
]
for i, (n, lab, eq, c) in enumerate(items):
    yi = y + i * 0.85
    box(s, 0.6, yi, 12.1, 0.7, fill=WHITE, border=c, bpt=2)
    box(s, 0.6, yi, 0.7, 0.7, fill=c, border=c, bpt=0)
    text_box(s, 0.6, yi + 0.10, 0.7, 0.5, n,
             size=18, bold=True, color=WHITE, align=CTR)
    text_box(s, 1.5, yi + 0.13, 2.6, 0.5, lab,
             size=15, bold=True, color=c)
    text_box(s, 4.2, yi + 0.13, 8.5, 0.5, eq,
             size=15, bold=True, color=BLACK, font='Consolas')

# What it does and what it doesn't
box(s, 0.6, 5.85, 5.95, 1.45, fill=GREEN_BG, border=GREEN)
text_box(s, 0.8, 5.95, 5.6, 0.5,
         "What LQT handles",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 6.40, 5.6, 0.95,
         ["• balancing the robot upright",
          "• tracking the shaped v_d",
          "• optimal Q/R tradeoff"],
         size=12, color=BLACK)

box(s, 6.75, 5.85, 5.95, 1.45, fill=RED_BG, border=RED)
text_box(s, 6.95, 5.95, 5.6, 0.5,
         "What LQT does NOT handle",
         size=14, bold=True, color=RED)
text_box(s, 6.95, 6.40, 5.6, 0.95,
         ["• caster-wheel kicks  → aux compensator (Sec 4.4)",
          "• model error / friction → DOB (Sec 4.5)"],
         size=12, color=BLACK)


# ----- save ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'lqt_explained.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  {len(prs.slides)} slides')
