"""
Teaching PPT: Pipeline of one neural network layer (Actor network in SAC)
For MECE 6397 anti-sloshing robot project (Choi et al. 2024)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=NAVY, line_w=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=ORANGE, weight=2.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = color
    return line


def header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, 0.4, 0.08, 12.5, 0.55, title,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.4, 0.78, 12.5, 0.4, subtitle,
                 size=14, color=GRAY)


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ---------------- Slide 1 — Title ----------------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

add_text(s, 0.6, 2.2, 12, 1.4,
         "Pipeline of a Neural Network Layer",
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.8, 12, 0.7,
         "How data flows through one layer of the SAC actor",
         size=24, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.6, 12, 0.6,
         "MECE 6397 — Anti-Sloshing Robot Project",
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)


# ---------------- Slide 2 — What is a layer? ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "What is a 'Layer'?",
           "One block that transforms its input into a new representation")

add_box(s, 0.8, 1.2, 11.7, 1.5, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.3, 11.3, 0.4,
         "A layer does ONE math operation + an activation function:",
         size=16, bold=True, color=NAVY)
add_text(s, 1.0, 1.7, 11.3, 0.5,
         "input vector  →  multiply by matrix W + bias b  →  apply nonlinearity  →  output vector",
         size=15, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 2.2, 11.3, 0.4,
         "z = W·x + b              h = ReLU(z)",
         size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 0.8, 2.95, 12, 0.4,
         "Your actor network has 3 such layers:",
         size=16, color=GRAY)

# Architecture diagram
add_box(s, 0.5, 3.5, 1.6, 1.0, fill=LIGHT, line=NAVY)
add_text(s, 0.5, 3.65, 1.6, 0.4, "obs", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.05, 1.6, 0.4, "(7)", size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_arrow(s, 2.2, 4.0, 3.0, 4.0)

add_box(s, 3.05, 3.5, 1.7, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 3.05, 3.6, 1.7, 0.4, "Layer 1", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 3.05, 4.0, 1.7, 0.4, "FC + ReLU", size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_arrow(s, 4.85, 4.0, 5.55, 4.0)

add_box(s, 5.6, 3.5, 1.7, 1.0, fill=LIGHT, line=NAVY)
add_text(s, 5.6, 3.65, 1.7, 0.4, "hidden", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 5.6, 4.05, 1.7, 0.4, "(256)", size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_arrow(s, 7.4, 4.0, 8.1, 4.0)

add_box(s, 8.15, 3.5, 1.7, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 8.15, 3.6, 1.7, 0.4, "Layer 2", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 8.15, 4.0, 1.7, 0.4, "FC + ReLU", size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_arrow(s, 9.95, 4.0, 10.65, 4.0)

add_box(s, 10.7, 3.5, 1.7, 1.0, fill=GREEN, line=GREEN)
add_text(s, 10.7, 3.6, 1.7, 0.4, "Layer 3", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 10.7, 4.0, 1.7, 0.4, "FC only", size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_arrow(s, 12.45, 4.0, 12.95, 4.0)
add_text(s, 11.2, 4.6, 2.0, 0.4, "→ action (1)", size=14, bold=True, color=NAVY)

add_box(s, 1.5, 5.4, 10.3, 1.7, fill=LIGHT, line=GREEN)
add_text(s, 1.7, 5.5, 10, 0.4, "Each layer:", size=16, bold=True, color=GREEN)
add_text(s, 1.7, 5.95, 10, 0.4,
         "• takes a vector in",
         size=14, color=NAVY)
add_text(s, 1.7, 6.25, 10, 0.4,
         "• multiplies by a matrix W + adds a bias b",
         size=14, color=NAVY)
add_text(s, 1.7, 6.55, 10, 0.4,
         "• applies an activation function (ReLU here)",
         size=14, color=NAVY)


# ---------------- Slide 3 — Input ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Walking Through Layer 1 — The Input",
           "What does Layer 1 receive?")

add_text(s, 0.8, 1.2, 12, 0.5,
         "The actor sees a 7-dimensional observation vector at every time step:",
         size=18, color=GRAY)

add_box(s, 0.8, 2.0, 11.7, 2.2, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 2.15, 11.3, 0.4,
         "obs  =  [ ψ ,  ψ̇ ,  x ,  ẋ ,  θ ,  θ̇ ,  v_ref ]",
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 2.8, 11.3, 0.4,
         "Example values:",
         size=15, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 3.2, 11.3, 0.5,
         "obs  =  [ 0.05,  0.1,  0.2,  0.3,  0.04,  0.2,  0.8 ]",
         size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 3.75, 11.3, 0.4,
         "↑  This is a vector of length 7 (7 numbers)",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Meaning of each entry
add_box(s, 0.8, 4.5, 11.7, 2.5, fill=LIGHT, line=ORANGE)
add_text(s, 1.0, 4.6, 11.3, 0.4,
         "What each number means:",
         size=16, bold=True, color=ORANGE)

meanings = [
    ("ψ = 0.05", "robot tilt angle (rad)"),
    ("ψ̇ = 0.1", "tilt angular velocity"),
    ("x = 0.2", "cart position (m)"),
    ("ẋ = 0.3", "cart velocity (m/s)"),
    ("θ = 0.04", "slosh angle (rad)"),
    ("θ̇ = 0.2", "slosh angular velocity"),
    ("v_ref = 0.8", "commanded velocity"),
]

y = 5.05
for i, (var, desc) in enumerate(meanings):
    col = i % 2
    row = i // 2
    xx = 1.2 + col * 5.7
    yy = y + row * 0.4
    add_text(s, xx, yy, 2.0, 0.4, var, size=13, bold=True, color=NAVY, font="Consolas")
    add_text(s, xx + 1.7, yy, 4.0, 0.4, desc, size=13, color=GRAY)


# ---------------- Slide 4 — Step A: Linear transformation ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Step A — Linear Transformation",
           "Multiply input by weight matrix W₁ + add bias b₁")

add_text(s, 0.8, 1.1, 12, 0.4,
         "Layer 1 has a weight matrix W₁ of size 256 × 7  (256 outputs, 7 inputs):",
         size=15, color=GRAY)

add_box(s, 0.8, 1.6, 11.7, 2.2, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.7, 11.3, 1.6,
         "         [  w₁,₁    w₁,₂   ...   w₁,₇  ]   ← row 1  (7 weights for output #1)\n"
         "         [  w₂,₁    w₂,₂   ...   w₂,₇  ]   ← row 2\n"
         "W₁  =   [   .         .       .     .    ]\n"
         "         [   .         .       .     .    ]\n"
         "         [w₂₅₆,₁  w₂₅₆,₂  ...  w₂₅₆,₇]   ← row 256",
         size=14, color=NAVY, font="Consolas")

add_text(s, 0.8, 4.0, 12, 0.4,
         "The layer computes:",
         size=15, color=GRAY)
add_box(s, 1.5, 4.4, 10.3, 0.8, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 4.5, 10.3, 0.6,
         "z  =  W₁ · obs  +  b₁",
         size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.8, 5.4, 12, 0.4,
         "Each output neuron i computes a weighted sum of all 7 inputs:",
         size=15, color=GRAY)
add_box(s, 1.5, 5.85, 10.3, 1.2, fill=LIGHT, line=GREEN)
add_text(s, 1.5, 5.95, 10.3, 0.5,
         "zᵢ  =  w_{i,1}·obs[1]  +  w_{i,2}·obs[2]  +  ...  +  w_{i,7}·obs[7]  +  bᵢ",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.5, 6.45, 10.3, 0.5,
         "= weighted sum of all 7 input features  +  bias",
         size=14, color=GREEN, align=PP_ALIGN.CENTER)


# ---------------- Slide 5 — Concrete math example ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Concrete Example — One Output Neuron",
           "What does neuron #1 actually compute?")

add_text(s, 0.8, 1.1, 12, 0.4,
         "Suppose neuron 1 has these weights and bias:",
         size=15, color=GRAY)
add_box(s, 0.8, 1.5, 11.7, 0.9, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.6, 11.3, 0.4,
         "w_{1,*}  =  [ 0.2,  −0.5,  0.1,  0.3,  1.2,  −0.8,  0.4 ]",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 2.0, 11.3, 0.4,
         "b_1  =  0.05",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")

add_box(s, 0.8, 2.65, 11.7, 3.5, fill=LIGHT, line=ORANGE)
add_text(s, 1.0, 2.75, 11.3, 0.4,
         "Compute the weighted sum:",
         size=15, bold=True, color=ORANGE)
add_text(s, 1.0, 3.15, 11.3, 2.9,
         "z₁  =     0.2  × 0.05      (ψ contribution)\n"
         "      +  (−0.5) × 0.1     (ψ̇ contribution)\n"
         "      +   0.1  × 0.2       (x contribution)\n"
         "      +   0.3  × 0.3       (ẋ contribution)\n"
         "      +   1.2  × 0.04      (θ contribution)   ← strong slosh weight\n"
         "      +  (−0.8) × 0.2     (θ̇ contribution)\n"
         "      +   0.4  × 0.8       (v_ref contribution)\n"
         "      +   0.05               (bias)",
         size=12, color=NAVY, font="Consolas")

add_box(s, 1.5, 6.3, 10.3, 0.9, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 6.4, 10.3, 0.7,
         "z₁  =  0.01 − 0.05 + 0.02 + 0.09 + 0.048 − 0.16 + 0.32 + 0.05  =  0.328",
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")


# ---------------- Slide 6 — All 256 neurons ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "256 Different Neurons, 256 Different 'Features'",
           "Each neuron looks at the input from a different angle")

add_text(s, 0.8, 1.1, 12, 0.4,
         "All 256 output neurons do the SAME calculation, but with DIFFERENT weights:",
         size=15, color=GRAY)

add_box(s, 0.8, 1.7, 11.7, 3.0, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.85, 11.3, 0.4,
         "z₁    =  0.328     (\"braking signal\")",
         size=15, color=NAVY, font="Consolas")
add_text(s, 1.0, 2.25, 11.3, 0.4,
         "z₂    =  −0.7      (\"forward push signal\")",
         size=15, color=NAVY, font="Consolas")
add_text(s, 1.0, 2.65, 11.3, 0.4,
         "z₃    =  1.4       (\"slosh imminent signal\")",
         size=15, color=NAVY, font="Consolas")
add_text(s, 1.0, 3.05, 11.3, 0.4,
         "z₄    =  0.0       (\"hold position\")",
         size=15, color=NAVY, font="Consolas")
add_text(s, 1.0, 3.45, 11.3, 0.4,
         "...",
         size=15, color=GRAY, font="Consolas")
add_text(s, 1.0, 3.85, 11.3, 0.4,
         "z₂₅₆ =  0.05      (\"don't move signal\")",
         size=15, color=NAVY, font="Consolas")
add_text(s, 1.0, 4.3, 11.3, 0.4,
         "Output of matrix multiply: vector of 256 numbers",
         size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_box(s, 0.8, 4.95, 11.7, 2.1, fill=LIGHT, line=GREEN)
add_text(s, 1.0, 5.05, 11.3, 0.4,
         "Each neuron extracts a different LEARNED FEATURE:",
         size=15, bold=True, color=GREEN)
add_text(s, 1.0, 5.5, 11.3, 0.4,
         "•  Neuron 1 might detect:  \"tilt + slosh in same direction → DANGER\"",
         size=13, color=NAVY)
add_text(s, 1.0, 5.85, 11.3, 0.4,
         "•  Neuron 2 might detect:  \"fast cart velocity → need to brake soon\"",
         size=13, color=NAVY)
add_text(s, 1.0, 6.2, 11.3, 0.4,
         "•  Neuron 3 might detect:  \"θ̇ > 0 and θ > 0 → slosh is accelerating\"",
         size=13, color=NAVY)
add_text(s, 1.0, 6.55, 11.3, 0.4,
         "These features are learned AUTOMATICALLY during training — you don't design them.",
         size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ---------------- Slide 7 — Step B: ReLU activation ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Step B — Activation Function (ReLU)",
           "Pass each output through max(0, x)")

add_box(s, 0.6, 1.2, 6.0, 2.6, fill=LIGHT, line=NAVY)
add_text(s, 0.8, 1.3, 5.6, 0.4,
         "ReLU = Rectified Linear Unit:",
         size=15, bold=True, color=NAVY)
add_text(s, 0.8, 1.8, 5.6, 0.5,
         "ReLU(x)  =  max(0, x)",
         size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 2.4, 5.6, 0.4,
         "= x  if x > 0",
         size=16, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 2.8, 5.6, 0.4,
         "= 0  if x ≤ 0",
         size=16, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 3.3, 5.6, 0.4,
         "Negatives become 0; positives pass through.",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ReLU graph
add_box(s, 6.9, 1.2, 6.0, 2.6, fill=WHITE, line=NAVY)
add_text(s, 7.1, 1.3, 5.6, 0.4, "Graph of ReLU:",
         size=14, bold=True, color=ORANGE)
add_text(s, 7.1, 1.8, 5.6, 1.9,
         "      |\n"
         "  2  ┤        ╱\n"
         "  1  ┤      ╱\n"
         "  0  ┤────╱─────────→\n"
         "      └────┴─────► input\n"
         "       −1   0   1   2",
         size=12, color=NAVY, font="Consolas")

add_text(s, 0.8, 4.0, 12, 0.4,
         "Apply to each of the 256 output values:",
         size=15, color=GRAY)
add_box(s, 0.8, 4.4, 11.7, 1.4, fill=LIGHT, line=ORANGE)
add_text(s, 1.0, 4.5, 11.3, 0.4,
         "Pre-ReLU:    [ 0.328,  −0.7,  1.4,  0.0,  ...,  0.05 ]",
         size=14, color=NAVY, font="Consolas")
add_text(s, 1.0, 4.95, 11.3, 0.4,
         "After ReLU:  [ 0.328,   0,   1.4,  0.0,  ...,  0.05 ]",
         size=14, bold=True, color=ORANGE, font="Consolas")
add_text(s, 1.0, 5.4, 11.3, 0.4,
         "↑ negative values became 0",
         size=12, color=GREEN)

# Why ReLU?
add_box(s, 0.8, 6.0, 11.7, 1.1, fill=LIGHT, line=GREEN)
add_text(s, 1.0, 6.05, 11.3, 0.4,
         "Why ReLU?",
         size=14, bold=True, color=GREEN)
add_text(s, 1.0, 6.4, 11.3, 0.4,
         "•  Adds NONLINEARITY (without it, stacked layers collapse to one linear map)",
         size=11, color=NAVY)
add_text(s, 1.0, 6.7, 11.3, 0.4,
         "•  Cheap — just one comparison.   •  Gradients flow cleanly.   •  Encourages sparse features.",
         size=11, color=NAVY)


# ---------------- Slide 8 — Layer 1 output ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Step C — Output of Layer 1",
           "256-dimensional vector becomes the input to Layer 2")

add_text(s, 0.8, 1.2, 12, 0.4,
         "Putting Steps A + B together:",
         size=16, color=GRAY)

add_box(s, 1.5, 1.8, 10.3, 1.4, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 1.95, 10.3, 0.5,
         "h₁  =  ReLU( W₁ · obs  +  b₁ )",
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.5, 2.6, 10.3, 0.4,
         "= 256-dimensional vector  (most positive, some zero)",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_box(s, 0.8, 3.5, 11.7, 1.3, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 3.6, 11.3, 0.4,
         "Example output:",
         size=15, bold=True, color=NAVY)
add_text(s, 1.0, 4.0, 11.3, 0.6,
         "h₁  =  [ 0.328,  0,  1.4,  0,  0.7,  0,  ...,  0.05,  0,  0.2 ]",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 4.45, 11.3, 0.4,
         "256 numbers — passed to Layer 2",
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_box(s, 0.8, 5.0, 11.7, 2.0, fill=LIGHT, line=GREEN)
add_text(s, 1.0, 5.1, 11.3, 0.4,
         "What just happened?",
         size=15, bold=True, color=GREEN)
add_text(s, 1.0, 5.5, 11.3, 0.4,
         "•  Started with 7 raw measurements (ψ, ψ̇, x, ẋ, θ, θ̇, v_ref)",
         size=13, color=NAVY)
add_text(s, 1.0, 5.85, 11.3, 0.4,
         "•  Created 256 LEARNED FEATURES (each a different combination of the 7 inputs)",
         size=13, color=NAVY)
add_text(s, 1.0, 6.2, 11.3, 0.4,
         "•  Filtered out negative features (ReLU) — keeping only \"active\" signals",
         size=13, color=NAVY)
add_text(s, 1.0, 6.55, 11.3, 0.4,
         "•  Result: a richer 256-dim representation that Layer 2 can build upon",
         size=13, bold=True, color=GREEN)


# ---------------- Slide 9 — Full actor pipeline ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Full Actor Pipeline — All 3 Layers",
           "obs (7) → 256 → 256 → action (1)")

# Compact diagram
y0 = 1.0

# Input
add_box(s, 4.5, y0, 4.3, 0.7, fill=LIGHT, line=NAVY)
add_text(s, 4.5, y0 + 0.15, 4.3, 0.4, "INPUT  obs  (shape: 7)",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_arrow(s, 6.65, y0 + 0.75, 6.65, y0 + 1.05, color=ORANGE, weight=3)

# Layer 1
y1 = y0 + 1.1
add_box(s, 1.0, y1, 11.3, 1.3, fill=LIGHT, line=ORANGE)
add_text(s, 1.2, y1 + 0.05, 11.0, 0.4, "LAYER 1 — Feature extraction",
         size=14, bold=True, color=ORANGE)
add_text(s, 1.2, y1 + 0.45, 11.0, 0.35,
         "z₁ = W₁·obs + b₁         W₁:(256×7)  b₁:(256)        Total params: 2,048",
         size=12, color=NAVY, font="Consolas")
add_text(s, 1.2, y1 + 0.78, 11.0, 0.35,
         "h₁ = ReLU(z₁)             shape: (256)",
         size=12, color=NAVY, font="Consolas")
add_arrow(s, 6.65, y1 + 1.35, 6.65, y1 + 1.65, color=ORANGE, weight=3)

# Layer 2
y2 = y1 + 1.7
add_box(s, 1.0, y2, 11.3, 1.3, fill=LIGHT, line=ORANGE)
add_text(s, 1.2, y2 + 0.05, 11.0, 0.4, "LAYER 2 — Feature combination",
         size=14, bold=True, color=ORANGE)
add_text(s, 1.2, y2 + 0.45, 11.0, 0.35,
         "z₂ = W₂·h₁ + b₂          W₂:(256×256)  b₂:(256)     Total params: 65,792",
         size=12, color=NAVY, font="Consolas")
add_text(s, 1.2, y2 + 0.78, 11.0, 0.35,
         "h₂ = ReLU(z₂)             shape: (256)",
         size=12, color=NAVY, font="Consolas")
add_arrow(s, 6.65, y2 + 1.35, 6.65, y2 + 1.65, color=GREEN, weight=3)

# Layer 3
y3 = y2 + 1.7
add_box(s, 1.0, y3, 11.3, 1.0, fill=LIGHT, line=GREEN)
add_text(s, 1.2, y3 + 0.05, 11.0, 0.4, "LAYER 3 — Output (no activation)",
         size=14, bold=True, color=GREEN)
add_text(s, 1.2, y3 + 0.4, 11.0, 0.35,
         "action = W₃·h₂ + b₃     W₃:(1×256)    b₃:(1)         Total params: 257",
         size=12, color=NAVY, font="Consolas")

# Output
add_arrow(s, 6.65, y3 + 1.05, 6.65, y3 + 1.3, color=NAVY, weight=3)


# ---------------- Slide 10 — Layer interpretation ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "What Each Layer Actually Does",
           "Physical interpretation of layer purpose")

layers = [
    ("Layer 1 — Feature Extraction",
     "Takes raw obs (7 numbers) and creates 256 different combinations.",
     "Each neuron learns to detect a specific pattern, e.g.:\n"
     "•  \"tilt + slosh in same direction → DANGER\"\n"
     "•  \"fast cart velocity → need to brake soon\"\n"
     "•  \"θ̇ > 0 AND θ > 0 → slosh accelerating\"",
     ORANGE),
    ("Layer 2 — Feature Combination",
     "Takes Layer 1's 256 features and combines them into 256 higher-level features.",
     "Each neuron learns abstract combinations, e.g.:\n"
     "•  \"DANGER + braking signal → strong negative action\"\n"
     "•  \"stable + tracking error → moderate forward push\"\n"
     "Has the most weights (65,792) — the \"deep thinking\" layer.",
     ORANGE),
    ("Layer 3 — Output Collapsing",
     "Collapses 256 high-level features into ONE number — the action.",
     "action  =  w₁·h₂[1] + w₂·h₂[2] + ... + w₂₅₆·h₂[256] + bias\n"
     "         =  weighted vote of all 256 high-level features\n"
     "No activation — output can be ±, then clipped to [−4, 4] m/s².",
     GREEN),
]

y = 1.1
for title, summary, body, color in layers:
    add_box(s, 0.6, y, 12.1, 1.95, fill=LIGHT, line=color)
    add_text(s, 0.85, y + 0.05, 11.7, 0.4, title,
             size=15, bold=True, color=color)
    add_text(s, 0.85, y + 0.45, 11.7, 0.4, summary,
             size=12, color=GRAY)
    add_text(s, 0.85, y + 0.85, 11.7, 1.1, body,
             size=11, color=NAVY)
    y += 2.05


# ---------------- Slide 11 — Numbers table ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Numbers to Remember",
           "Weight counts for each layer")

# Table header
add_box(s, 1.0, 1.3, 11.3, 0.55, fill=NAVY, line=NAVY)
header_cells = ["Quantity", "Layer 1", "Layer 2", "Layer 3"]
col_x = [1.0, 4.0, 6.7, 9.5]
col_w = [3.0, 2.7, 2.8, 2.8]
for i, hd in enumerate(header_cells):
    add_text(s, col_x[i], 1.4, col_w[i], 0.4, hd,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Input dim", "7", "256", "256"),
    ("Output dim", "256", "256", "1"),
    ("W shape", "(256, 7)", "(256, 256)", "(1, 256)"),
    ("Weights in W", "1,792", "65,536", "256"),
    ("Biases", "256", "256", "1"),
    ("Total params", "2,048", "65,792", "257"),
    ("Activation", "ReLU", "ReLU", "none"),
]

for r, row in enumerate(rows):
    yy = 1.85 + r * 0.5
    fill = LIGHT if r % 2 == 0 else WHITE
    add_box(s, 1.0, yy, 11.3, 0.5, fill=fill, line=NAVY, line_w=0.5)
    for c, cell in enumerate(row):
        is_total = (rows[r][0] == "Total params")
        col = ORANGE if (is_total and c > 0) else NAVY
        bold = is_total
        add_text(s, col_x[c], yy + 0.08, col_w[c], 0.35, cell,
                 size=13, bold=bold, color=col, align=PP_ALIGN.CENTER,
                 font="Consolas")

add_box(s, 1.0, 5.6, 11.3, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 5.7, 11.3, 0.4,
         "Total weights in actor:",
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.05, 11.3, 0.5,
         "2,048  +  65,792  +  257  =  68,097  ≈  67K weights",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_text(s, 1.0, 6.85, 11.3, 0.4,
         "Most of the weights live in Layer 2 (the 256×256 middle layer).",
         size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------- Slide 12 — Forward pass ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Forward Pass — One obs Becomes One Action",
           "What happens at every timestep during runtime")

add_box(s, 0.6, 1.3, 12.1, 1.5, fill=LIGHT, line=NAVY)
add_text(s, 0.8, 1.4, 11.7, 0.4, "At every time t:",
         size=15, bold=True, color=NAVY)
add_text(s, 0.8, 1.85, 11.7, 0.5,
         "obs(t)  →  Layer 1  →  h₁ (256)  →  Layer 2  →  h₂ (256)  →  Layer 3  →  action(t)",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 0.8, 2.35, 11.7, 0.4,
         "↑ feature extraction       ↑ feature combination       ↑ vote → action",
         size=12, color=ORANGE, align=PP_ALIGN.CENTER)

facts = [
    ("Computation time", "Microseconds on a CPU"),
    ("Multiplications", "~150,000 per forward pass (mostly Layer 2)"),
    ("Output", "ONE action value, clipped to [-4, 4] m/s²"),
    ("How often?", "Every 0.01 seconds (100 Hz control rate)"),
    ("During runtime", "ONLY forward pass — no learning anymore"),
    ("During training", "Forward pass + backward pass (next slide)"),
]

y = 3.1
for i, (k, v) in enumerate(facts):
    yy = y + i * 0.55
    add_box(s, 0.8, yy, 11.7, 0.5, fill=LIGHT, line=GREEN, line_w=0.8)
    add_text(s, 1.0, yy + 0.07, 4.0, 0.4, k, size=13, bold=True, color=GREEN)
    add_text(s, 5.2, yy + 0.07, 7.0, 0.4, v, size=13, color=NAVY)


# ---------------- Slide 13 — Backward pass ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Backward Pass — How the Network Learns",
           "Gradients flow BACKWARD through the layers")

add_text(s, 0.8, 1.1, 12, 0.4,
         "After receiving a reward, gradients flow from output back to input:",
         size=15, color=GRAY)

add_box(s, 0.6, 1.6, 12.1, 4.4, fill=LIGHT, line=NAVY)
add_text(s, 0.85, 1.7, 11.6, 4.3,
         "reward\n"
         "   │\n"
         "   ▼  loss = -Q(s, action)\n"
         "   │\n"
         "   ▼  ∂loss/∂action          ← Layer 3 output gradient\n"
         "   │\n"
         "   ▼  ∂loss/∂W₃              ← UPDATE Layer 3 weights\n"
         "   ▼  ∂loss/∂h₂              ← gradient passed back to Layer 2\n"
         "   │\n"
         "   ▼  ∂loss/∂W₂              ← UPDATE Layer 2 weights\n"
         "   ▼  ∂loss/∂h₁              ← gradient passed back to Layer 1\n"
         "   │\n"
         "   ▼  ∂loss/∂W₁              ← UPDATE Layer 1 weights",
         size=13, color=NAVY, font="Consolas")

add_box(s, 0.8, 6.2, 11.7, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 0.8, 6.3, 11.7, 0.4,
         "Update rule (gradient descent):",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 6.7, 11.7, 0.5,
         "W_new  =  W_old  −  learning_rate  ×  ∂loss/∂W",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")


# ---------------- Slide 14 — Visual summary ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Visual Summary — One Layer's Pipeline",
           "Input → Linear → ReLU → Output")

add_box(s, 3.0, 1.1, 7.3, 0.9, fill=LIGHT, line=NAVY)
add_text(s, 3.0, 1.2, 7.3, 0.4, "INPUT VECTOR  (7 numbers)",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 3.0, 1.55, 7.3, 0.4,
         "[ ψ,  ψ̇,  x,  ẋ,  θ,  θ̇,  v_ref ]",
         size=12, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_arrow(s, 6.65, 2.05, 6.65, 2.4, color=ORANGE, weight=3)

add_box(s, 3.5, 2.45, 6.3, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 3.5, 2.55, 6.3, 0.4, "Multiply by W₁  +  Add bias b₁",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 3.5, 2.95, 6.3, 0.4, "(256 × 7 matrix)  +  (256 vector)",
         size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_arrow(s, 6.65, 3.5, 6.65, 3.85, color=ORANGE, weight=3)

add_box(s, 3.0, 3.9, 7.3, 0.9, fill=LIGHT, line=NAVY)
add_text(s, 3.0, 4.0, 7.3, 0.4, "PRE-ACTIVATION  (256 numbers)",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 3.0, 4.35, 7.3, 0.4,
         "[ z₁,  z₂,  ...,  z₂₅₆ ]",
         size=12, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_arrow(s, 6.65, 4.85, 6.65, 5.2, color=GREEN, weight=3)

add_box(s, 3.5, 5.25, 6.3, 0.7, fill=GREEN, line=GREEN)
add_text(s, 3.5, 5.4, 6.3, 0.4, "Apply ReLU  =  max(0, x)",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_arrow(s, 6.65, 5.97, 6.65, 6.32, color=GREEN, weight=3)

add_box(s, 3.0, 6.4, 7.3, 0.9, fill=LIGHT, line=GREEN)
add_text(s, 3.0, 6.5, 7.3, 0.4, "POST-ACTIVATION  (256 numbers, sparse)",
         size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, 3.0, 6.85, 7.3, 0.4,
         "[ h₁,  h₂,  ...,  h₂₅₆ ]   →  input to next layer",
         size=12, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")


# ---------------- Slide 15 — Summary ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Summary — One Layer in 4 Sentences",
           "The takeaway")

points = [
    ("1.  Linear transformation",
     "Multiply input by weight matrix W, add bias b:   z = W·x + b",
     NAVY),
    ("2.  Nonlinearity",
     "Apply an activation like ReLU:   h = max(0, z)",
     ORANGE),
    ("3.  Output",
     "The resulting vector becomes the input to the next layer.",
     GREEN),
    ("4.  Learning",
     "Weights W and biases b are what get UPDATED during training — they store the\n"
     "agent's 'knowledge' about how to map states to actions.",
     PURPLE),
]

y = 1.2
for title, body, color in points:
    add_box(s, 0.6, y, 12.1, 1.15, fill=LIGHT, line=color)
    add_text(s, 0.85, y + 0.1, 11.7, 0.4, title,
             size=16, bold=True, color=color)
    add_text(s, 0.85, y + 0.5, 11.7, 0.6, body,
             size=13, color=NAVY)
    y += 1.25

add_box(s, 1.0, 6.3, 11.3, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 6.4, 11.3, 0.5,
         "The Big Idea",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.85, 11.3, 0.4,
         "One layer = (matrix multiply + bias) + (activation function).  Stack them → deep network.",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/neural_layer_explained.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
