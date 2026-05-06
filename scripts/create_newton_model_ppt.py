#!/usr/bin/env python3
"""
create_newton_model_ppt.py
Slide deck explaining the Newton derivation of the two-body
sloshing-robot model:
   - sloshing pendulum:  theta-double-dot equation
   - robot body:         psi state-space  q-dot = A q + B u
and why the two equations decouple in practice.

Run:    python3 scripts/create_newton_model_ppt.py
Output: results/newton_model_explained.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


# ----- palette ---------------------------------------------------------------
WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)
PURPLE_BG= rgb(0xF5, 0xEC, 0xFB)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


# ----- presentation ----------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def line(slide, x1, y1, x2, y2, color=BLACK, weight=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = conn.line._get_or_add_ln()
        ln.set('cap', 'flat')
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', 'dash')
    return conn


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def circle(slide, cx, cy, r, fill=ORANGE, border=BLACK, bpt=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - r), Inches(cy - r),
                                 Inches(2 * r), Inches(2 * r))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE,
        bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.3, 11.7, 1.4,
         "Newton's Laws → Two-Body Robot Model",
         size=48, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.7, 11.7, 0.7,
         "How the sloshing  θ̈  and robot-body  q̇ = Aq + Bu  equations "
         "actually come from physics",
         size=22, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 6.0, 11.7, 0.5,
         "Anti-Sloshing Serving Robot  ·  MECE 6397",
         size=18, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.5, 11.7, 0.5,
         "Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2 — Big picture: two coupled bodies
# =============================================================================
s = blank_slide()
title_bar(s, "Two rigid bodies, one input  u",
          "u (acceleration command) drives both ψ and θ — but they are written separately")

# Left — schematic
box(s, 0.5, 1.3, 5.8, 5.6, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 1.4, 5.4, 0.4, "The physical picture",
         size=18, bold=True, color=BLUE)

# robot body (tilted box)
import math
psi_deg = 12
def rotxy(x, y, ang_deg, cx, cy):
    a = math.radians(ang_deg)
    dx, dy = x - cx, y - cy
    return cx + dx*math.cos(a) - dy*math.sin(a), \
           cy + dx*math.sin(a) + dy*math.cos(a)

# wheels (axle) center
ax_x, ax_y = 3.4, 6.3
# robot body corners (un-tilted) — width 1.6, height 3.2, sitting on axle
bw, bh = 1.6, 3.2
corners = [(ax_x - bw/2, ax_y - bh),
           (ax_x + bw/2, ax_y - bh),
           (ax_x + bw/2, ax_y),
           (ax_x - bw/2, ax_y)]
rot = [rotxy(x, y, psi_deg, ax_x, ax_y) for x, y in corners]
# draw the four sides
for i in range(4):
    x1, y1 = rot[i]
    x2, y2 = rot[(i + 1) % 4]
    line(s, x1, y1, x2, y2, color=BLACK, weight=2.5)

# wheel
circle(s, ax_x, ax_y, 0.18, fill=GRAY, border=BLACK)
# ground
line(s, 0.7, 6.55, 6.1, 6.55, color=BLACK, weight=1.5)
# psi label and tilted axis
line(s, ax_x, ax_y, ax_x, ax_y - 1.4, color=GRAY, weight=1.2, dash=True)
text_box(s, ax_x - 0.95, ax_y - 1.6, 0.6, 0.4, "ψ", size=22, bold=True,
         color=BLUE)

# pendulum pivot near top of body
pv_x, pv_y = rotxy(ax_x, ax_y - bh + 0.3, psi_deg, ax_x, ax_y)
# pendulum bob — tilted by θ relative to body's local "down"
theta_deg = 25
# pendulum unit vector in body frame is (sin θ, cos θ) downward,
# then rotate by ψ
def body_down(dx_local, dy_local):
    a = math.radians(psi_deg)
    return dx_local*math.cos(a) - dy_local*math.sin(a), \
           dx_local*math.sin(a) + dy_local*math.cos(a)
pl = 1.6
pdx, pdy = body_down(math.sin(math.radians(theta_deg)) * pl,
                     math.cos(math.radians(theta_deg)) * pl)
bob_x, bob_y = pv_x + pdx, pv_y + pdy
# rod
line(s, pv_x, pv_y, bob_x, bob_y, color=ORANGE, weight=3.0)
# bob
circle(s, bob_x, bob_y, 0.18, fill=ORANGE, border=BLACK)
# theta dashed body-frame axis
ddx, ddy = body_down(0, pl)
line(s, pv_x, pv_y, pv_x + ddx, pv_y + ddy, color=GRAY, weight=1.0, dash=True)
text_box(s, bob_x + 0.05, bob_y - 0.35, 0.7, 0.4, "θ", size=20, bold=True,
         color=ORANGE)

# input arrow at base
arrow(s, ax_x + 0.5, ax_y + 0.05, ax_x + 1.5, ax_y + 0.05, color=RED, weight=3)
text_box(s, ax_x + 0.6, ax_y + 0.1, 0.9, 0.4, "u (=  ẍ)",
         size=14, bold=True, color=RED)

# legend
text_box(s, 0.7, 5.0, 5.4, 0.4, "ψ  =  body tilt  (inverted-pendulum)",
         size=12, color=BLUE, bold=True)
text_box(s, 0.7, 5.3, 5.4, 0.4, "θ  =  liquid-sloshing pendulum angle",
         size=12, color=ORANGE, bold=True)
text_box(s, 0.7, 5.6, 5.4, 0.4, "u  =  wheel-base linear acceleration",
         size=12, color=RED, bold=True)

# Right — coupling structure
box(s, 6.6, 1.3, 6.4, 5.6, fill=GREEN_BG, border=GREEN)
text_box(s, 6.8, 1.4, 6.0, 0.5, "Coupling structure",
         size=18, bold=True, color=GREEN)

text_box(s, 6.8, 2.0, 6.0, 0.5,
         "The control input  u  appears in BOTH equations:",
         size=14, color=BLACK)

text_box(s, 6.9, 2.55, 6.0, 0.5,
         "ψ:    q̇  =  A q  +  B u            ← linear, robot body",
         size=16, bold=True, color=BLUE, font='Consolas')

text_box(s, 6.9, 3.05, 6.0, 1.4,
         "θ:    θ̈  =  −(g/l) sin θ\n"
         "             +  (1/l) u cos θ\n"
         "             −  (b/(m l²)) θ̇          ← nonlinear, pendulum",
         size=16, bold=True, color=ORANGE, font='Consolas')

text_box(s, 6.8, 4.7, 6.0, 0.4,
         "Why two separate equations?",
         size=15, bold=True, color=GREEN)
text_box(s, 6.8, 5.05, 6.0, 1.8,
         ["• Body inertia (~42 kg·m²) ≫ liquid (~0.5 kg·m²)",
          "• Sloshing barely pushes back on the body",
          "• So we treat ψ as if θ didn't exist",
          "• And θ as a pendulum hanging in an accelerating frame",
          "• ONE-WAY COUPLING:   u → ψ ,   u → θ  (no θ → ψ)"],
         size=13, color=BLACK)


# =============================================================================
# SLIDE 3 — Free body diagram of the sloshing pendulum
# =============================================================================
s = blank_slide()
title_bar(s, "Step 1: Free body diagram of the liquid (pendulum model)",
          "Pivot accelerates with u → work in the pivot's non-inertial frame")

# Diagram on left
box(s, 0.5, 1.3, 6.5, 5.6, fill=YELLOW_BG, border=ORANGE)
# pivot
pv_x, pv_y = 3.0, 2.0
circle(s, pv_x, pv_y, 0.10, fill=BLACK, border=BLACK)
text_box(s, pv_x + 0.15, pv_y - 0.15, 1.5, 0.3, "pivot (in container)",
         size=11, color=GRAY)
# vertical reference (gravity-down)
line(s, pv_x, pv_y, pv_x, pv_y + 3.5, color=GRAY, weight=1.0, dash=True)
# rod (tilted by θ)
theta_deg = 28
pl_in = 2.6
rdx = pl_in * math.sin(math.radians(theta_deg))
rdy = pl_in * math.cos(math.radians(theta_deg))
bob_x, bob_y = pv_x + rdx, pv_y + rdy
line(s, pv_x, pv_y, bob_x, bob_y, color=ORANGE, weight=4.0)
circle(s, bob_x, bob_y, 0.30, fill=ORANGE, border=BLACK, bpt=2)
text_box(s, bob_x + 0.35, bob_y - 0.20, 0.6, 0.4, "m",
         size=18, bold=True, color=BLACK)
# theta arc label
text_box(s, pv_x + 0.10, pv_y + 0.55, 0.6, 0.4, "θ",
         size=20, bold=True, color=ORANGE)
# length label
text_box(s, pv_x + rdx/2 + 0.15, pv_y + rdy/2 - 0.20, 0.5, 0.4, "l",
         size=20, bold=True, color=BLACK)

# forces on the bob
# gravity mg (downward)
arrow(s, bob_x, bob_y + 0.30, bob_x, bob_y + 1.4, color=GREEN, weight=3)
text_box(s, bob_x + 0.10, bob_y + 0.95, 0.8, 0.4, "m g",
         size=16, bold=True, color=GREEN)
# pseudo-force −m u (opposite to pivot acceleration u)
# pivot accel is to the right, so pseudo-force on bob is to the left
arrow(s, bob_x - 0.30, bob_y, bob_x - 1.4, bob_y, color=RED, weight=3)
text_box(s, bob_x - 1.55, bob_y - 0.45, 0.9, 0.4, "− m u",
         size=16, bold=True, color=RED)
# damping (small curved arrow indicator) — represented by a label
text_box(s, bob_x - 0.5, bob_y + 0.45, 1.5, 0.4, "−b θ̇  (damping)",
         size=12, color=GRAY)

# axis
text_box(s, pv_x - 1.2, pv_y - 0.5, 4.0, 0.4,
         "(non-inertial pivot frame)", size=12, color=GRAY)

# Right — torque balance
box(s, 7.3, 1.3, 5.7, 5.6, fill=BLUE_BG, border=BLUE)
text_box(s, 7.5, 1.4, 5.4, 0.5,
         "Torque about the pivot   τ = I α",
         size=18, bold=True, color=BLUE)

text_box(s, 7.5, 2.05, 5.4, 0.4,
         "Moment of inertia of point-mass bob:",
         size=14, color=BLACK)
text_box(s, 7.7, 2.45, 5.2, 0.4,
         "I  =  m l²",
         size=18, bold=True, color=BLACK, font='Consolas')

text_box(s, 7.5, 3.05, 5.4, 0.4,
         "Sum of moments (CCW positive):",
         size=14, color=BLACK)
text_box(s, 7.7, 3.45, 5.2, 2.0,
         ["gravity:        − m g l sin θ",
          "pivot pseudo:  + m u l cos θ",
          "damping:       − b θ̇"],
         size=15, color=BLACK, font='Consolas')

text_box(s, 7.5, 5.40, 5.4, 0.4,
         "Newton's 2nd law (rotational):",
         size=14, bold=True, color=BLUE)
text_box(s, 7.5, 5.80, 5.4, 0.5,
         "m l² θ̈  =  − m g l sin θ  +  m u l cos θ  − b θ̇",
         size=15, bold=True, color=BLUE, font='Consolas')

text_box(s, 7.5, 6.40, 5.4, 0.4,
         "→ Slide 4 simplifies this.",
         size=13, color=GRAY)


# =============================================================================
# SLIDE 4 — Sloshing equation (final form)
# =============================================================================
s = blank_slide()
title_bar(s, "Step 2: Divide by m l²  →  the sloshing equation",
          "Same equation that appears in your model card")

# Top: starting point
box(s, 0.7, 1.3, 11.9, 1.0, fill=BLUE_BG, border=BLUE)
text_box(s, 0.9, 1.45, 11.6, 0.4,
         "Starting point  (from torque balance):",
         size=14, bold=True, color=BLUE)
text_box(s, 0.9, 1.85, 11.6, 0.4,
         "m l² θ̈   =   − m g l sin θ   +   m u l cos θ   −   b θ̇",
         size=18, bold=True, color=BLACK, font='Consolas')

# Divide step
text_box(s, 0.7, 2.55, 12.0, 0.4,
         "Divide every term by  m l²  :",
         size=14, color=GRAY)

# Final form
box(s, 0.7, 3.0, 11.9, 1.5, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 0.9, 3.15, 11.6, 0.4,
         "Sloshing pendulum equation",
         size=14, bold=True, color=GREEN)
text_box(s, 0.9, 3.55, 11.6, 0.7,
         "θ̈   =   − (g / l) sin θ   +   (1 / l) u cos θ   −   ( b / (m l²) ) θ̇",
         size=22, bold=True, color=GREEN, font='Consolas')

# annotation row
arrow(s, 2.0, 4.55, 2.0, 4.95, color=GREEN, weight=2)
arrow(s, 5.5, 4.55, 5.5, 4.95, color=ORANGE, weight=2)
arrow(s, 9.5, 4.55, 9.5, 4.95, color=RED, weight=2)

text_box(s, 0.9, 5.0, 3.2, 0.5, "gravity restoring",
         size=14, bold=True, color=GREEN, align=CTR)
text_box(s, 0.9, 5.4, 3.2, 0.7,
         "pulls liquid back\ntoward vertical",
         size=12, color=GRAY, align=CTR)

text_box(s, 4.2, 5.0, 3.0, 0.5, "robot acceleration",
         size=14, bold=True, color=ORANGE, align=CTR)
text_box(s, 4.2, 5.4, 3.0, 0.7,
         "this is what u shaping\ntries to control",
         size=12, color=GRAY, align=CTR)

text_box(s, 8.0, 5.0, 3.0, 0.5, "viscous damping",
         size=14, bold=True, color=RED, align=CTR)
text_box(s, 8.0, 5.4, 3.0, 0.7,
         "energy dissipation\n(small for water)",
         size=12, color=GRAY, align=CTR)

# small-angle box
box(s, 0.7, 6.3, 11.9, 0.95, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.9, 6.4, 11.6, 0.45,
         "Small-angle linearization  (sin θ ≈ θ ,  cos θ ≈ 1):",
         size=13, bold=True, color=ORANGE)
text_box(s, 0.9, 6.80, 11.6, 0.45,
         "θ̈  +  2 δ ω_f θ̇  +  ω_f² θ   =   ω_f² · ( u / g )       "
         "where  ω_f² = g/l ,   2 δ ω_f = b/(m l²)",
         size=14, bold=True, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 5 — Robot body free body diagram
# =============================================================================
s = blank_slide()
title_bar(s, "Step 3: Free body diagram of the robot body",
          "Inverted pendulum on accelerating wheels — same Newton recipe")

# Diagram on left
box(s, 0.5, 1.3, 6.5, 5.6, fill=YELLOW_BG, border=ORANGE)

# axle and ground
ax_x, ax_y = 3.0, 6.4
line(s, 0.7, ax_y + 0.18, 6.9, ax_y + 0.18, color=BLACK, weight=1.5)
circle(s, ax_x, ax_y, 0.20, fill=GRAY, border=BLACK)
text_box(s, ax_x - 0.95, ax_y + 0.10, 0.9, 0.4, "wheel",
         size=11, color=GRAY)

# vertical reference
line(s, ax_x, ax_y, ax_x, ax_y - 4.0, color=GRAY, weight=1.0, dash=True)

# tilted body axis
psi_deg = 18
bL = 3.4   # length to CoM
bdx = bL * math.sin(math.radians(psi_deg))
bdy = bL * math.cos(math.radians(psi_deg))
com_x, com_y = ax_x + bdx, ax_y - bdy
line(s, ax_x, ax_y, com_x, com_y, color=BLUE, weight=4.0)

# CoM marker
circle(s, com_x, com_y, 0.18, fill=BLUE, border=BLACK)
text_box(s, com_x + 0.20, com_y - 0.18, 1.0, 0.4, "M  (CoM)",
         size=14, bold=True, color=BLUE)

# psi label
text_box(s, ax_x + 0.10, ax_y - 1.2, 0.6, 0.4, "ψ",
         size=22, bold=True, color=BLUE)
# L label
text_box(s, ax_x + bdx/2 - 0.45, ax_y - bdy/2, 0.6, 0.4, "L",
         size=18, bold=True, color=BLACK)

# gravity at CoM
arrow(s, com_x, com_y + 0.20, com_x, com_y + 1.2, color=GREEN, weight=3)
text_box(s, com_x + 0.10, com_y + 0.7, 0.8, 0.4, "M g",
         size=15, bold=True, color=GREEN)

# pivot pseudo-force at CoM (−M u, since base accelerates by u)
arrow(s, com_x - 0.20, com_y, com_x - 1.3, com_y, color=RED, weight=3)
text_box(s, com_x - 1.55, com_y - 0.40, 1.0, 0.4, "− M u",
         size=15, bold=True, color=RED)

# u arrow at base
arrow(s, ax_x + 0.30, ax_y - 0.10, ax_x + 1.4, ax_y - 0.10,
      color=RED, weight=3)
text_box(s, ax_x + 0.40, ax_y - 0.5, 1.0, 0.4, "u (= ẍ)",
         size=14, bold=True, color=RED)

# Right — torque balance
box(s, 7.3, 1.3, 5.7, 5.6, fill=BLUE_BG, border=BLUE)
text_box(s, 7.5, 1.4, 5.4, 0.5,
         "Torque about the wheel axle",
         size=18, bold=True, color=BLUE)

text_box(s, 7.5, 2.05, 5.4, 0.4,
         "Body inertia about axle:",
         size=14, color=BLACK)
text_box(s, 7.7, 2.45, 5.2, 0.4,
         "I_p   (≈ M L² + I_cm)",
         size=18, bold=True, color=BLACK, font='Consolas')

text_box(s, 7.5, 3.10, 5.4, 0.4,
         "Sum of moments:",
         size=14, color=BLACK)
text_box(s, 7.7, 3.50, 5.2, 1.4,
         ["gravity:        + M g L sin ψ",
          "pivot pseudo:  − M L u cos ψ",
          "(unstable: gravity tips it FURTHER)"],
         size=15, color=BLACK, font='Consolas')

text_box(s, 7.5, 5.10, 5.4, 0.4,
         "Newton's 2nd law:",
         size=14, bold=True, color=BLUE)
text_box(s, 7.5, 5.50, 5.4, 0.5,
         "I_p ψ̈   =   M g L sin ψ  −  M L u cos ψ",
         size=16, bold=True, color=BLUE, font='Consolas')

text_box(s, 7.5, 6.20, 5.4, 0.5,
         "Linearize  (sin ψ ≈ ψ ,  cos ψ ≈ 1):",
         size=12, color=GRAY)
text_box(s, 7.5, 6.55, 5.4, 0.45,
         "I_p ψ̈   ≈   M g L ψ  −  M L u",
         size=15, bold=True, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 6 — State-space form
# =============================================================================
s = blank_slide()
title_bar(s, "Step 4: Pack the body equation into state-space  q̇ = A q + B u",
          "This is what the LQT controller actually uses")

# State definition
box(s, 0.6, 1.3, 12.1, 1.1, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "State vector",
         size=14, bold=True, color=BLUE)
text_box(s, 0.8, 1.80, 11.8, 0.5,
         "q   =   [ ψ ,   ψ̇ ,   x ,   ẋ ]ᵀ          "
         "(tilt, tilt-rate, position, velocity)",
         size=16, bold=True, color=BLACK, font='Consolas')

# A and B matrices side-by-side
box(s, 0.6, 2.6, 6.1, 4.0, fill=GREEN_BG, border=GREEN)
text_box(s, 0.8, 2.7, 5.7, 0.4, "Dynamics matrix  A  (4×4)",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 3.15, 5.7, 3.1,
         ["       ⌈   0          1     0    0   ⌉",
          " A  =  |  M g L/I_p   0     0    0   |",
          "       |   0          0     0    1   |",
          "       ⌊   0          0     0    0   ⌋",
          "",
          " row 2 row carries gravity tipping",
          " term M g L / I_p — that is what",
          " makes the robot UNSTABLE (positive",
          " eigenvalue) and forces us to use",
          " feedback control."],
         size=13, color=BLACK, font='Consolas')

box(s, 6.9, 2.6, 5.8, 4.0, fill=RED_BG, border=RED)
text_box(s, 7.1, 2.7, 5.4, 0.4, "Input matrix  B  (4×1)",
         size=14, bold=True, color=RED)
text_box(s, 7.1, 3.15, 5.4, 3.1,
         ["          ⌈     0      ⌉",
          "  B  =   |  − M L/I_p  |",
          "          |     0      |",
          "          ⌊     1      ⌋",
          "",
          " row 2: u tilts the body the",
          " opposite way to what you",
          " might expect (Segway physics).",
          " row 4: u IS the linear",
          " acceleration ẍ ."],
         size=13, color=BLACK, font='Consolas')

# Bottom — what u, ψ, x mean physically
box(s, 0.6, 6.7, 12.1, 0.65, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 6.78, 11.8, 0.5,
         "u  =  wheel acceleration command  →  appears in BOTH "
         "the body (via −ML/I_p) AND the sloshing (via cosθ/l).",
         size=13, bold=True, color=PURPLE)


# =============================================================================
# SLIDE 7 — Why decoupled
# =============================================================================
s = blank_slide()
title_bar(s, "Step 5: Why we treat ψ and θ as decoupled",
          "Newton would couple them — but the back-reaction is negligible")

# Top message
box(s, 0.6, 1.3, 12.1, 0.95, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.45, 11.8, 0.4,
         "Strict Newton: the sloshing liquid also pushes back on the body.",
         size=15, bold=True, color=BLACK)
text_box(s, 0.8, 1.85, 11.8, 0.4,
         "But that back-reaction is tiny because the liquid mass is ~1% "
         "of the body inertia.",
         size=14, color=GRAY)

# Two-column comparison
box(s, 0.6, 2.45, 6.1, 4.6, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 2.55, 5.7, 0.4, "Body inertia (about axle)",
         size=15, bold=True, color=BLUE)
text_box(s, 0.8, 3.0, 5.7, 1.2,
         ["I_p   ≈   M L²",
          "      ≈   30 kg · (1.2 m)²",
          "      ≈   42 kg · m²"],
         size=15, color=BLACK, font='Consolas')
text_box(s, 0.8, 4.5, 5.7, 0.4,
         "Liquid pendulum inertia",
         size=15, bold=True, color=ORANGE)
text_box(s, 0.8, 4.95, 5.7, 1.2,
         ["m l²   ≈   1.5 kg · (0.6 m)²",
          "       ≈   0.5 kg · m²"],
         size=15, color=BLACK, font='Consolas')
text_box(s, 0.8, 6.05, 5.7, 0.4,
         "Ratio:   m l²  /  I_p   ≈   1 %",
         size=16, bold=True, color=RED, font='Consolas')

box(s, 6.9, 2.45, 5.8, 4.6, fill=GREEN_BG, border=GREEN)
text_box(s, 7.1, 2.55, 5.4, 0.4,
         "Consequence",
         size=15, bold=True, color=GREEN)
text_box(s, 7.1, 3.0, 5.4, 4.0,
         ["• ψ-equation: ignore θ-back-reaction",
          "  → linear, time-invariant, easy to",
          "    design LQT/LQR for",
          "",
          "• θ-equation: treat pivot as a",
          "  prescribed input (driven by u),",
          "  pendulum doesn't fight back on",
          "  the body",
          "",
          "• ONE-WAY coupling:",
          "      u  →  ψ ",
          "      u  →  θ        (no  θ → ψ)",
          "",
          "• This is what lets us:",
          "  — design body LQT (Sec 4.3),",
          "  — design liquid shaper F_e(s)",
          "      (Sec 4.2)",
          "  in parallel and combine them."],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 8 — Summary
# =============================================================================
s = blank_slide()
title_bar(s, "Summary: the model in one page",
          "From Newton to the equations the controller actually sees")

# 3 big boxes
y0 = 1.3
box(s, 0.6, y0, 12.1, 1.45, fill=ORANGE_BG, border=ORANGE)
text_box(s, 0.8, y0 + 0.05, 11.8, 0.4, "Sloshing  (liquid pendulum):",
         size=15, bold=True, color=ORANGE)
text_box(s, 0.8, y0 + 0.45, 11.8, 0.5,
         "θ̈   =   − (g/l) sin θ   +   (1/l) u cos θ   −   "
         "( b / (m l²) ) θ̇",
         size=18, bold=True, color=BLACK, font='Consolas')
text_box(s, 0.8, y0 + 1.00, 11.8, 0.4,
         "→ derives the natural frequency  ω_f = √(g/l)   used to "
         "build  F_e(s)  and  T_e .",
         size=12, color=GRAY)

y1 = y0 + 1.6
box(s, 0.6, y1, 12.1, 1.45, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, y1 + 0.05, 11.8, 0.4, "Robot body  (inverted pendulum on wheels):",
         size=15, bold=True, color=BLUE)
text_box(s, 0.8, y1 + 0.45, 11.8, 0.5,
         "I_p ψ̈   =   M g L ψ   −   M L u                  "
         "(after small-angle linearization)",
         size=17, bold=True, color=BLACK, font='Consolas')
text_box(s, 0.8, y1 + 1.00, 11.8, 0.4,
         "→ packed into  q̇ = A q + B u  →  fed to the LQT controller "
         "(Sec 4.3).",
         size=12, color=GRAY)

y2 = y1 + 1.6
box(s, 0.6, y2, 12.1, 1.45, fill=GREEN_BG, border=GREEN)
text_box(s, 0.8, y2 + 0.05, 11.8, 0.4, "How the full SBSFC uses these models:",
         size=15, bold=True, color=GREEN)
text_box(s, 0.8, y2 + 0.45, 11.8, 1.0,
         ["• Liquid model  →  pick  ω_f , δ   →  build  F_e(s)  →  shape v_ref  (no sloshing)",
          "• Body model   →  design LQT gains K, K_track  →  balance the robot  (no tipping)",
          "• Aux compensator + DOB  →  cancel everything the models miss  (Secs 4.4–4.5)"],
         size=13, color=BLACK)

# Footnote
text_box(s, 0.6, 6.65, 12.1, 0.4,
         "All four control blocks (shaper, LQT, aux, DOB) start from the "
         "two Newton equations on this slide.",
         size=14, color=GRAY, align=CTR)


# ----- save ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'newton_model_explained.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  {len(prs.slides)} slides')
