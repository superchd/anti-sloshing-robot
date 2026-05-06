"""
Teaching PPT: RL Parameters Explained
For MECE 6397 anti-sloshing robot project (Choi et al. 2024)

Audience: professor / classmates
Goal: defend WHY each RL parameter has its value.

Structure:
  Slide 1 — Title
  Slide 2 — The 3 buckets (overview map)
  Slide 3 — Bucket 1: Problem definition (obs / action / sample time)
  Slide 4 — Bucket 2a: SAC core (γ, LR, τ, batch)
  Slide 5 — Bucket 2b: Memory & exploration (buffer, warm-start, entropy)
  Slide 6 — Bucket 2c: Network & training budget
  Slide 7 — Bucket 3: Reward engineering (50:20:2:0.05)
  Slide 8 — Parameter classification (measured / default / physics-tied / chosen)
  Slide 9 — One-sentence summary
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=NAVY, line_w=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, 0.4, 0.08, 12.5, 0.55, title,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.4, 0.78, 12.5, 0.4, subtitle,
                 size=14, color=GRAY)


def add_table(slide, x, y, col_widths, rows, *, header=True,
              header_fill=NAVY, header_color=WHITE,
              row_h=0.42, font_size=12, body_font="Calibri",
              first_col_bold=True):
    """Generic table renderer. rows is a list of tuples."""
    cur_y = y
    for r_idx, row in enumerate(rows):
        is_hdr = header and r_idx == 0
        fill = header_fill if is_hdr else (LIGHT if r_idx % 2 == 1 else WHITE)
        line = NAVY if not is_hdr else header_fill
        cur_x = x
        # background row
        add_box(slide, x, cur_y, sum(col_widths), row_h,
                fill=fill, line=line, line_w=0.5)
        for c_idx, cell in enumerate(row):
            color = header_color if is_hdr else NAVY
            bold = is_hdr or (first_col_bold and c_idx == 0)
            add_text(slide, cur_x, cur_y + 0.06, col_widths[c_idx], row_h - 0.1,
                     str(cell),
                     size=font_size, bold=bold, color=color,
                     align=PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT,
                     font=body_font if c_idx == 0 else "Consolas")
            cur_x += col_widths[c_idx]
        cur_y += row_h


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------------- Slide 1 — Title ----------------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

add_text(s, 0.6, 1.9, 12, 1.4,
         "RL Parameters",
         size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.3, 12, 0.8,
         "What every value means and why we chose it",
         size=22, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.4, 12, 0.6,
         "SAC controller for the anti-sloshing robot",
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 5.6, 12, 0.5,
         "MECE 6397 — Anti-Sloshing Robot Project",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ---------------- Slide 2 — Overview: 3 buckets ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "RL Parameters — Three Buckets",
           "Group every parameter into one of three categories")

add_text(s, 0.6, 1.0, 12.1, 0.5,
         "The story we'll tell:",
         size=16, bold=True, color=NAVY)

# Bucket 1
add_box(s, 0.6, 1.55, 12.1, 1.55, fill=LIGHT, line=ORANGE)
add_text(s, 0.85, 1.65, 11.6, 0.4,
         "Bucket 1  —  Problem definition",
         size=18, bold=True, color=ORANGE)
add_text(s, 0.85, 2.05, 11.6, 0.4,
         "What the agent sees and does",
         size=13, bold=True, color=NAVY)
add_text(s, 0.85, 2.40, 11.6, 0.6,
         "•  Observation: 7-D vector  [ψ, ψ̇, x, ẋ, θ, θ̇, v_ref]\n"
         "•  Action: 1 scalar (cart acceleration u, clipped to ±4 m/s²)\n"
         "•  Sample time: 10 ms  (100 Hz control)",
         size=12, color=NAVY, font="Consolas")

# Bucket 2
add_box(s, 0.6, 3.20, 12.1, 1.95, fill=LIGHT, line=ORANGE)
add_text(s, 0.85, 3.30, 11.6, 0.4,
         "Bucket 2  —  SAC algorithm hyperparameters",
         size=18, bold=True, color=ORANGE)
add_text(s, 0.85, 3.70, 11.6, 0.4,
         "How the agent learns",
         size=13, bold=True, color=NAVY)
add_text(s, 0.85, 4.05, 11.6, 1.05,
         "•  Core:    γ = 0.99   |   LR = 3e-4 (Adam)   |   τ = 5e-3   |   batch = 256\n"
         "•  Memory:  replay buffer 1M   |   warm-start = 1,000 steps\n"
         "•  Explore: target entropy = −1   (= −dim(action))\n"
         "•  Network: 2 hidden layers × 256 units, ReLU\n"
         "•  Budget:  600 episodes × 1,000 steps = 10 s each",
         size=12, color=NAVY, font="Consolas")

# Bucket 3
add_box(s, 0.6, 5.25, 12.1, 1.95, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 5.35, 11.6, 0.4,
         "Bucket 3  —  Reward engineering   (the actual design)",
         size=18, bold=True, color=GREEN)
add_text(s, 0.85, 5.75, 11.6, 0.4,
         "What the agent is rewarded for",
         size=13, bold=True, color=NAVY)
add_text(s, 0.85, 6.10, 11.6, 1.05,
         "•  r = − ( 50·θ²  +  20·ψ²  +  2·(ẋ − v_ref)²  +  0.05·Δu² )  −  100·fall\n"
         "•  Ratios 50 : 20 : 2 : 0.05  →  slosh > tilt > tracking > smoothness\n"
         "•  Terminal penalty 100 if  |ψ| > 20°  (robot fell)",
         size=12, color=NAVY, font="Consolas")


# ---------------- Slide 3 — Bucket 1: Problem definition ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Bucket 1 — Problem Definition",
           "What the agent sees and does (defined before any algorithm choice)")

# Diagram strip
add_box(s, 0.5, 1.10, 3.5, 0.8, fill=LIGHT, line=NAVY)
add_text(s, 0.5, 1.20, 3.5, 0.3, "obs (7-D)",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 1.50, 3.5, 0.3, "[ψ, ψ̇, x, ẋ, θ, θ̇, v_ref]",
         size=10, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")

add_box(s, 4.5, 1.10, 4.3, 0.8, fill=ORANGE, line=ORANGE)
add_text(s, 4.5, 1.20, 4.3, 0.3, "SAC actor",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 4.5, 1.50, 4.3, 0.3, "256 → 256 → 1",
         size=11, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_box(s, 9.3, 1.10, 3.5, 0.8, fill=LIGHT, line=NAVY)
add_text(s, 9.3, 1.20, 3.5, 0.3, "action (1-D)",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 9.3, 1.50, 3.5, 0.3, "u ∈ [−4, +4] m/s²",
         size=11, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")

# Table of parameters
col_widths = [3.3, 2.2, 7.0]
rows = [
    ("Parameter",            "Value",            "Engineering reason"),
    ("Observation dim",      "7",                "Minimum needed to be Markov"),
    ("Action dim",           "1",                "One scalar acceleration command"),
    ("Action bounds",        "±4 m/s²",          "Motor saturation limit (hardware)"),
    ("Sample time",          "10 ms",            "Matches plant integration step (100 Hz)"),
    ("Episode length",       "10 s = 1,000 steps", "Long enough for slosh to develop and decay"),
]
add_table(s, 0.5, 2.30, col_widths, rows,
          header_fill=NAVY, font_size=12, row_h=0.50)

# Take-away
add_box(s, 0.6, 5.6, 12.1, 1.6, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 5.70, 11.6, 0.4,
         "How to say it:",
         size=14, bold=True, color=GREEN)
add_text(s, 0.85, 6.05, 11.6, 1.05,
         "\"Before any algorithm, I had to define what the agent observes and outputs.\n"
         " 7 inputs (the system's full state plus the reference), 1 output (acceleration command),\n"
         " refreshed every 10 ms — exactly the same interface as the classical SBSFC controller.\"",
         size=12, color=NAVY)


# ---------------- Slide 4 — Bucket 2a: SAC core ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Bucket 2a — SAC Algorithm Core",
           "γ, learning rate, target smoothing, batch size — the four core knobs")

col_widths = [2.6, 1.7, 5.6, 2.8]
rows = [
    ("Parameter", "Value", "What it controls", "Why this value"),
    ("Discount  γ",
     "0.99",
     "Look-ahead horizon: 1/(1−γ) ≈ 100 steps ≈ 1 s",
     "Matches slosh time constant"),
    ("Learning rate",
     "3e-4",
     "Step size for Adam optimizer",
     "SAC paper default"),
    ("Target smooth  τ",
     "5e-3",
     "Slow target update:  Q̄ ← τQ + (1−τ)Q̄",
     "Stabilizes Q bootstrap"),
    ("Mini-batch size",
     "256",
     "Samples drawn per gradient step",
     "Standard SAC value"),
]
add_table(s, 0.5, 1.10, col_widths, rows,
          header_fill=NAVY, font_size=11, row_h=0.85)

# Highlight box
add_box(s, 0.6, 5.50, 12.1, 1.7, fill=LIGHT, line=ORANGE)
add_text(s, 0.85, 5.60, 11.6, 0.4,
         "Project-specific choice (not a default):",
         size=14, bold=True, color=ORANGE)
add_text(s, 0.85, 5.95, 11.6, 0.4,
         "γ = 0.99",
         size=20, bold=True, color=NAVY, font="Consolas")
add_text(s, 0.85, 6.40, 11.6, 0.7,
         "Discount factor sets how far the agent looks ahead.\n"
         "1/(1−0.99) ≈ 100 steps × 10 ms = 1 second  →  exactly the slosh oscillation period.",
         size=12, color=NAVY)


# ---------------- Slide 5 — Bucket 2b: Memory & exploration ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Bucket 2b — Memory & Exploration",
           "Replay buffer, warm-start, target entropy")

col_widths = [3.0, 1.8, 5.0, 2.9]
rows = [
    ("Parameter",            "Value",       "What it does",
     "Why this value"),
    ("Replay buffer length", "1,000,000",   "Stores past (s, a, r, s′) transitions",
     "≈ 3 hours of experience"),
    ("Warm-start steps",     "1,000",       "Random actions before learning starts",
     "Fills buffer with diverse data"),
    ("Target entropy  H̄",     "−1",          "Setpoint for policy stochasticity",
     "= −dim(action) heuristic"),
    ("Entropy LR",           "3e-4",        "How fast α adapts to hit H̄",
     "Same as actor/critic LR"),
]
add_table(s, 0.5, 1.10, col_widths, rows,
          header_fill=NAVY, font_size=11, row_h=0.75)

# Why entropy matters
add_box(s, 0.6, 5.10, 12.1, 2.10, fill=LIGHT, line=ORANGE)
add_text(s, 0.85, 5.20, 11.6, 0.4,
         "Why is target entropy = −1 the project-specific value?",
         size=14, bold=True, color=ORANGE)
add_text(s, 0.85, 5.60, 11.6, 1.5,
         "•  SAC explores by being stochastic:  π(a|s) = N(μ, σ).\n"
         "•  The temperature α auto-tunes σ so the policy entropy stays near H̄.\n"
         "•  Heuristic for continuous control:  H̄ = −dim(action).\n"
         "•  Our action is 1-D  →  H̄ = −1.\n"
         "•  Effect: actor stays moderately exploratory throughout training.",
         size=12, color=NAVY)


# ---------------- Slide 6 — Bucket 2c: Network & budget ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Bucket 2c — Network & Training Budget",
           "Capacity of the neural networks and how long we train")

col_widths = [3.2, 2.0, 7.7]
rows = [
    ("Parameter",                "Value",      "Why this value"),
    ("Hidden layers (actor)",    "2",          "Standard MLP depth for SAC"),
    ("Hidden units / layer",     "256",        "Default — enough capacity for a 7-D obs"),
    ("Activation",               "ReLU",       "Default — fast, non-saturating"),
    ("# critics",                "2  (+2 targets)", "Twin Q-trick; takes min(Q₁,Q₂) to fight overestimation"),
    ("Total actor weights",      "≈ 67,000",   "(7→256) + (256→256) + (256→1)"),
    ("Total critic weights",     "≈ 67,000 each",   "(8→256) + (256→256) + (256→1) per critic"),
    ("Max episodes",             "600",        "Empirical: convergence in 300–500 episodes"),
    ("Steps per episode",        "1,000",      "= 10 s at 100 Hz"),
    ("Stop criterion",           "avg reward ≥ −5", "Slosh ≈ 0.3°, ψ ≈ 0.5°  →  safe operation"),
]
add_table(s, 0.5, 1.10, col_widths, rows,
          header_fill=NAVY, font_size=12, row_h=0.55)


# ---------------- Slide 7 — Bucket 3: Reward engineering ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Bucket 3 — Reward Engineering",
           "The only knob the engineer actually designs")

# Reward equation
add_box(s, 0.6, 1.05, 12.1, 1.05, fill=ORANGE, line=ORANGE)
add_text(s, 0.6, 1.20, 12.1, 0.7,
         "r  =  −( 50·θ²  +  20·ψ²  +  2·(ẋ − v_ref)²  +  0.05·Δu² )  −  100·fall",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

# Coefficient table
col_widths = [1.3, 3.5, 3.0, 5.0]
rows = [
    ("Coef.", "Penalizes",                "Priority",        "Why this magnitude"),
    ("50",    "θ²   (slosh angle)",       "1st — primary",   "Largest weight: slosh is what we want to kill"),
    ("20",    "ψ²   (robot tilt)",        "2nd",             "Keep robot upright as a hard constraint"),
    ("2",     "(ẋ − v_ref)²",             "3rd",             "Track velocity, but not at the cost of slosh"),
    ("0.05",  "Δu²  (control jerk)",      "4th — soft",      "Suppress chatter without dominating"),
    ("100",   "fall  (|ψ| > 20°)",        "Terminal",        "Big enough to make falling unambiguously bad"),
]
add_table(s, 0.5, 2.30, col_widths, rows,
          header_fill=NAVY, font_size=11, row_h=0.55)

# Take-away
add_box(s, 0.6, 5.7, 12.1, 1.5, fill=LIGHT, line=GREEN)
add_text(s, 0.85, 5.80, 11.6, 0.4,
         "How to say it:",
         size=14, bold=True, color=GREEN)
add_text(s, 0.85, 6.15, 11.6, 1.0,
         "\"The relative ratios — 50 : 20 : 2 : 0.05 — encode the engineering priorities.\n"
         " Slosh is 25× more important than velocity tracking. The absolute magnitudes don't matter;\n"
         " the ratios do. This is the only place where I encoded *what the robot should care about*.\"",
         size=12, color=NAVY)


# ---------------- Slide 8 — Parameter classification ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Defending Each Parameter — Three Categories",
           "When the professor asks 'why this value?', point to the category")

col_widths = [3.3, 4.3, 4.9]
rows = [
    ("Type",            "Examples",                                 "How to defend"),
    ("Algorithmic default",
     "LR = 3e-4, batch = 256, τ = 5e-3, buffer = 1M, layers = 256",
     "\"I used the SAC paper's published value\""),
    ("Physics-tied",
     "γ = 0.99, sample time = 10 ms, target entropy = −1, episode = 10 s",
     "\"This value is locked by the system, not chosen freely\""),
    ("Engineering choice",
     "Reward weights 50 : 20 : 2 : 0.05, fall threshold 20°",
     "\"I designed these to encode my priorities\""),
]
add_table(s, 0.5, 1.10, col_widths, rows,
          header_fill=NAVY, font_size=12, row_h=1.00)

add_box(s, 0.6, 5.0, 12.1, 2.1, fill=LIGHT, line=ORANGE)
add_text(s, 0.85, 5.10, 11.6, 0.4,
         "Rule of thumb during Q&A:",
         size=14, bold=True, color=ORANGE)
add_text(s, 0.85, 5.50, 11.6, 1.5,
         "•  If asked about a default-type parameter →  cite the SAC paper.\n"
         "•  If asked about a physics-tied parameter →  cite the system property (slosh period, action dim).\n"
         "•  If asked about an engineering choice →  explain the priority encoded by the ratio.\n"
         "•  If you don't know →  say 'I would tune via grid search / sensitivity analysis.'  Never invent a reason.",
         size=12, color=NAVY)


# ---------------- Slide 9 — Summary ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Summary — Three Numbers That Define the RL Setup",
           "If you remember nothing else, remember these")

# Three big highlight boxes
boxes = [
    ("γ  =  0.99",
     "Discount factor",
     "100-step look-ahead = 1 second = slosh oscillation period",
     ORANGE),
    ("H̄  =  −1",
     "Target entropy",
     "= −dim(action).  Keeps policy stochastic for 1-D continuous control",
     PURPLE),
    ("50 : 20 : 2 : 0.05",
     "Reward priority ratio",
     "slosh > tilt > tracking > smoothness  (the actual engineering decision)",
     GREEN),
]

y = 1.10
for big, label, desc, color in boxes:
    add_box(s, 0.6, y, 12.1, 1.65, fill=LIGHT, line=color)
    add_text(s, 0.85, y + 0.10, 4.5, 0.7, big,
             size=28, bold=True, color=color, font="Consolas")
    add_text(s, 0.85, y + 0.85, 4.5, 0.4, label,
             size=13, bold=True, color=NAVY)
    add_text(s, 5.5, y + 0.45, 7.0, 1.0, desc,
             size=14, color=NAVY)
    y += 1.85

add_box(s, 0.6, 6.75, 12.1, 0.55, fill=ORANGE, line=ORANGE)
add_text(s, 0.6, 6.85, 12.1, 0.4,
         "Everything else is either a SAC-paper default or a hardware constant.",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/rl_parameters_explained.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
