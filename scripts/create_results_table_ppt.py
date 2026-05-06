"""
Academic-style results table — RL vs SBSFC vs Without SBSFC (LPF)
Mirrors the layout of Choi et al. 2024 Table 3, extended to 3 columns.

Scenario 1: Sudden Start and Stop  (T = 30 s, N = 30,000 steps)

Numbers from MATLAB run:
    RL    : mean|θ| = 0.2377°,  var = 0.1934°²,  max = 2.0809°
    SBSFC : mean|θ| = 0.3167°,  var = 0.1931°²,  max = 1.7120°
    LPF   : mean|θ| = 5.0659°,  var = 35.4129°², max = 11.4124°
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY   = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF4, 0xF1, 0xEC)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xC6, 0x28, 0x28)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_hline(slide, x, y, w, weight=1.5, color=BLACK):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y),
                                      Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = color
    return line


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ============================================================
# Slide 1 — 3-column academic table  (RL | SBSFC | Without SBSFC)
# ============================================================
s = prs.slides.add_slide(blank)

# ---- Title (top-left like the paper) ----
add_text(s, 0.6, 0.45, 12.0, 0.55,
         "Table 1",
         size=22, bold=True, color=BLACK)
add_text(s, 0.6, 0.95, 12.0, 0.55,
         "Comparison of pendulum swing angles: RL, SBSFC, and without SBSFC.",
         size=16, color=BLACK)

# ---- Top thick horizontal rule ----
table_x = 0.6
table_w = 12.0
y_rule_top = 1.65
add_hline(s, table_x, y_rule_top,         table_w, weight=2.0, color=BLACK)

# ---- Column header row ----
y_hdr = 1.75
add_text(s, table_x,        y_hdr, 3.6, 0.45, "Value",
         size=15, bold=False, color=BLACK, align=PP_ALIGN.LEFT)
add_text(s, table_x + 3.6,  y_hdr, 2.8, 0.45, "RL",
         size=15, bold=False, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 6.4,  y_hdr, 2.8, 0.45, "With SBSFC",
         size=15, bold=False, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 9.2,  y_hdr, 2.8, 0.45, "Without SBSFC",
         size=15, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

# ---- Thin horizontal rule under header ----
add_hline(s, table_x, y_hdr + 0.55, table_w, weight=1.0, color=BLACK)

# ---- Row 1: Mean |θ| ----
y_r1 = 2.65
add_text(s, table_x,        y_r1, 3.6, 0.45, "Mean ( |θ|  [degree] )",
         size=15, color=BLACK, align=PP_ALIGN.LEFT)
add_text(s, table_x + 3.6,  y_r1, 2.8, 0.45, "0.2377",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 6.4,  y_r1, 2.8, 0.45, "0.3167",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 9.2,  y_r1, 2.8, 0.45, "5.0659",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)

# ---- Row 2: Variance θ² ----
y_r2 = 3.50
add_text(s, table_x,        y_r2, 3.6, 0.45, "Variance ( θ  [degree²] )",
         size=15, color=BLACK, align=PP_ALIGN.LEFT)
add_text(s, table_x + 3.6,  y_r2, 2.8, 0.45, "0.1934",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 6.4,  y_r2, 2.8, 0.45, "0.1931",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 9.2,  y_r2, 2.8, 0.45, "35.4129",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)

# ---- Row 3: Max |θ| ----
y_r3 = 4.35
add_text(s, table_x,        y_r3, 3.6, 0.45, "Max ( |θ|  [degree] )",
         size=15, color=BLACK, align=PP_ALIGN.LEFT)
add_text(s, table_x + 3.6,  y_r3, 2.8, 0.45, "2.0809",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 6.4,  y_r3, 2.8, 0.45, "1.7120",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 9.2,  y_r3, 2.8, 0.45, "11.4124",
         size=15, color=BLACK, align=PP_ALIGN.CENTER)

# ---- Bottom thick horizontal rule ----
add_hline(s, table_x, 5.10, table_w, weight=2.0, color=BLACK)

# ---- Footnote (italic) ----
add_text(s, 0.6, 5.30, 12.0, 0.4,
         "*Results from Scenario 1 only (sudden start and stop, T = 30 s).",
         size=13, color=BLACK, italic=True)

# ---- Reduction summary block (green) ----
add_text(s, 0.6, 5.95, 12.0, 0.45,
         "Reduction vs. Without SBSFC (LPF baseline):",
         size=15, bold=True, color=GREEN)
add_text(s, 0.6, 6.40, 12.0, 0.45,
         "•  SBSFC :   Mean = 93.7 %     Variance = 99.5 %",
         size=14, bold=True, color=GREEN, font="Consolas")
add_text(s, 0.6, 6.80, 12.0, 0.45,
         "•  RL    :   Mean = 95.3 %     Variance = 99.5 %",
         size=14, bold=True, color=GREEN, font="Consolas")


# ============================================================
# Slide 2 — RL vs SBSFC head-to-head (the more interesting comparison)
# ============================================================
s = prs.slides.add_slide(blank)

add_text(s, 0.6, 0.45, 12.0, 0.55,
         "Table 2",
         size=22, bold=True, color=BLACK)
add_text(s, 0.6, 0.95, 12.0, 0.55,
         "RL vs. SBSFC head-to-head (LPF excluded for clarity).",
         size=16, color=BLACK)

table_x = 1.5
table_w = 10.3

# Top rule
add_hline(s, table_x, 1.75, table_w, weight=2.0, color=BLACK)

# Headers
y_hdr = 1.85
add_text(s, table_x,         y_hdr, 4.3, 0.45, "Value",
         size=16, color=BLACK, align=PP_ALIGN.LEFT)
add_text(s, table_x + 4.3,   y_hdr, 3.0, 0.45, "RL",
         size=16, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, table_x + 7.3,   y_hdr, 3.0, 0.45, "SBSFC",
         size=16, color=BLACK, align=PP_ALIGN.CENTER)

# Header rule
add_hline(s, table_x, 2.40, table_w, weight=1.0, color=BLACK)

rows = [
    ("Mean ( |θ|  [degree] )",       "0.2377",  "0.3167"),
    ("Variance ( θ  [degree²] )",    "0.1934",  "0.1931"),
    ("Max ( |θ|  [degree] )",        "2.0809",  "1.7120"),
    ("Mean ( |ψ|  [degree] )",       "—",       "3.0956"),
]

y_row = 2.55
row_h = 0.65
for value, rl, sbsfc in rows:
    add_text(s, table_x,         y_row, 4.3, row_h, value,
             size=15, color=BLACK, align=PP_ALIGN.LEFT)
    # winner highlighting (smaller is better, ignore "—")
    rl_better    = (rl != "—" and sbsfc != "—" and float(rl) < float(sbsfc))
    sbsfc_better = (rl != "—" and sbsfc != "—" and float(sbsfc) < float(rl))
    rl_color    = GREEN  if rl_better    else BLACK
    sbsfc_color = GREEN  if sbsfc_better else BLACK
    add_text(s, table_x + 4.3,   y_row, 3.0, row_h, rl,
             size=15, bold=rl_better, color=rl_color, align=PP_ALIGN.CENTER)
    add_text(s, table_x + 7.3,   y_row, 3.0, row_h, sbsfc,
             size=15, bold=sbsfc_better, color=sbsfc_color, align=PP_ALIGN.CENTER)
    y_row += row_h

# Bottom rule
add_hline(s, table_x, y_row + 0.05, table_w, weight=2.0, color=BLACK)

# Footnote
add_text(s, table_x, y_row + 0.20, table_w, 0.40,
         "*Bold green = winner.  Smaller is better for all metrics.",
         size=12, color=BLACK, italic=True)

# Take-away
add_text(s, 0.6, 6.10, 12.0, 0.45,
         "Take-away:",
         size=15, bold=True, color=NAVY)
add_text(s, 0.6, 6.50, 12.0, 0.45,
         "•  RL wins on average performance (mean) — 25 % lower than SBSFC.",
         size=13, color=NAVY)
add_text(s, 0.6, 6.85, 12.0, 0.45,
         "•  SBSFC wins on worst-case peak (max) — its shaper is analytically optimal.",
         size=13, color=NAVY)


# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/results_table.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
