"""
Teaching PPT: How to derive q̇ = A·q + B·u from the robot body Newton equation
For MECE 6397 anti-sloshing robot project (Choi et al. 2024)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------- Style constants ----------------
NAVY = RGBColor(0x1F, 0x3A, 0x68)
ORANGE = RGBColor(0xD8, 0x5A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=NAVY, line_w=1.25):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=ORANGE, weight=2.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = color
    return line


def header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, 0.4, 0.08, 12.5, 0.55, title,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 0.4, 0.78, 12.5, 0.4, subtitle,
                 size=14, color=GRAY)


# ============================================================
# Build deck
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------------- Slide 1 — Title ----------------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

add_text(s, 0.6, 2.2, 12, 1.4,
         "From Newton's Equation to State-Space Form",
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 3.8, 12, 0.7,
         "How to build q̇ = A·q + B·u step by step",
         size=24, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.6, 12, 0.6,
         "MECE 6397 — Anti-Sloshing Robot Project",
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)


# ---------------- Slide 2 — Goal ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Our Goal",
           "Convert one second-order ODE into a first-order matrix equation")

add_box(s, 0.8, 1.3, 11.7, 1.4, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.45, 11, 0.4, "We start with this physics equation:",
         size=18, color=GRAY)
add_text(s, 1.0, 1.85, 11, 0.7,
         "I_p · ψ̈  =  M·g·L · sin(ψ)  −  M·L · u · cos(ψ)",
         size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s, 0.8, 3.0, 12, 0.4, "And we want to end up here:",
         size=18, color=GRAY)
add_box(s, 0.8, 3.4, 11.7, 1.0, fill=ORANGE, line=ORANGE)
add_text(s, 0.8, 3.55, 11.7, 0.7,
         "q̇  =  A · q  +  B · u",
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.8, 4.7, 12, 0.4,
         "Why? Because controllers (LQR / LQT) need first-order matrix form.",
         size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_box(s, 1.5, 5.4, 10.3, 1.6, fill=LIGHT, line=GREEN)
add_text(s, 1.7, 5.5, 10, 0.4, "The Plan — 4 Steps:",
         size=18, bold=True, color=GREEN)
add_text(s, 1.7, 5.9, 10, 0.4,
         "1.  Linearize (sin → ψ, cos → 1)",
         size=15, color=NAVY)
add_text(s, 1.7, 6.2, 10, 0.4,
         "2.  Define state vector q = [ψ, ψ̇]ᵀ",
         size=15, color=NAVY)
add_text(s, 1.7, 6.5, 10, 0.4,
         "3.  Write q̇₁ and q̇₂ as separate first-order equations",
         size=15, color=NAVY)
add_text(s, 1.7, 6.8, 10, 0.4,
         "4.  Stack into matrix form → read off A and B",
         size=15, color=NAVY)


# ---------------- Slide 3 — Step 1: Linearize ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Step 1 — Linearize for Small Tilt",
           "The robot only tilts a few degrees, so we can simplify")

add_box(s, 0.6, 1.2, 6.0, 2.4, fill=LIGHT, line=NAVY)
add_text(s, 0.8, 1.3, 5.6, 0.4, "Original (nonlinear):",
         size=16, bold=True, color=NAVY)
add_text(s, 0.8, 1.85, 5.6, 0.7,
         "I_p·ψ̈ = MgL·sin(ψ) − ML·u·cos(ψ)",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 2.7, 5.6, 0.4,
         "Has sin(ψ) and cos(ψ) — hard to solve.",
         size=14, color=GRAY)

add_arrow(s, 6.7, 2.4, 7.5, 2.4, color=ORANGE, weight=3)

add_box(s, 7.6, 1.2, 5.2, 2.4, fill=LIGHT, line=GREEN)
add_text(s, 7.8, 1.3, 4.8, 0.4, "Small-angle trick:",
         size=16, bold=True, color=GREEN)
add_text(s, 7.8, 1.8, 4.8, 0.4, "sin(ψ) ≈ ψ",
         size=18, bold=True, color=NAVY)
add_text(s, 7.8, 2.2, 4.8, 0.4, "cos(ψ) ≈ 1",
         size=18, bold=True, color=NAVY)
add_text(s, 7.8, 2.7, 4.8, 0.4,
         "Valid for ψ < ~15°.  Robot tilt is small.",
         size=13, color=GRAY)

add_box(s, 1.5, 4.0, 10.3, 1.4, fill=ORANGE, line=ORANGE)
add_text(s, 1.5, 4.15, 10.3, 0.5, "Plug them in →",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 4.55, 10.3, 0.7,
         "I_p · ψ̈  =  MgL · ψ  −  ML · u",
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.8, 5.7, 12, 0.4,
         "Now divide both sides by I_p to isolate ψ̈:",
         size=16, color=GRAY)
add_box(s, 1.5, 6.1, 10.3, 1.0, fill=LIGHT, line=NAVY)
add_text(s, 1.5, 6.25, 10.3, 0.7,
         "ψ̈  =  (MgL / I_p) · ψ  −  (ML / I_p) · u",
         size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ---------------- Slide 4 — Step 2: Define state vector ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Step 2 — Define a State Vector",
           "Trade one 2nd-order ODE for two 1st-order ODEs")

add_text(s, 0.8, 1.2, 12, 0.5,
         "The current equation is SECOND-order (it has ψ̈).",
         size=18, color=GRAY)
add_text(s, 0.8, 1.65, 12, 0.5,
         "Controllers want FIRST-order. The trick: treat ψ̇ as its own variable.",
         size=18, color=GRAY)

add_box(s, 1.5, 2.4, 10.3, 1.7, fill=LIGHT, line=NAVY)
add_text(s, 1.7, 2.55, 10, 0.4, "Define the state vector q:",
         size=16, bold=True, color=NAVY)
add_text(s, 1.7, 3.05, 10, 0.9,
         "q  =  [ q₁ ]  =  [ ψ  ]      ← angle\n      [ q₂ ]      [ ψ̇ ]      ← angular velocity",
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_box(s, 0.8, 4.5, 5.8, 2.6, fill=LIGHT, line=ORANGE)
add_text(s, 1.0, 4.6, 5.4, 0.4, "Why two variables?",
         size=16, bold=True, color=ORANGE)
add_text(s, 1.0, 5.05, 5.4, 1.9,
         "A 2nd-order ODE describes a system with TWO\n"
         "pieces of state info you must track:\n\n"
         "  • where it is now  (ψ)\n"
         "  • how fast it's moving  (ψ̇)\n\n"
         "If you know both → you can predict the future.",
         size=13, color=NAVY)

add_box(s, 6.9, 4.5, 5.8, 2.6, fill=LIGHT, line=GREEN)
add_text(s, 7.1, 4.6, 5.4, 0.4, "Analogy:",
         size=16, bold=True, color=GREEN)
add_text(s, 7.1, 5.05, 5.4, 1.9,
         "Think of a car:\n\n"
         "  • position alone is not enough\n"
         "  • velocity alone is not enough\n"
         "  • together they FULLY describe the car\n\n"
         "Same idea: q = [ψ, ψ̇] fully describes the body.",
         size=13, color=NAVY)


# ---------------- Slide 5 — Step 3: Write q̇₁ and q̇₂ ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Step 3 — Write q̇ Component by Component",
           "Take time-derivative of each entry of q")

add_text(s, 0.8, 1.2, 12, 0.5,
         "We need q̇ — that means take time-derivative of each row of q.",
         size=18, color=GRAY)

# q̇₁
add_box(s, 0.6, 1.9, 6.0, 2.3, fill=LIGHT, line=NAVY)
add_text(s, 0.8, 2.0, 5.6, 0.4, "First entry: q̇₁",
         size=16, bold=True, color=NAVY)
add_text(s, 0.8, 2.45, 5.6, 0.5,
         "q₁ = ψ ,  so  q̇₁ = ψ̇",
         size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 3.0, 5.6, 0.4,
         "And ψ̇ is exactly q₂ by our definition.",
         size=14, color=GRAY)
add_box(s, 0.9, 3.5, 5.6, 0.65, fill=ORANGE, line=ORANGE)
add_text(s, 0.9, 3.55, 5.6, 0.55,
         "q̇₁  =  q₂",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# q̇₂
add_box(s, 6.8, 1.9, 6.0, 2.3, fill=LIGHT, line=NAVY)
add_text(s, 7.0, 2.0, 5.6, 0.4, "Second entry: q̇₂",
         size=16, bold=True, color=NAVY)
add_text(s, 7.0, 2.45, 5.6, 0.5,
         "q₂ = ψ̇ ,  so  q̇₂ = ψ̈",
         size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 7.0, 3.0, 5.6, 0.4,
         "And ψ̈ comes from Step 1's linearized eq.",
         size=14, color=GRAY)
add_box(s, 7.1, 3.5, 5.6, 0.65, fill=ORANGE, line=ORANGE)
add_text(s, 7.1, 3.55, 5.6, 0.55,
         "q̇₂  =  (MgL/I_p)·ψ − (ML/I_p)·u",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Substitute ψ = q₁
add_box(s, 1.5, 4.6, 10.3, 1.0, fill=LIGHT, line=GREEN)
add_text(s, 1.5, 4.7, 10.3, 0.4,
         "Replace ψ with q₁ (same thing):",
         size=15, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
add_text(s, 1.5, 5.1, 10.3, 0.5,
         "q̇₂  =  (MgL/I_p)·q₁  −  (ML/I_p)·u",
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Two equations together
add_text(s, 0.8, 5.9, 12, 0.4, "Now we have two first-order equations:",
         size=16, bold=True, color=NAVY)
add_box(s, 1.5, 6.3, 10.3, 1.0, fill=LIGHT, line=NAVY)
add_text(s, 1.5, 6.35, 10.3, 0.5,
         "q̇₁  =  0·q₁  +  1·q₂  +  0·u",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 6.78, 10.3, 0.5,
         "q̇₂  =  (MgL/I_p)·q₁  +  0·q₂  +  (−ML/I_p)·u",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ---------------- Slide 6 — Step 4: Stack into matrix ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Step 4 — Stack Into Matrix Form",
           "The two scalar equations become one matrix equation")

add_text(s, 0.8, 1.1, 12, 0.4,
         "Look at the two equations side-by-side and notice the pattern:",
         size=16, color=GRAY)

add_box(s, 1.0, 1.6, 11.3, 1.5, fill=LIGHT, line=NAVY)
add_text(s, 1.0, 1.7, 11.3, 0.5,
         "q̇₁  =     0    ·q₁  +    1    ·q₂   +     0     ·u",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 2.15, 11.3, 0.5,
         "q̇₂  =  (MgL/I_p)·q₁  +    0    ·q₂   +  (−ML/I_p)·u",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 1.0, 2.65, 11.3, 0.4,
         "↑ coefficients of q₁, q₂  go into A      ↑ coefficient of u goes into B",
         size=12, color=ORANGE, align=PP_ALIGN.CENTER)

# Big matrix equation
add_box(s, 1.0, 3.4, 11.3, 2.4, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 3.55, 11.3, 0.4,
         "Stack them into matrix form:",
         size=15, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
add_text(s, 1.0, 4.0, 11.3, 1.7,
         "[ q̇₁ ]     [   0           1   ] [ q₁ ]     [    0    ]\n"
         "[ q̇₂ ]  =  [ MgL/I_p       0   ] [ q₂ ]  +  [ −ML/I_p ] · u",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

# Final form
add_box(s, 1.5, 6.1, 10.3, 1.1, fill=LIGHT, line=GREEN)
add_text(s, 1.5, 6.2, 10.3, 0.4,
         "Which we write compactly as:",
         size=14, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
add_text(s, 1.5, 6.55, 10.3, 0.6,
         "q̇  =  A · q  +  B · u",
         size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ---------------- Slide 7 — Read off A and B ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "What Are A and B?",
           "Just bookkeeping matrices holding the coefficients")

add_box(s, 0.6, 1.3, 6.1, 2.6, fill=LIGHT, line=NAVY)
add_text(s, 0.8, 1.4, 5.7, 0.4,
         "A — system matrix",
         size=18, bold=True, color=NAVY)
add_text(s, 0.8, 1.85, 5.7, 0.4,
         "How current state drives future state",
         size=12, color=GRAY)
add_text(s, 0.8, 2.3, 5.7, 1.2,
         "A  =  [   0           1   ]\n      [ MgL/I_p       0   ]",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 0.8, 3.4, 5.7, 0.4,
         "Top row: ψ̇ = ψ̇  (kinematics)\nBottom row: dynamics from physics",
         size=11, color=GRAY)

add_box(s, 6.9, 1.3, 6.1, 2.6, fill=LIGHT, line=ORANGE)
add_text(s, 7.1, 1.4, 5.7, 0.4,
         "B — input matrix",
         size=18, bold=True, color=ORANGE)
add_text(s, 7.1, 1.85, 5.7, 0.4,
         "How control input u pushes the state",
         size=12, color=GRAY)
add_text(s, 7.1, 2.3, 5.7, 1.2,
         "B  =  [    0    ]\n      [ −ML/I_p ]",
         size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, 7.1, 3.4, 5.7, 0.4,
         "u doesn't directly change ψ (top = 0)\nu changes ψ̇ via torque (bottom)",
         size=11, color=GRAY)

add_box(s, 0.8, 4.2, 12.2, 2.9, fill=LIGHT, line=GREEN)
add_text(s, 1.0, 4.3, 11.8, 0.4,
         "Reading the matrices physically:",
         size=16, bold=True, color=GREEN)
add_text(s, 1.0, 4.75, 11.8, 0.4,
         "•  A[1,2] = 1  →  \"the rate of change of angle equals angular velocity\"  (definition)",
         size=14, color=NAVY)
add_text(s, 1.0, 5.15, 11.8, 0.4,
         "•  A[2,1] = MgL/I_p  →  \"gravity creates angular acceleration proportional to tilt\"",
         size=14, color=NAVY)
add_text(s, 1.0, 5.55, 11.8, 0.4,
         "•  B[2] = −ML/I_p  →  \"acceleration command u creates angular acceleration\"",
         size=14, color=NAVY)
add_text(s, 1.0, 6.05, 11.8, 0.4,
         "Every entry has a physical meaning — A and B are NOT mysterious.",
         size=14, color=GREEN, bold=True)
add_text(s, 1.0, 6.5, 11.8, 0.5,
         "They just store the coefficients we already wrote in Step 3.",
         size=14, color=GRAY)


# ---------------- Slide 8 — Why bother? ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Why Convert to State-Space?",
           "What this form unlocks")

reasons = [
    ("1.  LQR / LQT need it",
     "Optimal controllers are derived assuming q̇ = Aq + Bu form.\n"
     "K = −R⁻¹·Bᵀ·P comes from the Algebraic Riccati Equation, which uses A and B directly."),
    ("2.  Easy to solve numerically",
     "First-order systems integrate cleanly with ode45, Euler, RK4, etc.\n"
     "MATLAB / Simulink expect this form for state-space blocks."),
    ("3.  Generalizes to many states",
     "Add more variables (cart position, slosh angle, ...) → just bigger A and B.\n"
     "The math stays exactly the same."),
    ("4.  Stability / controllability tests",
     "Eigenvalues of A → stability. Rank of [B AB A²B ...] → controllability.\n"
     "These tools all assume the q̇ = Aq + Bu form."),
]

y0 = 1.3
for i, (title, body) in enumerate(reasons):
    add_box(s, 0.6, y0 + i * 1.4, 12.1, 1.25, fill=LIGHT, line=NAVY)
    add_text(s, 0.8, y0 + i * 1.4 + 0.05, 11.7, 0.4, title,
             size=16, bold=True, color=ORANGE)
    add_text(s, 0.8, y0 + i * 1.4 + 0.45, 11.7, 0.8, body,
             size=13, color=NAVY)


# ---------------- Slide 9 — Summary ----------------
s = prs.slides.add_slide(blank)
header_bar(s, "Summary — The 4-Step Recipe",
           "Use this for ANY second-order system (mass-spring, pendulum, RLC...)")

steps = [
    ("Step 1 — Linearize",
     "sin(ψ)→ψ , cos(ψ)→1 → I_p ψ̈ = MgL·ψ − ML·u",
     NAVY),
    ("Step 2 — Define state",
     "q = [ψ, ψ̇]ᵀ — pack position + velocity into one vector",
     ORANGE),
    ("Step 3 — Write q̇ component-wise",
     "q̇₁ = q₂  ,   q̇₂ = (MgL/I_p)·q₁ − (ML/I_p)·u",
     GREEN),
    ("Step 4 — Stack into matrix",
     "Read off A and B from the coefficients → q̇ = A·q + B·u",
     RED),
]

y0 = 1.2
for i, (title, body, color) in enumerate(steps):
    add_box(s, 0.6, y0 + i * 1.0, 12.1, 0.85, fill=LIGHT, line=color)
    add_text(s, 0.85, y0 + i * 1.0 + 0.05, 4.0, 0.5, title,
             size=16, bold=True, color=color)
    add_text(s, 4.95, y0 + i * 1.0 + 0.05, 7.5, 0.8, body,
             size=13, color=NAVY, font="Consolas")

add_box(s, 1.0, 5.6, 11.3, 1.6, fill=ORANGE, line=ORANGE)
add_text(s, 1.0, 5.75, 11.3, 0.5,
         "The Big Idea",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.2, 11.3, 0.9,
         "A and B are just bookkeeping. They store the coefficients\n"
         "of the linearized ODE in a form that controllers can use.",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
out = "/Users/hyundae/MATLAB-Drive/Project/results/state_space_derivation_explained.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
