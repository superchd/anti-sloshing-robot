#!/usr/bin/env python3
"""
create_aux_compensator_ppt.py
Teaching deck for Section 4.4 of Choi et al. (2024):
"Auxiliary compensator for balancing a robot"

Walks students through:
   why caster wheels create discontinuous disturbances
   -> Eq. (4):  u^L_t  =  K_c  sgn( d psi^L / dt )
   -> noise problem with derivative -> low-pass filter
   -> what it does in time, when it fires, role vs DOB

Run:    python3 scripts/create_aux_compensator_ppt.py
Output: results/aux_compensator_explained.pptx
"""
import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)
PURPLE_BG= rgb(0xF5, 0xEC, 0xFB)
TEAL     = rgb(0x00, 0x7A, 0x86)
TEAL_BG  = rgb(0xE0, 0xF2, 0xF1)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def line_(slide, x1, y1, x2, y2, color=BLACK, weight=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = conn.line._get_or_add_ln()
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', 'dash')
    return conn


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def circle(slide, cx, cy, r, fill=ORANGE, border=BLACK, bpt=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - r), Inches(cy - r),
                                 Inches(2 * r), Inches(2 * r))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE,
        bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.1, 11.7, 1.1,
         "Auxiliary Compensator",
         size=48, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.3, 11.7, 1.4,
         "Catching the casters before they tip the robot\n"
         "u^L_t  =  K_c  sgn( d ψ^L / dt )",
         size=22, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 5.5, 11.7, 0.5,
         "Choi et al. 2024 · Section 4.4",
         size=20, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.5, 11.7, 0.5,
         "MECE 6397   ·   Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2 — Why we need it (caster picture)
# =============================================================================
s = blank_slide()
title_bar(s, "Why an extra controller? — caster wheels make sudden kicks",
          "LQT alone is too slow for sharp, discontinuous disturbances")

# Left — caster picture
box(s, 0.5, 1.3, 6.2, 5.6, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.7, 1.4, 5.8, 0.5, "What casters do",
         size=16, bold=True, color=ORANGE)

# ground
line_(s, 0.7, 6.45, 6.4, 6.45, color=BLACK, weight=2)
# robot base
box(s, 1.6, 5.45, 4.2, 1.0, fill=LTGRAY, border=BLACK, bpt=2, rounded=False)
# main drive wheel
circle(s, 3.7, 6.55, 0.25, fill=GRAY, border=BLACK)
# casters - small wheels at front and back
circle(s, 1.9, 6.55, 0.15, fill=ORANGE, border=BLACK)
circle(s, 5.5, 6.55, 0.15, fill=ORANGE, border=BLACK)
text_box(s, 1.4, 6.75, 1.1, 0.4, "caster", size=11, bold=True, color=ORANGE)
text_box(s, 5.0, 6.75, 1.2, 0.4, "caster", size=11, bold=True, color=ORANGE)

# tilted obstacle (bump)
line_(s, 5.20, 6.45, 5.80, 6.30, color=BLACK, weight=2)

# pitch indicator
text_box(s, 2.5, 4.0, 3.0, 0.5, "small ψ (pitch)",
         size=14, bold=True, color=BLUE, align=CTR)

# annotations
text_box(s, 0.7, 1.95, 5.8, 0.4,
         "When the front (or rear) caster hits a tiny bump:",
         size=12, color=BLACK)
text_box(s, 0.7, 2.30, 5.8, 1.7,
         ["• body suddenly rocks forward/backward",
          "• ψ̇ flips sign within milliseconds",
          "• it is a transient, discontinuous disturbance",
          "• acts on the body, not on the liquid",
          "• LQT is smooth — too slow to catch this"],
         size=12, color=BLACK)

# Right — why LQT alone fails
box(s, 6.9, 1.3, 6.0, 5.6, fill=RED_BG, border=RED)
text_box(s, 7.1, 1.4, 5.6, 0.5,
         "Why LQT alone is not enough",
         size=16, bold=True, color=RED)
text_box(s, 7.1, 1.95, 5.6, 5.0,
         ["• LQT is OPTIMAL for the smooth nominal model.",
          "",
          "• But casters create:",
          "    – impulse-like torques",
          "    – discontinuous (jump) disturbances",
          "    – durations of only a few ms",
          "",
          "• Riccati-tuned gain K is bounded and smooth,",
          "  so it cannot react fast enough.",
          "",
          "• We need an extra term that fires the moment",
          "  ψ̇ flips sign — that is exactly when a caster",
          "  has just kicked the robot."],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 3 — Eq. (4)
# =============================================================================
s = blank_slide()
title_bar(s, "Equation (4): the auxiliary compensator term",
          "A small but immediate kick in the right direction")

# Big formula
box(s, 0.6, 1.3, 12.1, 1.5, fill=ORANGE_BG, border=ORANGE, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Auxiliary control law",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.85, 11.8, 0.7,
         "u^L_t   =   K_c   ·   sgn(  d ψ^L / dt  )",
         size=26, bold=True, color=BLACK, font='Consolas', align=CTR)

# 3 piece breakdown
box(s, 0.5, 3.0, 4.0, 4.0, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 3.10, 3.7, 0.5,
         "K_c   →   gain",
         size=14, bold=True, color=BLUE, font='Consolas')
text_box(s, 0.7, 3.55, 3.7, 3.3,
         ["• constant scalar to tune",
          "• determines KICK SIZE",
          "• too small → still drifts",
          "• too large → chatters",
          "  (controller buzzes)",
          "",
          "• in your code:",
          "    p.K_c = 1.5"],
         size=12, color=BLACK)

box(s, 4.65, 3.0, 4.0, 4.0, fill=PURPLE_BG, border=PURPLE)
text_box(s, 4.85, 3.10, 3.7, 0.5,
         "sgn( · )   →   sign",
         size=14, bold=True, color=PURPLE, font='Consolas')
text_box(s, 4.85, 3.55, 3.7, 3.4,
         ["sgn(x) =  +1   if  x > 0",
          "sgn(x) =   0   if  x = 0",
          "sgn(x) = −1   if  x < 0",
          "",
          "→ output is ALWAYS",
          "  ±K_c (or zero)",
          "",
          "→ never grows",
          "  with the magnitude",
          "  of the input"],
         size=12, color=BLACK, font='Consolas')

box(s, 8.8, 3.0, 4.0, 4.0, fill=GREEN_BG, border=GREEN)
text_box(s, 9.0, 3.10, 3.7, 0.5,
         "d ψ^L / dt   →   trigger",
         size=13, bold=True, color=GREEN, font='Consolas')
text_box(s, 9.0, 3.55, 3.7, 3.4,
         ["• ψ^L  =  pitch as seen",
          "  through a low-pass",
          "  filter (≈ 0.8 ψ).",
          "",
          "• its derivative tells",
          "  us the DIRECTION of",
          "  the body's pitch",
          "  motion right now.",
          "",
          "• when it flips sign",
          "  → caster kicked!"],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 4 — derivative noise + LPF fix
# =============================================================================
s = blank_slide()
title_bar(s, "Practical issue: derivative is noisy → use a low-pass filter",
          "What the paper says: 'a low-pass filter is employed'")

# top message
box(s, 0.6, 1.3, 12.1, 1.05, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Why a raw derivative is dangerous",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.80, 11.8, 0.5,
         "Numerical differentiation of a noisy IMU signal AMPLIFIES noise. "
         "A jitter of 0.001 rad becomes a spike of 0.001/dt — huge.",
         size=13, color=BLACK)

# Steps shown as a chain
y = 2.55
items = [
    ("1.", "Read pitch ψ from IMU",
     "      raw, slightly noisy",      BLUE),
    ("2.", "Apply a soft proxy ψ^L = α ψ_LPF + (1−α) ψ_prev",
     "      low-pass smoothing",       PURPLE),
    ("3.", "Compute  d ψ^L / dt    via finite difference",
     "      smaller derivative noise", ORANGE),
    ("4.", "Take sgn( · ) → ±1 only",
     "      kills noise amplitude",   RED),
    ("5.", "Multiply by K_c → final u_aux",
     "      bounded ±K_c",            GREEN),
]
for i, (n, lab, sub, c) in enumerate(items):
    yi = y + i * 0.85
    box(s, 0.6, yi, 12.1, 0.75, fill=WHITE, border=c, bpt=2)
    box(s, 0.6, yi, 0.75, 0.75, fill=c, border=c, bpt=0)
    text_box(s, 0.6, yi + 0.10, 0.75, 0.55, n,
             size=20, bold=True, color=WHITE, align=CTR)
    text_box(s, 1.6, yi + 0.10, 8.0, 0.5, lab,
             size=15, bold=True, color=c)
    text_box(s, 9.5, yi + 0.18, 3.1, 0.5, sub,
             size=11, color=GRAY, align=LFT)


# =============================================================================
# SLIDE 5 — Time-domain story
# =============================================================================
s = blank_slide()
title_bar(s, "What happens in time — a caster kick scenario",
          "u_aux fires only when the pitch motion changes direction")

# Three rows: ψ, dψ/dt, u_aux
def panel(slide, y_top, label, color):
    box(slide, 0.6, y_top, 12.1, 1.6, fill=BLUE_BG, border=color)
    text_box(slide, 0.8, y_top + 0.05, 4.0, 0.4, label,
             size=14, bold=True, color=color)
    # axes
    ax_l = 4.5
    ax_t = y_top + 0.20
    ax_w = 7.8
    ax_h = 1.30
    cy = ax_t + ax_h * 0.6
    line_(slide, ax_l, cy, ax_l + ax_w, cy, color=BLACK, weight=1)
    line_(slide, ax_l, ax_t + 0.05, ax_l, ax_t + ax_h - 0.05,
          color=BLACK, weight=1)
    return ax_l, ax_t, ax_w, ax_h, cy

# Row 1 — pitch ψ
y0 = 1.3
ax_l, ax_t, ax_w, ax_h, cy = panel(s, y0, "ψ  (pitch angle)", BLUE)
# baseline + bump
n = 80
prev = (ax_l, cy)
for i in range(1, n + 1):
    x = ax_l + (i / n) * ax_w
    # small wobble + spike around fraction 0.4
    t = i / n
    if 0.40 <= t <= 0.55:
        # rapid up-then-down
        local = (t - 0.40) / 0.15
        y = cy - 0.45 * math.sin(local * math.pi)
    else:
        y = cy + 0.05 * math.sin(t * 14)
    line_(s, prev[0], prev[1], x, y, color=BLUE, weight=2.5)
    prev = (x, y)

text_box(s, ax_l + 0.40 * ax_w, cy + 0.10, 2.0, 0.4,
         "← caster bump", size=11, bold=True, color=ORANGE)

# Row 2 — dψ/dt
y0 = 3.0
ax_l, ax_t, ax_w, ax_h, cy = panel(s, y0, "d ψ^L / dt  (filtered rate)", PURPLE)
prev = (ax_l, cy)
for i in range(1, n + 1):
    t = i / n
    x = ax_l + t * ax_w
    if 0.40 <= t <= 0.55:
        # sharp positive then negative
        local = (t - 0.40) / 0.15
        y = cy - 0.45 * math.cos(local * math.pi)
    else:
        y = cy + 0.04 * math.cos(t * 14)
    line_(s, prev[0], prev[1], x, y, color=PURPLE, weight=2.5)
    prev = (x, y)

# zero crossing markers
text_box(s, ax_l + 0.475 * ax_w - 0.45, cy - 0.55, 1.5, 0.4,
         "sign flip!", size=11, bold=True, color=RED)

# Row 3 — u_aux (sgn)
y0 = 4.7
ax_l, ax_t, ax_w, ax_h, cy = panel(s, y0, "u_aux  =  K_c · sgn( · )", GREEN)
# square pulses: 0 outside, +K_c then -K_c around bump
prev = (ax_l, cy)
for i in range(1, n + 1):
    t = i / n
    x = ax_l + t * ax_w
    if 0.40 <= t < 0.475:
        y = cy - 0.45
    elif 0.475 <= t <= 0.55:
        y = cy + 0.45
    else:
        y = cy
    line_(s, prev[0], prev[1], x, y, color=GREEN, weight=3)
    prev = (x, y)

# explanation strip below
box(s, 0.6, 6.55, 12.1, 0.85, fill=GREEN_BG, border=GREEN)
text_box(s, 0.8, 6.65, 11.8, 0.7,
         "Outside the caster event,  d ψ^L / dt ≈ 0  →  u_aux = 0 .   "
         "During the kick, u_aux jumps to ±K_c — a fast, "
         "bounded shove that the smooth LQT cannot deliver.",
         size=12, color=BLACK)


# =============================================================================
# SLIDE 6 — Block diagram (where it sits)
# =============================================================================
s = blank_slide()
title_bar(s, "Where  u_aux  sits in the full SBSFC loop",
          "It is added on top of LQT (Eq. 3); DOB compensates the rest")

# block diagram horizontal
y_top = 1.7
# v_d
box(s, 0.5, y_top, 1.7, 0.9, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.5, y_top + 0.18, 1.7, 0.6, "v_d",
         size=18, bold=True, color=BLACK, align=CTR)
arrow(s, 2.2, y_top + 0.45, 2.55, y_top + 0.45, color=GRAY, weight=2)

# F_e
box(s, 2.55, y_top, 1.7, 0.9, fill=PURPLE_BG, border=PURPLE)
text_box(s, 2.55, y_top + 0.18, 1.7, 0.6, "F_e(s)",
         size=18, bold=True, color=PURPLE, align=CTR)
arrow(s, 4.25, y_top + 0.45, 4.6, y_top + 0.45, color=GRAY, weight=2)

# LQT
box(s, 4.6, y_top, 1.7, 0.9, fill=BLUE_BG, border=BLUE)
text_box(s, 4.6, y_top + 0.18, 1.7, 0.6, "LQT",
         size=18, bold=True, color=BLUE, align=CTR)
arrow(s, 6.30, y_top + 0.45, 7.05, y_top + 0.45, color=GRAY, weight=2)

# Sum node
box(s, 7.05, y_top, 1.0, 0.9, fill=WHITE, border=BLACK, bpt=2)
text_box(s, 7.05, y_top + 0.18, 1.0, 0.6, "+",
         size=22, bold=True, color=BLACK, align=CTR)
arrow(s, 8.05, y_top + 0.45, 8.40, y_top + 0.45, color=GRAY, weight=2)

# Plant
box(s, 8.4, y_top, 2.0, 0.9, fill=ORANGE_BG, border=ORANGE)
text_box(s, 8.4, y_top + 0.18, 2.0, 0.6, "Plant (robot)",
         size=14, bold=True, color=ORANGE, align=CTR)
arrow(s, 10.4, y_top + 0.45, 10.8, y_top + 0.45, color=GRAY, weight=2)

# state q
box(s, 10.8, y_top, 1.5, 0.9, fill=WHITE, border=BLACK)
text_box(s, 10.8, y_top + 0.20, 1.5, 0.6, "q (state)",
         size=14, bold=True, color=BLACK, align=CTR)

# u_aux feed coming from below into sum
box(s, 4.6, y_top + 2.0, 3.4, 1.1, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 4.7, y_top + 2.07, 3.2, 0.4,
         "u_aux  =  K_c sgn(d ψ^L / dt)",
         size=14, bold=True, color=GREEN, align=CTR, font='Consolas')
text_box(s, 4.7, y_top + 2.50, 3.2, 0.5,
         "fires on caster kicks",
         size=12, color=GRAY, align=CTR)
arrow(s, 7.55, y_top + 2.0, 7.55, y_top + 0.95, color=GREEN, weight=2.5)

# DOB feed into sum
box(s, 8.5, y_top + 2.0, 3.4, 1.1, fill=RED_BG, border=RED, bpt=2.5)
text_box(s, 8.6, y_top + 2.07, 3.2, 0.4,
         "u_dob  =  − d̂ / m",
         size=14, bold=True, color=RED, align=CTR, font='Consolas')
text_box(s, 8.6, y_top + 2.50, 3.2, 0.5,
         "cancels small continuous drift",
         size=12, color=GRAY, align=CTR)
arrow(s, 7.85, y_top + 2.0, 7.75, y_top + 0.95, color=RED, weight=2.5)

# state feedback path
arrow(s, 11.5, y_top + 0.95, 11.5, y_top + 4.4, color=GRAY, weight=2)
arrow(s, 11.5, y_top + 4.4, 5.4, y_top + 4.4, color=GRAY, weight=2)
arrow(s, 5.4, y_top + 4.4, 5.4, y_top + 0.95, color=GRAY, weight=2)
text_box(s, 9.0, y_top + 4.45, 2.5, 0.4, "feedback q",
         size=12, color=GRAY)

# IMU feed for u_aux
text_box(s, 5.6, y_top + 3.2, 4.0, 0.5,
         "ψ from IMU  →  d ψ^L / dt",
         size=11, color=GRAY, font='Consolas')

# bottom — total
box(s, 0.6, 6.55, 12.1, 0.8, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 6.65, 11.8, 0.5,
         "u_total  =  u_LQT  +  u_aux  +  u_DOB    "
         "(this is exactly what controller.m computes)",
         size=14, bold=True, color=BLUE, font='Consolas', align=CTR)


# =============================================================================
# SLIDE 7 — Aux vs DOB role split
# =============================================================================
s = blank_slide()
title_bar(s, "Role split: aux compensator vs. DOB",
          "Two robustness blocks, two different jobs")

# Headers
box(s, 0.5, 1.3, 6.2, 5.6, fill=GREEN_BG, border=GREEN)
text_box(s, 0.7, 1.4, 5.8, 0.5,
         "Auxiliary compensator  (Sec 4.4)",
         size=16, bold=True, color=GREEN)
text_box(s, 0.7, 1.95, 5.8, 5.0,
         ["Handles:",
          "  • caster wheel kicks",
          "  • sudden, transient disturbances",
          "  • small-amplitude impulses",
          "",
          "Form:    u^L = K_c · sgn(dψ^L/dt)",
          "",
          "Strength:",
          "  • REACTS in one timestep",
          "  • bounded output (±K_c)",
          "  • simple to tune (one gain)",
          "",
          "Weakness:",
          "  • only fires on direction CHANGE",
          "  • cannot estimate magnitude",
          "  • can chatter if K_c too big"],
         size=12, color=BLACK)

box(s, 6.9, 1.3, 6.0, 5.6, fill=RED_BG, border=RED)
text_box(s, 7.1, 1.4, 5.6, 0.5,
         "Disturbance Observer  (Sec 4.5)",
         size=16, bold=True, color=RED)
text_box(s, 7.1, 1.95, 5.6, 5.0,
         ["Handles:",
          "  • model uncertainty",
          "  • friction / floor slope",
          "  • slow continuous drifts",
          "",
          "Form:    d̂̇ = η (d − d̂),  u_dob = −d̂/m",
          "",
          "Strength:",
          "  • estimates magnitude precisely",
          "  • smooth correction over time",
          "  • cancels biases LQT cannot",
          "",
          "Weakness:",
          "  • bandwidth-limited by η",
          "  • not fast enough for caster kicks"],
         size=12, color=BLACK)

# Bottom line
box(s, 0.5, 7.0, 12.4, 0.4, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.6, 7.0, 12.2, 0.4,
         "Together: aux for sharp transients, DOB for slow biases — "
         "they cover what LQT cannot.",
         size=12, bold=True, color=PURPLE, align=CTR)


# =============================================================================
# SLIDE 8 — In your code
# =============================================================================
s = blank_slide()
title_bar(s, "How it shows up in your controller.m",
          "controller.m, lines ~54-59")

# code block
box(s, 0.6, 1.3, 12.1, 3.3, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Auxiliary compensator block (case 'sbsfc'):",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.85, 11.8, 2.7,
         ["    psi_L = q(1) * 0.8;",
          "    dpsi_L = (psi_L - ctrl_state.psi_L_prev) / p.dt;",
          "    ctrl_state.psi_L_prev = psi_L;",
          "",
          "    ctrl_state.psi_L_filtered = ...",
          "        ctrl_state.alpha_lpf * dpsi_L + ...",
          "        (1 - ctrl_state.alpha_lpf) * ctrl_state.psi_L_filtered;",
          "",
          "    u_aux = p.K_c * sign(ctrl_state.psi_L_filtered);"],
         size=13, color=BLACK, font='Consolas')

# annotations panel
box(s, 0.6, 4.75, 5.95, 2.5, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 4.85, 5.6, 0.5, "Match with the paper",
         size=14, bold=True, color=BLUE)
text_box(s, 0.8, 5.30, 5.6, 1.9,
         ["• psi_L                = ψ^L",
          "• dpsi_L              = d ψ^L / dt",
          "• psi_L_filtered      = LPF version",
          "• alpha_lpf           = filter weight",
          "• p.K_c               = K_c",
          "• sign( · )           = sgn( · )"],
         size=12, color=BLACK, font='Consolas')

box(s, 6.75, 4.75, 5.95, 2.5, fill=GREEN_BG, border=GREEN)
text_box(s, 6.95, 4.85, 5.6, 0.5, "What you can tune",
         size=14, bold=True, color=GREEN)
text_box(s, 6.95, 5.30, 5.6, 1.9,
         ["• p.K_c          → kick magnitude",
          "  (start ~1, raise if drifts after",
          "   bump, lower if chatter)",
          "",
          "• ctrl_state.alpha_lpf → filter speed",
          "  (smaller = smoother, slower)"],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 9 — Summary
# =============================================================================
s = blank_slide()
title_bar(s, "Summary — auxiliary compensator on one page",
          "Small, fast, bounded — built to react")

# 5 bullets style
y = 1.4
items = [
    ("Trigger",  "fires when  d ψ^L / dt  changes sign (caster bump)",  BLUE),
    ("Formula",  "u^L  =  K_c · sgn( d ψ^L / dt )    (Eq. 4)",          ORANGE),
    ("Why sgn",  "bounded output ±K_c  →  noise stays harmless",        PURPLE),
    ("Why LPF",  "raw derivative is noisy; LPF gives clean rate",       RED),
    ("Role",     "transient disturbances only; DOB handles the rest",   GREEN),
]
for i, (k, v, c) in enumerate(items):
    yi = y + i * 0.85
    box(s, 0.6, yi, 12.1, 0.7, fill=WHITE, border=c, bpt=2.0)
    box(s, 0.6, yi, 0.7, 0.7, fill=c, border=c, bpt=0)
    # leading dot
    text_box(s, 0.6, yi + 0.10, 0.7, 0.5, "•",
             size=24, bold=True, color=WHITE, align=CTR)
    text_box(s, 1.5, yi + 0.13, 2.5, 0.5, k,
             size=15, bold=True, color=c)
    text_box(s, 4.2, yi + 0.13, 8.5, 0.5, v,
             size=14, color=BLACK)

# closing line
box(s, 0.6, 6.05, 12.1, 1.25, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 6.15, 11.8, 0.5,
         "Big picture",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 6.55, 11.8, 0.7,
         "Aux compensator is the 'reflex': fast, sharp, bounded. "
         "LQT keeps the robot upright on average, DOB cleans up slow biases, "
         "and aux catches sudden caster kicks before they grow.",
         size=13, color=BLACK)


# ----- save ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'aux_compensator_explained.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  {len(prs.slides)} slides')
