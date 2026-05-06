#!/usr/bin/env python3
"""
create_reference_shaping_progress_ppt.py
Teaching deck that walks students through the FULL progression of
reference shaping for the anti-sloshing robot:

   Newton's law (FBD)
    -> small-angle linearization
    -> Laplace transform
    -> transfer function G(s)
    -> shaper F_e(s) (pole cancellation)
    -> apply F_e(s) to v_ref

Run:    python3 scripts/create_reference_shaping_progress_ppt.py
Output: results/reference_shaping_progress.pptx
"""
import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)
PURPLE_BG= rgb(0xF5, 0xEC, 0xFB)
TEAL     = rgb(0x00, 0x7A, 0x86)
TEAL_BG  = rgb(0xE0, 0xF2, 0xF1)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def line_(slide, x1, y1, x2, y2, color=BLACK, weight=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = conn.line._get_or_add_ln()
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', 'dash')
    return conn


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def circle(slide, cx, cy, r, fill=ORANGE, border=BLACK, bpt=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - r), Inches(cy - r),
                                 Inches(2 * r), Inches(2 * r))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE,
        bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


def step_chip(slide, l, t, n, label, color):
    """Small numbered chip used in roadmap-style slides."""
    box(slide, l, t, 0.55, 0.55, fill=color, border=color, bpt=0)
    text_box(slide, l, t + 0.04, 0.55, 0.5, str(n), size=20, bold=True,
             color=WHITE, align=CTR)
    text_box(slide, l + 0.65, t + 0.07, 4.0, 0.5, label, size=15, bold=True,
             color=color)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.1, 11.7, 1.1,
         "Reference Shaping, Step by Step",
         size=46, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.3, 11.7, 1.4,
         "From Newton's law on a pendulum  →  Laplace transform  →  G(s)  "
         "→  the shaping filter F_e(s)",
         size=22, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 5.5, 11.7, 0.5,
         "A 6-step walkthrough for students",
         size=20, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.5, 11.7, 0.5,
         "MECE 6397   ·   Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2 — Roadmap
# =============================================================================
s = blank_slide()
title_bar(s, "Roadmap: 6 steps from physics to a working shaper",
          "Each box becomes one slide")

steps = [
    (1, "Newton's law on the pendulum (FBD)",      ORANGE),
    (2, "Small-angle linearization",                ORANGE),
    (3, "Laplace transform (time -> s-domain)",     RED),
    (4, "Transfer function   G(s)",                 RED),
    (5, "Design  F_e(s)  to cancel G(s) poles",     PURPLE),
    (6, "Apply  F_e(s)  to v_ref  (smooth ramp)",   GREEN),
]

# Centered roadmap
y = 1.5
for i, (n, lab, c) in enumerate(steps):
    yi = y + i * 0.85
    box(s, 1.4, yi, 10.5, 0.7, fill=WHITE, border=c, bpt=2.5)
    box(s, 1.4, yi, 0.7, 0.7, fill=c, border=c, bpt=0)
    text_box(s, 1.4, yi + 0.10, 0.7, 0.5, str(n),
             size=22, bold=True, color=WHITE, align=CTR)
    text_box(s, 2.3, yi + 0.13, 9.0, 0.5, lab,
             size=18, bold=True, color=c, align=LFT)
    if i < len(steps) - 1:
        # arrow down
        arrow(s, 6.65, yi + 0.72, 6.65, yi + 0.83,
              color=GRAY, weight=1.5)

text_box(s, 0.6, 6.85, 12.1, 0.5,
         "We will follow this exact order. Skip none — each step depends "
         "on the one above.",
         size=13, color=GRAY, align=CTR)


# =============================================================================
# SLIDE 3 — Step 1: Newton FBD on the pendulum
# =============================================================================
s = blank_slide()
title_bar(s, "Step 1: Newton's law on the sloshing pendulum",
          "Sum of moments about the pivot equals  I·θ̈")

# left FBD
box(s, 0.5, 1.3, 6.2, 5.6, fill=YELLOW_BG, border=ORANGE)
pv_x, pv_y = 3.0, 2.2
circle(s, pv_x, pv_y, 0.10, fill=BLACK, border=BLACK)
text_box(s, pv_x + 0.2, pv_y - 0.18, 1.6, 0.3, "pivot", size=11, color=GRAY)
# vertical reference
line_(s, pv_x, pv_y, pv_x, pv_y + 3.5, color=GRAY, weight=1.0, dash=True)
# rod
theta_deg = 28
pl_in = 2.6
rdx = pl_in * math.sin(math.radians(theta_deg))
rdy = pl_in * math.cos(math.radians(theta_deg))
bob_x, bob_y = pv_x + rdx, pv_y + rdy
line_(s, pv_x, pv_y, bob_x, bob_y, color=ORANGE, weight=4.0)
circle(s, bob_x, bob_y, 0.30, fill=ORANGE, border=BLACK, bpt=2)
text_box(s, bob_x + 0.35, bob_y - 0.20, 0.6, 0.4, "m",
         size=18, bold=True, color=BLACK)
text_box(s, pv_x + 0.10, pv_y + 0.55, 0.6, 0.4, "θ",
         size=20, bold=True, color=ORANGE)
text_box(s, pv_x + rdx/2 + 0.15, pv_y + rdy/2 - 0.20, 0.5, 0.4, "l",
         size=20, bold=True, color=BLACK)
# gravity
arrow(s, bob_x, bob_y + 0.30, bob_x, bob_y + 1.4, color=GREEN, weight=3)
text_box(s, bob_x + 0.10, bob_y + 0.95, 0.8, 0.4, "m g",
         size=16, bold=True, color=GREEN)
# pseudo-force
arrow(s, bob_x - 0.30, bob_y, bob_x - 1.4, bob_y, color=RED, weight=3)
text_box(s, bob_x - 1.55, bob_y - 0.45, 0.9, 0.4, "− m u",
         size=16, bold=True, color=RED)
text_box(s, bob_x - 0.5, bob_y + 0.45, 1.5, 0.4, "−b θ̇  (damping)",
         size=12, color=GRAY)

# right derivation
box(s, 7.0, 1.3, 6.0, 5.6, fill=BLUE_BG, border=BLUE)
text_box(s, 7.2, 1.4, 5.7, 0.5,
         "Sum of moments about the pivot",
         size=16, bold=True, color=BLUE)
text_box(s, 7.2, 2.0, 5.7, 0.5, "Inertia of a point-mass bob:",
         size=13, color=BLACK)
text_box(s, 7.4, 2.4, 5.4, 0.5, "I = m l²",
         size=18, bold=True, color=BLACK, font='Consolas')

text_box(s, 7.2, 3.05, 5.7, 0.5, "Three torques (CCW positive):",
         size=13, color=BLACK)
text_box(s, 7.4, 3.45, 5.4, 1.5,
         ["gravity:        − m g l sin θ",
          "pivot pseudo:  + m u l cos θ",
          "damping:       − b θ̇"],
         size=14, color=BLACK, font='Consolas')

text_box(s, 7.2, 5.05, 5.7, 0.4, "Newton's 2nd law (rotational):",
         size=14, bold=True, color=BLUE)
box(s, 7.2, 5.45, 5.7, 0.85, fill=WHITE, border=BLUE, bpt=2)
text_box(s, 7.3, 5.55, 5.5, 0.7,
         "m l² θ̈  =  −m g l sin θ  +  m u l cos θ  −  b θ̇",
         size=14, bold=True, color=BLACK, font='Consolas')

text_box(s, 7.2, 6.45, 5.7, 0.4,
         "→ this is still NONLINEAR (sin and cos).",
         size=13, color=GRAY)


# =============================================================================
# SLIDE 4 — Step 2: small-angle linearization
# =============================================================================
s = blank_slide()
title_bar(s, "Step 2: Linearize for small angles",
          "Trick: when θ is small, replace sinθ → θ and cosθ → 1")

# Top — the approximation
box(s, 0.6, 1.3, 12.1, 1.4, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Why we can do this:",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.80, 11.8, 0.85,
         ["• Liquid sloshing in a serving robot is small (a few degrees).",
          "• For |θ| < 10°: sinθ ≈ θ  (error < 0.5%) , cosθ ≈ 1.",
          "• A mildly curved equation becomes a clean LINEAR one — and "
          "linear equations have closed-form Laplace solutions."],
         size=12, color=BLACK)

# Step-by-step substitution
text_box(s, 0.6, 2.85, 12.1, 0.5,
         "Apply the approximation to the Newton equation:",
         size=15, bold=True, color=BLUE)

box(s, 0.6, 3.35, 12.1, 0.8, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 3.45, 11.8, 0.6,
         "Before (Step 1):   m l² θ̈  =  −m g l sin θ  +  m u l cos θ  −  b θ̇",
         size=15, color=BLACK, font='Consolas')

box(s, 0.6, 4.25, 12.1, 0.8, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 4.35, 11.8, 0.6,
         "After substitution:  m l² θ̈  =  −m g l · θ  +  m u l · 1  −  b θ̇",
         size=15, color=BLACK, font='Consolas')

# Divide by m l²
text_box(s, 0.6, 5.20, 12.1, 0.5,
         "Now divide every term by  m l²  and rearrange:",
         size=13, color=GRAY)

box(s, 0.6, 5.65, 12.1, 1.2, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 0.8, 5.75, 11.8, 0.4,
         "Linearized sloshing ODE",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 6.20, 11.8, 0.6,
         "θ̈   +   2 δ ω_f θ̇   +   ω_f² θ   =   ω_f² · ( u / g )",
         size=20, bold=True, color=GREEN, font='Consolas')

text_box(s, 0.6, 6.95, 12.1, 0.4,
         "where  ω_f² = g/l   and   2 δ ω_f = b/(m l²).   Clean linear "
         "2nd-order ODE — ready for Laplace.",
         size=12, color=GRAY, align=CTR)


# =============================================================================
# SLIDE 5 — Step 3: Laplace transform
# =============================================================================
s = blank_slide()
title_bar(s, "Step 3: Laplace transform — what it does and why",
          "Convert ODE in time  →  algebra in s")

# Two-column setup
box(s, 0.5, 1.3, 6.2, 5.7, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.7, 1.40, 5.8, 0.5,
         "What Laplace replaces",
         size=18, bold=True, color=PURPLE)
text_box(s, 0.7, 1.95, 5.8, 4.0,
         ["The Laplace transform is a substitution rule:",
          "",
          "    f(t)         →    F(s)",
          "    df/dt        →    s · F(s)",
          "    d²f/dt²      →    s² · F(s)",
          "",
          "(assuming zero initial conditions)",
          "",
          "Differential equations become",
          "ALGEBRAIC equations in s.",
          "",
          "We can:",
          "  • factor like polynomials,",
          "  • read off poles & zeros,",
          "  • design filters by canceling them."],
         size=13, color=BLACK)

# Right — apply it
box(s, 7.0, 1.3, 6.0, 5.7, fill=BLUE_BG, border=BLUE)
text_box(s, 7.2, 1.40, 5.6, 0.5,
         "Apply it to our linearized ODE",
         size=18, bold=True, color=BLUE)

text_box(s, 7.2, 1.95, 5.6, 0.5, "Starting equation:",
         size=13, color=BLACK)
text_box(s, 7.2, 2.35, 5.6, 0.5,
         "θ̈ + 2δω_f θ̇ + ω_f² θ  =  ω_f² (u/g)",
         size=14, bold=True, color=BLACK, font='Consolas')

text_box(s, 7.2, 3.05, 5.6, 0.5, "Substitute  d/dt → s :",
         size=13, color=BLACK)
text_box(s, 7.2, 3.45, 5.6, 0.6,
         "s² Θ(s) + 2δω_f s Θ(s) + ω_f² Θ(s)\n"
         "                                    =  ω_f² U(s)/g",
         size=13, bold=True, color=BLACK, font='Consolas')

text_box(s, 7.2, 4.40, 5.6, 0.5, "Factor Θ(s):",
         size=13, color=BLACK)
text_box(s, 7.2, 4.80, 5.6, 0.6,
         "( s² + 2δω_f s + ω_f² ) Θ(s)  =  ω_f² U(s)/g",
         size=14, bold=True, color=BLACK, font='Consolas')

text_box(s, 7.2, 5.55, 5.6, 0.5, "Solve for Θ(s)/U(s):",
         size=13, bold=True, color=BLUE)
box(s, 7.2, 5.95, 5.6, 0.95, fill=WHITE, border=BLUE, bpt=2.5)
text_box(s, 7.3, 6.05, 5.4, 0.8,
         "Θ(s)/(U(s)/g)  =  ω_f² / (s² + 2δω_f s + ω_f²)",
         size=14, bold=True, color=BLUE, font='Consolas', align=CTR)


# =============================================================================
# SLIDE 6 — Step 4: G(s) and its poles
# =============================================================================
s = blank_slide()
title_bar(s, "Step 4: The plant transfer function  G(s)  and its poles",
          "G(s) describes how the input u shakes the liquid")

# Big formula
box(s, 0.6, 1.3, 12.1, 1.4, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4, "Plant transfer function",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 1.80, 11.8, 0.8,
         "G(s)  =  ω_f²  /  ( s²  +  2 δ ω_f s  +  ω_f² )",
         size=24, bold=True, color=GREEN, align=CTR, font='Consolas')

# Left — pole formula
box(s, 0.5, 2.95, 6.2, 4.1, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 3.05, 5.8, 0.5,
         "Poles  =  roots of denominator",
         size=16, bold=True, color=BLUE)
text_box(s, 0.7, 3.55, 5.8, 0.5,
         "s²  +  2 δ ω_f s  +  ω_f²  =  0",
         size=14, bold=True, color=BLACK, font='Consolas')
text_box(s, 0.7, 4.10, 5.8, 0.5,
         "Use the quadratic formula:",
         size=13, color=BLACK)
text_box(s, 0.7, 4.50, 5.8, 0.7,
         "s_{1,2}  =  −δ ω_f  ±  j ω_d ,",
         size=15, bold=True, color=BLACK, font='Consolas')
text_box(s, 0.7, 5.10, 5.8, 0.5,
         "where  ω_d = ω_f √(1 − δ²)  (damped frequency).",
         size=13, color=GRAY)
text_box(s, 0.7, 5.65, 5.8, 0.5,
         "Two complex-conjugate poles:",
         size=13, color=BLACK)
text_box(s, 0.7, 6.05, 5.8, 0.95,
         "• Real part  −δ ω_f   →   damping",
         size=14, color=BLUE, font='Consolas')
text_box(s, 0.7, 6.40, 5.8, 0.95,
         "• Imag part  ± j ω_d  →   oscillation",
         size=14, color=RED, font='Consolas')

# Right — pole map (cartoon)
box(s, 7.0, 2.95, 6.0, 4.1, fill=YELLOW_BG, border=ORANGE)
text_box(s, 7.2, 3.05, 5.6, 0.5,
         "Pole map (s-plane)",
         size=16, bold=True, color=ORANGE)

# axes in the box
ax_l, ax_t, ax_w, ax_h = 7.4, 3.7, 5.3, 3.1
cx, cy = ax_l + ax_w / 2, ax_t + ax_h / 2
# real axis (horizontal)
line_(s, ax_l + 0.2, cy, ax_l + ax_w - 0.2, cy, color=BLACK, weight=1.5)
# imaginary axis (vertical)
line_(s, cx, ax_t + 0.2, cx, ax_t + ax_h - 0.2, color=BLACK, weight=1.5)
text_box(s, ax_l + ax_w - 0.5, cy + 0.05, 0.5, 0.4, "Re",
         size=12, color=BLACK)
text_box(s, cx + 0.10, ax_t + 0.10, 0.5, 0.4, "Im",
         size=12, color=BLACK)

# the two poles (left half-plane, complex-conjugate)
px, py_top, py_bot = cx - 1.0, cy - 0.95, cy + 0.95
# x markers
xsz = 0.20
line_(s, px - xsz, py_top - xsz, px + xsz, py_top + xsz, color=BLUE, weight=4)
line_(s, px - xsz, py_top + xsz, px + xsz, py_top - xsz, color=BLUE, weight=4)
line_(s, px - xsz, py_bot - xsz, px + xsz, py_bot + xsz, color=BLUE, weight=4)
line_(s, px - xsz, py_bot + xsz, px + xsz, py_bot - xsz, color=BLUE, weight=4)

text_box(s, px + 0.25, py_top - 0.20, 2.4, 0.4,
         "−δω_f + j ω_d", size=11, bold=True, color=BLUE, font='Consolas')
text_box(s, px + 0.25, py_bot - 0.20, 2.4, 0.4,
         "−δω_f − j ω_d", size=11, bold=True, color=BLUE, font='Consolas')
text_box(s, ax_l + 0.2, ax_t + ax_h - 0.4, 5.0, 0.4,
         "← these poles cause the ringing!",
         size=12, color=RED, bold=True)


# =============================================================================
# SLIDE 7 — Step 5: design F_e(s)
# =============================================================================
s = blank_slide()
title_bar(s, "Step 5: Design  F_e(s)  to CANCEL G(s)'s poles",
          "Pole-zero cancellation — kill the ringing at the source")

# Idea panel
box(s, 0.6, 1.3, 12.1, 1.0, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Idea:",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 1.80, 11.8, 0.4,
         "Put a filter F_e(s) BEFORE G(s). If the zeros of F_e match "
         "the poles of G, they kill each other → no ringing.",
         size=14, color=BLACK)

# left — F_e structure
box(s, 0.5, 2.45, 6.2, 4.5, fill=RED_BG, border=RED)
text_box(s, 0.7, 2.55, 5.8, 0.5,
         "F_e(s) — exponential shaper",
         size=16, bold=True, color=RED)
text_box(s, 0.7, 3.10, 5.8, 1.0,
         "F_e(s)  =  (1 − e^{−sT_e})  /  ( T_e (s − μ) )",
         size=15, bold=True, color=BLACK, font='Consolas')

text_box(s, 0.7, 4.10, 5.8, 0.5, "Recipe (Choi et al., Sec. 4.2):",
         size=14, bold=True, color=RED)
text_box(s, 0.7, 4.55, 5.8, 1.7,
         ["• T_e   =   2π / ω_d        (one damped oscillation period)",
          "• μ    =   −δ ω_f           (matches G's real-part decay)",
          "",
          "These two choices place F_e's zeros exactly at",
          "  s = −δω_f ± j n ω_d   for  n = ±1, ±2, …"],
         size=12, color=BLACK, font='Consolas')

text_box(s, 0.7, 6.30, 5.8, 0.5,
         "→ F_e zeros at n = ±1 land on G(s)'s poles.",
         size=13, bold=True, color=RED)

# right — block diagram & cancellation visual
box(s, 7.0, 2.45, 6.0, 4.5, fill=GREEN_BG, border=GREEN)
text_box(s, 7.2, 2.55, 5.6, 0.5,
         "Pole-zero cancellation",
         size=16, bold=True, color=GREEN)

# blocks
box(s, 7.4, 3.20, 1.3, 0.7, fill=WHITE, border=GREEN)
text_box(s, 7.4, 3.30, 1.3, 0.5, "v_ref",
         size=14, bold=True, color=BLACK, align=CTR)
arrow(s, 8.7, 3.55, 9.0, 3.55, color=GREEN, weight=2)
box(s, 9.0, 3.20, 1.4, 0.7, fill=WHITE, border=RED)
text_box(s, 9.0, 3.30, 1.4, 0.5, "F_e(s)",
         size=14, bold=True, color=RED, align=CTR)
arrow(s, 10.4, 3.55, 10.8, 3.55, color=GREEN, weight=2)
box(s, 10.8, 3.20, 1.4, 0.7, fill=WHITE, border=BLUE)
text_box(s, 10.8, 3.30, 1.4, 0.5, "G(s)",
         size=14, bold=True, color=BLUE, align=CTR)
arrow(s, 12.2, 3.55, 12.6, 3.55, color=GREEN, weight=2)
text_box(s, 12.6, 3.30, 0.6, 0.5, "θ",
         size=14, bold=True, color=BLACK)

text_box(s, 7.2, 4.20, 5.6, 0.5,
         "Combined transfer:",
         size=13, color=BLACK)
text_box(s, 7.2, 4.60, 5.6, 0.5,
         "F_e(s) · G(s)  →  poles & zeros CANCEL  →  no ringing",
         size=13, bold=True, color=GREEN, font='Consolas')

text_box(s, 7.2, 5.30, 5.6, 1.5,
         ["• G(s)'s poles at  −δω_f ± j ω_d  are EXACTLY",
          "  the n=1 zeros of F_e(s).",
          "• They cancel out algebraically.",
          "• The resonance the liquid would amplify is",
          "  removed BEFORE it ever reaches the liquid."],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 8 — Step 6: apply F_e(s) to v_ref
# =============================================================================
s = blank_slide()
title_bar(s, "Step 6: Apply  F_e(s)  to the velocity reference",
          "What the user sees: a step command becomes a smooth ramp")

# Inverse-Laplace explanation
box(s, 0.6, 1.3, 12.1, 1.0, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "The shaped signal in time:",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.80, 11.8, 0.4,
         "If v_ref is a step from 0 to V, then F_e(s)·v_ref(t) gives an "
         "exponential ramp that finishes after exactly  T_e  seconds.",
         size=13, color=BLACK)

# Left — closed-form formula
box(s, 0.5, 2.45, 6.2, 4.5, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 2.55, 5.8, 0.5,
         "Closed-form shaped step",
         size=16, bold=True, color=BLUE)
text_box(s, 0.7, 3.05, 5.8, 1.6,
         "             ⌈  V · ( 1 − e^{μ t} ) / ( 1 − e^{μ T_e} ),\n"
         "v_shaped(t) =                                    0 ≤ t ≤ T_e\n"
         "             |  V,                                t > T_e\n"
         "             ⌊  0,                                t < 0",
         size=11, color=BLACK, font='Consolas')

text_box(s, 0.7, 4.80, 5.8, 0.5, "Read it as:",
         size=13, bold=True, color=BLUE)
text_box(s, 0.7, 5.20, 5.8, 1.7,
         ["• 0 → T_e: smooth exponential rise to V",
          "• after T_e: stays at V forever",
          "• never overshoots V",
          "• has no high-frequency content above 1/T_e",
          "  → won't shake the liquid"],
         size=12, color=BLACK)

# Right — sketch
box(s, 7.0, 2.45, 6.0, 4.5, fill=GREEN_BG, border=GREEN)
text_box(s, 7.2, 2.55, 5.6, 0.5,
         "What the signal looks like",
         size=16, bold=True, color=GREEN)

# axes
ax_l, ax_t, ax_w, ax_h = 7.4, 3.10, 5.5, 3.55
# x-axis
line_(s, ax_l + 0.5, ax_t + ax_h - 0.6, ax_l + ax_w - 0.2,
      ax_t + ax_h - 0.6, color=BLACK, weight=1.5)
# y-axis
line_(s, ax_l + 0.5, ax_t + 0.2, ax_l + 0.5, ax_t + ax_h - 0.6,
      color=BLACK, weight=1.5)
text_box(s, ax_l + ax_w - 0.5, ax_t + ax_h - 0.55, 0.5, 0.4, "t",
         size=12, color=BLACK)
text_box(s, ax_l + 0.20, ax_t + 0.0, 0.5, 0.4, "v",
         size=12, color=BLACK)

# raw step (red, vertical jump)
y_top = ax_t + 0.4
y_bot = ax_t + ax_h - 0.6
x0 = ax_l + 0.9
line_(s, x0, y_bot, x0, y_top, color=RED, weight=3)
line_(s, x0, y_top, ax_l + ax_w - 0.4, y_top, color=RED, weight=3)
text_box(s, ax_l + ax_w - 2.4, y_top - 0.4, 2.0, 0.4,
         "raw step v_ref", size=11, bold=True, color=RED)

# shaped ramp (green)
# build a rough exponential rise polyline
import math as _m
n_pts = 12
T_e_x = 1.6  # inches along the x-axis representing T_e
mu_param = -0.6
prev_x, prev_y = x0, y_bot
for i in range(1, n_pts + 1):
    frac = i / n_pts
    t = frac * T_e_x
    # normalized rise (approx)
    val = (1 - _m.exp(mu_param * t)) / (1 - _m.exp(mu_param * T_e_x))
    x = x0 + t
    y = y_bot - val * (y_bot - y_top)
    line_(s, prev_x, prev_y, x, y, color=GREEN, weight=3.5)
    prev_x, prev_y = x, y
# hold flat after T_e
line_(s, prev_x, prev_y, ax_l + ax_w - 0.4, prev_y, color=GREEN, weight=3.5)
text_box(s, ax_l + 1.0, y_top + 0.55, 4.0, 0.4,
         "shaped v_ref  ↑  (with F_e)",
         size=11, bold=True, color=GREEN)

# T_e marker
line_(s, x0 + T_e_x, y_bot, x0 + T_e_x, y_bot - 0.15,
      color=GRAY, weight=1.5)
text_box(s, x0 + T_e_x - 0.3, y_bot + 0.05, 1.0, 0.35,
         "T_e", size=12, bold=True, color=GRAY, align=CTR)


# =============================================================================
# SLIDE 9 — Full pipeline summary
# =============================================================================
s = blank_slide()
title_bar(s, "Putting it all together — the full pipeline",
          "From physics to a working anti-sloshing reference")

# Big horizontal pipeline
y_top = 1.4
box_w, box_h = 1.6, 1.0
xs = [0.6, 2.5, 4.4, 6.3, 8.2, 10.1, 12.0]
labels = [
    ("Newton\nFBD",          ORANGE,  ORANGE_BG),
    ("Linearize\n(small θ)", ORANGE,  ORANGE_BG),
    ("Laplace\nd/dt → s",    RED,     RED_BG),
    ("Plant\nG(s)",          BLUE,    BLUE_BG),
    ("Shaper\nF_e(s)",       PURPLE,  PURPLE_BG),
    ("Apply to\nv_ref",      GREEN,   GREEN_BG),
    ("No\nsloshing!",        TEAL,    TEAL_BG),
]
for i, (lab, c, bg) in enumerate(labels):
    box(s, xs[i] - 0.1, y_top, box_w, box_h, fill=bg, border=c, bpt=2.0)
    text_box(s, xs[i] - 0.1, y_top + 0.12, box_w, box_h - 0.2,
             lab, size=12, bold=True, color=c, align=CTR)
    if i < len(labels) - 1:
        arrow(s, xs[i] + box_w - 0.1, y_top + box_h / 2,
              xs[i + 1] - 0.1, y_top + box_h / 2,
              color=GRAY, weight=2)

# Three "what we got" boxes
y_b = 2.9
box(s, 0.6, y_b, 4.0, 4.1, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, y_b + 0.10, 3.7, 0.5,
         "Physics → math",
         size=15, bold=True, color=BLUE)
text_box(s, 0.8, y_b + 0.55, 3.7, 3.4,
         ["1. Newton FBD on the",
          "   liquid pendulum.",
          "",
          "2. Linearize for",
          "   small angles.",
          "",
          "3. Get a clean LTI",
          "   2nd-order ODE."],
         size=12, color=BLACK)

box(s, 4.7, y_b, 4.0, 4.1, fill=PURPLE_BG, border=PURPLE)
text_box(s, 4.9, y_b + 0.10, 3.7, 0.5,
         "Math → design",
         size=15, bold=True, color=PURPLE)
text_box(s, 4.9, y_b + 0.55, 3.7, 3.4,
         ["4. Laplace turns the",
          "   ODE into G(s).",
          "",
          "5. Read off poles",
          "   −δω_f ± j ω_d .",
          "",
          "6. Build F_e(s) whose",
          "   zeros land on those",
          "   poles."],
         size=12, color=BLACK)

box(s, 8.8, y_b, 4.0, 4.1, fill=GREEN_BG, border=GREEN)
text_box(s, 9.0, y_b + 0.10, 3.7, 0.5,
         "Design → behavior",
         size=15, bold=True, color=GREEN)
text_box(s, 9.0, y_b + 0.55, 3.7, 3.4,
         ["7. Shape v_ref through",
          "   F_e(s) — step becomes",
          "   smooth ramp of length T_e.",
          "",
          "8. Liquid sees no",
          "   resonant input.",
          "",
          "9. No overshoot,",
          "   no ringing,",
          "   no spilling."],
         size=12, color=BLACK)


# ----- save ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'reference_shaping_progress.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  {len(prs.slides)} slides')
