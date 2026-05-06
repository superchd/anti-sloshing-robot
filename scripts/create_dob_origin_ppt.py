#!/usr/bin/env python3
"""
create_dob_origin_ppt.py
Slide deck answering the question:

   "Did the author invent Eq. (6)?  --  No, it's a classic recipe."

Walks through:
  - the short answer (No, [13-15])
  - where each piece of Eq. (6) comes from
  - the three goals the design satisfies
  - the author's actual contribution
  - why Eq. (6) is needed at all

Run:    python3 scripts/create_dob_origin_ppt.py
Output: results/dob_origin_explained.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)
PURPLE_BG= rgb(0xF5, 0xEC, 0xFB)
TEAL     = rgb(0x00, 0x7A, 0x86)
TEAL_BG  = rgb(0xE0, 0xF2, 0xF1)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE,
        bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.1, 11.7, 1.1,
         "Did the author invent Eq. (6)?",
         size=46, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.3, 11.7, 1.4,
         "No — it's a classic DOB recipe from the 1980s.\n"
         "The novelty is the WAY they use it.",
         size=22, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 5.5, 11.7, 0.5,
         "Choi et al. 2024 · Section 4.5",
         size=20, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.5, 11.7, 0.5,
         "MECE 6397   ·   Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2 — Short answer
# =============================================================================
s = blank_slide()
title_bar(s, "Short answer:  NO  —  the authors did not invent Eq. (6)",
          "It is a standard DOB design from control theory")

# Big "NO" callout
box(s, 2.0, 1.5, 9.3, 1.6, fill=RED_BG, border=RED, bpt=3)
text_box(s, 2.0, 1.65, 9.3, 0.7,
         "NO",
         size=72, bold=True, color=RED, align=CTR)
text_box(s, 2.0, 2.55, 9.3, 0.5,
         "(they're using it, not inventing it)",
         size=18, color=GRAY, align=CTR)

# Citation evidence
box(s, 0.6, 3.45, 12.1, 1.35, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 3.55, 11.8, 0.4,
         "The paper itself tells you so",
         size=15, bold=True, color=ORANGE)
text_box(s, 0.8, 3.95, 11.8, 0.85,
         "Right before Eq. (6) the paper says:\n"
         "    \"...the following DOB is considered [13–15]\"",
         size=15, color=BLACK, font='Consolas')

# History note
box(s, 0.6, 4.95, 12.1, 2.0, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 5.05, 11.8, 0.4,
         "Where it really comes from",
         size=15, bold=True, color=BLUE)
text_box(s, 0.8, 5.45, 11.8, 1.5,
         ["• DOB framework dates back to the 1980s (Ohnishi, Umeno, Hori)",
          "• Same family as: Kalman filter, Luenberger observer,",
          "  exponential moving average, simple PI integrator",
          "• Used in robotics, motors, drones, hard disks for ~40 years"],
         size=13, color=BLACK)


# =============================================================================
# SLIDE 3 — Where each piece comes from
# =============================================================================
s = blank_slide()
title_bar(s, "Where each piece of Eq. (6) comes from",
          "Two recycled ideas glued together")

# Big formula
box(s, 0.6, 1.3, 12.1, 1.2, fill=PURPLE_BG, border=PURPLE, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.5,
         "d̂̇_t   =   η  (  d_t  −  d̂_t  )    =    η ( q̇_t − A q_t − B u_t  −  d̂_t )",
         size=18, bold=True, color=BLACK, font='Consolas', align=CTR)
text_box(s, 0.8, 1.95, 11.8, 0.4,
         "Two pieces. Neither is new.",
         size=13, color=GRAY, align=CTR)

# Two-column breakdown
box(s, 0.5, 2.65, 6.2, 4.5, fill=ORANGE_BG, border=ORANGE)
text_box(s, 0.7, 2.75, 5.8, 0.5,
         "Piece 1:  the underbraced part",
         size=15, bold=True, color=ORANGE)
text_box(s, 0.7, 3.20, 5.8, 0.5,
         "d_t  =  q̇_t  −  A q_t  −  B u_t",
         size=15, bold=True, color=BLACK, font='Consolas')
text_box(s, 0.7, 3.75, 5.8, 3.3,
         ["From Eq. (1) of the paper:",
          "    q̇_t  =  A q_t + B u_t + d_t",
          "",
          "Just rearrange.",
          "",
          "Not a design choice —",
          "literally the DEFINITION of d flipped:",
          "",
          "  'disturbance = what the model",
          "                  cannot explain.'"],
         size=12, color=BLACK)

box(s, 6.9, 2.65, 6.0, 4.5, fill=GREEN_BG, border=GREEN)
text_box(s, 7.1, 2.75, 5.6, 0.5,
         "Piece 2:  the outer chasing part",
         size=15, bold=True, color=GREEN)
text_box(s, 7.1, 3.20, 5.6, 0.5,
         "d̂̇  =  η ( d  −  d̂ )",
         size=15, bold=True, color=BLACK, font='Consolas')
text_box(s, 7.1, 3.75, 5.6, 3.3,
         ["A first-order low-pass filter:",
          "  • d̂ < d  →  d̂ rises",
          "  • d̂ > d  →  d̂ falls",
          "  • speed = η",
          "",
          "Same shape as:",
          "  • thermostat",
          "  • Kalman filter (simplest case)",
          "  • PI controller integrator",
          "  • exponential moving average"],
         size=12, color=BLACK)


# =============================================================================
# SLIDE 4 — Three goals it satisfies
# =============================================================================
s = blank_slide()
title_bar(s, "Why this specific design? — three goals",
          "Minimum complexity that still works")

# 3 goal cards
y = 1.5
goals = [
    ("Track  d  accurately",
     "The  (d − d̂)  feedback drives the error to zero in steady state.",
     BLUE),
    ("Smooth out noise",
     "Low-pass behaviour — η sets the bandwidth (high η = fast, low η = smooth).",
     PURPLE),
    ("Stay simple — 1 knob",
     "Only one parameter, η, to tune. No covariance matrices, no gain scheduling.",
     GREEN),
]
for i, (g, why, c) in enumerate(goals):
    yi = y + i * 1.4
    box(s, 0.6, yi, 12.1, 1.2, fill=WHITE, border=c, bpt=2.5)
    box(s, 0.6, yi, 1.2, 1.2, fill=c, border=c, bpt=0)
    text_box(s, 0.6, yi + 0.30, 1.2, 0.6, str(i + 1),
             size=42, bold=True, color=WHITE, align=CTR)
    text_box(s, 2.0, yi + 0.20, 10.5, 0.5, g,
             size=18, bold=True, color=c)
    text_box(s, 2.0, yi + 0.65, 10.5, 0.5, why,
             size=13, color=BLACK)

# bottom note
box(s, 0.6, 6.0, 12.1, 1.2, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 6.10, 11.8, 0.4,
         "Why pick THIS over a Kalman filter or sliding-mode observer?",
         size=13, bold=True, color=ORANGE)
text_box(s, 0.8, 6.50, 11.8, 0.7,
         "Minimum complexity, proven convergence (Appendix A), and good "
         "enough for the disturbances this robot actually faces.",
         size=12, color=BLACK)


# =============================================================================
# SLIDE 5 — Author's actual contribution
# =============================================================================
s = blank_slide()
title_bar(s, "So what IS the author's contribution?",
          "Not Eq. (6) itself — the way they combine it")

# Big idea
box(s, 0.6, 1.3, 12.1, 1.2, fill=PURPLE_BG, border=PURPLE, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Analogy",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 1.80, 11.8, 0.7,
         "A chef didn't invent salt. The recipe that puts salt + butter "
         "+ lemon in just the right way is the contribution.",
         size=13, color=BLACK)

# Three numbered contributions
y = 2.75
contribs = [
    ("Apply DOB to a sloshing robot",
     "Not done before in this exact food-serving / sloshing context.", ORANGE),
    ("Pair it with the auxiliary compensator",
     "Aux  → big sudden kicks (caster).   DOB  → small continuous drift.\n"
     "Two robustness blocks splitting the work by timescale.", GREEN),
    ("Combine everything with LQT + F_e(s) shaping",
     "All four pieces glued into one framework (Eq. 8 / Fig. 1).\n"
     "That integrated design is the paper's real novelty.", BLUE),
]
for i, (head, body, c) in enumerate(contribs):
    yi = y + i * 1.35
    box(s, 0.6, yi, 12.1, 1.15, fill=WHITE, border=c, bpt=2.0)
    box(s, 0.6, yi, 0.95, 1.15, fill=c, border=c, bpt=0)
    text_box(s, 0.6, yi + 0.27, 0.95, 0.6, str(i + 1),
             size=36, bold=True, color=WHITE, align=CTR)
    text_box(s, 1.75, yi + 0.12, 10.7, 0.5, head,
             size=15, bold=True, color=c)
    text_box(s, 1.75, yi + 0.55, 10.7, 0.55, body,
             size=12, color=BLACK)


# =============================================================================
# SLIDE 6 — Why make Eq. (6) at all?
# =============================================================================
s = blank_slide()
title_bar(s, "Why bother with Eq. (6) in this paper at all?",
          "Because LQT alone trusts a model that the real robot doesn't obey")

# left: what LQT assumes
box(s, 0.5, 1.3, 6.2, 5.6, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 1.4, 5.8, 0.5,
         "LQT (Sec 4.3) assumes...",
         size=15, bold=True, color=BLUE)
text_box(s, 0.7, 1.95, 5.8, 5.0,
         ["• A clean linear model:",
          "    q̇  =  A q  +  B u",
          "",
          "• No uncertainties.",
          "",
          "• No external impulses.",
          "",
          "• No friction or model mismatch.",
          "",
          "Riccati gives an OPTIMAL gain K",
          "for THIS idealized system."],
         size=13, color=BLACK)

# right: what the real robot has
box(s, 6.9, 1.3, 6.0, 5.6, fill=RED_BG, border=RED)
text_box(s, 7.1, 1.4, 5.6, 0.5,
         "Real robot has...",
         size=15, bold=True, color=RED)
text_box(s, 7.1, 1.95, 5.6, 5.0,
         ["• Caster wheel kicks",
          "• Friction (caster, drive)",
          "• Floor irregularities, slope",
          "• Modeling error  ΔU  (Eq. 5)",
          "• Payload changes (food added)",
          "",
          "→  These all violate the LQT model.",
          "→  Without compensation:",
          "      drift, instability, spilling."],
         size=13, color=BLACK)

# bottom takeaway
box(s, 0.6, 7.0, 12.1, 0.4, fill=GREEN_BG, border=GREEN)
text_box(s, 0.6, 7.0, 12.1, 0.4,
         "DOB's whole purpose:  estimate that 'extra stuff' d  →  cancel it  "
         "→  LQT sees a clean model again.",
         size=12, bold=True, color=GREEN, align=CTR)


# =============================================================================
# SLIDE 7 — TL;DR
# =============================================================================
s = blank_slide()
title_bar(s, "TL;DR — five lines",
          "What to remember when you cite Eq. (6)")

y = 1.5
items = [
    ("Eq. (6) is borrowed",
     "not invented — references [13–15] in the paper",                BLUE),
    ("Underbraced part",
     "rearranged Eq. (1) → 'disturbance is what the model can't explain'",
                                                                       ORANGE),
    ("Outer part",
     "standard low-pass tracker → 'smoothly chase the truth'",         PURPLE),
    ("Author's contribution",
     "combining DOB with aux compensator + LQT + shaping",             RED),
    ("Why use it",
     "make the messy real robot behave like the clean nominal model",  GREEN),
]
for i, (k, v, c) in enumerate(items):
    yi = y + i * 0.95
    box(s, 0.6, yi, 12.1, 0.8, fill=WHITE, border=c, bpt=2.0)
    box(s, 0.6, yi, 3.0, 0.8, fill=c, border=c, bpt=0)
    text_box(s, 0.7, yi + 0.20, 2.9, 0.5, k,
             size=14, bold=True, color=WHITE, align=CTR)
    text_box(s, 3.7, yi + 0.20, 9.0, 0.5, v,
             size=13, color=BLACK)


# ----- save ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'dob_origin_explained.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  {len(prs.slides)} slides')
