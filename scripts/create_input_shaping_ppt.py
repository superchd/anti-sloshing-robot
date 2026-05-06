#!/usr/bin/env python3
"""
create_input_shaping_ppt.py
Plain-English slide deck explaining F_e(s) (the reference-shaping filter)
and T_e (its duration parameter) from Choi et al. (2024), Section 4.2.

Run:    python3 scripts/create_input_shaping_ppt.py
Output: results/input_shaping_explained.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


# ----- palette ---------------------------------------------------------------
WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


# ----- presentation ----------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # add arrow head via xml
    from pptx.oxml.ns import qn
    from lxml import etree
    line = conn.line._get_or_add_ln()
    head = etree.SubElement(line, qn('a:headEnd'))
    tail = etree.SubElement(line, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE, bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


# =============================================================================
# SLIDE 1: Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.4, 11.7, 1.2,
         "Input Shaping for Anti-Sloshing",
         size=54, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.7, 11.7, 0.8,
         r"What  F_e(s)  and  T_e  actually do",
         size=30, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 5.8, 11.7, 0.5,
         "Based on Choi et al. 2024, Section 4.2",
         size=18, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.4, 11.7, 0.5,
         "MECE 6397  ·  Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2: The Problem
# =============================================================================
s = blank_slide()
title_bar(s, "The problem: step commands excite sloshing",
          "Instant velocity changes kick the liquid at every frequency")

# Left: step command + sloshing result
text_box(s, 0.8, 1.4, 5.5, 0.5, "What you command:", size=18, bold=True)
box(s, 0.8, 2.0, 5.5, 1.6, fill=YELLOW_BG, border=ORANGE)
text_box(s, 1.0, 2.25, 5.1, 0.5,
         "v_ref  =  step from 0 → 0.6 m/s", size=18, bold=True)
text_box(s, 1.0, 2.85, 5.1, 0.6,
         "A sudden jump at t=0.\nContains energy at ALL frequencies.",
         size=15, color=GRAY)

text_box(s, 0.8, 3.9, 5.5, 0.5, "What the liquid does:", size=18, bold=True,
         color=RED)
box(s, 0.8, 4.5, 5.5, 2.3, fill=RED_BG, border=RED)
text_box(s, 1.0, 4.75, 5.1, 2.0,
         ["• Huge overshoot (up to 70%)",
          "• Rings for 5+ seconds",
          "• Spills over the edge",
          "• Robot rocks / tips"],
         size=16, color=BLACK)

# Right: the insight
text_box(s, 7.2, 1.4, 5.6, 0.5, "Why does this happen?", size=18, bold=True,
         color=PURPLE)
box(s, 7.2, 2.0, 5.6, 4.8, fill=rgb(0xF5, 0xEC, 0xFB), border=PURPLE)
text_box(s, 7.5, 2.2, 5.0, 0.5,
         "The liquid has a natural frequency ω_f.",
         size=16, bold=True)
text_box(s, 7.5, 2.75, 5.0, 0.5,
         "ω_f ≈ 9.92 rad/s  (about 1.58 Hz)",
         size=15, color=GRAY)
text_box(s, 7.5, 3.4, 5.0, 0.5,
         "A step signal contains energy at EVERY frequency —",
         size=15)
text_box(s, 7.5, 3.85, 5.0, 0.5,
         "including ω_f. That energy drives the resonance.",
         size=15)
text_box(s, 7.5, 4.7, 5.0, 0.5, "The fix:", size=18, bold=True, color=GREEN)
text_box(s, 7.5, 5.15, 5.0, 1.5,
         ["Reshape the reference BEFORE it",
          "reaches the controller, so it has",
          "no energy at ω_f."],
         size=16, color=BLACK)


# =============================================================================
# SLIDE 3: Model the liquid as a 2nd-order system
# =============================================================================
s = blank_slide()
title_bar(s, "Step 1: model the liquid as a pendulum",
          "Sloshing = mass + restoring force + damping  =  2nd-order system")

# Big equation center-top
box(s, 2.0, 1.3, 9.3, 1.4, fill=BLUE_BG, border=BLUE)
text_box(s, 2.0, 1.45, 9.3, 0.45,
         "Transfer function of the sloshing dynamics:",
         size=14, color=GRAY, align=CTR)
text_box(s, 2.0, 1.85, 9.3, 0.9,
         "G(s)  =  ω_f²  /  (s² + 2 δ ω_f s + ω_f²)",
         size=30, bold=True, align=CTR, color=BLUE,
         font='Cambria Math')

# Three pieces underneath
y = 3.2
for i, (term, name, desc, col) in enumerate([
    ("ω_f", "natural frequency",
     "How fast the liquid oscillates on its own.\nSet by cup size + gravity.", BLUE),
    ("δ", "damping ratio",
     "How fast the oscillation dies out.\nWater has small δ → rings a lot.", ORANGE),
    ("s²", "inertia",
     "The liquid has mass — it can't change\nvelocity instantly.", GREEN),
]):
    x = 0.7 + i * 4.3
    box(s, x, y, 3.9, 3.3, fill=WHITE, border=col, bpt=2)
    text_box(s, x, y + 0.2, 3.9, 0.8, term, size=42, bold=True, color=col,
             align=CTR, font='Cambria Math')
    text_box(s, x, y + 1.3, 3.9, 0.5, name, size=18, bold=True, align=CTR)
    text_box(s, x, y + 2.0, 3.9, 1.2, desc, size=13, color=GRAY, align=CTR)

text_box(s, 0.8, 6.8, 11.7, 0.5,
         "Same math as a pendulum, a mass on a spring, or an RLC circuit.",
         size=14, color=GRAY, align=CTR)


# =============================================================================
# SLIDE 4: What F_e(s) is
# =============================================================================
s = blank_slide()
title_bar(s, "F_e(s):  a filter that smooths the command",
          "It sits in FRONT of the controller — pure feedforward, no sensors needed")

# Pipeline diagram
y = 1.7
text_box(s, 0.5, y, 2.3, 0.5, "v_ref (step)", size=16, bold=True, align=CTR,
         color=BLUE)
box(s, 3.0, y, 2.2, 0.9, fill=RED_BG, border=RED, bpt=2.5)
text_box(s, 3.0, y + 0.12, 2.2, 0.6, "F_e(s)", size=28, bold=True, align=CTR,
         color=RED, font='Cambria Math')
text_box(s, 5.4, y, 2.3, 0.5, "v_d (smooth)", size=16, bold=True, align=CTR,
         color=ORANGE)
box(s, 7.9, y, 3.1, 0.9, fill=BLUE_BG, border=BLUE, bpt=2.5)
text_box(s, 7.9, y + 0.12, 3.1, 0.6, "LQT + aux + DOB",
         size=20, bold=True, align=CTR, color=BLUE)
text_box(s, 11.2, y, 1.7, 0.5, "robot", size=16, bold=True, align=CTR,
         color=GREEN)

arrow(s, 2.8, y + 0.45, 3.0, y + 0.45)
arrow(s, 5.2, y + 0.45, 5.4, y + 0.45)
arrow(s, 7.6, y + 0.45, 7.9, y + 0.45)
arrow(s, 11.0, y + 0.45, 11.2, y + 0.45)

text_box(s, 3.0, y + 1.1, 2.2, 0.4, "shaper",
         size=12, color=RED, align=CTR)
text_box(s, 7.9, y + 1.1, 3.1, 0.4, "feedback controller",
         size=12, color=BLUE, align=CTR)

# The filter formula
box(s, 1.0, 3.3, 11.3, 1.3, fill=YELLOW_BG, border=ORANGE)
text_box(s, 1.0, 3.4, 11.3, 0.4,
         "The exponential filter formula (Eq. 2 in the paper):",
         size=13, color=GRAY, align=CTR)
text_box(s, 1.0, 3.75, 11.3, 0.8,
         "F_e(s) = [ μ / (e^(μT_e) − 1) ]  ·  [ (1 − e^(μT_e) · e^(−T_e·s)) / (s − μ) ]",
         size=22, bold=True, align=CTR, font='Cambria Math')

# What does it do
text_box(s, 0.8, 4.9, 11.7, 0.5, "What it actually does (time domain):",
         size=18, bold=True)
box(s, 0.8, 5.5, 5.8, 1.7, fill=RED_BG, border=RED)
text_box(s, 1.0, 5.65, 5.6, 0.5, "Input: step command", size=14, bold=True,
         color=GRAY)
text_box(s, 1.0, 6.05, 5.6, 1.0,
         "─┐\n │\n─┘",
         size=22, color=RED, font='Consolas')

box(s, 7.0, 5.5, 5.8, 1.7, fill=GREEN_BG, border=GREEN)
text_box(s, 7.2, 5.65, 5.6, 0.5, "Output: smooth exponential ramp",
         size=14, bold=True, color=GRAY)
text_box(s, 7.2, 6.05, 5.6, 1.0,
         "       ────\n     ╱\n____╱",
         size=22, color=GREEN, font='Consolas')


# =============================================================================
# SLIDE 5: What T_e is
# =============================================================================
s = blank_slide()
title_bar(s, "T_e:  how long the ramp takes",
          "One number, with one magic value that makes everything work")

# Definition box
box(s, 0.8, 1.4, 11.7, 1.4, fill=BLUE_BG, border=BLUE)
text_box(s, 1.0, 1.55, 11.3, 0.5,
         "T_e  =  the duration over which F_e(s) stretches the step command.",
         size=20, bold=True, align=CTR)
text_box(s, 1.0, 2.1, 11.3, 0.6,
         "If T_e = 0.5 s, then a step at t=0 becomes a smooth ramp reaching full value at t=0.5 s.",
         size=15, color=GRAY, align=CTR)

# The magic formula
box(s, 0.8, 3.1, 11.7, 1.6, fill=YELLOW_BG, border=ORANGE, bpt=3)
text_box(s, 1.0, 3.2, 11.3, 0.5,
         "The paper's design choice:",
         size=15, color=GRAY, align=CTR)
text_box(s, 1.0, 3.65, 11.3, 0.6,
         "T_e  =  2π / (ω_f · √(1 − δ²))",
         size=30, bold=True, align=CTR, color=ORANGE, font='Cambria Math')
text_box(s, 1.0, 4.3, 11.3, 0.5,
         "→  T_e  =  one full period of the liquid's natural sloshing",
         size=17, bold=True, align=CTR, color=RED)

# Numerical example
text_box(s, 0.8, 5.0, 11.7, 0.5, "Plugging in the paper's numbers:",
         size=17, bold=True)
box(s, 0.8, 5.6, 5.7, 1.5, fill=WHITE, border=BLUE, bpt=2)
text_box(s, 1.0, 5.75, 5.3, 1.3,
         ["ω_f  =  9.922 rad/s",
          "δ    =  0.05",
          "──────────────────",
          "T_e  ≈  0.63 s"],
         size=16, font='Consolas')

box(s, 6.8, 5.6, 5.9, 1.5, fill=GREEN_BG, border=GREEN, bpt=2)
text_box(s, 7.0, 5.75, 5.5, 1.3,
         ["Meaning:",
          "Every step command becomes",
          "a smooth ramp over ~0.63 s.",
          "That's one slosh period."],
         size=15)


# =============================================================================
# SLIDE 6: Why it works — pole-zero cancellation
# =============================================================================
s = blank_slide()
title_bar(s, "Why T_e = 2π/ω_d works: pole–zero cancellation",
          "The filter's zeros sit exactly on the plant's resonance")

# Left: intuition
box(s, 0.5, 1.4, 6.0, 5.8, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 1.55, 5.6, 0.5, "The liquid is like a bell",
         size=20, bold=True, color=BLUE)
text_box(s, 0.7, 2.1, 5.6, 5.0,
         ["• Every system has poles — frequencies",
          "  where it naturally wants to ring.",
          "",
          "• Zeros are the opposite — frequencies",
          "  the system refuses to respond to.",
          "",
          "• If you place a ZERO exactly on top of",
          "  a POLE, they CANCEL. The ringing",
          "  mode is mathematically deleted.",
          "",
          "• That's what T_e = 2π/ω_d does —",
          "  creates zeros right where the",
          "  plant's resonance poles live."],
         size=14, color=BLACK)

# Right: pole-zero sketch
box(s, 7.0, 1.4, 5.9, 5.8, fill=WHITE, border=GRAY)

# axes
cx, cy = 9.95, 4.3
ax_len = 2.2

arrow(s, cx - ax_len, cy, cx + ax_len, cy, color=BLACK, weight=1.5)
arrow(s, cx, cy + ax_len, cx, cy - ax_len, color=BLACK, weight=1.5)
text_box(s, cx + ax_len + 0.05, cy - 0.2, 0.4, 0.3, "Re", size=11, color=GRAY)
text_box(s, cx - 0.35, cy - ax_len - 0.3, 0.4, 0.3, "Im", size=11, color=GRAY)

# plant poles as X
px = cx - 0.5
for py in [cy - 1.0, cy + 1.0]:
    shp1 = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                  Inches(px - 0.12), Inches(py - 0.12),
                                  Inches(px + 0.12), Inches(py + 0.12))
    shp1.line.color.rgb = BLUE; shp1.line.width = Pt(3.5)
    shp2 = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                  Inches(px - 0.12), Inches(py + 0.12),
                                  Inches(px + 0.12), Inches(py - 0.12))
    shp2.line.color.rgb = BLUE; shp2.line.width = Pt(3.5)

# Fe zeros as circles
for n_off in [-2.0, -1.0, 0.0, 1.0, 2.0]:
    py = cy + n_off * 1.0
    zshp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(px - 0.15), Inches(py - 0.15),
                              Inches(0.3), Inches(0.3))
    zshp.fill.background()
    zshp.line.color.rgb = RED
    zshp.line.width = Pt(2.5)

# Vertical dashed guide
guide = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                               Inches(px), Inches(cy - ax_len + 0.1),
                               Inches(px), Inches(cy + ax_len - 0.1))
guide.line.color.rgb = RED
guide.line.width = Pt(1)
guide.line.dash_style = 7  # dash

# Labels
text_box(s, 7.1, 2.5, 2.5, 0.4, "× plant poles",
         size=13, bold=True, color=BLUE)
text_box(s, 7.1, 2.9, 2.5, 0.4, "○ F_e(s) zeros",
         size=13, bold=True, color=RED)

# callout
text_box(s, 10.6, 2.95, 2.3, 1.0,
         ["n=1 zeros land",
          "exactly on the",
          "plant poles.",
          "→ CANCEL"],
         size=11, color=BLACK, bold=True)


# =============================================================================
# SLIDE 7: The result
# =============================================================================
s = blank_slide()
title_bar(s, "The result: same target, no sloshing",
          "Fig. 2 in the paper — step response of G(s) with three different inputs")

# Insert the diagram we already generated
diag = os.path.join(os.path.dirname(__file__), '..', 'results',
                    'input_shaping_diagram.png')
if os.path.exists(diag):
    s.shapes.add_picture(diag, Inches(0.5), Inches(1.2),
                         width=Inches(8.5), height=Inches(5.3))

# Right-side takeaway
box(s, 9.3, 1.3, 3.7, 5.3, fill=GREEN_BG, border=GREEN)
text_box(s, 9.5, 1.45, 3.3, 0.5, "Takeaways",
         size=22, bold=True, color=GREEN)
text_box(s, 9.5, 2.05, 3.3, 4.5,
         ["Raw step:",
          "  huge overshoot, rings",
          "  for 5+ seconds.",
          "",
          "Generic LPF:",
          "  better, but still",
          "  oscillates.",
          "",
          "Input shaping (F_e):",
          "  no overshoot,",
          "  no ringing,",
          "  settles in ~1 s.",
          "",
          "All three reach 1.0 —",
          "only F_e gets there",
          "WITHOUT exciting the",
          "sloshing mode."],
         size=13, color=BLACK)

# Bottom quote box
box(s, 0.5, 6.7, 12.3, 0.65, fill=BLUE_BG, border=BLUE)
text_box(s, 0.5, 6.78, 12.3, 0.5,
         "Key insight:  smoothing helps, but only input shaping places a zero ON the resonance pole — so the mode is eliminated, not just attenuated.",
         size=14, bold=True, color=BLUE, align=CTR)


# =============================================================================
# SLIDE 8: Summary
# =============================================================================
s = blank_slide()
title_bar(s, "Summary: F_e and T_e in one slide")

# Two big columns
box(s, 0.7, 1.4, 5.9, 5.5, fill=RED_BG, border=RED, bpt=2.5)
text_box(s, 0.7, 1.6, 5.9, 0.7,
         "F_e(s)",
         size=44, bold=True, color=RED, align=CTR, font='Cambria Math')
text_box(s, 0.9, 2.4, 5.5, 0.5,
         "A function (a transfer function).",
         size=16, bold=True, align=CTR)
text_box(s, 0.9, 2.95, 5.5, 3.8,
         ["What it is:",
          "  a feedforward filter",
          "",
          "Inputs:  step command v_ref",
          "Outputs: smooth ramp v_d",
          "",
          "Where it lives:",
          "  BEFORE the LQT + DOB.",
          "  No sensors. No feedback.",
          "",
          "Design goal:",
          "  place zeros on the plant's",
          "  resonance poles."],
         size=13, color=BLACK)

box(s, 6.8, 1.4, 5.9, 5.5, fill=ORANGE_BG, border=ORANGE, bpt=2.5)
text_box(s, 6.8, 1.6, 5.9, 0.7,
         "T_e",
         size=44, bold=True, color=ORANGE, align=CTR, font='Cambria Math')
text_box(s, 7.0, 2.4, 5.5, 0.5,
         "A number (in seconds).",
         size=16, bold=True, align=CTR)
text_box(s, 7.0, 2.95, 5.5, 3.8,
         ["What it is:",
          "  duration of the ramp",
          "",
          "Value:  2π / (ω_f·√(1−δ²))",
          "        = one slosh period",
          "",
          "Meaning:",
          "  stretch the step command",
          "  over one full period of",
          "  the liquid's resonance.",
          "",
          "Paper's number:",
          "  T_e ≈ 0.63 s"],
         size=13, color=BLACK)

# Bottom line
box(s, 0.7, 7.0, 12.0, 0.4, fill=BLUE, border=BLUE, rounded=False)
text_box(s, 0.7, 7.02, 12.0, 0.4,
         "F_e(s) is the machine.   T_e is the dial on the machine.",
         size=15, bold=True, color=WHITE, align=CTR)


# ----- save -----------------------------------------------------------------
OUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'input_shaping_explained.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
