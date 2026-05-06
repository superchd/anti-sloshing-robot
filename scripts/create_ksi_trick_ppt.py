#!/usr/bin/env python3
"""
create_ksi_trick_ppt.py
Slide deck explaining WHY Eq. (7) introduces the helper variable

    ksi_t  =  d_hat_t  -  eta * q_t

and how that lets us avoid differentiating q.

Run:    python3 scripts/create_ksi_trick_ppt.py
Output: results/ksi_trick_explained.pptx
"""
import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE


def rgb(r, g, b): return RGBColor(r, g, b)


WHITE    = rgb(0xFF, 0xFF, 0xFF)
BLACK    = rgb(0x1A, 0x1A, 0x1A)
GRAY     = rgb(0x55, 0x55, 0x55)
LTGRAY   = rgb(0xBB, 0xBB, 0xBB)
BLUE     = rgb(0x15, 0x65, 0xC0)
BLUE_BG  = rgb(0xE3, 0xF2, 0xFD)
RED      = rgb(0xC0, 0x39, 0x2B)
RED_BG   = rgb(0xFD, 0xEC, 0xEA)
GREEN    = rgb(0x2E, 0x7D, 0x32)
GREEN_BG = rgb(0xE8, 0xF5, 0xE9)
ORANGE   = rgb(0xE6, 0x7E, 0x22)
ORANGE_BG= rgb(0xFD, 0xF2, 0xE5)
YELLOW_BG= rgb(0xFF, 0xF6, 0xCC)
PURPLE   = rgb(0x6A, 0x1B, 0x9A)
PURPLE_BG= rgb(0xF5, 0xEC, 0xFB)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def text_box(slide, l, t, w, h, text, size=18, bold=False, color=BLACK,
             align=LFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def box(slide, l, t, w, h, fill=BLUE_BG, border=BLUE, bpt=1.8, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(bpt)
    shp.shadow.inherit = False
    return shp


def line_(slide, x1, y1, x2, y2, color=BLACK, weight=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = conn.line._get_or_add_ln()
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', 'dash')
    return conn


def arrow(slide, x1, y1, x2, y2, color=BLACK, weight=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn


def title_bar(slide, title, subtitle=None):
    box(slide, 0.0, 0.0, 13.333, 0.95, fill=BLUE, border=BLUE,
        bpt=0, rounded=False)
    text_box(slide, 0.4, 0.15, 13.0, 0.7, title, size=28, bold=True,
             color=WHITE, align=LFT)
    if subtitle:
        text_box(slide, 0.4, 0.55, 13.0, 0.4, subtitle, size=14,
                 color=rgb(0xDD, 0xE8, 0xF7), align=LFT)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = blank_slide()
box(s, 0, 0, 13.333, 7.5, fill=BLUE, border=BLUE, bpt=0, rounded=False)
text_box(s, 0.8, 2.1, 11.7, 1.1,
         "Why introduce  ξ ?",
         size=52, bold=True, color=WHITE, align=CTR)
text_box(s, 0.8, 3.3, 11.7, 1.4,
         "The derivative-avoidance trick — turning Eq. (6) into Eq. (7)\n"
         "so we use  q  only, never  q̇",
         size=22, color=rgb(0xE3, 0xF2, 0xFD), align=CTR)
text_box(s, 0.8, 5.5, 11.7, 0.5,
         "Choi et al. 2024 · Section 4.5 · Eq. (7)",
         size=20, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)
text_box(s, 0.8, 6.5, 11.7, 0.5,
         "MECE 6397   ·   Hyundae Cha",
         size=16, color=rgb(0xBB, 0xDE, 0xFB), align=CTR)


# =============================================================================
# SLIDE 2 — The problem
# =============================================================================
s = blank_slide()
title_bar(s, "Step 1: The problem with Eq. (6)",
          "It needs  q̇  — and computing  q̇  on real hardware is dangerous")

# The bad equation
box(s, 0.6, 1.3, 12.1, 1.4, fill=RED_BG, border=RED, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Eq. (6) — what it asks of us",
         size=14, bold=True, color=RED)
text_box(s, 0.8, 1.85, 11.8, 0.7,
         "d̂̇_t   =   η  (  q̇_t   −   A q_t   −   B u_t   −   d̂_t  )",
         size=22, bold=True, color=BLACK, font='Consolas', align=CTR)

# spotlight on q̇
box(s, 0.6, 2.85, 12.1, 0.8, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 2.95, 11.8, 0.6,
         "See that  q̇  sitting inside?  That's bad news on real hardware.",
         size=15, bold=True, color=ORANGE, align=CTR)

# Why it's bad — left
box(s, 0.5, 3.85, 6.2, 3.3, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 3.95, 5.8, 0.5,
         "Why  q̇  is a problem",
         size=15, bold=True, color=BLUE)
text_box(s, 0.7, 4.45, 5.8, 2.6,
         ["• Sensors give you  q  (smooth).",
          "  No 'velocity sensor' for perfect q̇.",
          "",
          "• To get q̇, you must compute:",
          "       q̇ ≈ ( q_now − q_prev ) / dt",
          "",
          "• That is NUMERICAL DIFFERENTIATION."],
         size=12, color=BLACK, font='Consolas')

# Right — noise amplification
box(s, 6.9, 3.85, 6.0, 3.3, fill=RED_BG, border=RED)
text_box(s, 7.1, 3.95, 5.6, 0.5,
         "Numerical derivatives amplify noise",
         size=15, bold=True, color=RED)
text_box(s, 7.1, 4.45, 5.6, 2.6,
         ["• tiny sensor jitter:  0.001 rad",
          "• sample time:         dt = 0.001 s",
          "",
          "• jitter through derivative:",
          "       0.001 / 0.001  =  1.0 rad/s",
          "",
          "• Garbage q̇  →  garbage d̂  →  bad ctrl."],
         size=12, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 3 — The trick (definition of ξ)
# =============================================================================
s = blank_slide()
title_bar(s, "Step 2: The trick — define a helper variable  ξ",
          "Just a renaming. Nothing magic yet.")

# Big definition
box(s, 0.6, 1.3, 12.1, 1.5, fill=PURPLE_BG, border=PURPLE, bpt=2.5)
text_box(s, 0.8, 1.40, 11.8, 0.4,
         "Definition (paper, just before Eq. 7)",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 1.85, 11.8, 0.7,
         "ξ_t   ≡   d̂_t   −   η  q_t",
         size=28, bold=True, color=BLACK, font='Consolas', align=CTR)

# What this is
box(s, 0.6, 3.0, 12.1, 1.0, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 3.10, 11.8, 0.4,
         "Read it as:",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 3.50, 11.8, 0.4,
         "'Instead of tracking d̂ directly, I'll track this combination "
         "d̂ − η q.'",
         size=14, color=BLACK)

# Why THIS combo?
box(s, 0.6, 4.15, 12.1, 1.05, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 4.25, 11.8, 0.4,
         "Why subtract  η q  specifically?",
         size=14, bold=True, color=BLUE)
text_box(s, 0.8, 4.65, 11.8, 0.5,
         "Because — when we differentiate ξ — the η q̇ term will EXACTLY "
         "cancel the bad q̇ inside Eq. (6). Watch.",
         size=13, color=BLACK)

# Forward to next slide
box(s, 0.6, 5.40, 12.1, 1.7, fill=GREEN_BG, border=GREEN)
text_box(s, 0.8, 5.50, 11.8, 0.4,
         "What we will do on the next slide:",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 5.95, 11.8, 1.05,
         ["1.  Take  ξ̇  =  d̂̇  −  η q̇",
          "2.  Substitute Eq. (6) for d̂̇",
          "3.  Watch the  q̇  terms cancel"],
         size=14, color=BLACK, font='Consolas')


# =============================================================================
# SLIDE 4 — The cancellation (math chain)
# =============================================================================
s = blank_slide()
title_bar(s, "Step 3: Differentiate ξ — and watch  q̇  cancel",
          "Three lines of algebra, that's all")

# Chain of equations
y = 1.4
ks = [
    ("Differentiate the definition of ξ:",
     "ξ̇   =   d̂̇   −   η q̇",
     PURPLE),
    ("Substitute Eq. (6) for d̂̇ :",
     "ξ̇   =   η ( q̇  −  A q  −  B u  −  d̂ )   −   η q̇",
     BLUE),
    ("Distribute η on the right side:",
     "ξ̇   =   η q̇   −   η A q   −   η B u   −   η d̂   −   η q̇",
     ORANGE),
    ("The two  η q̇  terms cancel exactly:",
     "ξ̇   =   η ( − A q  −  B u )   −   η d̂",
     GREEN),
]
for i, (note, eq, c) in enumerate(ks):
    yi = y + i * 1.15
    box(s, 0.6, yi, 12.1, 1.0, fill=WHITE, border=c, bpt=2)
    text_box(s, 0.8, yi + 0.05, 11.8, 0.4, note,
             size=13, bold=True, color=c)
    text_box(s, 0.8, yi + 0.45, 11.8, 0.5, eq,
             size=15, bold=True, color=BLACK, font='Consolas', align=CTR)

# celebration line
box(s, 0.6, 6.10, 12.1, 1.1, fill=GREEN_BG, border=GREEN, bpt=2.5)
text_box(s, 0.8, 6.20, 11.8, 0.4,
         "Result — Eq. (7)",
         size=14, bold=True, color=GREEN)
text_box(s, 0.8, 6.60, 11.8, 0.5,
         "ξ̇   =   η ( − A q  −  B u )   −   η d̂           "
         "✓  q̇  is GONE.  Only  q  remains.",
         size=15, bold=True, color=GREEN, font='Consolas', align=CTR)


# =============================================================================
# SLIDE 5 — Online use (two-step recipe)
# =============================================================================
s = blank_slide()
title_bar(s, "Step 4: What we actually compute every timestep",
          "Two steps. q only — never q̇.")

# Step 1
box(s, 0.6, 1.4, 12.1, 1.85, fill=BLUE_BG, border=BLUE)
text_box(s, 0.8, 1.50, 1.0, 0.6, "1.",
         size=36, bold=True, color=BLUE, align=CTR)
text_box(s, 1.8, 1.50, 10.7, 0.5,
         "INTEGRATE  ξ  (smooth, never differentiates):",
         size=15, bold=True, color=BLUE)
text_box(s, 1.8, 1.95, 10.7, 0.6,
         "ξ_new   =   ξ_old   +   dt  ·  η ( − A q  −  B u  −  d̂ )",
         size=16, bold=True, color=BLACK, font='Consolas')
text_box(s, 1.8, 2.65, 10.7, 0.5,
         "(integration is the OPPOSITE of differentiation — it smooths, "
         "doesn't amplify, noise)",
         size=12, color=GRAY)

# Step 2
box(s, 0.6, 3.45, 12.1, 1.85, fill=GREEN_BG, border=GREEN)
text_box(s, 0.8, 3.55, 1.0, 0.6, "2.",
         size=36, bold=True, color=GREEN, align=CTR)
text_box(s, 1.8, 3.55, 10.7, 0.5,
         "RECOVER  d̂  from ξ and q:",
         size=15, bold=True, color=GREEN)
text_box(s, 1.8, 4.00, 10.7, 0.6,
         "d̂_t   =   ξ_t   +   η  q_t",
         size=18, bold=True, color=BLACK, font='Consolas')
text_box(s, 1.8, 4.70, 10.7, 0.5,
         "(uses position q, never the noisy derivative)",
         size=12, color=GRAY)

# the punch line
box(s, 0.6, 5.55, 12.1, 1.55, fill=YELLOW_BG, border=ORANGE, bpt=2.5)
text_box(s, 0.8, 5.65, 11.8, 0.4,
         "What we got:",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 6.05, 11.8, 1.0,
         ["• Same disturbance estimate as Eq. (6)",
          "• Without ever computing  q̇",
          "• Survives noisy real-world sensors."],
         size=14, color=BLACK)


# =============================================================================
# SLIDE 6 — Picture: integration vs differentiation
# =============================================================================
s = blank_slide()
title_bar(s, "Big picture — integration smooths, differentiation explodes",
          "Why the  ξ  trick is worth doing")

# Comparison table
box(s, 0.6, 1.4, 12.1, 4.0, fill=WHITE, border=BLACK, bpt=1.5)

# header row
box(s, 0.6, 1.4, 4.0, 0.65, fill=BLUE, border=BLUE, bpt=0)
box(s, 4.6, 1.4, 5.0, 0.65, fill=BLUE, border=BLUE, bpt=0)
box(s, 9.6, 1.4, 3.1, 0.65, fill=BLUE, border=BLUE, bpt=0)
text_box(s, 0.6, 1.50, 4.0, 0.5, "Method",
         size=16, bold=True, color=WHITE, align=CTR)
text_box(s, 4.6, 1.50, 5.0, 0.5, "What we do to get d̂",
         size=16, bold=True, color=WHITE, align=CTR)
text_box(s, 9.6, 1.50, 3.1, 0.5, "Noise behavior",
         size=16, bold=True, color=WHITE, align=CTR)

# row 1
box(s, 0.6, 2.05, 4.0, 1.6, fill=RED_BG, border=RED, bpt=1)
box(s, 4.6, 2.05, 5.0, 1.6, fill=RED_BG, border=RED, bpt=1)
box(s, 9.6, 2.05, 3.1, 1.6, fill=RED_BG, border=RED, bpt=1)
text_box(s, 0.7, 2.50, 3.8, 0.7,
         "Eq. (6) directly",
         size=15, bold=True, color=RED)
text_box(s, 4.7, 2.20, 4.8, 1.4,
         ["DIFFERENTIATE q to get q̇,",
          "then plug into",
          "d̂̇ = η(q̇ − Aq − Bu − d̂)"],
         size=12, color=BLACK, font='Consolas')
text_box(s, 9.7, 2.45, 2.9, 0.7,
         "EXPLODES",
         size=18, bold=True, color=RED, align=CTR)
text_box(s, 9.7, 3.00, 2.9, 0.5,
         "tiny jitter → huge spike",
         size=11, color=GRAY, align=CTR)

# row 2
box(s, 0.6, 3.65, 4.0, 1.6, fill=GREEN_BG, border=GREEN, bpt=1)
box(s, 4.6, 3.65, 5.0, 1.6, fill=GREEN_BG, border=GREEN, bpt=1)
box(s, 9.6, 3.65, 3.1, 1.6, fill=GREEN_BG, border=GREEN, bpt=1)
text_box(s, 0.7, 4.10, 3.8, 0.7,
         "Eq. (7) with ξ",
         size=15, bold=True, color=GREEN)
text_box(s, 4.7, 3.80, 4.8, 1.4,
         ["INTEGRATE ξ,",
          "then  d̂ = ξ + η q",
          "(uses q only)"],
         size=12, color=BLACK, font='Consolas')
text_box(s, 9.7, 4.05, 2.9, 0.7,
         "STAYS CLEAN",
         size=18, bold=True, color=GREEN, align=CTR)
text_box(s, 9.7, 4.60, 2.9, 0.5,
         "integration smooths noise",
         size=11, color=GRAY, align=CTR)

# bottom takeaway
box(s, 0.6, 5.65, 12.1, 1.45, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 5.75, 11.8, 0.4,
         "The whole point of the trick:",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 6.15, 11.8, 0.95,
         "Convert a 'differentiate-the-noisy-signal' problem  →  "
         "into an 'integrate-the-clean-signal' problem.",
         size=14, color=BLACK, align=CTR)


# =============================================================================
# SLIDE 7 — But your code uses q̇?
# =============================================================================
s = blank_slide()
title_bar(s, "But wait — your controller.m DOES use  q̇ ?",
          "Yes, and that's fine for a clean simulation. Real hardware would use ξ.")

# Code line
box(s, 0.6, 1.4, 12.1, 1.4, fill=YELLOW_BG, border=ORANGE)
text_box(s, 0.8, 1.50, 11.8, 0.4,
         "controller.m, line 70",
         size=14, bold=True, color=ORANGE)
text_box(s, 0.8, 1.95, 11.8, 0.7,
         "residual = (q(2) − ctrl_state.q2_prev) / p.dt − q_dot_est(2);",
         size=14, color=BLACK, font='Consolas')
text_box(s, 0.8, 2.45, 11.8, 0.4,
         "← that  (q(2) − q2_prev) / dt   IS numerical differentiation of q.",
         size=12, color=RED)

# why it's OK in sim
box(s, 0.5, 3.0, 6.2, 4.0, fill=BLUE_BG, border=BLUE)
text_box(s, 0.7, 3.10, 5.8, 0.5,
         "Why it works in your sim",
         size=15, bold=True, color=BLUE)
text_box(s, 0.7, 3.60, 5.8, 3.3,
         ["• MATLAB integration is clean",
          "  (no encoder noise)",
          "",
          "• q is computed from the ODE,",
          "  not measured",
          "",
          "• So q̇ via finite difference",
          "  has tiny numerical error",
          "",
          "• Eq. (6) form is fine here"],
         size=13, color=BLACK)

# why ξ on hardware
box(s, 6.9, 3.0, 6.0, 4.0, fill=GREEN_BG, border=GREEN)
text_box(s, 7.1, 3.10, 5.6, 0.5,
         "Why ξ on real hardware",
         size=15, bold=True, color=GREEN)
text_box(s, 7.1, 3.60, 5.6, 3.3,
         ["• IMU / encoder noise is real",
          "",
          "• Numerical q̇ would chatter",
          "",
          "• Use Eq. (7):",
          "    1. integrate ξ",
          "    2. d̂ = ξ + η q",
          "",
          "• That's the version that",
          "  survives a real serving robot"],
         size=13, color=BLACK)


# =============================================================================
# SLIDE 8 — TL;DR
# =============================================================================
s = blank_slide()
title_bar(s, "TL;DR — four lines",
          "What to remember about ξ")

y = 1.6
items = [
    ("The problem",
     "q̇ in Eq. (6) is dangerous — numerical differentiation amplifies noise",
     RED),
    ("The trick",
     "ξ ≡ d̂ − η q   chosen so the bad  q̇  cancels in  ξ̇",
     PURPLE),
    ("Online use",
     "integrate ξ (safe), then add η q to recover d̂",
     BLUE),
    ("End result",
     "final algorithm uses ONLY  q,  never  q̇   ✓",
     GREEN),
]
for i, (k, v, c) in enumerate(items):
    yi = y + i * 1.05
    box(s, 0.6, yi, 12.1, 0.85, fill=WHITE, border=c, bpt=2.0)
    box(s, 0.6, yi, 2.7, 0.85, fill=c, border=c, bpt=0)
    text_box(s, 0.7, yi + 0.22, 2.5, 0.5, k,
             size=15, bold=True, color=WHITE, align=CTR)
    text_box(s, 3.5, yi + 0.22, 9.0, 0.5, v,
             size=14, color=BLACK)

# closing line
box(s, 0.6, 6.0, 12.1, 1.2, fill=PURPLE_BG, border=PURPLE)
text_box(s, 0.8, 6.10, 11.8, 0.5,
         "One sentence",
         size=14, bold=True, color=PURPLE)
text_box(s, 0.8, 6.55, 11.8, 0.6,
         "ξ is the single algebraic trick that makes the DOB usable on "
         "real hardware.",
         size=14, color=BLACK)


# ----- save ------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, '..', 'results')
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, 'ksi_trick_explained.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'  {len(prs.slides)} slides')
