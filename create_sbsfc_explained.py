#!/usr/bin/env python3
"""
create_sbsfc_explained.py
SBSFC block diagram with plain-English callout annotations for presentation.
Explains each block so first-time viewers can understand.

Run:    python3 create_sbsfc_explained.py
Output: results/sbsfc_explained.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

def rgb(r, g, b): return RGBColor(r, g, b)

BLACK     = rgb(0x1A, 0x1A, 0x1A)
WHITE     = rgb(0xFF, 0xFF, 0xFF)
BLUE      = rgb(0x26, 0x66, 0xCF)
BLUE_BG   = rgb(0xEB, 0xF3, 0xFF)
BLUE_FG   = rgb(0x1A, 0x50, 0xAA)
YELLOW    = rgb(0xCC, 0x99, 0x00)
YELLOW_BG = rgb(0xFF, 0xFB, 0xE6)
YELLOW_FG = rgb(0x99, 0x77, 0x00)
RED       = rgb(0xCC, 0x22, 0x22)
RED_DASH  = rgb(0xDD, 0x44, 0x44)
GRAY      = rgb(0x55, 0x55, 0x55)
LTGRAY    = rgb(0xAA, 0xAA, 0xAA)
GREEN     = rgb(0x22, 0x88, 0x22)
ORANGE    = rgb(0xEE, 0x77, 0x00)
MAGENTA   = rgb(0xBB, 0x22, 0x88)
BOX_BG    = rgb(0xFA, 0xFA, 0xFA)
BOX_BDR   = rgb(0x33, 0x33, 0x33)
SIG_COLOR = rgb(0x22, 0x22, 0x22)
BLUE_SIG  = rgb(0x22, 0x55, 0xCC)
YEL_SIG   = rgb(0xBB, 0x88, 0x00)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT
RGT = PP_ALIGN.RIGHT

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1: SBSFC Block Diagram with Callout Annotations
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background; bf = bg.fill; bf.solid(); bf.fore_color.rgb = WHITE

# ─── Helpers (bound to current slide) ────────────────────────────────────────

def box(l, t, w, h, fill=BOX_BG, border=BOX_BDR, bpt=1.5, dash=False, rounded=False):
    kind = 5 if rounded else 1
    s = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(bpt)
    if dash:
        ln = s.line._ln
        for old in ln.findall(qn('a:prstDash')): ln.remove(old)
        etree.SubElement(ln, qn('a:prstDash')).set('val', 'lgDash')
    return s

def write(shape, rows, anchor='ctr'):
    tf = shape.text_frame; tf.word_wrap = True
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', anchor)
        bp.set('lIns', '36000'); bp.set('rIns', '36000')
        bp.set('tIns', '18000'); bp.set('bIns', '18000')
    tf.clear()
    for i, (s, sz, bd, co, al) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if not s: continue
        r = p.add_run(); r.text = s
        r.font.size = Pt(sz); r.font.bold = bd; r.font.color.rgb = co

def label(l, t, w, h, text, sz, bold, color, al=CTR):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return tb

def mlabel(l, t, w, h, text, sz, bold, color, al=LFT):
    """Multi-line text box with word wrap."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return tb

def fwd_arrow(x1, y1, x2, y2, color=SIG_COLOR, pt=2.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(pt)
    ln = c.line._ln
    for tag in [qn('a:headEnd'), qn('a:tailEnd')]:
        for el in ln.findall(tag): ln.remove(el)
    etree.SubElement(ln, qn('a:headEnd')).set('type', 'none')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')

def line(x1, y1, x2, y2, color=SIG_COLOR, pt=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(pt)

def dash_line(x1, y1, x2, y2, color=LTGRAY, pt=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(pt)
    ln = c.line._ln
    for old in ln.findall(qn('a:prstDash')): ln.remove(old)
    etree.SubElement(ln, qn('a:prstDash')).set('val', 'lgDash')

def summing_junction(cx, cy, r=0.12):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r), Inches(2*r), Inches(2*r))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BOX_BDR; s.line.width = Pt(1.5)
    return s

def callout_box(l, t, w, h, num, title, desc, num_color, bg_color, border_color):
    """Numbered callout with title and description."""
    # Background box
    b = box(l, t, w, h, bg_color, border_color, bpt=2.0, rounded=True)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(l - 0.12), Inches(t - 0.12), Inches(0.32), Inches(0.32))
    circ.fill.solid(); circ.fill.fore_color.rgb = num_color
    circ.line.color.rgb = num_color; circ.line.width = Pt(0)
    tf = circ.text_frame
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')
        bp.set('lIns', '0'); bp.set('rIns', '0')
        bp.set('tIns', '0'); bp.set('bIns', '0')
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = num
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    # Text content
    write(b, [
        (title, 10, True, border_color, LFT),
        (desc, 9, False, GRAY, LFT),
    ], anchor='ctr')
    return b

# ─── Slide title ─────────────────────────────────────────────────────────────
title_bg = box(0, 0, 13.333, 0.55, rgb(0x1A, 0x3A, 0x5C), rgb(0x1A, 0x3A, 0x5C), 0)
write(title_bg, [
    ("SBSFC Controller — How It Works", 18, True, WHITE, CTR)
], anchor='ctr')
label(0.3, 0.12, 3.0, 0.30, "Self-Balancing Slosh-Free Control", 9, False, rgb(0x88, 0xBB, 0xEE), LFT)

# ═══════════════════════════════════════════════════════════════════════════════
#  BLOCK DIAGRAM (center of slide, slightly compact)
# ═══════════════════════════════════════════════════════════════════════════════

# Offsets: diagram in x: 0.3–8.5, y: 0.7–4.0
DY = 0.20  # shift everything down a bit

# SBSFC outer boundary (red dashed)
box(0.3, 0.65+DY, 8.0, 3.25, WHITE, RED_DASH, bpt=2.5, dash=True)
label(0.35, 3.55+DY, 4.0, 0.30,
      "[Self-balancing slosh-free control (SBSFC)]", 8, True, BLACK, LFT)

# Upper body Controller group (blue dashed)
box(0.5, 0.90+DY, 5.40, 1.75, BLUE_BG, BLUE, bpt=2.0, dash=True)
label(1.7, 0.72+DY, 2.5, 0.22, "Upper body Controller", 9, True, BLUE_FG)

# Lower body Controller group (yellow dashed)
box(4.60, 2.85+DY, 2.50, 0.85, YELLOW_BG, YELLOW, bpt=2.0, dash=True)
label(5.5, 3.42+DY, 1.8, 0.22, "Lower body Controller", 8, True, YELLOW_FG)

# ─── Blocks ──────────────────────────────────────────────────────────────────

# Reference Shaping
rs = box(1.0, 1.30+DY, 1.25, 0.45, BOX_BG, BOX_BDR, 1.5)
write(rs, [("Reference", 8, True, BLACK, CTR),
           ("Shaping",   8, True, BLACK, CTR)])

# Upper body controller
uc = box(2.90, 1.30+DY, 1.40, 0.45, BOX_BG, BOX_BDR, 1.5)
write(uc, [("Upper body", 8, True, BLACK, CTR),
           ("controller", 8, True, BLACK, CTR)])

# Disturbance observer
do = box(3.00, 2.15+DY, 1.40, 0.42, BOX_BG, BOX_BDR, 1.5)
write(do, [("Disturbance", 8, True, BLACK, CTR),
           ("observer",    8, True, BLACK, CTR)])

# Compensator
comp = box(4.80, 2.95+DY, 1.55, 0.55, BOX_BG, YELLOW, 1.5)
write(comp, [("Compensator", 8, True, BLACK, CTR),
             ("sgn(\u03c8\u1d38)", 9, True, YELLOW_FG, CTR)])

# State vector q_t
label(1.60, 1.88+DY, 0.95, 0.65,
      "q = [\u03c8, \u03c8\u0307, x, x\u0307]", 8, False, BLACK, CTR)

# ─── Signal labels ───────────────────────────────────────────────────────────
label(0.15, 1.32+DY, 0.55, 0.22, "v_ref", 9, True, BLACK)
label(2.28, 1.13+DY, 0.45, 0.22, "v_d", 9, True, BLACK)
label(4.60, 1.10+DY, 0.55, 0.22, "u\u209c\u1d58", 9, True, BLUE_FG)
label(4.70, 2.18+DY, 0.55, 0.22, "d\u0302\u209c", 9, True, BLUE_FG)
label(6.45, 3.05+DY, 0.50, 0.22, "u\u209c\u1d38", 9, True, YELLOW_FG)
label(7.10, 0.50+DY, 0.45, 0.22, "d\u209c", 10, True, RED)
label(7.70, 1.28+DY, 0.85, 0.25, "u\u209c = x\u0308\u209c", 9, True, BLACK)

# ─── Summing junctions ───────────────────────────────────────────────────────
sj1 = summing_junction(6.10, 1.52+DY, 0.12)
sj2 = summing_junction(7.20, 1.52+DY, 0.12)
label(5.90, 1.25+DY, 0.15, 0.15, "+", 8, True, BLACK)
label(5.90, 1.60+DY, 0.15, 0.15, "\u2013", 9, True, RED)
label(7.00, 1.25+DY, 0.15, 0.15, "+", 8, True, BLACK)
label(7.30, 1.25+DY, 0.15, 0.15, "+", 8, True, RED)

# ─── Signal flow arrows ─────────────────────────────────────────────────────
fwd_arrow(0.60, 1.52+DY, 1.00, 1.52+DY, BLACK, 1.8)
fwd_arrow(2.25, 1.52+DY, 2.90, 1.52+DY, BLACK, 1.8)
fwd_arrow(4.30, 1.52+DY, 5.98, 1.52+DY, BLUE_SIG, 1.8)
# q_t feedback
line(2.50, 1.52+DY, 2.50, 1.88+DY, BLACK, 1.2)
fwd_arrow(2.50, 1.88+DY, 2.90, 1.65+DY, GRAY, 1.0)
fwd_arrow(2.60, 2.36+DY, 3.00, 2.36+DY, BLACK, 1.2)
# DOB output
fwd_arrow(4.40, 2.36+DY, 6.10, 2.36+DY, BLUE_SIG, 1.2)
line(6.10, 2.36+DY, 6.10, 1.64+DY, BLUE_SIG, 1.2)
# Compensator output
fwd_arrow(6.35, 3.22+DY, 6.70, 3.22+DY, YEL_SIG, 1.2)
line(6.70, 3.22+DY, 6.70, 1.64+DY, YEL_SIG, 1.2)
# Sum1 → Sum2
fwd_arrow(6.22, 1.52+DY, 7.08, 1.52+DY, BLACK, 1.8)
# d_t → Sum2
fwd_arrow(7.20, 0.80+DY, 7.20, 1.40+DY, RED, 1.8)
# Sum2 → output
fwd_arrow(7.32, 1.52+DY, 8.30, 1.52+DY, BLACK, 2.5)

# ═══════════════════════════════════════════════════════════════════════════════
#  CALLOUT BOXES — Right side
# ═══════════════════════════════════════════════════════════════════════════════

CB_X = 8.75   # callout boxes start x
CB_W = 4.30   # callout box width

# ① Reference Shaping
callout_box(CB_X, 0.65, CB_W, 0.70,
    "\u2460", "Smooth the speed command",
    "Filters v_ref to avoid exciting the liquid's\nnatural sloshing frequency (\u03c9_f = 9.9 rad/s)",
    rgb(0x44, 0x44, 0x44), rgb(0xF4, 0xF4, 0xF8), rgb(0x44, 0x44, 0x44))

# Dashed connector from callout to block
dash_line(CB_X, 1.00, 2.25, 1.52+DY, rgb(0x88, 0x88, 0x88), 1.0)

# ② Upper body controller (LQT)
callout_box(CB_X, 1.50, CB_W, 0.70,
    "\u2461", "Balance the robot (LQT)",
    "State feedback: reads [\u03c8, \u03c8\u0307, x, x\u0307] and computes\nacceleration to keep robot upright + track speed",
    BLUE_FG, rgb(0xED, 0xF3, 0xFF), BLUE_FG)

dash_line(CB_X, 1.85, 4.30, 1.52+DY, rgb(0x66, 0x88, 0xCC), 1.0)

# ③ Disturbance Observer (DOB)
callout_box(CB_X, 2.35, CB_W, 0.70,
    "\u2462", "Cancel disturbances (DOB)",
    "Detects unexpected forces (bumps, pushes)\nand subtracts them (\u2013 sign at junction)",
    BLUE_FG, rgb(0xED, 0xF3, 0xFF), BLUE_FG)

dash_line(CB_X, 2.70, 4.40, 2.36+DY, rgb(0x66, 0x88, 0xCC), 1.0)

# ④ Compensator
callout_box(CB_X, 3.20, CB_W, 0.70,
    "\u2463", "Damp residual wobble",
    "Bang-bang correction: watches lower body tilt\n\u03c8\u1d38 and quickly damps remaining oscillation",
    YELLOW_FG, rgb(0xFF, 0xFB, 0xEC), YELLOW_FG)

dash_line(CB_X, 3.55, 6.35, 3.22+DY, rgb(0xBB, 0x99, 0x33), 1.0)

# ═══════════════════════════════════════════════════════════════════════════════
#  BOTTOM: ψ vs θ explanation with robot side-view sketch
# ═══════════════════════════════════════════════════════════════════════════════

# Divider line
line(0.3, 4.25, 13.0, 4.25, rgb(0xDD, 0xDD, 0xDD), 1.0)

# Section title
label(0.3, 4.35, 5.0, 0.30,
      "Two Angles on One Robot — \u03c8 vs \u03b8", 14, True, BLACK, LFT)

# ── ψ explanation box (magenta) ──
psi_box = box(0.3, 4.80, 3.80, 1.50, rgb(0xFD, 0xF0, 0xF8), MAGENTA, 2.0, rounded=True)
write(psi_box, [
    ("\u03c8 = Robot Body Tilt", 13, True, MAGENTA, LFT),
    ("", 4, False, BLACK, LFT),
    ("Measured at: robot base (whole body)", 9, False, GRAY, LFT),
    ("Stability: UNSTABLE (inverted pendulum)", 9, False, RED, LFT),
    ("Controlled by: LQT state feedback (\u2461)", 9, False, GRAY, LFT),
    ("Goal: keep \u03c8 \u2248 0\u00b0 (don't fall over)", 9, True, BLACK, LFT),
], anchor='t')

# ── θ explanation box (orange) ──
th_box = box(4.40, 4.80, 3.80, 1.50, rgb(0xFF, 0xF5, 0xE8), ORANGE, 2.0, rounded=True)
write(th_box, [
    ("\u03b8 = Liquid Sloshing Angle", 13, True, ORANGE, LFT),
    ("", 4, False, BLACK, LFT),
    ("Measured at: tray (pendulum swing)", 9, False, GRAY, LFT),
    ("Stability: STABLE (gravity restores)", 9, False, GREEN, LFT),
    ("Controlled by: Input Shaping (\u2460)", 9, False, GRAY, LFT),
    ("Goal: minimize \u03b8 (don't spill liquid)", 9, True, BLACK, LFT),
], anchor='t')

# ── Coupling explanation box ──
coup_box = box(8.50, 4.80, 4.50, 1.50, rgb(0xF4, 0xF4, 0xF8), rgb(0x55, 0x55, 0x55), 1.5, rounded=True)
write(coup_box, [
    ("How they connect:", 12, True, BLACK, LFT),
    ("", 4, False, BLACK, LFT),
    ("u (acceleration) drives BOTH \u03c8 and \u03b8", 10, True, BLACK, LFT),
    ("but they are computed SEPARATELY:", 9, False, GRAY, LFT),
    ("", 3, False, BLACK, LFT),
    ("\u03c8:  q\u0307 = A\u00b7q + B\u00b7u          (linear)", 9, False, BLUE_FG, LFT),
    ("\u03b8:  \u03b8\u0308 = -(g/l)sin\u03b8 + (u/l)cos\u03b8  (nonlinear)", 9, False, ORANGE, LFT),
    ("", 3, False, BLACK, LFT),
    ("\u03b8 does NOT affect \u03c8 (0.5 kg vs 42 kg)", 9, True, RED, LFT),
], anchor='t')

# ── Simple robot side-view sketch showing both angles ──
# Robot body (center at 2.2, 7.0)
RX = 2.0; RY = 6.60

# Ground
line(0.5, RY + 0.55, 4.5, RY + 0.55, GRAY, 2.0)

# Wheels
for wx in [RX - 0.25, RX + 0.25]:
    wh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(wx - 0.08), Inches(RY + 0.38), Inches(0.16), Inches(0.16))
    wh.fill.solid(); wh.fill.fore_color.rgb = rgb(0x44, 0x44, 0x44)
    wh.line.color.rgb = BLACK; wh.line.width = Pt(1)

# Robot body (slightly tilted to show ψ)
body = box(RX - 0.18, RY - 0.35, 0.36, 0.72, rgb(0xDD, 0xDD, 0xE5), GRAY, 1.5, rounded=True)

# Tray
tray = box(RX - 0.35, RY - 0.42, 0.70, 0.06, rgb(0x88, 0xCC, 0x88), rgb(0x44, 0x88, 0x44), 1.5)

# ψ arc (magenta) — at robot base
label(RX + 0.35, RY + 0.15, 0.80, 0.25, "\u03c8 (body tilt)", 8, True, MAGENTA, LFT)
line(RX, RY + 0.35, RX, RY - 0.10, MAGENTA, 2.0)          # vertical ref
dash_line(RX, RY + 0.35, RX - 0.08, RY - 0.10, LTGRAY, 1.0)  # tilted body hint
fwd_arrow(RX + 0.38, RY + 0.28, RX + 0.05, RY + 0.20, MAGENTA, 1.0)

# Pendulum hanging from tray
pend_top_x = RX; pend_top_y = RY - 0.36
pend_bot_x = RX + 0.10; pend_bot_y = RY - 0.06
line(pend_top_x, pend_top_y, pend_bot_x, pend_bot_y, rgb(0x88, 0x88, 0x88), 2.0)

# Pendulum bob
bob = slide.shapes.add_shape(MSO_SHAPE.OVAL,
    Inches(pend_bot_x - 0.04), Inches(pend_bot_y - 0.04),
    Inches(0.08), Inches(0.08))
bob.fill.solid(); bob.fill.fore_color.rgb = ORANGE
bob.line.color.rgb = ORANGE; bob.line.width = Pt(1)

# θ arc (orange) — at tray pivot
label(RX + 0.35, RY - 0.50, 0.95, 0.25, "\u03b8 (sloshing)", 8, True, ORANGE, LFT)
# Dashed vertical ref from tray
dash_line(pend_top_x, pend_top_y, pend_top_x, pend_top_y + 0.32, rgb(0xEE, 0xAA, 0x55), 1.0)
fwd_arrow(RX + 0.38, RY - 0.44, RX + 0.08, RY - 0.22, ORANGE, 1.0)

# ── Analogy at the very bottom ──
analogy_bg = box(0.3, 7.05, 12.7, 0.38, rgb(0xF0, 0xF8, 0xF0), GREEN, 1.5, rounded=True)
write(analogy_bg, [
    ("Analogy:  You're on a bus holding coffee.  "
     "\u03c8 = how much YOU lean  |  "
     "\u03b8 = how much the COFFEE sloshes  |  "
     "Your leaning causes sloshing, but sloshing doesn't make you fall.",
     9, False, rgb(0x22, 0x66, 0x22), CTR)
], anchor='ctr')


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2: Signal Flow Summary (simple left-to-right)
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg2 = slide2.background; bf2 = bg2.fill; bf2.solid(); bf2.fore_color.rgb = WHITE

# Rebind helpers to slide2
def box2(l, t, w, h, fill, border, bpt=1.5, rounded=False):
    kind = 5 if rounded else 1
    s = slide2.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(bpt)
    return s

def write2(shape, rows, anchor='ctr'):
    tf = shape.text_frame; tf.word_wrap = True
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', anchor)
        bp.set('lIns', '45000'); bp.set('rIns', '45000')
        bp.set('tIns', '27000'); bp.set('bIns', '27000')
    tf.clear()
    for i, (s, sz, bd, co, al) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if not s: continue
        r = p.add_run(); r.text = s
        r.font.size = Pt(sz); r.font.bold = bd; r.font.color.rgb = co

def label2(l, t, w, h, text, sz, bold, color, al=CTR):
    tb = slide2.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return tb

def arrow2(x1, y1, x2, y2, color=SIG_COLOR, pt=3.0):
    c = slide2.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(pt)
    ln = c.line._ln
    for tag in [qn('a:headEnd'), qn('a:tailEnd')]:
        for el in ln.findall(tag): ln.remove(el)
    etree.SubElement(ln, qn('a:headEnd')).set('type', 'none')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow'); te.set('w', 'lg'); te.set('len', 'lg')

# Title
t2 = box2(0, 0, 13.333, 0.55, rgb(0x1A, 0x3A, 0x5C), rgb(0x1A, 0x3A, 0x5C))
write2(t2, [("Signal Flow:  How the Controller Processes One Command", 18, True, WHITE, CTR)])

# ── Main flow: 4 big numbered boxes with arrows ──
BH = 1.60  # box height
BY = 1.20  # box Y
BW = 2.40  # box width
gap = 0.45

# Box positions
x1 = 0.40
x2 = x1 + BW + gap
x3 = x2 + BW + gap
x4 = x3 + BW + gap

# ① Smooth
b1 = box2(x1, BY, BW, BH, rgb(0xF4, 0xF4, 0xF8), rgb(0x44, 0x44, 0x44), 2.0, True)
write2(b1, [
    ("\u2460 Smooth", 16, True, rgb(0x44, 0x44, 0x44), CTR),
    ("", 5, False, BLACK, CTR),
    ("Reference Shaping", 11, True, BLACK, CTR),
    ("", 4, False, BLACK, CTR),
    ('"Go 0.5 m/s" \u2192 gradual ramp', 10, False, GRAY, CTR),
    ("Avoids exciting liquid's", 10, False, GRAY, CTR),
    ("natural frequency", 10, False, GRAY, CTR),
], anchor='ctr')

arrow2(x1 + BW + 0.02, BY + BH/2, x2 - 0.02, BY + BH/2, BLACK, 3.0)

# ② Balance
b2 = box2(x2, BY, BW, BH, BLUE_BG, BLUE_FG, 2.0, True)
write2(b2, [
    ("\u2461 Balance", 16, True, BLUE_FG, CTR),
    ("", 5, False, BLACK, CTR),
    ("LQT Controller", 11, True, BLACK, CTR),
    ("", 4, False, BLACK, CTR),
    ("Reads: \u03c8, \u03c8\u0307, x, x\u0307", 10, False, GRAY, CTR),
    ("Keeps robot upright", 10, False, GRAY, CTR),
    ("while tracking speed", 10, False, GRAY, CTR),
], anchor='ctr')

arrow2(x2 + BW + 0.02, BY + BH/2, x3 - 0.02, BY + BH/2, BLACK, 3.0)

# ③ Cancel
b3 = box2(x3, BY, BW, BH, BLUE_BG, BLUE_FG, 2.0, True)
write2(b3, [
    ("\u2462 Cancel", 16, True, BLUE_FG, CTR),
    ("", 5, False, BLACK, CTR),
    ("Disturbance Observer", 11, True, BLACK, CTR),
    ("", 4, False, BLACK, CTR),
    ("Detects: bumps, pushes", 10, False, GRAY, CTR),
    ("Estimates force d\u0302\u209c", 10, False, GRAY, CTR),
    ("and subtracts it out", 10, False, GRAY, CTR),
], anchor='ctr')

arrow2(x3 + BW + 0.02, BY + BH/2, x4 - 0.02, BY + BH/2, BLACK, 3.0)

# ④ Damp
b4 = box2(x4, BY, BW, BH, YELLOW_BG, YELLOW_FG, 2.0, True)
write2(b4, [
    ("\u2463 Damp", 16, True, YELLOW_FG, CTR),
    ("", 5, False, BLACK, CTR),
    ("Compensator", 11, True, BLACK, CTR),
    ("", 4, False, BLACK, CTR),
    ("Watches lower body tilt", 10, False, GRAY, CTR),
    ("Quick bang-bang correction", 10, False, GRAY, CTR),
    ("for residual wobble", 10, False, GRAY, CTR),
], anchor='ctr')

# ── Result arrow ──
result_y = BY + BH + 0.35
arrow2(x4 + BW/2, BY + BH + 0.05, x4 + BW/2, result_y + 0.05, BLACK, 3.0)

# Result box
rb = box2(x1, result_y, x4 + BW - x1, 0.55, rgb(0xE8, 0xF8, 0xE8), GREEN, 2.5, True)
write2(rb, [
    ("u\u209c = \u2460 + \u2462 + \u2463  \u2192  acceleration sent to wheels  \u2192  robot moves WITHOUT spilling",
     13, True, rgb(0x15, 0x60, 0x28), CTR)
], anchor='ctr')

# ── Bottom section: what does u affect? ──
sec_y = result_y + 0.85

label2(0.4, sec_y, 12.5, 0.35,
       "What does the acceleration u do to the robot?", 14, True, BLACK, LFT)

# Two paths
path_y = sec_y + 0.50
path_h = 1.80

# Path 1: u → ψ
p1 = box2(0.5, path_y, 5.80, path_h, rgb(0xFD, 0xF0, 0xF8), MAGENTA, 2.0, True)
write2(p1, [
    ("u  \u2192  \u03c8 (Robot Body Tilt)", 14, True, MAGENTA, LFT),
    ("", 4, False, BLACK, LFT),
    ("Equation:  q\u0307 = A\u00b7q + B\u00b7u   (linear state-space)", 11, False, BLACK, LFT),
    ("", 3, False, BLACK, LFT),
    ("When u increases: robot tilts forward (\u03c8 increases)", 10, False, GRAY, LFT),
    ("LQT immediately corrects it back toward \u03c8 = 0\u00b0", 10, False, GRAY, LFT),
    ("Response: FAST (controlled directly by feedback)", 10, True, MAGENTA, LFT),
], anchor='t')

# Path 2: u → θ
p2 = box2(6.70, path_y, 6.20, path_h, rgb(0xFF, 0xF5, 0xE8), ORANGE, 2.0, True)
write2(p2, [
    ("u  \u2192  \u03b8 (Liquid Sloshing)", 14, True, ORANGE, LFT),
    ("", 4, False, BLACK, LFT),
    ("Equation:  \u03b8\u0308 = -(g/l)sin\u03b8 + (u/l)cos\u03b8 - damping", 11, False, BLACK, LFT),
    ("                                          (nonlinear pendulum)", 9, False, GRAY, LFT),
    ("When u increases: liquid swings backward (inertia)", 10, False, GRAY, LFT),
    ("Then oscillates at natural freq \u03c9_f = 9.9 rad/s", 10, False, GRAY, LFT),
    ("Response: SLOW (oscillatory, hard to stop once started)", 10, True, ORANGE, LFT),
], anchor='t')

# Key insight at bottom
key_y = path_y + path_h + 0.20
kb = box2(0.5, key_y, 12.4, 0.55, rgb(0xFE, 0xF0, 0xF0), RED, 2.0, True)
write2(kb, [
    ("KEY INSIGHT:  Same input u, different responses.  "
     "That's why we need Input Shaping (\u2460) — "
     "it prevents u from exciting \u03b8's natural frequency, "
     "even though \u03c8 is already handled by LQT (\u2461).",
     11, True, rgb(0x88, 0x22, 0x22), CTR)
], anchor='ctr')


# ═══════════════════════════════════════════════════════════════════════════════
#  Save
# ═══════════════════════════════════════════════════════════════════════════════
out_dir = os.path.join(os.path.dirname(__file__), 'results')
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'sbsfc_explained.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
print("2 slides:")
print("  Slide 1: Block diagram + numbered callouts + ψ vs θ explanation")
print("  Slide 2: Signal flow summary + how u affects ψ and θ differently")
