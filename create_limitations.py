#!/usr/bin/env python3
"""
create_limitations.py
Single slide summarizing the current project limitations and the
two planned next steps: RL controller + expanded scenario testing.

Run:    python3 create_limitations.py
Output: results/limitations.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def rgb(r, g, b): return RGBColor(r, g, b)

BLACK   = rgb(0x1A, 0x1A, 0x1A)
WHITE   = rgb(0xFF, 0xFF, 0xFF)
GRAY    = rgb(0x66, 0x66, 0x66)
LTGRAY  = rgb(0xDD, 0xDD, 0xDD)

RED_FG  = rgb(0xA8, 0x2A, 0x2A)
RED_BG  = rgb(0xFC, 0xEC, 0xEC)
RED_BRD = rgb(0xCC, 0x44, 0x44)

GRN_FG  = rgb(0x1E, 0x6E, 0x2E)
GRN_BG  = rgb(0xEB, 0xF7, 0xEE)
GRN_BRD = rgb(0x2E, 0x99, 0x4A)

BLU_FG  = rgb(0x1A, 0x50, 0xAA)
BLU_BG  = rgb(0xEB, 0xF3, 0xFF)
BLU_BRD = rgb(0x26, 0x66, 0xCF)

GOLD_FG  = rgb(0x99, 0x66, 0x00)
GOLD_BG  = rgb(0xFF, 0xF7, 0xE0)
GOLD_BRD = rgb(0xCC, 0x99, 0x00)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def text(x, y, w, h, s, fsize=11, color=BLACK, bold=False, italic=False,
         align=LFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.size = Pt(fsize); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color


def panel(x, y, w, h, bg, brd):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = bg
    shp.line.color.rgb = brd; shp.line.width = Pt(2)
    shp.shadow.inherit = False


def icon_badge(x, y, d, symbol, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = symbol
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE


# ============================================================
# TITLE
# ============================================================
text(0.4, 0.20, 12.5, 0.55,
     "Current Limitations and Next Steps",
     fsize=28, bold=True, color=BLACK, align=CTR)

text(0.4, 0.80, 12.5, 0.35,
     "What this project demonstrated — and what remains to be done",
     fsize=13, italic=True, color=GRAY, align=CTR)

# ============================================================
# LEFT PANEL — LIMITATIONS (red)
# ============================================================
panel(0.40, 1.35, 6.20, 5.75, RED_BG, RED_BRD)
icon_badge(0.60, 1.55, 0.55, "!", RED_BRD)
text(1.30, 1.55, 5.00, 0.55,
     "Current Limitations",
     fsize=22, bold=True, color=RED_FG, align=LFT)

limitations = [
    ("Hand-tuned classical controller",
     "SBSFC gains (K, K_c, η, Q, R) are set manually. Good for one plant, "
     "but re-tuning is needed when liquid, container, or robot change."),
    ("Known sloshing frequency assumed",
     "The input shaper needs ω_f exactly. Real liquids change with fill level, "
     "temperature, and container — a fixed notch mis-tunes quickly."),
    ("Only a few synthetic scenarios",
     "We tested 5 idealized cases (step, bump, push, rough terrain). Real "
     "delivery paths are richer — curves, obstacles, traffic."),
    ("1-D motion only",
     "Current controller assumes straight-line motion. Real serving robots turn, "
     "weave around tables, and navigate in 2-D."),
    ("No hardware validation yet",
     "All results are in simulation. Real wheel friction, motor delay, IMU noise, "
     "and liquid physics may differ from our pendulum model."),
]

yy = 2.35
for title_txt, body in limitations:
    # bullet dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.60), Inches(yy + 0.12),
                                 Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = RED_FG
    dot.line.color.rgb = RED_FG
    dot.shadow.inherit = False

    text(0.85, yy, 5.55, 0.33, title_txt,
         fsize=13, bold=True, color=RED_FG)
    text(0.85, yy + 0.35, 5.55, 0.55, body,
         fsize=11, color=BLACK)
    yy += 0.92

# ============================================================
# RIGHT PANEL — NEXT STEPS (green)
# ============================================================
panel(6.75, 1.35, 6.20, 5.75, GRN_BG, GRN_BRD)
icon_badge(6.95, 1.55, 0.55, "→", GRN_BRD)
text(7.65, 1.55, 5.00, 0.55,
     "Next Steps",
     fsize=22, bold=True, color=GRN_FG, align=LFT)

# Sub-section 1 — RL controller
sub1_y = 2.35
text(7.00, sub1_y, 5.5, 0.40,
     "① Add a Reinforcement Learning controller",
     fsize=15, bold=True, color=BLU_FG)

rl_points = [
    "Learn the control policy directly from simulation — no manual gain tuning.",
    "Adapt automatically to different liquids, fill levels, and containers.",
    "Can optimize for BOTH sloshing and tracking simultaneously, beyond what hand-tuning achieves.",
    "Compare against SBSFC as baseline → quantify how much improvement RL gives.",
]
yy = sub1_y + 0.42
for pt in rl_points:
    text(7.20, yy, 0.18, 0.25, "•", fsize=13, bold=True, color=BLU_FG)
    text(7.40, yy, 5.30, 0.55, pt, fsize=10.5, color=BLACK)
    yy += 0.40

# Sub-section 2 — Expanded scenarios
sub2_y = sub1_y + 2.55
text(7.00, sub2_y, 5.5, 0.40,
     "② Expand the scenario test set",
     fsize=15, bold=True, color=GOLD_FG)

scen_points = [
    "Add realistic restaurant paths: curves, aisles, 90° turns around tables.",
    "Vary liquid properties: fill level, viscosity, container shape.",
    "Stress tests: adversarial disturbances, multi-robot interactions.",
    "Move toward 2-D motion — currently all tests are 1-D straight-line.",
]
yy = sub2_y + 0.42
for pt in scen_points:
    text(7.20, yy, 0.18, 0.25, "•", fsize=13, bold=True, color=GOLD_FG)
    text(7.40, yy, 5.30, 0.55, pt, fsize=10.5, color=BLACK)
    yy += 0.40

# ============================================================
# FOOTER — Closing statement
# ============================================================
footer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.40), Inches(7.18),
                                Inches(12.55), Inches(0.28))
footer.fill.solid(); footer.fill.fore_color.rgb = BLU_BG
footer.line.color.rgb = BLU_BRD; footer.line.width = Pt(1)
footer.shadow.inherit = False
tf = footer.text_frame
tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
p = tf.paragraphs[0]; p.alignment = CTR
r = p.add_run()
r.text = ("SBSFC proves the concept in simulation  →  "
          "RL + richer scenarios are the path to a real-world serving robot.")
r.font.size = Pt(11); r.font.bold = True; r.font.italic = True
r.font.color.rgb = BLU_FG

# ============================================================
# SAVE
# ============================================================
os.makedirs("results", exist_ok=True)
out = "results/limitations.pptx"
prs.save(out)
print(f"Saved: {out}")
