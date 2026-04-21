#!/usr/bin/env python3
"""
create_feedback_components.py
Single-slide explainer for the three feedback pieces of SBSFC with
explicit reasoning for WHY each block is needed.

Run:    python3 create_feedback_components.py
Output: results/feedback_components.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

def rgb(r, g, b): return RGBColor(r, g, b)

BLACK  = rgb(0x1A, 0x1A, 0x1A)
WHITE  = rgb(0xFF, 0xFF, 0xFF)
GRAY   = rgb(0x66, 0x66, 0x66)
LTGRAY = rgb(0xDD, 0xDD, 0xDD)

BLUE_FG  = rgb(0x1A, 0x50, 0xAA)
BLUE_BG  = rgb(0xEB, 0xF3, 0xFF)
BLUE_BRD = rgb(0x26, 0x66, 0xCF)

GRN_FG   = rgb(0x1E, 0x6E, 0x2E)
GRN_BG   = rgb(0xEB, 0xF7, 0xEE)
GRN_BRD  = rgb(0x2E, 0x99, 0x4A)

RED_FG   = rgb(0xA8, 0x1E, 0x1E)
RED_BG   = rgb(0xFC, 0xEC, 0xEC)
RED_BRD  = rgb(0xCC, 0x33, 0x33)

PROBLEM_BG = rgb(0xFF, 0xEC, 0xEC)
PROBLEM_BR = rgb(0xCC, 0x66, 0x66)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# ----- Title -----
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.20), Inches(12.33), Inches(0.55))
tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
p = tf.paragraphs[0]; p.alignment = CTR
r = p.add_run(); r.text = "Why Three Feedback Blocks?  LQT + Compensator + DOB"
r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BLACK

# ----- Subtitle -----
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.78), Inches(12.33), Inches(0.35))
tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
p = tf.paragraphs[0]; p.alignment = CTR
r = p.add_run(); r.text = "Each block fixes a failure mode the previous one cannot"
r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = GRAY


def add_column(x_in, order, title, subtitle, problem_title, problem_text,
               eq, fix_text, fg, bg, brd):
    col_w = 4.0
    col_h = 5.0
    y = 1.25

    # background panel
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x_in), Inches(y),
                                   Inches(col_w), Inches(col_h))
    panel.fill.solid(); panel.fill.fore_color.rgb = bg
    panel.line.color.rgb = brd; panel.line.width = Pt(2)
    panel.shadow.inherit = False

    # order badge (circle "1", "2", "3")
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(x_in + 0.12), Inches(y + 0.12),
                                   Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = fg
    badge.line.color.rgb = fg
    badge.shadow.inherit = False
    tf = badge.text_frame; tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = str(order)
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE

    # title
    tb = slide.shapes.add_textbox(Inches(x_in + 0.7), Inches(y + 0.12),
                                  Inches(col_w - 0.8), Inches(0.5))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = fg

    # subtitle
    tb = slide.shapes.add_textbox(Inches(x_in + 0.7), Inches(y + 0.52),
                                  Inches(col_w - 0.8), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = subtitle
    r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = GRAY

    # "Problem it solves" header
    tb = slide.shapes.add_textbox(Inches(x_in + 0.25), Inches(y + 0.95),
                                  Inches(col_w - 0.5), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = "⚠  Problem it solves"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = rgb(0xB0, 0x44, 0x44)

    # Problem box
    prob_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x_in + 0.25), Inches(y + 1.25),
                                      Inches(col_w - 0.5), Inches(1.1))
    prob_box.fill.solid(); prob_box.fill.fore_color.rgb = PROBLEM_BG
    prob_box.line.color.rgb = PROBLEM_BR; prob_box.line.width = Pt(1)
    prob_box.shadow.inherit = False
    tf = prob_box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = problem_title
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = rgb(0x80, 0x30, 0x30)
    p2 = tf.add_paragraph(); p2.alignment = LFT
    r = p2.add_run(); r.text = problem_text
    r.font.size = Pt(11); r.font.color.rgb = BLACK

    # "How it fixes" header
    tb = slide.shapes.add_textbox(Inches(x_in + 0.25), Inches(y + 2.5),
                                  Inches(col_w - 0.5), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = "✓  How it fixes"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = fg

    # equation
    eq_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x_in + 0.25), Inches(y + 2.8),
                                    Inches(col_w - 0.5), Inches(0.55))
    eq_box.fill.solid(); eq_box.fill.fore_color.rgb = WHITE
    eq_box.line.color.rgb = brd; eq_box.line.width = Pt(1)
    eq_box.shadow.inherit = False
    tf = eq_box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = CTR
    r = p.add_run(); r.text = eq
    r.font.size = Pt(13); r.font.name = "Courier New"; r.font.bold = True
    r.font.color.rgb = BLACK

    # fix explanation
    tb = slide.shapes.add_textbox(Inches(x_in + 0.25), Inches(y + 3.45),
                                  Inches(col_w - 0.5), Inches(1.4))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = LFT
    r = p.add_run(); r.text = fix_text
    r.font.size = Pt(11); r.font.color.rgb = BLACK


# ----- Column 1: LQT -----
add_column(
    x_in=0.33, order=1,
    title="LQT",
    subtitle="Linear Quadratic Tracker",
    problem_title="Without any controller:",
    problem_text="Two-wheel robot is open-loop UNSTABLE — it just falls over. And we need to track a velocity command.",
    eq="u_lqt = K·q − K_track·q_d",
    fix_text="Uses full state q = [ψ, ψ̇, x, ẋ] to simultaneously balance the body AND track the shaped reference. Gain K is optimal for Q = penalty on ψ, ẋ and R = penalty on effort.",
    fg=BLUE_FG, bg=BLUE_BG, brd=BLUE_BRD,
)

# ----- Column 2: Compensator -----
add_column(
    x_in=4.67, order=2,
    title="Compensator",
    subtitle="Nonlinear auxiliary term",
    problem_title="LQT alone leaves:",
    problem_text="A slow pitch drift. LQT is based on the LINEAR model, but sin(ψ) ≠ ψ at real angles. The mismatch produces a tiny bias the LQT can't remove (it has no integral action).",
    eq="u_aux = K_c · sign(ψ̇_L,filt)",
    fix_text="Adds a small nonlinear nudge (K_c = 0.001) driven by the low-passed pitch-rate trend. Pushes the body back toward upright whenever drift appears. Linear LQT can't do this — sign() is nonlinear.",
    fg=GRN_FG, bg=GRN_BG, brd=GRN_BRD,
)

# ----- Column 3: DOB -----
add_column(
    x_in=9.00, order=3,
    title="DOB",
    subtitle="Disturbance Observer",
    problem_title="LQT + Compensator still fail on:",
    problem_text="External disturbances — someone pushes the robot, it hits a bump, the floor tilts. These forces are NOT in the model, so feedback alone reacts only AFTER the tilt builds up.",
    eq="u_dob = − d̂ / m",
    fix_text="Estimates the unknown force by comparing measured ψ̈ with the model's prediction — any mismatch must be external. Then feeds −d̂/m forward to cancel it BEFORE it grows. Pure feedback can't do this.",
    fg=RED_FG, bg=RED_BG, brd=RED_BRD,
)

# ----- Coverage matrix at the bottom -----
mat_y = 6.5
tb = slide.shapes.add_textbox(Inches(0.5), Inches(mat_y),
                              Inches(12.33), Inches(0.35))
tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
p = tf.paragraphs[0]; p.alignment = CTR
r = p.add_run(); r.text = "Coverage — remove any block and a failure mode returns:"
r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = GRAY

# mini table (manual using shapes)
col_x = [1.0, 5.0, 7.5, 10.0, 12.5]
headers = ["Disturbance type", "LQT", "Compensator", "DOB"]
rows = [
    ("Balance (body falls over)",   "✓",  "",   ""),
    ("Velocity tracking",           "✓",  "",   ""),
    ("Slow nonlinear pitch drift",  "✗",  "✓",  ""),
    ("External push / bump",        "✗",  "✗",  "✓"),
]

# Header row
hy = 6.85
for i, h in enumerate(headers):
    w = col_x[i+1] - col_x[i]
    tb = slide.shapes.add_textbox(Inches(col_x[i]), Inches(hy),
                                  Inches(w), Inches(0.25))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = LFT if i == 0 else CTR
    r = p.add_run(); r.text = h
    r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = (BLUE_FG if i==1 else GRN_FG if i==2 else RED_FG if i==3 else BLACK)

# Data rows
for ri, row in enumerate(rows):
    ry = hy + 0.05 + 0.24 * (ri + 1)
    for i, cell in enumerate(row):
        w = col_x[i+1] - col_x[i]
        tb = slide.shapes.add_textbox(Inches(col_x[i]), Inches(ry),
                                      Inches(w), Inches(0.22))
        tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = LFT if i == 0 else CTR
        r = p.add_run(); r.text = cell
        r.font.size = Pt(10)
        if cell == "✓":
            r.font.color.rgb = rgb(0x1E, 0x88, 0x1E); r.font.bold = True
        elif cell == "✗":
            r.font.color.rgb = rgb(0xCC, 0x33, 0x33); r.font.bold = True
        else:
            r.font.color.rgb = BLACK

# ----- Save -----
os.makedirs("results", exist_ok=True)
out = "results/feedback_components.pptx"
prs.save(out)
print(f"Saved: {out}")
