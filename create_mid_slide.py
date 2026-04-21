#!/usr/bin/env python3
"""
create_mid_slide.py
One-slide summary for the mid-project update (3-minute talk).
Four quadrants:
    (A) Objective + architecture in text
    (B) Table 3 headline numbers
    (C) Acceleration + FFT figure (why it works)
    (D) Roadblocks -> next steps

Run:    python3 create_mid_slide.py
Output: results/mid_slide.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def rgb(r, g, b): return RGBColor(r, g, b)

BLACK  = rgb(0x1A, 0x1A, 0x1A)
WHITE  = rgb(0xFF, 0xFF, 0xFF)
GRAY   = rgb(0x66, 0x66, 0x66)

BLU_FG = rgb(0x1A, 0x50, 0xAA); BLU_BG = rgb(0xEB, 0xF3, 0xFF); BLU_BRD = rgb(0x26, 0x66, 0xCF)
GRN_FG = rgb(0x1E, 0x6E, 0x2E); GRN_BG = rgb(0xEB, 0xF7, 0xEE); GRN_BRD = rgb(0x2E, 0x99, 0x4A)
RED_FG = rgb(0xA8, 0x2A, 0x2A); RED_BG = rgb(0xFC, 0xEC, 0xEC); RED_BRD = rgb(0xCC, 0x44, 0x44)
ORG_FG = rgb(0xB5, 0x55, 0x00); ORG_BG = rgb(0xFF, 0xF2, 0xE0); ORG_BRD = rgb(0xE0, 0x85, 0x20)

CTR = PP_ALIGN.CENTER
LFT = PP_ALIGN.LEFT

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def text(x, y, w, h, s, fsize=11, color=BLACK, bold=False, italic=False, align=LFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.size = Pt(fsize); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color


def panel(x, y, w, h, bg, brd, title, fg):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = bg
    shp.line.color.rgb = brd; shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    text(x + 0.15, y + 0.08, w - 0.3, 0.33, title,
         fsize=13, bold=True, color=fg, align=LFT)


# ---------------- TITLE ----------------
text(0.4, 0.15, 12.5, 0.50,
     "Anti-Sloshing Controller for a Food-Serving Robot  —  Mid-Project Update",
     fsize=22, bold=True, color=BLACK, align=CTR)
text(0.4, 0.65, 12.5, 0.30,
     "Hyundae Cha   |   MECE 6397   |   Extending Choi et al. 2024 (SBSFC)  →  RL + richer scenarios",
     fsize=12, italic=True, color=GRAY, align=CTR)

# ============================================================
# (A) Objective + architecture  — top-left
# ============================================================
panel(0.35, 1.05, 6.30, 2.90, BLU_BG, BLU_BRD,
      "(A)  Problem & Approach", BLU_FG)

text(0.55, 1.50, 6.0, 0.32,
     "Sloshing is driven by ACCELERATION, not velocity.",
     fsize=12, bold=True, color=BLU_FG)
text(0.55, 1.82, 6.0, 0.60,
     "A low-pass filter (LPF) smooths velocity but leaves a broadband "
     "acceleration spike at every step — which excites the liquid's "
     "resonance \u03C9_f.",
     fsize=10.5, color=BLACK)

text(0.55, 2.55, 6.0, 0.32,
     "SBSFC (reproduced in MATLAB):",
     fsize=12, bold=True, color=BLU_FG)

arch_lines = [
    ("Input shaper  F_c(s)", "feed-forward — zeros cancel sloshing poles"),
    ("LQT gain  K",          "state feedback on [x, \u03C8, \u1E8B, \u03C8\u0307]"),
    ("DOB + compensator",    "observe \u03C8\u0308 residual + sign(\u03C8\u0307) correction"),
]
yy = 2.90
for sym, desc in arch_lines:
    text(0.70, yy, 2.3, 0.28, "\u2022 " + sym,
         fsize=10.5, bold=True, color=BLU_FG)
    text(3.00, yy, 3.55, 0.28, desc, fsize=10.5, color=BLACK)
    yy += 0.28

text(0.55, 3.75, 6.0, 0.25,
     "Plant: two-wheel self-balancing robot + pendulum-equivalent slosh "
     "(\u03C9_f = 1.29 Hz).",
     fsize=9.5, italic=True, color=GRAY)

# ============================================================
# (B) Table 3 headline — top-right
# ============================================================
panel(6.75, 1.05, 6.25, 2.90, GRN_BG, GRN_BRD,
      "(B)  Preliminary Results  —  Scenario 1 (sudden start / stop)", GRN_FG)

# Build a small visual table
tbl_x, tbl_y = 6.95, 1.55
col_w = [2.70, 1.25, 1.25, 1.00]
headers = ["Metric", "SBSFC", "LPF", "Reduction"]

# Header row
cx = tbl_x
for i, hd in enumerate(headers):
    text(cx, tbl_y, col_w[i], 0.30, hd,
         fsize=11, bold=True, color=GRN_FG,
         align=CTR if i > 0 else LFT)
    cx += col_w[i]

# separator line (thin shape)
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(tbl_x), Inches(tbl_y + 0.32),
                             Inches(sum(col_w)), Inches(0.02))
sep.fill.solid(); sep.fill.fore_color.rgb = GRN_BRD
sep.line.fill.background(); sep.shadow.inherit = False

rows = [
    ("Mean  |\u03B8|   [deg]",        "0.317",  "5.066",  "93.7 %"),
    ("Variance  \u03B8   [deg\u00B2]", "0.193",  "35.413", "99.5 %"),
    ("Peak  |a|   [m/s\u00B2]",        "0.20",   "1.14",   "5.7\u00D7 smaller"),
]
ry = tbl_y + 0.42
for row in rows:
    cx = tbl_x
    for i, val in enumerate(row):
        text(cx, ry, col_w[i], 0.30, val,
             fsize=11,
             bold=(i == 3),
             color=(GRN_FG if i == 3 else BLACK),
             align=CTR if i > 0 else LFT)
        cx += col_w[i]
    ry += 0.38

# Highlight box
hl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(6.95), Inches(3.30),
                            Inches(5.85), Inches(0.55))
hl.fill.solid(); hl.fill.fore_color.rgb = WHITE
hl.line.color.rgb = GRN_BRD; hl.line.width = Pt(1.2)
hl.shadow.inherit = False
text(7.05, 3.38, 5.70, 0.40,
     "Sloshing angle reduced 93.7 %  —  spectral notch exactly at \u03C9_f = 1.29 Hz",
     fsize=11, bold=True, color=GRN_FG, align=CTR)

# ============================================================
# (C) Mechanism figure — bottom-left
# ============================================================
panel(0.35, 4.05, 6.30, 3.10, ORG_BG, ORG_BRD,
      "(C)  Why it works  —  acceleration & FFT", ORG_FG)

fig_path = "results/velocity_chain.png"
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(0.50), Inches(4.48),
                             width=Inches(6.00))
else:
    text(0.55, 5.30, 6.0, 0.5,
         "(velocity_chain.png not found — run plot_velocity_chain.m)",
         fsize=11, italic=True, color=GRAY)

text(0.55, 6.90, 6.0, 0.22,
     "Velocities look identical  \u2192  accelerations differ 5.7\u00D7  \u2192  "
     "SBSFC notches FFT at \u03C9_f.",
     fsize=9.5, italic=True, color=ORG_FG)

# ============================================================
# (D) Roadblocks -> Next steps — bottom-right
# ============================================================
panel(6.75, 4.05, 6.25, 3.10, RED_BG, RED_BRD,
      "(D)  Roadblocks  \u2192  Plan for remainder", RED_FG)

# Left column: roadblocks
text(6.95, 4.48, 2.90, 0.30,
     "Current limits",
     fsize=12, bold=True, color=RED_FG)
limits = [
    "Gains hand-tuned per plant",
    "\u03C9_f assumed known & fixed",
    "Only 1-D synthetic scenarios",
    "No hardware validation yet",
]
yy = 4.80
for it in limits:
    text(7.00, yy, 0.20, 0.28, "\u2022", fsize=11, bold=True, color=RED_FG)
    text(7.18, yy, 2.70, 0.28, it, fsize=10, color=BLACK)
    yy += 0.32

# Arrow divider
arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                             Inches(9.80), Inches(5.30),
                             Inches(0.35), Inches(0.40))
arr.fill.solid(); arr.fill.fore_color.rgb = GRN_BRD
arr.line.color.rgb = GRN_BRD
arr.shadow.inherit = False

# Right column: plan
text(10.25, 4.48, 2.75, 0.30,
     "Plan (rest of semester)",
     fsize=12, bold=True, color=GRN_FG)
plans = [
    "RL policy (PPO / SAC) \u2014 no manual tuning",
    "Domain randomization over \u03C9_f & fill",
    "Curved / 2-D paths, table turns",
    "Baseline: SBSFC  vs  RL on same plant",
]
yy = 4.80
for it in plans:
    text(10.30, yy, 0.20, 0.28, "\u2022", fsize=11, bold=True, color=GRN_FG)
    text(10.48, yy, 2.55, 0.28, it, fsize=10, color=BLACK)
    yy += 0.32

# Timeline strip
tl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(6.95), Inches(6.55),
                            Inches(5.85), Inches(0.50))
tl.fill.solid(); tl.fill.fore_color.rgb = WHITE
tl.line.color.rgb = BLU_BRD; tl.line.width = Pt(1)
tl.shadow.inherit = False
text(7.05, 6.62, 5.70, 0.38,
     "Timeline:  RL pipeline (late Apr)  \u2192  PPO vs SBSFC (early May)  "
     "\u2192  randomized-\u03C9_f eval (final report)",
     fsize=9.5, bold=True, italic=True, color=BLU_FG, align=CTR)

# ============================================================
# FOOTER
# ============================================================
text(0.4, 7.22, 12.5, 0.22,
     "Reference: Choi et al., \"Self-Balancing Slosh-Free Control of a "
     "Two-Wheeled Food-Serving Robot,\" Mechatronics, 2024.",
     fsize=9, italic=True, color=GRAY, align=CTR)

# ============================================================
# SAVE
# ============================================================
os.makedirs("results", exist_ok=True)
out = "results/mid_slide.pptx"
prs.save(out)
print(f"Saved: {out}")
